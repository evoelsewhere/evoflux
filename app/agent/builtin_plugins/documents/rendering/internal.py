"""Document-plugin renderers used by artifact QA and preview pipelines.

The renderers intentionally consume the same OOXML models that the authoring
engines write.  They do not shell out to LibreOffice, Poppler, Chromium, or any
other machine-wide application, so desktop and server builds behave the same.
"""

from __future__ import annotations

from dataclasses import dataclass
from io import BytesIO
import math
from pathlib import Path
import re
from typing import Any
from xml.etree import ElementTree

from PIL import Image, ImageChops, ImageColor, ImageDraw, ImageFont
import pypdfium2 as pdfium

from app.agent.builtin_plugins.documents.rendering.xlsx_formula import (
    FormulaEvaluation,
    evaluate_workbook_formulas,
    format_computed_value,
)


_FONT_ROOT = (
    Path(__file__).resolve().parents[3]
    / "builtin_skills"
    / "canvas-design"
    / "canvas-fonts"
)

_FONT_ALIASES = {
    "arial": "WorkSans",
    "helvetica": "WorkSans",
    "helveticaneue": "WorkSans",
    "sansserif": "WorkSans",
    "systemui": "WorkSans",
    "applesystem": "WorkSans",
    "timesnewroman": "Lora",
    "times": "Lora",
    "serif": "Lora",
    "couriernew": "JetBrainsMono",
    "courier": "JetBrainsMono",
    "monospace": "JetBrainsMono",
}
_FONT_FALLBACK_FAMILY = "WorkSans"


def _font(
    size: int,
    *,
    bold: bool = False,
    italic: bool = False,
    family: str | None = None,
) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    safe_family = re.sub(r"[^A-Za-z0-9]", "", family or _FONT_FALLBACK_FAMILY)
    safe_family = _FONT_ALIASES.get(safe_family.casefold(), safe_family)
    suffix = (
        "BoldItalic"
        if bold and italic
        else "Bold"
        if bold
        else "Italic"
        if italic
        else "Regular"
    )
    candidate = _FONT_ROOT / f"{safe_family}-{suffix}.ttf"
    if not candidate.is_file():
        fallback = (
            f"{_FONT_FALLBACK_FAMILY}-BoldItalic.ttf"
            if bold and italic
            else f"{_FONT_FALLBACK_FAMILY}-Bold.ttf"
            if bold
            else f"{_FONT_FALLBACK_FAMILY}-Italic.ttf"
            if italic
            else f"{_FONT_FALLBACK_FAMILY}-Regular.ttf"
        )
        candidate = _FONT_ROOT / fallback
    try:
        return ImageFont.truetype(str(candidate), max(8, size))
    except OSError:
        return ImageFont.load_default()


def _color(value: Any, default: str = "#ffffff") -> str:
    text = str(value or "").strip().lstrip("#")
    if len(text) == 8:
        text = text[-6:]
    if len(text) == 6:
        try:
            ImageColor.getrgb(f"#{text}")
            return f"#{text}"
        except ValueError:
            pass
    return default


def _wrapped_lines(
    draw: ImageDraw.ImageDraw,
    value: str,
    font: ImageFont.FreeTypeFont | ImageFont.ImageFont,
    width: int,
) -> list[str]:
    if width <= 1:
        return [value]
    lines: list[str] = []
    for paragraph in value.splitlines() or [""]:
        words = paragraph.split()
        if not words:
            lines.append("")
            continue
        current = words[0]
        for word in words[1:]:
            candidate = f"{current} {word}"
            if draw.textbbox((0, 0), candidate, font=font)[2] <= width:
                current = candidate
            else:
                lines.append(current)
                current = word
        lines.append(current)
    return lines


def _draw_text_box(
    draw: ImageDraw.ImageDraw,
    box: tuple[int, int, int, int],
    value: str,
    *,
    color: str = "#111827",
    size: int = 18,
    bold: bool = False,
    align: str = "left",
) -> None:
    left, top, right, bottom = box
    padding = max(3, size // 5)
    font = _font(size, bold=bold)
    lines = _wrapped_lines(draw, value, font, max(1, right - left - 2 * padding))
    line_height = max(size + 3, draw.textbbox((0, 0), "Ag", font=font)[3] + 3)
    y = top + padding
    for line in lines:
        if y + line_height > bottom:
            break
        text_width = draw.textbbox((0, 0), line, font=font)[2]
        if align == "center":
            x = left + max(padding, (right - left - text_width) // 2)
        elif align == "right":
            x = right - padding - text_width
        else:
            x = left + padding
        draw.text((x, y), line, fill=color, font=font)
        y += line_height


def count_pdf_pages(source: Path) -> int:
    """Return a PDF page count without rasterizing any page."""

    document = pdfium.PdfDocument(str(source))
    try:
        return len(document)
    finally:
        document.close()


def render_pdf_pages(
    source: Path,
    render_dir: Path,
    *,
    dpi: int = 144,
    max_pages: int | None = None,
    max_total_bytes: int | None = None,
    max_pixels_per_page: int | None = None,
) -> list[Path]:
    """Rasterize PDF pages with optional resource bounds.

    Artifact QA leaves the limits unset so it can inspect every page. The
    interactive preview supplies strict limits to keep untrusted documents
    from exhausting sidecar or WebView memory.
    """

    render_dir.mkdir(parents=True, exist_ok=True)
    document = pdfium.PdfDocument(str(source))
    outputs: list[Path] = []
    scale = dpi / 72
    total_bytes = 0
    try:
        page_count = len(document)
        if max_pages is not None:
            page_count = min(page_count, max(0, max_pages))
        for index in range(page_count):
            page = document[index]
            try:
                render_scale = scale
                if max_pixels_per_page is not None:
                    width, height = page.get_size()
                    pixels = width * height * render_scale * render_scale
                    if not math.isfinite(pixels) or pixels <= 0:
                        raise ValueError("PDF page has invalid dimensions")
                    if pixels > max_pixels_per_page:
                        render_scale *= math.sqrt(max_pixels_per_page / pixels)
                bitmap = page.render(scale=render_scale)
                try:
                    image = bitmap.to_pil()
                    destination = render_dir / f"page-{index + 1:03d}.png"
                    try:
                        image.save(destination)
                    finally:
                        image.close()
                finally:
                    bitmap.close()
            finally:
                page.close()
            output_bytes = destination.stat().st_size
            if (
                max_total_bytes is not None
                and total_bytes + output_bytes > max_total_bytes
            ):
                destination.unlink(missing_ok=True)
                break
            total_bytes += output_bytes
            outputs.append(destination)
    finally:
        document.close()
    return outputs


def render_docx_pages(source: Path, render_dir: Path, *, dpi: int = 144) -> list[Path]:
    """Render a deterministic semantic DOCX preview with Pillow.

    This is deliberately an OOXML-model renderer, not an Office emulator.  It
    catches overflow, missing content, table density, and pagination regressions
    without requiring a separately installed office suite.
    """
    from docx import Document

    document = Document(str(source))
    width = int(8.27 * dpi)
    height = int(11.69 * dpi)
    margin = int(0.65 * dpi)
    pages: list[Image.Image] = []
    image = Image.new("RGB", (width, height), "white")
    draw = ImageDraw.Draw(image)
    y = margin

    def new_page() -> None:
        nonlocal image, draw, y
        pages.append(image)
        image = Image.new("RGB", (width, height), "white")
        draw = ImageDraw.Draw(image)
        y = margin

    def ensure(space: int) -> None:
        if y + space > height - margin:
            new_page()

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            y += 10
            continue
        style_name = (paragraph.style.name if paragraph.style else "").lower()
        heading = "heading" in style_name or "title" in style_name
        size = 30 if "title" in style_name else 22 if heading else 14
        font = _font(size, bold=heading)
        lines = _wrapped_lines(draw, text, font, width - 2 * margin)
        line_height = size + 6
        ensure(max(line_height, len(lines) * line_height + 8))
        for line in lines:
            draw.text((margin, y), line, fill="#172033", font=font)
            y += line_height
        y += 8

    for table in document.tables:
        rows = len(table.rows)
        columns = max((len(row.cells) for row in table.rows), default=0)
        if rows == 0 or columns == 0:
            continue
        cell_width = (width - 2 * margin) // columns
        row_height = 42
        ensure(min(rows, 5) * row_height + 16)
        for row_index, row in enumerate(table.rows):
            ensure(row_height)
            for column_index, cell in enumerate(row.cells):
                box = (
                    margin + column_index * cell_width,
                    y,
                    margin + (column_index + 1) * cell_width,
                    y + row_height,
                )
                fill = "#e8eef7" if row_index == 0 else "#ffffff"
                draw.rectangle(box, fill=fill, outline="#aab4c3", width=1)
                _draw_text_box(
                    draw,
                    box,
                    cell.text,
                    size=11,
                    bold=row_index == 0,
                )
            y += row_height
        y += 16

    pages.append(image)
    render_dir.mkdir(parents=True, exist_ok=True)
    outputs: list[Path] = []
    for index, page in enumerate(pages, start=1):
        destination = render_dir / f"page-{index:03d}.png"
        page.save(destination)
        outputs.append(destination)
    return outputs


def _xlsx_cell_color(cell: Any) -> str:
    fill = getattr(cell, "fill", None)
    foreground = getattr(fill, "fgColor", None)
    if getattr(fill, "fill_type", None) == "solid":
        return _color(getattr(foreground, "rgb", None), "#ffffff")
    return "#ffffff"


@dataclass(frozen=True)
class XlsxRenderEvidence:
    paths: list[Path]
    chart_count: int
    rendered_chart_count: int


def _xlsx_column_width(sheet: Any, column: int) -> int:
    from openpyxl.utils import get_column_letter

    return max(
        42,
        min(
            280,
            int(
                (sheet.column_dimensions[get_column_letter(column)].width or 8.43) * 7
                + 12
            ),
        ),
    )


def _xlsx_row_height(sheet: Any, row: int) -> int:
    return max(24, min(120, int((sheet.row_dimensions[row].height or 18) * 1.34)))


def _xlsx_anchor_box(
    sheet: Any,
    drawing: Any,
    column_widths: list[int],
    row_heights: list[int],
) -> tuple[int, int, int, int]:
    marker = drawing.anchor._from
    left = 1 + sum(column_widths[: int(marker.col)])
    top = 1 + sum(row_heights[: int(marker.row)])
    left += int(float(marker.colOff or 0) / 914400 * 96)
    top += int(float(marker.rowOff or 0) / 914400 * 96)
    extent = getattr(drawing.anchor, "ext", None)
    if extent is not None:
        width = int(float(extent.cx) / 914400 * 96)
        height = int(float(extent.cy) / 914400 * 96)
    else:
        width = int(float(getattr(drawing, "width", 320) or 320))
        height = int(float(getattr(drawing, "height", 180) or 180))
    return left, top, left + max(width, 1), top + max(height, 1)


def _xlsx_source_formula(series: Any, branch: str) -> str:
    source = getattr(series, branch, None)
    if source is None:
        return ""
    reference = getattr(source, "numRef", None) or getattr(source, "strRef", None)
    return str(getattr(reference, "f", "") or "")


def _xlsx_series_values(
    chart: Any, evaluation: FormulaEvaluation, current_sheet: str
) -> list[list[float]]:
    rows: list[list[float]] = []
    for series in chart.ser:
        source = getattr(series, "val", None) or getattr(series, "yVal", None)
        literal = getattr(source, "numLit", None) if source is not None else None
        raw: list[Any]
        if literal is not None:
            raw = [point.v for point in getattr(literal, "pt", ())]
        else:
            formula = _xlsx_source_formula(
                series, "val" if getattr(series, "val", None) is not None else "yVal"
            )
            if not formula:
                continue
            try:
                raw = evaluation.reference_values(formula, current_sheet)
            except ValueError:
                continue
        values: list[float] = []
        for value in raw:
            try:
                values.append(float(value))
            except (TypeError, ValueError):
                values.append(0.0)
        if values:
            rows.append(values)
    return rows


def _render_xlsx_chart(
    draw: ImageDraw.ImageDraw,
    box: tuple[int, int, int, int],
    chart: Any,
    evaluation: FormulaEvaluation,
    sheet_name: str,
) -> bool:
    left, top, right, bottom = box
    if right - left < 60 or bottom - top < 50:
        return False
    values = _xlsx_series_values(chart, evaluation, sheet_name)
    if not values:
        return False
    draw.rounded_rectangle(box, radius=6, fill="#ffffff", outline="#c7ced8", width=1)
    plot = (left + 38, top + 32, right - 16, bottom - 28)
    plot_left, plot_top, plot_right, plot_bottom = plot
    if plot_right <= plot_left or plot_bottom <= plot_top:
        return False
    palette = ("#2563eb", "#0f766e", "#f59e0b", "#dc2626", "#7c3aed")
    flat = [value for row in values for value in row]
    maximum = max(max((abs(value) for value in flat), default=0.0), 1.0)
    chart_type = type(chart).__name__.upper()
    draw.line((plot_left, plot_bottom, plot_right, plot_bottom), fill="#94a3b8")
    draw.line((plot_left, plot_top, plot_left, plot_bottom), fill="#94a3b8")
    if "PIE" in chart_type or "DOUGHNUT" in chart_type:
        first = [abs(value) for value in values[0]]
        total = sum(first) or 1
        diameter = max(20, min(plot_right - plot_left, plot_bottom - plot_top))
        pie_box = (
            plot_left + (plot_right - plot_left - diameter) // 2,
            plot_top + (plot_bottom - plot_top - diameter) // 2,
            plot_left + (plot_right - plot_left + diameter) // 2,
            plot_top + (plot_bottom - plot_top + diameter) // 2,
        )
        angle = -90.0
        for index, value in enumerate(first):
            next_angle = angle + 360 * value / total
            draw.pieslice(
                pie_box,
                start=angle,
                end=next_angle,
                fill=palette[index % len(palette)],
            )
            angle = next_angle
        if "DOUGHNUT" in chart_type:
            inset = diameter // 3
            draw.ellipse(
                (
                    pie_box[0] + inset,
                    pie_box[1] + inset,
                    pie_box[2] - inset,
                    pie_box[3] - inset,
                ),
                fill="#ffffff",
            )
        return True
    if "LINE" in chart_type or "AREA" in chart_type or "SCATTER" in chart_type:
        for series_index, row in enumerate(values):
            points = [
                (
                    int(
                        plot_left
                        + index * (plot_right - plot_left) / max(len(row) - 1, 1)
                    ),
                    int(plot_bottom - value / maximum * (plot_bottom - plot_top)),
                )
                for index, value in enumerate(row)
            ]
            if "AREA" in chart_type and points:
                draw.polygon(
                    [
                        (points[0][0], plot_bottom),
                        *points,
                        (points[-1][0], plot_bottom),
                    ],
                    fill=palette[series_index % len(palette)] + "55",
                )
            if len(points) > 1:
                draw.line(points, fill=palette[series_index % len(palette)], width=3)
            for x, y in points:
                draw.ellipse(
                    (x - 3, y - 3, x + 3, y + 3),
                    fill=palette[series_index % len(palette)],
                )
        return True
    count = max(len(row) for row in values)
    group_width = max(1.0, (plot_right - plot_left) / max(count, 1))
    bar_width = max(2.0, group_width * 0.72 / max(len(values), 1))
    for series_index, row in enumerate(values):
        for index, value in enumerate(row):
            height = abs(value) / maximum * (plot_bottom - plot_top)
            x = plot_left + index * group_width + series_index * bar_width
            y = plot_bottom - height if value >= 0 else plot_bottom
            draw.rectangle(
                (int(x), int(y), int(x + bar_width - 1), int(y + height)),
                fill=palette[series_index % len(palette)],
            )
    return True


def render_xlsx_workbook_with_evidence(
    workbook: Any,
    render_dir: Path,
    *,
    evaluation: FormulaEvaluation | None = None,
) -> XlsxRenderEvidence:
    """Render worksheets with evaluated formulas and native chart drawings."""

    evaluation = evaluation or evaluate_workbook_formulas(workbook)
    render_dir.mkdir(parents=True, exist_ok=True)
    outputs: list[Path] = []
    chart_count = 0
    rendered_chart_count = 0
    for sheet_index, sheet in enumerate(workbook.worksheets, start=1):
        max_row = max(1, min(sheet.max_row, 250))
        max_column = max(1, min(sheet.max_column, 60))
        chart_markers = [
            chart.anchor._from
            for chart in sheet._charts  # noqa: SLF001
        ]
        layout_column_count = max(
            [max_column, *(int(marker.col) + 1 for marker in chart_markers)]
        )
        layout_row_count = max(
            [max_row, *(int(marker.row) + 1 for marker in chart_markers)]
        )
        column_widths = [
            _xlsx_column_width(sheet, column)
            for column in range(1, layout_column_count + 1)
        ]
        row_heights = [
            _xlsx_row_height(sheet, row) for row in range(1, layout_row_count + 1)
        ]
        chart_boxes = [
            _xlsx_anchor_box(sheet, chart, column_widths, row_heights)
            for chart in sheet._charts  # noqa: SLF001
        ]
        chart_count += len(chart_boxes)
        width = min(
            5000,
            max([sum(column_widths) + 2, *(box[2] + 2 for box in chart_boxes)]),
        )
        height = min(
            7000,
            max([sum(row_heights) + 2, *(box[3] + 2 for box in chart_boxes)]),
        )
        image = Image.new("RGB", (width, height), "white")
        draw = ImageDraw.Draw(image)
        y = 1
        for row in range(1, max_row + 1):
            x = 1
            row_height = row_heights[row - 1]
            for column in range(1, max_column + 1):
                column_width = column_widths[column - 1]
                if x >= width:
                    break
                right = min(width - 1, x + column_width)
                bottom = min(height - 1, y + row_height)
                cell = sheet.cell(row, column)
                draw.rectangle(
                    (x, y, right, bottom),
                    fill=_xlsx_cell_color(cell),
                    outline="#c7ced8",
                )
                computed = evaluation.display_value(
                    sheet.title, cell.coordinate, cell.value
                )
                value = format_computed_value(computed, str(cell.number_format or ""))
                font = getattr(cell, "font", None)
                alignment = getattr(cell, "alignment", None)
                _draw_text_box(
                    draw,
                    (x, y, right, bottom),
                    value,
                    color=_color(
                        getattr(getattr(font, "color", None), "rgb", None), "#1f2937"
                    ),
                    size=max(9, min(18, int(getattr(font, "sz", None) or 11))),
                    bold=bool(getattr(font, "bold", False)),
                    align=str(getattr(alignment, "horizontal", None) or "left"),
                )
                x += column_width
            y += row_height
            if y >= height:
                break
        for chart, box in zip(sheet._charts, chart_boxes, strict=True):  # noqa: SLF001
            if box[2] > width or box[3] > height:
                continue
            if _render_xlsx_chart(draw, box, chart, evaluation, sheet.title):
                rendered_chart_count += 1
        destination = render_dir / f"sheet-{sheet_index:03d}.png"
        image.save(destination)
        outputs.append(destination)
    return XlsxRenderEvidence(
        paths=outputs,
        chart_count=chart_count,
        rendered_chart_count=rendered_chart_count,
    )


def render_xlsx_workbook(workbook: Any, render_dir: Path) -> list[Path]:
    """Backward-compatible path-only wrapper around the evidence renderer."""

    return render_xlsx_workbook_with_evidence(workbook, render_dir).paths


def render_xlsx_file(source: Path, render_dir: Path) -> list[Path]:
    from openpyxl import load_workbook

    workbook = load_workbook(source, data_only=False)
    try:
        return render_xlsx_workbook(workbook, render_dir)
    finally:
        workbook.close()


_PPTX_THEME_NAMES = {
    "ACCENT_1": "accent1",
    "ACCENT_2": "accent2",
    "ACCENT_3": "accent3",
    "ACCENT_4": "accent4",
    "ACCENT_5": "accent5",
    "ACCENT_6": "accent6",
    "BACKGROUND_1": "lt1",
    "BACKGROUND_2": "lt2",
    "DARK_1": "dk1",
    "DARK_2": "dk2",
    "FOLLOWED_HYPERLINK": "folHlink",
    "HYPERLINK": "hlink",
    "LIGHT_1": "lt1",
    "LIGHT_2": "lt2",
    "TEXT_1": "dk1",
    "TEXT_2": "dk2",
}


def _pptx_theme_colors(owner: Any) -> dict[str, str]:
    """Return the active master theme's color scheme as RGB hex values."""

    try:
        master = owner.slide_masters[0]
    except (AttributeError, IndexError, TypeError):
        master = owner
    theme_part = next(
        (
            relationship.target_part
            for relationship in master.part.rels.values()
            if relationship.reltype.endswith("/theme")
        ),
        None,
    )
    if theme_part is None:
        return {}
    try:
        root = ElementTree.fromstring(theme_part.blob)
    except (ElementTree.ParseError, AttributeError, TypeError):
        return {}
    scheme = next(
        (
            element
            for element in root.iter()
            if element.tag.rsplit("}", 1)[-1] == "clrScheme"
        ),
        None,
    )
    if scheme is None:
        return {}
    colors: dict[str, str] = {}
    for wrapper in scheme:
        name = wrapper.tag.rsplit("}", 1)[-1]
        if not len(wrapper):
            continue
        color = wrapper[0]
        color_name = color.tag.rsplit("}", 1)[-1]
        value = color.get("val") if color_name == "srgbClr" else color.get("lastClr")
        if value and re.fullmatch(r"[0-9A-Fa-f]{6}", value):
            colors[name] = f"#{value.upper()}"
    return colors


def _pptx_apply_brightness(color: str, brightness: float) -> str:
    try:
        channels = ImageColor.getrgb(color)
        red, green, blue = channels[:3]
    except ValueError:
        return color
    amount = max(-1.0, min(1.0, brightness))
    if amount >= 0:
        channels = tuple(
            round(channel + (255 - channel) * amount) for channel in (red, green, blue)
        )
    else:
        channels = tuple(
            round(channel * (1 + amount)) for channel in (red, green, blue)
        )
    return "#" + "".join(f"{channel:02X}" for channel in channels)


def _pptx_rgb(
    value: Any, default: str, theme_colors: dict[str, str] | None = None
) -> str:
    try:
        rgb = value.rgb
    except (AttributeError, TypeError, ValueError):
        rgb = None
    color = _color(rgb, "") if rgb is not None else ""
    if not color and theme_colors:
        try:
            theme_name = value.theme_color.name
        except (AttributeError, TypeError, ValueError):
            theme_name = ""
        color = theme_colors.get(_PPTX_THEME_NAMES.get(theme_name, ""), "")
    if not color:
        return default
    try:
        brightness = float(value.brightness or 0)
    except (AttributeError, TypeError, ValueError):
        brightness = 0
    return _pptx_apply_brightness(color, brightness)


def _pptx_xml_flag(element: Any, name: str, *, default: bool = True) -> bool:
    if element is None:
        return default
    value = element.get(name)
    if value is None:
        return default
    return str(value).strip().lower() not in {"0", "false", "off", "no"}


def _pptx_has_explicit_background(owner: Any) -> bool:
    try:
        return bool(
            owner._element.xpath(  # noqa: SLF001
                "./*[local-name()='cSld']/*[local-name()='bg']"
            )
        )
    except (AttributeError, TypeError, ValueError):
        return False


def _pptx_background(slide: Any, theme_colors: dict[str, str]) -> str:
    owners = [slide]
    try:
        if slide.follow_master_background:
            owners.extend((slide.slide_layout, slide.slide_layout.slide_master))
    except (AttributeError, TypeError, ValueError):
        pass
    for owner in owners:
        if not _pptx_has_explicit_background(owner):
            continue
        try:
            return _pptx_rgb(owner.background.fill.fore_color, "#ffffff", theme_colors)
        except (AttributeError, TypeError, ValueError):
            continue
    return "#ffffff"


def _pptx_layered_shapes(slide: Any) -> list[Any]:
    """Return visible master, layout, and slide shapes in paint order.

    Master and layout placeholders describe inherited geometry and styles; they
    are not separate slideshow paint layers. Rendering only the concrete slide
    placeholder prevents template prompts and footer tokens from being doubled.
    """

    layout = slide.slide_layout
    master = layout.slide_master
    shapes: list[Any] = []
    show_master_shapes = _pptx_xml_flag(slide._element, "showMasterSp") and (  # noqa: SLF001
        _pptx_xml_flag(layout._element, "showMasterSp")  # noqa: SLF001
    )
    if show_master_shapes:
        for shape in master.shapes:
            if not getattr(shape, "is_placeholder", False):
                shapes.append(shape)
    for shape in layout.shapes:
        if not getattr(shape, "is_placeholder", False):
            shapes.append(shape)
    shapes.extend(slide.shapes)
    return shapes


@dataclass(frozen=True)
class _PptxRunStyle:
    size: int
    color: str
    family: str | None
    bold: bool
    italic: bool
    underline: bool


@dataclass
class _PptxTextLine:
    items: list[tuple[str, _PptxRunStyle]]
    width: float
    height: int


def _pptx_length_pixels(value: Any, scale: float) -> float:
    try:
        return max(0.0, float(value) * scale)
    except (TypeError, ValueError):
        return 0.0


def _pptx_run_style(
    run: Any,
    paragraph: Any,
    *,
    font_scale: float,
    theme_colors: dict[str, str],
) -> _PptxRunStyle:
    run_font = getattr(run, "font", None)
    paragraph_font = getattr(paragraph, "font", None)

    def inherited(name: str, default: Any) -> Any:
        run_value = getattr(run_font, name, None)
        if run_value is not None:
            return run_value
        paragraph_value = getattr(paragraph_font, name, None)
        return default if paragraph_value is None else paragraph_value

    size_value = inherited("size", None)
    try:
        points = float(size_value.pt)
    except (AttributeError, TypeError, ValueError):
        points = 18.0
    size = max(8, min(160, round(points * font_scale)))
    color_value = getattr(run_font, "color", None)
    if color_value is None or getattr(color_value, "type", None) is None:
        color_value = getattr(paragraph_font, "color", None)
    return _PptxRunStyle(
        size=size,
        color=_pptx_rgb(color_value, "#111827", theme_colors),
        family=inherited("name", None),
        bold=bool(inherited("bold", False)),
        italic=bool(inherited("italic", False)),
        underline=bool(inherited("underline", False)),
    )


def _pptx_text_width(
    draw: ImageDraw.ImageDraw, text: str, style: _PptxRunStyle
) -> float:
    return float(
        draw.textlength(
            text,
            font=_font(
                style.size,
                bold=style.bold,
                italic=style.italic,
                family=style.family,
            ),
        )
    )


def _pptx_line_height(style: _PptxRunStyle) -> int:
    font = _font(
        style.size,
        bold=style.bold,
        italic=style.italic,
        family=style.family,
    )
    bounds = font.getbbox("Ag")
    return round(max(style.size + 2, bounds[3] - bounds[1] + 3))


def _pptx_wrap_runs(
    draw: ImageDraw.ImageDraw,
    runs: list[tuple[str, _PptxRunStyle]],
    width: int,
) -> list[_PptxTextLine]:
    maximum = max(1, width)
    lines: list[_PptxTextLine] = []
    items: list[tuple[str, _PptxRunStyle]] = []
    line_width = 0.0
    line_height = 0

    def append_line(*, force: bool = False) -> None:
        nonlocal items, line_width, line_height
        while items and items[-1][0].isspace():
            text, style = items.pop()
            line_width -= _pptx_text_width(draw, text, style)
        if items or force:
            lines.append(
                _PptxTextLine(
                    items=list(items),
                    width=max(0.0, line_width),
                    height=max(12, line_height),
                )
            )
        items = []
        line_width = 0.0
        line_height = 0

    def append_piece(piece: str, style: _PptxRunStyle) -> None:
        nonlocal line_width, line_height
        if not piece:
            return
        piece_width = _pptx_text_width(draw, piece, style)
        items.append((piece, style))
        line_width += piece_width
        line_height = max(line_height, _pptx_line_height(style))

    for text, style in runs:
        for token in re.findall(r"\v|\n|[^\S\v\n]+|[^\s\v\n]+", text):
            if token in {"\v", "\n"}:
                append_line(force=True)
                continue
            if token.isspace() and not items:
                continue
            token_width = _pptx_text_width(draw, token, style)
            if items and line_width + token_width > maximum:
                append_line()
                if token.isspace():
                    continue
            if token_width <= maximum:
                append_piece(token, style)
                continue
            for character in token:
                character_width = _pptx_text_width(draw, character, style)
                if items and line_width + character_width > maximum:
                    append_line()
                append_piece(character, style)
    append_line(force=not lines)
    return lines


def _pptx_bullet_marker(
    paragraph: Any, counters: dict[tuple[int, str], int]
) -> str | None:
    properties = paragraph._p.pPr  # noqa: SLF001
    if properties is None:
        return None
    if properties.xpath("./*[local-name()='buNone']"):
        return None
    characters = properties.xpath("./*[local-name()='buChar']")
    if characters:
        return characters[0].get("char") or "•"
    automatic = properties.xpath("./*[local-name()='buAutoNum']")
    if not automatic:
        return None
    node = automatic[0]
    kind = node.get("type") or "arabicPeriod"
    level = int(getattr(paragraph, "level", 0) or 0)
    key = (level, kind)
    start = int(node.get("startAt") or 1)
    value = counters.get(key, start)
    counters[key] = value + 1
    if kind.startswith("alphaLc"):
        label = chr(ord("a") + (value - 1) % 26)
    elif kind.startswith("alphaUc"):
        label = chr(ord("A") + (value - 1) % 26)
    elif kind.startswith("romanLc"):
        label = str(value)
    elif kind.startswith("romanUc"):
        label = str(value)
    else:
        label = str(value)
    if kind.endswith("ParenBoth"):
        return f"({label})"
    if kind.endswith("ParenR"):
        return f"{label})"
    return f"{label}."


def _pptx_paragraph_alignment(paragraph: Any) -> str:
    try:
        name = paragraph.alignment.name
    except (AttributeError, TypeError, ValueError):
        return "left"
    if name in {"CENTER", "DISTRIBUTE", "THAI_DISTRIBUTE"}:
        return "center"
    if name == "RIGHT":
        return "right"
    return "left"


def _pptx_draw_line(
    draw: ImageDraw.ImageDraw,
    line: _PptxTextLine,
    *,
    x: float,
    y: float,
) -> None:
    cursor = x
    for text, style in line.items:
        font = _font(
            style.size,
            bold=style.bold,
            italic=style.italic,
            family=style.family,
        )
        draw.text((cursor, y), text, fill=style.color, font=font)
        text_width = _pptx_text_width(draw, text, style)
        if style.underline and text.strip():
            underline_y = y + _pptx_line_height(style) - 2
            draw.line(
                (cursor, underline_y, cursor + text_width, underline_y),
                fill=style.color,
                width=max(1, style.size // 14),
            )
        cursor += text_width


def _pptx_draw_text_frame(
    layer: Image.Image,
    shape: Any,
    *,
    x_scale: float,
    y_scale: float,
    theme_colors: dict[str, str],
) -> None:
    draw = ImageDraw.Draw(layer)
    frame = shape.text_frame
    width, height = layer.size
    left = _pptx_length_pixels(frame.margin_left, x_scale)
    right = width - _pptx_length_pixels(frame.margin_right, x_scale)
    top = _pptx_length_pixels(frame.margin_top, y_scale)
    bottom = height - _pptx_length_pixels(frame.margin_bottom, y_scale)
    content_width = max(1, round(right - left))
    font_scale = x_scale * 914_400 / 72
    counters: dict[tuple[int, str], int] = {}
    laid_out: list[
        tuple[
            Any,
            list[_PptxTextLine],
            str | None,
            _PptxRunStyle,
            float,
            float,
            float,
            float,
            float,
        ]
    ] = []
    total_height = 0.0
    for paragraph in frame.paragraphs:
        runs = list(paragraph.runs)
        if runs:
            styled_runs = [
                (
                    run.text,
                    _pptx_run_style(
                        run,
                        paragraph,
                        font_scale=font_scale,
                        theme_colors=theme_colors,
                    ),
                )
                for run in runs
            ]
        else:
            synthetic = type("SyntheticRun", (), {"font": None})()
            styled_runs = [
                (
                    paragraph.text,
                    _pptx_run_style(
                        synthetic,
                        paragraph,
                        font_scale=font_scale,
                        theme_colors=theme_colors,
                    ),
                )
            ]
        marker_style = styled_runs[0][1]
        marker = _pptx_bullet_marker(paragraph, counters)
        properties = paragraph._p.pPr  # noqa: SLF001
        margin_left = 0.0
        indent = 0.0
        if properties is not None:
            try:
                margin_left = float(properties.get("marL") or 0) * x_scale
                indent = float(properties.get("indent") or 0) * x_scale
            except (TypeError, ValueError):
                margin_left = 0.0
                indent = 0.0
        text_left = max(0.0, margin_left)
        first_left = max(0.0, margin_left + (0 if marker else indent))
        available = max(1, round(content_width - min(text_left, first_left)))
        lines = _pptx_wrap_runs(draw, styled_runs, available)
        try:
            line_spacing = float(paragraph.line_spacing or 1.0)
        except (TypeError, ValueError):
            line_spacing = 1.0
        if line_spacing > 10:
            # Absolute DrawingML line spacing is expressed in EMUs.
            absolute_line_height = line_spacing * y_scale
            ratio = 1.0
        else:
            absolute_line_height = 0.0
            ratio = max(0.5, min(4.0, line_spacing))
        for line in lines:
            line.height = max(
                line.height, round(line.height * ratio), round(absolute_line_height)
            )
        before = _pptx_length_pixels(paragraph.space_before, y_scale)
        after = _pptx_length_pixels(paragraph.space_after, y_scale)
        paragraph_height = sum(line.height for line in lines)
        laid_out.append(
            (
                paragraph,
                lines,
                marker,
                marker_style,
                text_left,
                first_left,
                indent,
                before,
                after,
            )
        )
        total_height += before + paragraph_height + after
    try:
        anchor_name = frame.vertical_anchor.name
    except (AttributeError, TypeError, ValueError):
        anchor_name = "TOP"
    if anchor_name == "MIDDLE":
        y = top + max(0.0, (bottom - top - total_height) / 2)
    elif anchor_name == "BOTTOM":
        y = max(top, bottom - total_height)
    else:
        y = top
    for (
        paragraph,
        lines,
        marker,
        marker_style,
        text_left,
        first_left,
        indent,
        before,
        after,
    ) in laid_out:
        y += before
        alignment = _pptx_paragraph_alignment(paragraph)
        for line_index, line in enumerate(lines):
            base_left = text_left if marker or line_index else first_left
            available = max(1.0, content_width - base_left)
            if alignment == "center":
                line_x = left + base_left + max(0.0, (available - line.width) / 2)
            elif alignment == "right":
                line_x = left + base_left + max(0.0, available - line.width)
            else:
                line_x = left + base_left
            if marker and line_index == 0:
                marker_width = _pptx_text_width(draw, marker, marker_style)
                marker_x = left + max(0.0, text_left + min(indent, -marker_width - 3))
                _pptx_draw_line(
                    draw,
                    _PptxTextLine(
                        [(marker, marker_style)],
                        marker_width,
                        _pptx_line_height(marker_style),
                    ),
                    x=marker_x,
                    y=y,
                )
            _pptx_draw_line(draw, line, x=line_x, y=y)
            y += line.height
            if y > bottom:
                return
        y += after
        if y > bottom:
            return


def _pptx_transform_image(image: Image.Image, shape: Any) -> Image.Image:
    transform = shape._element.spPr.xfrm  # noqa: SLF001
    if _pptx_xml_flag(transform, "flipH", default=False):
        image = image.transpose(Image.Transpose.FLIP_LEFT_RIGHT)
    if _pptx_xml_flag(transform, "flipV", default=False):
        image = image.transpose(Image.Transpose.FLIP_TOP_BOTTOM)
    try:
        rotation = float(shape.rotation or 0)
    except (AttributeError, TypeError, ValueError):
        rotation = 0
    if rotation % 360:
        image = image.rotate(-rotation, expand=True, resample=Image.Resampling.BICUBIC)
    return image


def _pptx_paste_centered(
    canvas: Image.Image, overlay: Image.Image, box: tuple[int, int, int, int]
) -> None:
    center_x = (box[0] + box[2]) / 2
    center_y = (box[1] + box[3]) / 2
    left = round(center_x - overlay.width / 2)
    top = round(center_y - overlay.height / 2)
    canvas.paste(overlay, (left, top), overlay)


def _pptx_picture_layer(shape: Any, width: int, height: int) -> Image.Image:
    with Image.open(BytesIO(shape.image.blob)) as source:
        picture = source.convert("RGBA")
    crop_left = max(0.0, min(0.99, float(shape.crop_left or 0)))
    crop_right = max(0.0, min(0.99, float(shape.crop_right or 0)))
    crop_top = max(0.0, min(0.99, float(shape.crop_top or 0)))
    crop_bottom = max(0.0, min(0.99, float(shape.crop_bottom or 0)))
    if crop_left + crop_right < 0.999 and crop_top + crop_bottom < 0.999:
        left = round(picture.width * crop_left)
        top = round(picture.height * crop_top)
        right = max(left + 1, round(picture.width * (1 - crop_right)))
        bottom = max(top + 1, round(picture.height * (1 - crop_bottom)))
        picture = picture.crop((left, top, right, bottom))
    picture = picture.resize((max(1, width), max(1, height)), Image.Resampling.LANCZOS)
    try:
        geometry = shape._element.spPr.prstGeom.get("prst")  # noqa: SLF001
    except (AttributeError, TypeError, ValueError):
        geometry = ""
    if geometry == "ellipse":
        mask = Image.new("L", picture.size, 0)
        ImageDraw.Draw(mask).ellipse(
            (0, 0, picture.width - 1, picture.height - 1), fill=255
        )
        alpha = picture.getchannel("A")
        picture.putalpha(ImageChops.multiply(alpha, mask))
    return _pptx_transform_image(picture, shape)


def _render_chart(
    draw: ImageDraw.ImageDraw, box: tuple[int, int, int, int], chart: Any
) -> None:
    left, top, right, bottom = box
    draw.rectangle(box, fill="#ffffff", outline="#c6cfdb", width=1)
    try:
        series = list(chart.series)
        values = [list(item.values) for item in series]
    except (AttributeError, TypeError, ValueError):
        return
    flat = [float(value or 0) for row in values for value in row]
    if not flat:
        return
    maximum = max(max(flat), 1)
    count = max(len(row) for row in values)
    group_width = max(1, (right - left - 24) // max(1, count))
    palette = ("#2563eb", "#14b8a6", "#f59e0b", "#ef4444", "#8b5cf6")
    for series_index, row in enumerate(values):
        bar_width = max(2, group_width // max(1, len(values)) - 2)
        for index, value in enumerate(row):
            magnitude = max(0.0, float(value or 0)) / maximum
            x = left + 12 + index * group_width + series_index * (bar_width + 2)
            y = bottom - 12 - int((bottom - top - 28) * magnitude)
            draw.rectangle(
                (x, y, x + bar_width, bottom - 12),
                fill=palette[series_index % len(palette)],
            )


def render_pptx_pages(
    source: Path, render_dir: Path, *, width: int = 1280
) -> list[Path]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    presentation = Presentation(str(source))
    slide_width = int(presentation.slide_width or 1)
    slide_height = int(presentation.slide_height or 1)
    height = max(1, int(width * slide_height / slide_width))
    x_scale = width / slide_width
    y_scale = height / slide_height
    render_dir.mkdir(parents=True, exist_ok=True)
    outputs: list[Path] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        theme_colors = _pptx_theme_colors(slide.slide_layout.slide_master)
        background = _pptx_background(slide, theme_colors)
        image = Image.new("RGB", (width, height), background)
        draw = ImageDraw.Draw(image)
        for shape in _pptx_layered_shapes(slide):
            box = (
                int(shape.left * x_scale),
                int(shape.top * y_scale),
                int((shape.left + shape.width) * x_scale),
                int((shape.top + shape.height) * y_scale),
            )
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or hasattr(shape, "image"):
                try:
                    picture = _pptx_picture_layer(
                        shape,
                        max(1, box[2] - box[0]),
                        max(1, box[3] - box[1]),
                    )
                    _pptx_paste_centered(image, picture, box)
                except (AttributeError, OSError, TypeError, ValueError):
                    draw.rectangle(box, fill="#e5e7eb", outline="#94a3b8")
                continue
            if getattr(shape, "has_table", False):
                table = shape.table
                row_heights = [max(1, int(row.height * y_scale)) for row in table.rows]
                column_widths = [
                    max(1, int(column.width * x_scale)) for column in table.columns
                ]
                row_top = box[1]
                for row_index, row in enumerate(table.rows):
                    column_left = box[0]
                    for column_index, cell in enumerate(row.cells):
                        if getattr(cell, "is_spanned", False):
                            column_left += column_widths[column_index]
                            continue
                        cell_width = sum(
                            column_widths[
                                column_index : column_index
                                + int(getattr(cell, "span_width", 1) or 1)
                            ]
                        )
                        cell_height = sum(
                            row_heights[
                                row_index : row_index
                                + int(getattr(cell, "span_height", 1) or 1)
                            ]
                        )
                        cell_box = (
                            column_left,
                            row_top,
                            column_left + cell_width,
                            row_top + cell_height,
                        )
                        fill = "#e8eef7" if row_index == 0 else "#ffffff"
                        try:
                            if cell.fill.type is not None:
                                fill = _pptx_rgb(
                                    cell.fill.fore_color, fill, theme_colors
                                )
                        except (AttributeError, TypeError, ValueError):
                            pass
                        draw.rectangle(
                            cell_box,
                            fill=fill,
                            outline="#94a3b8",
                        )
                        text_layer = Image.new(
                            "RGBA",
                            (
                                max(1, cell_box[2] - cell_box[0]),
                                max(1, cell_box[3] - cell_box[1]),
                            ),
                            (0, 0, 0, 0),
                        )
                        cell_shape = type(
                            "TableCellText",
                            (),
                            {"text_frame": cell.text_frame},
                        )()
                        _pptx_draw_text_frame(
                            text_layer,
                            cell_shape,
                            x_scale=x_scale,
                            y_scale=y_scale,
                            theme_colors=theme_colors,
                        )
                        image.paste(
                            text_layer,
                            (cell_box[0], cell_box[1]),
                            text_layer,
                        )
                        column_left += column_widths[column_index]
                    row_top += row_heights[row_index]
                continue
            if getattr(shape, "has_chart", False):
                _render_chart(draw, box, shape.chart)
                continue
            fill: str | None = None
            line: str | None = None
            try:
                if shape.fill.type is not None:
                    fill = _pptx_rgb(shape.fill.fore_color, "#ffffff", theme_colors)
            except (AttributeError, TypeError, ValueError):
                pass
            try:
                if shape.line.fill.type is not None:
                    line = _pptx_rgb(shape.line.color, "#94a3b8", theme_colors)
            except (AttributeError, TypeError, ValueError):
                pass
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                draw.rectangle(box, fill=fill, outline=line, width=1 if line else 0)
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                text_layer = Image.new(
                    "RGBA",
                    (
                        max(1, box[2] - box[0]),
                        max(1, box[3] - box[1]),
                    ),
                    (0, 0, 0, 0),
                )
                _pptx_draw_text_frame(
                    text_layer,
                    shape,
                    x_scale=x_scale,
                    y_scale=y_scale,
                    theme_colors=theme_colors,
                )
                text_layer = _pptx_transform_image(text_layer, shape)
                _pptx_paste_centered(image, text_layer, box)
        destination = render_dir / f"slide-{slide_index:03d}.png"
        image.save(destination)
        outputs.append(destination)
    return outputs


__all__ = [
    "count_pdf_pages",
    "render_docx_pages",
    "render_pdf_pages",
    "render_pptx_pages",
    "render_xlsx_file",
    "render_xlsx_workbook",
]
