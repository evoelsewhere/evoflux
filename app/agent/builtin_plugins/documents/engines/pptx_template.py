"""Plugin-owned PPTX template engine using direct OOXML cloning."""

from __future__ import annotations

import asyncio
from copy import deepcopy
from dataclasses import dataclass, field
import hashlib
import json
from pathlib import Path
import re
from typing import Annotated, Any, Literal

from pydantic import BaseModel, ConfigDict, Field, model_validator

from app.agent.builtin_plugins.documents.rendering.internal import render_pptx_pages
from app.agent.builtin_plugins.documents.rendering.runtime import file_sha256

pptx_sha256 = file_sha256
MAX_TEMPLATE_SLIDES = 80
MAX_EDITS_PER_SLIDE = 160
MAX_PLACEHOLDER_FILLS_PER_SLIDE = 32


class TemplateObjectEdit(BaseModel):
    model_config = ConfigDict(extra="forbid")
    operation: Literal[
        "set_text",
        "replace_text",
        "replace_image",
        "set_table_cell",
        "set_chart_series",
    ]
    target_id: str = Field(pattern=r"^(?:sh|im|tb|ch)/[1-9][0-9]*/[1-9][0-9]*$")
    text: str | None = Field(default=None, max_length=20_000)
    find: str | None = Field(default=None, min_length=1, max_length=4_000)
    replace: str | None = Field(default=None, max_length=20_000)
    asset_path: str | None = Field(default=None, max_length=2_000)
    alt: str | None = Field(default=None, max_length=1_000)
    row: int | None = Field(default=None, ge=0, le=500)
    column: int | None = Field(default=None, ge=0, le=500)
    series_index: int | None = Field(default=None, ge=0, le=100)
    values: list[float] | None = Field(default=None, min_length=1, max_length=2_000)

    @model_validator(mode="after")
    def validate_operation_fields(self) -> TemplateObjectEdit:
        required = {
            "set_text": self.text is not None,
            "replace_text": self.find is not None and self.replace is not None,
            "replace_image": bool(self.asset_path),
            "set_table_cell": self.row is not None
            and self.column is not None
            and self.text is not None,
            "set_chart_series": self.series_index is not None
            and self.values is not None,
        }
        if not required[self.operation]:
            raise ValueError(f"{self.operation} is missing required fields")
        return self


class TemplateTextPlaceholderFill(BaseModel):
    model_config = ConfigDict(extra="forbid")
    operation: Literal["set_text"]
    placeholder_idx: int = Field(ge=0, le=65_535)
    placeholder_type: Literal[
        "TITLE",
        "CENTER_TITLE",
        "SUBTITLE",
        "BODY",
        "OBJECT",
        "VERTICAL_TITLE",
        "VERTICAL_BODY",
        "VERTICAL_OBJECT",
    ]
    text: str = Field(max_length=20_000)


class TemplateImagePlaceholderFill(BaseModel):
    model_config = ConfigDict(extra="forbid")
    operation: Literal["set_image"]
    placeholder_idx: int = Field(ge=0, le=65_535)
    placeholder_type: Literal["PICTURE"]
    asset_path: str = Field(min_length=1, max_length=2_000)
    alt: str | None = Field(default=None, max_length=1_000)

    @model_validator(mode="after")
    def validate_raster_asset(self) -> TemplateImagePlaceholderFill:
        if Path(self.asset_path).suffix.lower() not in {
            ".png",
            ".jpg",
            ".jpeg",
            ".gif",
            ".webp",
        }:
            raise ValueError("set_image requires a supported raster image asset")
        return self


TemplatePlaceholderFill = Annotated[
    TemplateTextPlaceholderFill | TemplateImagePlaceholderFill,
    Field(discriminator="operation"),
]


class TemplateSlidePlan(BaseModel):
    model_config = ConfigDict(extra="forbid")
    output_slide: int = Field(ge=1, le=MAX_TEMPLATE_SLIDES)
    source_slide: int = Field(ge=1, le=MAX_TEMPLATE_SLIDES)
    narrative_role: str = Field(min_length=1, max_length=240)
    reuse_mode: Literal["duplicate-slide", "use-layout"] = "duplicate-slide"
    edits: list[TemplateObjectEdit] = Field(
        default_factory=list, max_length=MAX_EDITS_PER_SLIDE
    )
    placeholder_fills: list[TemplatePlaceholderFill] = Field(
        default_factory=list, max_length=MAX_PLACEHOLDER_FILLS_PER_SLIDE
    )
    speaker_notes: str | None = Field(default=None, max_length=40_000)

    @model_validator(mode="after")
    def validate_reuse_mode(self) -> TemplateSlidePlan:
        if self.reuse_mode == "duplicate-slide":
            if self.placeholder_fills:
                raise ValueError(
                    "duplicate-slide uses edits; placeholder_fills require use-layout"
                )
            return self
        if self.edits:
            raise ValueError("use-layout uses placeholder_fills; edits are not allowed")
        if not self.placeholder_fills:
            raise ValueError("use-layout requires at least one placeholder fill")
        targets = [
            (fill.placeholder_idx, fill.placeholder_type)
            for fill in self.placeholder_fills
        ]
        if len(targets) != len(set(targets)):
            raise ValueError("placeholder_fills must target each placeholder once")
        return self


class TemplateDeckProject(BaseModel):
    model_config = ConfigDict(extra="forbid")
    schema_version: Literal[2] = 2
    title: str = Field(min_length=1, max_length=240)
    source_sha256: str = Field(pattern=r"^[0-9a-f]{64}$")
    template_confirmed: Literal[True]
    output_slides: list[TemplateSlidePlan] = Field(
        min_length=1, max_length=MAX_TEMPLATE_SLIDES
    )
    omitted_source_slides: list[int] = Field(default_factory=list)

    @model_validator(mode="after")
    def validate_slide_map(self) -> TemplateDeckProject:
        expected = list(range(1, len(self.output_slides) + 1))
        actual = [slide.output_slide for slide in self.output_slides]
        if actual != expected:
            raise ValueError("output_slides must be ordered sequentially from 1")
        if len(self.omitted_source_slides) != len(set(self.omitted_source_slides)):
            raise ValueError("omitted_source_slides must not contain duplicates")
        return self


@dataclass
class TemplatePipelineResult:
    action: str
    source_pptx: Path
    work_dir: Path
    manifest_path: Path | None = None
    output: Path | None = None
    previews: list[Path] = field(default_factory=list)
    layout_paths: list[Path] = field(default_factory=list)
    slide_count: int = 0
    issues: list[dict[str, Any]] = field(default_factory=list)
    metadata: dict[str, Any] = field(default_factory=dict)

    @property
    def passed(self) -> bool:
        return not any(issue.get("severity") == "error" for issue in self.issues)


def load_template_project(path: Path) -> TemplateDeckProject:
    return TemplateDeckProject.model_validate_json(path.read_text(encoding="utf-8"))


def load_template_manifest(path: Path) -> dict[str, Any]:
    value = json.loads(path.read_text(encoding="utf-8"))
    if not isinstance(value, dict):
        raise ValueError("template manifest must be an object")
    return value


def _is_picture_shape(shape: Any) -> bool:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return True
    # A populated PowerPoint picture placeholder remains shape_type=PLACEHOLDER
    # in python-pptx even though its OOXML element is a native ``p:pic``.
    return str(getattr(shape._element, "tag", "")).rsplit("}", 1)[-1] == "pic"


def _record_kind(shape: Any) -> tuple[str, str]:
    shape_type = str(getattr(getattr(shape, "shape_type", None), "name", ""))
    if shape_type == "GROUP":
        return "gr", "group"
    if getattr(shape, "has_table", False):
        return "tb", "table"
    if getattr(shape, "has_chart", False):
        return "ch", "chart"
    if _is_picture_shape(shape):
        return "im", "image"
    return "sh", "textbox" if getattr(shape, "has_text_frame", False) else "shape"


def _placeholder_metadata(shape: Any) -> dict[str, Any] | None:
    if not getattr(shape, "is_placeholder", False):
        return None
    try:
        placeholder_type = shape.placeholder_format.type
        type_name = str(getattr(placeholder_type, "name", placeholder_type))
        index = int(shape.placeholder_format.idx)
    except (AttributeError, TypeError, ValueError):
        return None
    metadata: dict[str, Any] = {"idx": index, "type": type_name}
    try:
        nodes = shape._element.xpath(".//*[local-name()='ph']")
    except (AttributeError, TypeError, ValueError):
        nodes = []
    if nodes:
        for source_name, output_name in (("orient", "orientation"), ("sz", "size")):
            value = nodes[0].get(source_name)
            if value:
                metadata[output_name] = str(value)
    return metadata


def _picture_alt_text(shape: Any) -> str:
    try:
        properties = shape._element.xpath(
            "./*[local-name()='nvPicPr']/*[local-name()='cNvPr']"
        )
    except (AttributeError, TypeError, ValueError):
        properties = []
    if not properties:
        return ""
    return str(properties[0].get("descr") or properties[0].get("title") or "")


def _picture_crop(shape: Any) -> dict[str, float]:
    crop: dict[str, float] = {}
    for name in ("left", "right", "top", "bottom"):
        try:
            crop[name] = float(getattr(shape, f"crop_{name}") or 0)
        except (AttributeError, TypeError, ValueError):
            crop[name] = 0.0
    return crop


def _picture_geometry(shape: Any) -> str:
    try:
        nodes = shape._element.xpath(
            "./*[local-name()='spPr']/*[local-name()='prstGeom']"
        )
    except (AttributeError, TypeError, ValueError):
        nodes = []
    return str(nodes[0].get("prst") or "rect") if nodes else "rect"


def _picture_transform(shape: Any) -> dict[str, Any]:
    try:
        nodes = shape._element.xpath("./*[local-name()='spPr']/*[local-name()='xfrm']")
    except (AttributeError, TypeError, ValueError):
        nodes = []
    transform = nodes[0] if nodes else None
    return {
        "rotation": float(getattr(shape, "rotation", 0) or 0),
        "flipHorizontal": bool(transform is not None and transform.get("flipH") == "1"),
        "flipVertical": bool(transform is not None and transform.get("flipV") == "1"),
    }


def _image_metadata(shape: Any) -> dict[str, Any]:
    image = shape.image
    pixel_width, pixel_height = image.size
    dpi_x, dpi_y = image.dpi
    crop = _picture_crop(shape)
    relation_id = shape._pic.blip_rId  # noqa: SLF001 - inspect native image lineage
    related_part = shape.part.related_part(relation_id)
    frame_width_inches = max(float(shape.width) / 914_400, 0.0001)
    frame_height_inches = max(float(shape.height) / 914_400, 0.0001)
    visible_width = max(0.0, 1 - crop["left"] - crop["right"])
    visible_height = max(0.0, 1 - crop["top"] - crop["bottom"])
    return {
        "relationshipId": relation_id,
        "partName": str(related_part.partname),
        "contentType": str(related_part.content_type),
        "sha256": hashlib.sha256(image.blob).hexdigest(),
        "filename": str(image.filename),
        "pixelWidth": int(pixel_width),
        "pixelHeight": int(pixel_height),
        "dpiX": int(dpi_x),
        "dpiY": int(dpi_y),
        "effectiveDpiX": round(pixel_width * visible_width / frame_width_inches, 2),
        "effectiveDpiY": round(pixel_height * visible_height / frame_height_inches, 2),
        "crop": crop,
        "altText": _picture_alt_text(shape),
        "geometry": _picture_geometry(shape),
        **_picture_transform(shape),
    }


def _shape_record(
    shape: Any,
    slide_number: int,
    *,
    layer: str = "slide",
    z_order: int | None = None,
    addressable: bool = True,
) -> dict[str, Any]:
    prefix, kind = _record_kind(shape)
    record: dict[str, Any] = {
        "id": f"{prefix}/{slide_number}/{shape.shape_id}",
        "slide": slide_number,
        "layer": layer,
        "shapeId": shape.shape_id,
        "kind": kind,
        "addressable": addressable and kind in {"textbox", "image", "table", "chart"},
        "name": shape.name,
        "left": int(shape.left),
        "top": int(shape.top),
        "width": int(shape.width),
        "height": int(shape.height),
    }
    if z_order is not None:
        record["zOrder"] = z_order
    placeholder = _placeholder_metadata(shape)
    if placeholder is not None:
        record["placeholder"] = placeholder
    if getattr(shape, "has_text_frame", False):
        record["text"] = shape.text
    if getattr(shape, "has_table", False):
        record.update(
            {"rows": len(shape.table.rows), "columns": len(shape.table.columns)}
        )
    if getattr(shape, "has_chart", False):
        record["seriesCount"] = len(shape.chart.series)
    if kind == "image":
        record["image"] = _image_metadata(shape)
    if kind == "group":
        record["children"] = [
            _shape_record(
                child,
                slide_number,
                layer=layer,
                z_order=child_z_order,
                addressable=False,
            )
            for child_z_order, child in enumerate(shape.shapes)
        ]
    return record


def _inspect_records(presentation: Any) -> list[dict[str, Any]]:
    return [
        _shape_record(shape, slide_number, z_order=z_order)
        for slide_number, slide in enumerate(presentation.slides, start=1)
        for z_order, shape in enumerate(slide.shapes)
    ]


def _container_name(container: Any) -> str:
    name = str(getattr(container, "name", "") or "")
    if name:
        return name
    try:
        return str(container._element.cSld.get("name") or "")  # noqa: SLF001
    except (AttributeError, TypeError, ValueError):
        return ""


def _inherited_container_record(
    container: Any, *, slide_number: int, layer: Literal["layout", "master"]
) -> dict[str, Any]:
    records = [
        _shape_record(shape, slide_number, layer=layer, z_order=z_order)
        for z_order, shape in enumerate(container.shapes)
    ]
    return {
        "partName": str(container.part.partname),
        "sha256": hashlib.sha256(container.part.blob).hexdigest(),
        "name": _container_name(container),
        "placeholders": [
            record for record in records if isinstance(record.get("placeholder"), dict)
        ],
        "records": records,
    }


def _slide_lineage(slide: Any, slide_number: int) -> dict[str, Any]:
    layout = slide.slide_layout
    master = layout.slide_master
    slide_placeholders = [
        {
            "shapeId": int(shape.shape_id),
            **(metadata or {}),
        }
        for shape in slide.shapes
        if (metadata := _placeholder_metadata(shape)) is not None
    ]
    return {
        "slide": slide_number,
        "slideId": int(slide.slide_id),
        "layout": _inherited_container_record(
            layout, slide_number=slide_number, layer="layout"
        ),
        "master": _inherited_container_record(
            master, slide_number=slide_number, layer="master"
        ),
        "slidePlaceholders": slide_placeholders,
    }


def _presentation_lineage(presentation: Any) -> list[dict[str, Any]]:
    return [
        _slide_lineage(slide, slide_number)
        for slide_number, slide in enumerate(presentation.slides, start=1)
    ]


def _manifest_lineage_by_slide(manifest: dict[str, Any]) -> dict[int, dict[str, Any]]:
    values = manifest.get("slideLineage")
    if not isinstance(values, list):
        return {}
    lineage: dict[int, dict[str, Any]] = {}
    for value in values:
        if not isinstance(value, dict):
            continue
        try:
            slide_number = int(value.get("slide", 0))
        except (TypeError, ValueError):
            continue
        if slide_number > 0:
            lineage[slide_number] = value
    return lineage


def _layout_placeholder_records(lineage: dict[str, Any]) -> list[dict[str, Any]]:
    layout = lineage.get("layout")
    values = layout.get("placeholders") if isinstance(layout, dict) else None
    return [value for value in values or [] if isinstance(value, dict)]


def _validate_layout_placeholder_fill(
    fill: TemplatePlaceholderFill,
    *,
    lineage: dict[str, Any],
    source_slide: int,
) -> None:
    records = _layout_placeholder_records(lineage)
    indexed: list[dict[str, Any]] = []
    exact: list[dict[str, Any]] = []
    for record in records:
        placeholder = record.get("placeholder")
        if not isinstance(placeholder, dict):
            continue
        try:
            placeholder_idx = int(placeholder.get("idx", -1))
        except (TypeError, ValueError):
            continue
        if placeholder_idx != fill.placeholder_idx:
            continue
        indexed.append(record)
        if str(placeholder.get("type") or "") == fill.placeholder_type:
            exact.append(record)
    if len(exact) == 1:
        return
    if len(exact) > 1:
        raise ValueError(
            f"source slide {source_slide} layout contains duplicate placeholder "
            f"idx/type {fill.placeholder_idx}/{fill.placeholder_type}"
        )
    if indexed:
        actual_types = sorted(
            {str(record.get("placeholder", {}).get("type") or "") for record in indexed}
        )
        raise ValueError(
            f"source slide {source_slide} layout placeholder idx "
            f"{fill.placeholder_idx} has type {actual_types}, not "
            f"{fill.placeholder_type}"
        )
    raise ValueError(
        f"source slide {source_slide} layout has no placeholder idx "
        f"{fill.placeholder_idx}"
    )


def validate_template_project(
    project: TemplateDeckProject, manifest: dict[str, Any], *, source_pptx: Path
) -> dict[str, Any]:
    actual_digest = file_sha256(source_pptx)
    expected_digest = str(manifest.get("sourceSha256") or "")
    if actual_digest != project.source_sha256 or actual_digest != expected_digest:
        raise ValueError("source PPTX changed after inspection; inspect it again")
    source_slide_count = int(manifest.get("slideCount", 0))
    used = {slide.source_slide for slide in project.output_slides}
    if any(slide < 1 or slide > source_slide_count for slide in used):
        raise ValueError("source_slide is outside the inspected presentation")
    expected_omissions = set(range(1, source_slide_count + 1)) - used
    if set(project.omitted_source_slides) != expected_omissions:
        raise ValueError(
            "omitted_source_slides must explicitly list every unused source slide"
        )
    records = {str(record.get("id")): record for record in manifest.get("records", [])}
    lineage_by_slide = _manifest_lineage_by_slide(manifest)
    operation_kind = {
        "set_text": {"textbox"},
        "replace_text": {"textbox"},
        "replace_image": {"image"},
        "set_table_cell": {"table"},
        "set_chart_series": {"chart"},
    }
    for slide in project.output_slides:
        if slide.reuse_mode == "use-layout":
            lineage = lineage_by_slide.get(slide.source_slide)
            if lineage is None:
                raise ValueError(
                    "use-layout requires slideLineage from a fresh template inspection"
                )
            layout = lineage.get("layout")
            master = lineage.get("master")
            if not isinstance(layout, dict) or not layout.get("partName"):
                raise ValueError(
                    f"source slide {slide.source_slide} has no inspected layout lineage"
                )
            if not isinstance(master, dict) or not master.get("partName"):
                raise ValueError(
                    f"source slide {slide.source_slide} has no inspected master lineage"
                )
            for fill in slide.placeholder_fills:
                _validate_layout_placeholder_fill(
                    fill,
                    lineage=lineage,
                    source_slide=slide.source_slide,
                )
        for edit in slide.edits:
            record = records.get(edit.target_id)
            if record is None:
                raise ValueError(
                    f"target_id was not found in inspect manifest: {edit.target_id}"
                )
            if record.get("addressable") is False:
                raise ValueError(f"target {edit.target_id} is not directly addressable")
            if int(record.get("slide", 0)) != slide.source_slide:
                raise ValueError(
                    f"target {edit.target_id} does not belong to source slide {slide.source_slide}"
                )
            kind = str(record.get("kind"))
            if kind not in operation_kind[edit.operation]:
                expected = next(iter(operation_kind[edit.operation]))
                raise ValueError(f"{edit.operation} requires a {expected} target")
    return {
        "valid": True,
        "source_slide_count": source_slide_count,
        "output_slide_count": len(project.output_slides),
        "edit_count": sum(len(slide.edits) for slide in project.output_slides),
        "placeholder_fill_count": sum(
            len(slide.placeholder_fills) for slide in project.output_slides
        ),
        "use_layout_slide_count": sum(
            1 for slide in project.output_slides if slide.reuse_mode == "use-layout"
        ),
        "preserve_only_slide_count": sum(
            1
            for slide in project.output_slides
            if not slide.edits and not slide.placeholder_fills
        ),
    }


def template_catalog() -> dict[str, Any]:
    return {
        "workflow": "direct-openxml-pptx-template",
        "style_behavior": {
            "uploaded_template_is_style_confirmation": True,
            "ask_style_question": False,
            "ambiguous_upload": "Ask whether the PPTX is a visual template or only a content source.",
        },
        "invariants": [
            "Duplicate source slide XML or create a new slide from its inspected native layout.",
            "Preserve masters, layouts, themes, transitions, timing, and untouched objects.",
            "Edit only stable shape IDs declared by inspect.",
            "Fill use-layout placeholders only by inspected placeholder idx and type.",
            "Render every result with the bundled internal OOXML renderer.",
        ],
        "supported_edits": {
            "set_text": "Replace text while retaining the first run style.",
            "replace_text": "Replace a substring while retaining surrounding style.",
            "replace_image": "Swap image bytes while preserving frame geometry.",
            "set_table_cell": "Update one native table cell.",
            "set_chart_series": "Update one native chart series.",
            "speaker_notes": "Set notes on the output slide.",
        },
        "supported_placeholder_fills": {
            "set_text": "Fill an inspected native text placeholder on a source slide layout.",
            "set_image": "Fill an inspected native PICTURE placeholder with a local raster image.",
        },
        "project_json_schema": TemplateDeckProject.model_json_schema(),
    }


def _write_inspect_artifacts(source: Path, work_dir: Path) -> TemplatePipelineResult:
    from pptx import Presentation

    presentation = Presentation(str(source))
    records = _inspect_records(presentation)
    slide_lineage = _presentation_lineage(presentation)
    previews = render_pptx_pages(source, work_dir / "previews")
    layout_paths: list[Path] = []
    for slide_number in range(1, len(presentation.slides) + 1):
        layout_path = (
            work_dir / "layouts" / f"source-slide-{slide_number:03d}.layout.json"
        )
        layout_path.parent.mkdir(parents=True, exist_ok=True)
        layout_path.write_text(
            json.dumps(
                {
                    "slide": slide_number,
                    "lineage": slide_lineage[slide_number - 1],
                    "records": [
                        record for record in records if record["slide"] == slide_number
                    ],
                },
                ensure_ascii=False,
                indent=2,
            )
            + "\n",
            encoding="utf-8",
        )
        layout_paths.append(layout_path)
    manifest_path = work_dir / "template-manifest.json"
    manifest = {
        "schemaVersion": 2,
        "engine": "evoflux-direct-openxml",
        "sourcePath": str(source.resolve()),
        "sourceSha256": file_sha256(source),
        "slideCount": len(presentation.slides),
        "records": records,
        "slideLineage": slide_lineage,
        "slideArtifacts": [
            {
                "slide": index,
                "previewPath": str(previews[index - 1]),
                "layoutPath": str(layout_paths[index - 1]),
            }
            for index in range(1, len(presentation.slides) + 1)
        ],
    }
    manifest_path.write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2) + "\n", encoding="utf-8"
    )
    return TemplatePipelineResult(
        action="inspect",
        source_pptx=source,
        work_dir=work_dir,
        manifest_path=manifest_path,
        previews=previews,
        layout_paths=layout_paths,
        slide_count=len(presentation.slides),
        metadata={"engine": "evoflux-direct-openxml"},
    )


def _rewrite_relationship_ids(element: Any, mapping: dict[str, str]) -> None:
    from pptx.oxml.ns import qn

    attributes = (qn("r:id"), qn("r:embed"), qn("r:link"))
    for node in element.iter():
        for attribute in attributes:
            old = node.get(attribute)
            if old in mapping:
                node.set(attribute, mapping[old])


def _clone_partname(source_part: Any, used_partnames: set[str]) -> Any:
    from pptx.opc.packuri import PackURI

    template = getattr(source_part.__class__, "partname_template", None)
    if not template:
        source_name = str(source_part.partname)
        if source_name.startswith("/ppt/media/"):
            template = f"/ppt/media/image%d.{source_part.partname.ext}"
        else:
            path = Path(source_name)
            stem = re.sub(r"\d+$", "", path.stem) or "part"
            template = f"{path.parent.as_posix()}/{stem}%d{path.suffix}"
            if not template.startswith("/"):
                template = f"/{template}"
    index = 1
    while template % index in used_partnames:
        index += 1
    value = template % index
    used_partnames.add(value)
    return PackURI(value)


def _clone_related_part(
    source_part: Any,
    *,
    package: Any,
    clones: dict[int, Any],
    used_partnames: set[str],
) -> Any:
    existing = clones.get(id(source_part))
    if existing is not None:
        return existing
    clone = source_part.__class__.load(
        _clone_partname(source_part, used_partnames),
        source_part.content_type,
        package,
        source_part.blob,
    )
    clones[id(source_part)] = clone
    relationship_mapping: dict[str, str] = {}
    for relation in source_part.rels.values():
        target = (
            relation.target_ref
            if relation.is_external
            else _clone_related_part(
                relation.target_part,
                package=package,
                clones=clones,
                used_partnames=used_partnames,
            )
        )
        relationship_mapping[relation.rId] = clone.rels._add_relationship(  # noqa: SLF001
            relation.reltype,
            target,
            relation.is_external,
        )
    if hasattr(clone, "_element"):
        _rewrite_relationship_ids(clone._element, relationship_mapping)  # noqa: SLF001
    return clone


def _notes_text(slide: Any) -> str | None:
    try:
        if not slide.has_notes_slide:
            return None
        return str(slide.notes_slide.notes_text_frame.text)
    except (AttributeError, TypeError, ValueError):
        return None


def _copy_notes_text(source_slide: Any, target_slide: Any) -> None:
    value = _notes_text(source_slide)
    if value is not None:
        target_slide.notes_slide.notes_text_frame.text = value


def _duplicate_slide(presentation: Any, source_slide: Any) -> Any:
    from pptx.shapes.shapetree import SlideShapes

    target = presentation.slides.add_slide(source_slide.slide_layout)
    duplicate = deepcopy(source_slide._element)  # noqa: SLF001 - OOXML-preserving clone
    mapping: dict[str, str] = {}
    cloned_parts: dict[int, Any] = {}
    used_partnames = {
        str(part.partname) for part in presentation.part.package.iter_parts()
    }
    target_layout_rel = next(
        rel for rel in target.part.rels.values() if rel.reltype.endswith("/slideLayout")
    )
    for relation in source_slide.part.rels.values():
        if relation.reltype.endswith("/slideLayout"):
            mapping[relation.rId] = target_layout_rel.rId
        elif relation.reltype.endswith("/notesSlide"):
            continue
        else:
            clone_relation_target = relation.reltype.endswith(("/chart", "/image"))
            relation_target = (
                relation.target_ref
                if relation.is_external
                else _clone_related_part(
                    relation.target_part,
                    package=presentation.part.package,
                    clones=cloned_parts,
                    used_partnames=used_partnames,
                )
                if clone_relation_target
                else relation.target_part
            )
            mapping[relation.rId] = target.part.rels._add_relationship(  # noqa: SLF001
                relation.reltype,
                relation_target,
                relation.is_external,  # noqa: SLF001
            )
    _rewrite_relationship_ids(duplicate, mapping)
    target.part._element = duplicate  # noqa: SLF001
    target._element = duplicate  # noqa: SLF001
    # ``Slide.shapes`` is cached when python-pptx creates the blank target.
    # Point that proxy at the cloned shape tree as well as replacing the XML.
    target.__dict__["shapes"] = SlideShapes(duplicate.cSld.spTree, target)
    _copy_notes_text(source_slide, target)
    return target


def _remove_original_slides(presentation: Any, count: int) -> None:
    slide_ids = presentation.slides._sldIdLst  # noqa: SLF001
    for slide_id in list(slide_ids)[:count]:
        presentation.part.drop_rel(slide_id.rId)
        slide_ids.remove(slide_id)


def _shape_by_id(slide: Any, shape_id: int) -> Any:
    for shape in slide.shapes:
        if shape.shape_id == shape_id:
            return shape
    raise ValueError(f"shape id {shape_id} is missing from cloned slide")


def _placeholder_by_identity(
    slide: Any, placeholder_idx: int, placeholder_type: str
) -> Any:
    indexed: list[Any] = []
    exact: list[Any] = []
    for shape in slide.placeholders:
        metadata = _placeholder_metadata(shape)
        if metadata is None or int(metadata["idx"]) != placeholder_idx:
            continue
        indexed.append(shape)
        if str(metadata["type"]) == placeholder_type:
            exact.append(shape)
    if len(exact) == 1:
        return exact[0]
    if len(exact) > 1:
        raise ValueError(
            f"slide contains duplicate placeholder idx/type "
            f"{placeholder_idx}/{placeholder_type}"
        )
    if indexed:
        actual_types = sorted(
            {
                str((_placeholder_metadata(shape) or {}).get("type") or "")
                for shape in indexed
            }
        )
        raise ValueError(
            f"slide placeholder idx {placeholder_idx} has type {actual_types}, "
            f"not {placeholder_type}"
        )
    raise ValueError(f"slide has no placeholder idx {placeholder_idx}")


def _set_text_preserving_style(shape: Any, value: str) -> None:
    runs = [run for paragraph in shape.text_frame.paragraphs for run in paragraph.runs]
    if not runs:
        shape.text = value
        return
    runs[0].text = value
    for run in runs[1:]:
        run.text = ""


def _replace_text_preserving_style(shape: Any, find: str, replace: str) -> None:
    """Replace text without collapsing the existing paragraph/run structure.

    Replacement characters inherit the formatting of the run containing the
    first matched character. Text outside each match stays assigned to its
    original run, including matches that span adjacent rich-text runs.
    """

    paragraphs = list(shape.text_frame.paragraphs)
    expected_matches = shape.text.count(find)
    paragraph_matches = sum(
        "".join(run.text for run in paragraph.runs).count(find)
        for paragraph in paragraphs
    )
    if expected_matches != paragraph_matches:
        raise ValueError(
            "replace_text cannot preserve formatting when the match crosses a "
            "paragraph boundary"
        )

    for paragraph in paragraphs:
        runs = list(paragraph.runs)
        if not runs:
            continue
        original = "".join(run.text for run in runs)
        matches: list[tuple[int, int]] = []
        cursor = 0
        while (start := original.find(find, cursor)) >= 0:
            matches.append((start, start + len(find)))
            cursor = start + len(find)
        if not matches:
            continue

        bounds: list[tuple[int, int]] = []
        offset = 0
        for run in runs:
            end = offset + len(run.text)
            bounds.append((offset, end))
            offset = end
        output = ["" for _ in runs]

        def append_original(start: int, end: int) -> None:
            for index, (run_start, run_end) in enumerate(bounds):
                overlap_start = max(start, run_start)
                overlap_end = min(end, run_end)
                if overlap_start < overlap_end:
                    output[index] += original[overlap_start:overlap_end]

        def run_index_at(position: int) -> int:
            for index, (run_start, run_end) in enumerate(bounds):
                if run_start <= position < run_end:
                    return index
            return max(0, len(runs) - 1)

        cursor = 0
        for start, end in matches:
            append_original(cursor, start)
            output[run_index_at(start)] += replace
            cursor = end
        append_original(cursor, len(original))
        for run, value in zip(runs, output, strict=True):
            run.text = value


def _set_picture_alt_text(shape: Any, value: str) -> None:
    try:
        properties = shape._element.xpath(
            "./*[local-name()='nvPicPr']/*[local-name()='cNvPr']"
        )
    except (AttributeError, TypeError, ValueError):
        properties = []
    if not properties:
        raise ValueError(f"image shape {shape.shape_id} has no native picture metadata")
    properties[0].set("descr", value)


def _replace_image(
    slide: Any, shape: Any, source: Path, *, alt_text: str | None = None
) -> None:
    old_relation_id = shape._element.blipFill.blip.rEmbed  # noqa: SLF001
    image_part, relation_id = slide.part.get_or_add_image_part(str(source))
    del image_part
    if old_relation_id != relation_id:
        # Drop the prior image relationship only when this shape is its sole
        # consumer; ``drop_rel`` keeps it when another shape still references it.
        slide.part.drop_rel(old_relation_id)
    shape._element.blipFill.blip.rEmbed = relation_id  # noqa: SLF001
    if alt_text is not None:
        _set_picture_alt_text(shape, alt_text)


def _replace_chart_series(shape: Any, series_index: int, values: list[float]) -> None:
    from pptx.chart.data import CategoryChartData

    chart = shape.chart
    series_values = [list(series.values) for series in chart.series]
    if series_index >= len(series_values):
        raise ValueError(f"chart series index {series_index} is out of range")
    if len(values) != len(series_values[series_index]):
        raise ValueError("replacement chart series must preserve category count")
    series_values[series_index] = values
    try:
        categories = [str(category) for category in chart.plots[0].categories]
    except (AttributeError, TypeError):
        categories = [str(index + 1) for index in range(len(values))]
    data = CategoryChartData()
    data.categories = categories
    for index, series in enumerate(chart.series):
        data.add_series(str(series.name or f"Series {index + 1}"), series_values[index])
    chart.replace_data(data)


def _project_asset(project_dir: Path, asset_path: str) -> Path:
    path = Path(asset_path)
    candidate = (path if path.is_absolute() else project_dir / path).resolve(
        strict=False
    )
    try:
        candidate.relative_to(project_dir.resolve())
    except ValueError as exc:
        raise ValueError(
            "replacement image must stay inside the project directory"
        ) from exc
    if not candidate.is_file():
        raise FileNotFoundError(f"replacement image does not exist: {candidate}")
    if candidate.suffix.lower() not in {".png", ".jpg", ".jpeg", ".gif", ".webp"}:
        raise ValueError("replacement image must be a supported raster image")
    return candidate


def _apply_placeholder_fill(
    slide: Any, fill: TemplatePlaceholderFill, project_dir: Path
) -> None:
    placeholder = _placeholder_by_identity(
        slide,
        fill.placeholder_idx,
        fill.placeholder_type,
    )
    if isinstance(fill, TemplateTextPlaceholderFill):
        _set_text_preserving_style(placeholder, fill.text)
        return
    path = _project_asset(project_dir, fill.asset_path)
    insert_picture = getattr(placeholder, "insert_picture", None)
    if not callable(insert_picture):
        raise ValueError(
            f"placeholder {fill.placeholder_idx}/{fill.placeholder_type} "
            "does not support native image insertion"
        )
    picture = insert_picture(str(path))
    if fill.alt is not None:
        _set_picture_alt_text(picture, fill.alt)


def _apply_edit(slide: Any, edit: TemplateObjectEdit, project_dir: Path) -> None:
    shape_id = int(edit.target_id.rsplit("/", 1)[1])
    shape = _shape_by_id(slide, shape_id)
    if edit.operation == "set_text":
        _set_text_preserving_style(shape, edit.text or "")
    elif edit.operation == "replace_text":
        current = shape.text
        if edit.find not in current:
            raise ValueError(f"target {edit.target_id} does not contain {edit.find!r}")
        _replace_text_preserving_style(shape, edit.find or "", edit.replace or "")
    elif edit.operation == "replace_image":
        path = _project_asset(project_dir, edit.asset_path or "")
        _replace_image(slide, shape, path, alt_text=edit.alt)
    elif edit.operation == "set_table_cell":
        assert edit.row is not None and edit.column is not None
        if edit.row >= len(shape.table.rows) or edit.column >= len(shape.table.columns):
            raise ValueError("table cell index is out of range")
        shape.table.cell(edit.row, edit.column).text = edit.text or ""
    else:
        assert edit.series_index is not None and edit.values is not None
        _replace_chart_series(shape, edit.series_index, edit.values)


def _placeholder_identities(
    values: list[dict[str, Any]], *, nested: bool
) -> list[dict[str, Any]]:
    identities: list[dict[str, Any]] = []
    for value in values:
        placeholder = value.get("placeholder") if nested else value
        if not isinstance(placeholder, dict):
            continue
        try:
            idx = int(placeholder.get("idx", -1))
        except (TypeError, ValueError):
            continue
        placeholder_type = str(placeholder.get("type") or "")
        if idx < 0 or not placeholder_type:
            continue
        identities.append({"idx": idx, "type": placeholder_type})
    return sorted(identities, key=lambda item: (item["idx"], item["type"]))


def _lineage_signature(
    value: dict[str, Any], *, reuse_mode: str, expected: bool
) -> dict[str, Any]:
    signature: dict[str, Any] = {
        "layout": {
            "partName": value["layout"]["partName"],
            "sha256": value["layout"]["sha256"],
            "placeholders": [
                record.get("placeholder") for record in value["layout"]["placeholders"]
            ],
        },
        "master": {
            "partName": value["master"]["partName"],
            "sha256": value["master"]["sha256"],
            "placeholders": [
                record.get("placeholder") for record in value["master"]["placeholders"]
            ],
        },
    }
    if reuse_mode == "duplicate-slide":
        signature["slidePlaceholders"] = value["slidePlaceholders"]
        return signature
    if expected:
        cloneable = [
            record
            for record in value["layout"]["placeholders"]
            if str(record.get("placeholder", {}).get("type") or "")
            not in {"DATE", "FOOTER", "SLIDE_NUMBER"}
        ]
        signature["slidePlaceholders"] = _placeholder_identities(cloneable, nested=True)
    else:
        signature["slidePlaceholders"] = _placeholder_identities(
            value["slidePlaceholders"], nested=False
        )
    return signature


def _verify_placeholder_fills(
    plan: TemplateSlidePlan, output_slide: Any, output_index: int
) -> list[dict[str, Any]]:
    issues: list[dict[str, Any]] = []
    for fill in plan.placeholder_fills:
        try:
            placeholder = _placeholder_by_identity(
                output_slide,
                fill.placeholder_idx,
                fill.placeholder_type,
            )
        except ValueError as exc:
            issues.append(
                {
                    "severity": "error",
                    "code": "template-placeholder-fill-drift",
                    "message": str(exc),
                    "slide": output_index,
                    "sourceSlide": plan.source_slide,
                    "placeholderIdx": fill.placeholder_idx,
                    "placeholderType": fill.placeholder_type,
                }
            )
            continue
        if isinstance(fill, TemplateTextPlaceholderFill) and (
            str(placeholder.text) != fill.text
        ):
            issues.append(
                {
                    "severity": "error",
                    "code": "template-placeholder-fill-drift",
                    "message": (
                        f"text placeholder {fill.placeholder_idx}/"
                        f"{fill.placeholder_type} changed after reopen"
                    ),
                    "slide": output_index,
                    "sourceSlide": plan.source_slide,
                    "placeholderIdx": fill.placeholder_idx,
                    "placeholderType": fill.placeholder_type,
                }
            )
        elif isinstance(fill, TemplateImagePlaceholderFill) and (
            not _is_picture_shape(placeholder)
            or (fill.alt is not None and _picture_alt_text(placeholder) != fill.alt)
        ):
            issues.append(
                {
                    "severity": "error",
                    "code": "template-placeholder-fill-drift",
                    "message": (
                        f"image placeholder {fill.placeholder_idx}/"
                        f"{fill.placeholder_type} changed after reopen"
                    ),
                    "slide": output_index,
                    "sourceSlide": plan.source_slide,
                    "placeholderIdx": fill.placeholder_idx,
                    "placeholderType": fill.placeholder_type,
                }
            )
    return issues


def _verify_template_lineage(
    source_lineage: list[dict[str, Any]],
    output_lineage: list[dict[str, Any]],
    project: TemplateDeckProject,
    output_presentation: Any,
    source_notes: list[str | None],
) -> list[dict[str, Any]]:
    issues: list[dict[str, Any]] = []
    if len(output_lineage) != len(project.output_slides):
        return [
            {
                "severity": "error",
                "code": "template-slide-count-drift",
                "message": "output slide count differs from the validated template map",
                "expected": len(project.output_slides),
                "actual": len(output_lineage),
            }
        ]
    for output_index, (plan, actual) in enumerate(
        zip(project.output_slides, output_lineage, strict=True), start=1
    ):
        expected = source_lineage[plan.source_slide - 1]
        expected_signature = _lineage_signature(
            expected, reuse_mode=plan.reuse_mode, expected=True
        )
        actual_signature = _lineage_signature(
            actual, reuse_mode=plan.reuse_mode, expected=False
        )
        if actual_signature != expected_signature:
            issues.append(
                {
                    "severity": "error",
                    "code": "template-lineage-drift",
                    "message": (
                        f"output slide {output_index} no longer matches source slide "
                        f"{plan.source_slide} master/layout/placeholder lineage"
                    ),
                    "slide": output_index,
                    "sourceSlide": plan.source_slide,
                    "expected": expected_signature,
                    "actual": actual_signature,
                }
            )
        expected_notes = (
            plan.speaker_notes
            if plan.speaker_notes is not None
            else source_notes[plan.source_slide - 1]
            if plan.reuse_mode == "duplicate-slide"
            else None
        )
        actual_notes = _notes_text(output_presentation.slides[output_index - 1])
        if expected_notes != actual_notes:
            issues.append(
                {
                    "severity": "error",
                    "code": "template-notes-drift",
                    "message": f"speaker notes changed on output slide {output_index}",
                    "slide": output_index,
                    "sourceSlide": plan.source_slide,
                }
            )
        if plan.reuse_mode == "use-layout":
            issues.extend(
                _verify_placeholder_fills(
                    plan,
                    output_presentation.slides[output_index - 1],
                    output_index,
                )
            )
    return issues


def _build_template(
    source: Path,
    project: TemplateDeckProject,
    project_dir: Path,
    output: Path,
    work_dir: Path,
    *,
    publish: bool,
) -> TemplatePipelineResult:
    from pptx import Presentation

    presentation = Presentation(str(source))
    originals = list(presentation.slides)
    source_lineage = _presentation_lineage(presentation)
    source_notes = [_notes_text(slide) for slide in originals]
    output_slides: list[Any] = []
    for plan in project.output_slides:
        source_slide = originals[plan.source_slide - 1]
        output_slides.append(
            _duplicate_slide(presentation, source_slide)
            if plan.reuse_mode == "duplicate-slide"
            else presentation.slides.add_slide(source_slide.slide_layout)
        )
    _remove_original_slides(presentation, len(originals))
    for plan, slide in zip(project.output_slides, output_slides, strict=True):
        if plan.reuse_mode == "duplicate-slide":
            for edit in plan.edits:
                _apply_edit(slide, edit, project_dir=project_dir)
        else:
            for fill in plan.placeholder_fills:
                _apply_placeholder_fill(slide, fill, project_dir=project_dir)
        if plan.speaker_notes is not None:
            slide.notes_slide.notes_text_frame.text = plan.speaker_notes
    candidate = output if publish else work_dir / "template-preview.pptx"
    candidate.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(candidate))
    previews = render_pptx_pages(candidate, work_dir / "previews")
    layout_paths: list[Path] = []
    output_presentation = Presentation(str(candidate))
    records = _inspect_records(output_presentation)
    output_lineage = _presentation_lineage(output_presentation)
    issues = _verify_template_lineage(
        source_lineage,
        output_lineage,
        project,
        output_presentation,
        source_notes,
    )
    lineage_verified = not any(
        issue.get("code") in {"template-slide-count-drift", "template-lineage-drift"}
        for issue in issues
    )
    notes_verified = not any(
        issue.get("code") == "template-notes-drift" for issue in issues
    )
    for slide_number in range(1, len(project.output_slides) + 1):
        path = work_dir / "layouts" / f"output-slide-{slide_number:03d}.layout.json"
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_text(
            json.dumps(
                {
                    "slide": slide_number,
                    "lineage": output_lineage[slide_number - 1],
                    "records": [
                        record for record in records if record["slide"] == slide_number
                    ],
                },
                ensure_ascii=False,
                indent=2,
            )
            + "\n",
            encoding="utf-8",
        )
        layout_paths.append(path)
    manifest_path = work_dir / "template-output-manifest.json"
    manifest_path.write_text(
        json.dumps(
            {
                "schemaVersion": 2,
                "engine": "evoflux-direct-openxml",
                "sourceSha256": project.source_sha256,
                "slideCount": len(project.output_slides),
                "records": records,
                "slideLineage": output_lineage,
                "slideSources": [
                    {
                        "outputSlide": plan.output_slide,
                        "sourceSlide": plan.source_slide,
                        "reuseMode": plan.reuse_mode,
                        "layoutPartName": output_lineage[index]["layout"]["partName"],
                    }
                    for index, plan in enumerate(project.output_slides)
                ],
                "lineageVerified": lineage_verified,
                "notesVerified": notes_verified,
            },
            ensure_ascii=False,
            indent=2,
        )
        + "\n",
        encoding="utf-8",
    )
    return TemplatePipelineResult(
        action="compose" if publish else "render",
        source_pptx=source,
        work_dir=work_dir,
        manifest_path=manifest_path,
        output=output if publish else None,
        previews=previews,
        layout_paths=layout_paths,
        slide_count=len(project.output_slides),
        issues=issues,
        metadata={
            "engine": "evoflux-direct-openxml",
            "preserved_master_layouts": lineage_verified,
            "lineage_verified": lineage_verified,
            "notes_verified": notes_verified,
            "use_layout_slide_count": sum(
                1 for plan in project.output_slides if plan.reuse_mode == "use-layout"
            ),
            "placeholder_fill_count": sum(
                len(plan.placeholder_fills) for plan in project.output_slides
            ),
        },
    )


async def inspect_pptx_template(
    source_pptx: Path, *, workspace_root: Path, work_dir: Path
) -> TemplatePipelineResult:
    del workspace_root
    return await asyncio.to_thread(_write_inspect_artifacts, source_pptx, work_dir)


async def render_pptx_template(
    source_pptx: Path,
    project_path: Path,
    manifest_path: Path,
    *,
    workspace_root: Path,
    work_dir: Path,
) -> TemplatePipelineResult:
    del workspace_root
    project = load_template_project(project_path)
    validate_template_project(
        project, load_template_manifest(manifest_path), source_pptx=source_pptx
    )
    return await asyncio.to_thread(
        _build_template,
        source_pptx,
        project,
        project_path.parent.resolve(),
        work_dir / "template-preview.pptx",
        work_dir,
        publish=False,
    )


async def compose_pptx_template(
    source_pptx: Path,
    project_path: Path,
    manifest_path: Path,
    output: Path,
    *,
    workspace_root: Path,
    work_dir: Path,
) -> TemplatePipelineResult:
    del workspace_root
    project = load_template_project(project_path)
    validate_template_project(
        project, load_template_manifest(manifest_path), source_pptx=source_pptx
    )
    try:
        return await asyncio.to_thread(
            _build_template,
            source_pptx,
            project,
            project_path.parent.resolve(),
            output,
            work_dir,
            publish=True,
        )
    except Exception:
        output.unlink(missing_ok=True)
        raise


__all__ = [
    "TemplateDeckProject",
    "TemplateImagePlaceholderFill",
    "TemplateObjectEdit",
    "TemplatePlaceholderFill",
    "TemplatePipelineResult",
    "TemplateSlidePlan",
    "TemplateTextPlaceholderFill",
    "compose_pptx_template",
    "inspect_pptx_template",
    "load_template_manifest",
    "load_template_project",
    "pptx_sha256",
    "render_pptx_template",
    "template_catalog",
    "validate_template_project",
]
