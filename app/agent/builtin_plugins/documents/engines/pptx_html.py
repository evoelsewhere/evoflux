"""Plugin-owned HTML/Tailwind PowerPoint authoring engine."""

from __future__ import annotations

import asyncio
import base64
import binascii
from dataclasses import dataclass, field
import json
from pathlib import Path
import re
from statistics import median
from typing import Any, Literal

from PIL import Image, ImageChops, ImageStat
from pydantic import BaseModel, ConfigDict, Field, field_validator, model_validator

from app.agent.builtin_plugins.documents.engines.html_slide_broker import (
    get_html_slide_render_broker,
)
from app.agent.builtin_plugins.documents.rendering.internal import render_pptx_pages

MAX_SLIDES = 80
MAX_ASSETS_PER_SLIDE = 80
MAX_HTML_BYTES = 2_000_000
MAX_CSS_BYTES = 2_000_000
MAX_ASSET_BYTES = 20_000_000
MAX_TOTAL_ASSET_BYTES = 60_000_000
PER_SLIDE_SIMILARITY_MIN = 0.90
DECK_MEDIAN_SIMILARITY_MIN = 0.95
_PLACEHOLDER = re.compile(r"\b(?:lorem ipsum|todo|tbd|click to add)\b", re.I)
_UNSAFE_HTML = re.compile(
    r"<\s*/?\s*(?:script|iframe|object|embed|video|audio|canvas|form|input|button|textarea|select|html|head|body|meta|link|base|style|template)\b",
    re.I,
)
_EVENT_HANDLER = re.compile(r"(?:\s|/)on[a-z0-9_-]+\s*=", re.I)
_UNSAFE_URL = re.compile(
    r"(?:javascript:|vbscript:|data:text/html|https?://|(?<!asset:)//)", re.I
)
_UNSAFE_CSS = re.compile(
    r"(?:</?style\b|@import|expression\s*\(|-moz-binding|behavior\s*:)", re.I
)
_ASSET_TOKEN = re.compile(r"asset://([A-Za-z][A-Za-z0-9_-]{0,63})")
_ASSET_SUFFIXES = {
    ".png",
    ".jpg",
    ".jpeg",
    ".gif",
    ".webp",
    ".svg",
    ".woff",
    ".woff2",
    ".ttf",
    ".otf",
}
_MIME_TYPES = {
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".gif": "image/gif",
    ".webp": "image/webp",
    ".svg": "image/svg+xml",
    ".woff": "font/woff",
    ".woff2": "font/woff2",
    ".ttf": "font/ttf",
    ".otf": "font/otf",
}


class HtmlSlidePlan(BaseModel):
    model_config = ConfigDict(extra="forbid")

    id: str = Field(pattern=r"^[A-Za-z][A-Za-z0-9_-]{0,79}$")
    html_path: str = Field(min_length=1, max_length=2000)
    style_paths: list[str] = Field(default_factory=list, max_length=8)
    assets: dict[str, str] = Field(default_factory=dict)
    speaker_notes: str | None = Field(default=None, max_length=40_000)

    @field_validator("assets")
    @classmethod
    def validate_assets(cls, value: dict[str, str]) -> dict[str, str]:
        if len(value) > MAX_ASSETS_PER_SLIDE:
            raise ValueError(
                f"a slide may reference at most {MAX_ASSETS_PER_SLIDE} assets"
            )
        for key, path in value.items():
            if not re.fullmatch(r"[A-Za-z][A-Za-z0-9_-]{0,63}", key):
                raise ValueError(f"invalid asset key: {key}")
            if not path or len(path) > 2000:
                raise ValueError(f"invalid asset path for {key}")
        return value


class HtmlPptxProject(BaseModel):
    model_config = ConfigDict(extra="forbid")

    schema_version: Literal[7] = 7
    title: str = Field(min_length=1, max_length=240)
    width: int = Field(default=1280, ge=640, le=3840)
    height: int = Field(default=720, ge=360, le=2160)
    slides: list[HtmlSlidePlan] = Field(min_length=1, max_length=MAX_SLIDES)

    @model_validator(mode="after")
    def unique_slides(self) -> HtmlPptxProject:
        ids = [slide.id for slide in self.slides]
        if len(ids) != len(set(ids)):
            raise ValueError("slide ids must be unique")
        return self


class NativeTextRun(BaseModel):
    model_config = ConfigDict(extra="forbid")

    text: str = Field(max_length=20_000)
    font_family: str | None = Field(default=None, max_length=500)
    font_size: float | None = Field(default=None, ge=1, le=512)
    bold: bool = False
    italic: bool = False
    underline: bool = False
    color: str | None = Field(default=None, max_length=64)
    letter_spacing: float = Field(default=0, ge=-32, le=128)


class NativeTextBullet(BaseModel):
    model_config = ConfigDict(extra="forbid")

    kind: Literal["bullet", "number"]
    marker: str | None = Field(default=None, max_length=8)
    level: int = Field(default=0, ge=0, le=8)
    start: int | None = Field(default=None, ge=1, le=32_767)


class NativeTextParagraph(BaseModel):
    model_config = ConfigDict(extra="forbid")

    runs: list[NativeTextRun] = Field(min_length=1, max_length=1_000)
    bullet: NativeTextBullet | None = None
    level: int = Field(default=0, ge=0, le=8)
    text_align: str | None = Field(default=None, max_length=32)


class NativeTextPadding(BaseModel):
    model_config = ConfigDict(extra="forbid")

    left: float = Field(default=0, ge=0, le=512)
    right: float = Field(default=0, ge=0, le=512)
    top: float = Field(default=0, ge=0, le=512)
    bottom: float = Field(default=0, ge=0, le=512)


class NativeTextElement(BaseModel):
    model_config = ConfigDict(extra="forbid")

    kind: Literal["text"]
    name: str = Field(min_length=1, max_length=160)
    role: str | None = Field(default=None, max_length=80)
    x: float
    y: float
    width: float = Field(gt=0)
    height: float = Field(gt=0)
    padding: NativeTextPadding = Field(default_factory=NativeTextPadding)
    text_align: str = "left"
    vertical_align: str = "top"
    line_height_ratio: float = Field(default=1.2, ge=0.5, le=5)
    rotation: float = Field(default=0, ge=-360, le=360)
    paragraphs: list[NativeTextParagraph] = Field(
        default_factory=list, max_length=1_000
    )
    # Version-6 compatibility. New renderers emit paragraphs/runs.
    text: str | None = Field(default=None, max_length=40_000)
    font_family: str | None = Field(default=None, max_length=500)
    font_size: float | None = Field(default=None, ge=1, le=512)
    bold: bool = False
    italic: bool = False
    underline: bool = False
    color: str | None = Field(default=None, max_length=64)

    @model_validator(mode="after")
    def has_text_content(self) -> NativeTextElement:
        if not self.paragraphs and self.text is None:
            raise ValueError("native text requires paragraphs or legacy text")
        return self


class NativeImageElement(BaseModel):
    model_config = ConfigDict(extra="forbid")

    kind: Literal["image"]
    name: str = Field(min_length=1, max_length=160)
    x: float
    y: float
    width: float = Field(gt=0)
    height: float = Field(gt=0)
    asset_id: str = Field(min_length=1, max_length=64)
    alt: str | None = Field(default=None, max_length=1_000)


class FlattenedTextRecord(BaseModel):
    model_config = ConfigDict(extra="forbid")

    name: str = Field(min_length=1, max_length=160)
    reason: str = Field(min_length=1, max_length=500)
    characters: int = Field(default=0, ge=0)


class TextCoverage(BaseModel):
    model_config = ConfigDict(extra="forbid")

    visible_blocks: int = Field(ge=0)
    visible_characters: int = Field(ge=0)
    native_blocks: int = Field(ge=0)
    native_characters: int = Field(ge=0)
    flattened: list[FlattenedTextRecord] = Field(default_factory=list)

    @model_validator(mode="after")
    def consistent_counts(self) -> TextCoverage:
        if self.native_blocks > self.visible_blocks:
            raise ValueError("native text blocks exceed visible text blocks")
        if self.native_characters > self.visible_characters:
            raise ValueError("native text characters exceed visible text characters")
        return self


class HtmlRenderResponse(BaseModel):
    model_config = ConfigDict(extra="forbid")

    preview_png_base64: str = Field(min_length=1)
    shell_png_base64: str = Field(min_length=1)
    editable_elements: list[NativeTextElement | NativeImageElement] = Field(
        default_factory=list, max_length=5_000
    )
    text_coverage: TextCoverage
    issues: list[dict[str, Any]] = Field(default_factory=list, max_length=5_000)


@dataclass
class HtmlPptxPipelineResult:
    action: str
    work_dir: Path
    output: Path | None = None
    previews: list[Path] = field(default_factory=list)
    layout_paths: list[Path] = field(default_factory=list)
    manifest_path: Path | None = None
    issues: list[dict[str, Any]] = field(default_factory=list)
    metadata: dict[str, Any] = field(default_factory=dict)

    @property
    def passed(self) -> bool:
        return not any(issue.get("severity") == "error" for issue in self.issues)


def load_html_pptx_project(path: Path) -> HtmlPptxProject:
    return HtmlPptxProject.model_validate_json(path.read_text(encoding="utf-8"))


def _project_file(project_dir: Path, value: str, *, label: str) -> Path:
    candidate = (project_dir / value).resolve(strict=False)
    try:
        candidate.relative_to(project_dir)
    except (ValueError, binascii.Error) as exc:
        raise ValueError(f"{label} must stay inside the project directory") from exc
    if not candidate.is_file():
        raise FileNotFoundError(f"{label} does not exist: {candidate}")
    return candidate


def _read_limited(path: Path, limit: int, *, label: str) -> str:
    if path.stat().st_size > limit:
        raise ValueError(f"{label} exceeds {limit} bytes")
    return path.read_text(encoding="utf-8")


def _validate_markup(html: str, css: str, *, slide_id: str) -> None:
    if _UNSAFE_HTML.search(html):
        raise ValueError(f"slide {slide_id} contains an unsafe HTML element")
    if _EVENT_HANDLER.search(html):
        raise ValueError(f"slide {slide_id} contains an inline event handler")
    if _UNSAFE_URL.search(html) or _UNSAFE_URL.search(css):
        raise ValueError(f"slide {slide_id} contains a network or executable URL")
    if _UNSAFE_CSS.search(css):
        raise ValueError(f"slide {slide_id} contains unsafe CSS")
    if "data-slide-root" not in html:
        raise ValueError(f"slide {slide_id} must contain one data-slide-root element")


def _slide_sources(
    slide: HtmlSlidePlan, project_dir: Path
) -> tuple[str, str, dict[str, Path]]:
    html_path = _project_file(
        project_dir, slide.html_path, label=f"slide {slide.id} html_path"
    )
    html = _read_limited(html_path, MAX_HTML_BYTES, label=f"slide {slide.id} HTML")
    css_parts: list[str] = []
    for value in slide.style_paths:
        path = _project_file(project_dir, value, label=f"slide {slide.id} style path")
        css_parts.append(
            _read_limited(path, MAX_CSS_BYTES, label=f"slide {slide.id} CSS")
        )
    css = "\n".join(css_parts)
    if len(css.encode("utf-8")) > MAX_CSS_BYTES:
        raise ValueError(f"slide {slide.id} CSS exceeds {MAX_CSS_BYTES} bytes")
    assets: dict[str, Path] = {}
    total_asset_bytes = 0
    for key, value in slide.assets.items():
        path = _project_file(project_dir, value, label=f"slide {slide.id} asset {key}")
        if path.suffix.lower() not in _ASSET_SUFFIXES:
            raise ValueError(f"unsupported slide asset type: {path.suffix}")
        if path.stat().st_size > MAX_ASSET_BYTES:
            raise ValueError(
                f"slide asset exceeds {MAX_ASSET_BYTES} bytes: {path.name}"
            )
        total_asset_bytes += path.stat().st_size
        if total_asset_bytes > MAX_TOTAL_ASSET_BYTES:
            raise ValueError(
                f"slide {slide.id} assets exceed {MAX_TOTAL_ASSET_BYTES} bytes"
            )
        assets[key] = path
    _validate_markup(html, css, slide_id=slide.id)
    referenced = set(_ASSET_TOKEN.findall(html + "\n" + css))
    missing = sorted(referenced - set(assets))
    if missing:
        raise ValueError(
            f"slide {slide.id} references undeclared assets: {', '.join(missing)}"
        )
    return html, css, assets


def validate_html_pptx_project(
    project: HtmlPptxProject, project_path: Path
) -> dict[str, Any]:
    project_dir = project_path.parent.resolve()
    placeholder_slides: list[str] = []
    for slide in project.slides:
        html, _css, _assets = _slide_sources(slide, project_dir)
        if _PLACEHOLDER.search(html):
            placeholder_slides.append(slide.id)
    if placeholder_slides:
        raise ValueError(
            "unresolved placeholder text in slides: " + ", ".join(placeholder_slides)
        )
    return {
        "valid": True,
        "schema_version": project.schema_version,
        "slide_count": len(project.slides),
        "canvas": {"width": project.width, "height": project.height, "unit": "px"},
        "rendering": "desktop-webview-html-css",
        "representation": "html-shell-editable-text",
        "editable_kinds": ["text", "simple-raster-image"],
        "quality_policy": {
            "evidence": "runtime-render-only",
            "per_slide_similarity_min": PER_SLIDE_SIMILARITY_MIN,
            "deck_median_similarity_min": DECK_MEDIAN_SIMILARITY_MIN,
        },
    }


def html_pptx_catalog() -> dict[str, Any]:
    return {
        "workflow": "html-shell-editable-text-pptx",
        "schema_version": 7,
        "canvas": {"width": 1280, "height": 720, "unit": "CSS px"},
        "project": {
            "required": ["schema_version", "title", "slides"],
            "slide_required": ["id", "html_path"],
            "slide_optional": ["style_paths", "assets", "speaker_notes"],
        },
        "html_contract": {
            "root": "exactly one element with data-slide-root",
            "asset_urls": "asset://<declared-key>",
            "text": {
                "default": "ordinary visible HTML text becomes editable PowerPoint text",
                "art_opt_out": 'data-pptx-text-mode="art"',
                "compatibility_hint": 'data-pptx-editable="text"',
            },
            "image": {
                "optional": (
                    'data-pptx-editable="image" data-pptx-asset="<declared-key>"'
                ),
                "constraint": (
                    "declared raster with fill sizing and no crop, mask, radius, "
                    "filter, transform, opacity, or blend effect"
                ),
            },
            "scripts": False,
            "network": False,
        },
        "runtime": {
            "renderer": "connected EvoFlux desktop WebView",
            "tailwind": "build-time curated utility runtime plus project-local CSS",
            "headless_fallback": False,
        },
        "editability": (
            "HTML/CSS is the visual source of truth. The WebView removes eligible "
            "glyphs from a high-resolution shell and returns native text manifests; "
            "images, charts, gradients, icons, and decoration remain in the shell."
        ),
        "quality_policy": {
            "source": "runtime render evidence, never an author-authored score",
            "per_slide_similarity_min": PER_SLIDE_SIMILARITY_MIN,
            "deck_median_similarity_min": DECK_MEDIAN_SIMILARITY_MIN,
            "requires_text_coverage": True,
        },
    }


def _inline_assets(
    html: str, css: str, assets: dict[str, Path]
) -> tuple[str, str, dict[str, dict[str, str]]]:
    metadata: dict[str, dict[str, str]] = {}
    for key, path in assets.items():
        mime = _MIME_TYPES[path.suffix.lower()]
        encoded = base64.b64encode(path.read_bytes()).decode("ascii")
        data_url = f"data:{mime};base64,{encoded}"
        token = f"asset://{key}"
        html = html.replace(token, data_url)
        css = css.replace(token, data_url)
        metadata[key] = {
            "mime_type": mime,
            "suffix": path.suffix.lower(),
        }
    return html, css, metadata


def _decode_png(value: str, destination: Path, *, label: str) -> Path:
    try:
        data = base64.b64decode(value, validate=True)
    except ValueError as exc:
        raise ValueError(f"{label} is not valid base64") from exc
    if len(data) > 40_000_000:
        raise ValueError(f"{label} exceeds 40 MB")
    destination.parent.mkdir(parents=True, exist_ok=True)
    destination.write_bytes(data)
    try:
        with Image.open(destination) as image:
            if image.format != "PNG":
                raise ValueError(f"{label} must be PNG")
            image.verify()
    except Exception:
        destination.unlink(missing_ok=True)
        raise
    return destination


def _preview_quality(path: Path, *, width: int, height: int) -> dict[str, Any]:
    with Image.open(path).convert("RGB") as image:
        actual_width, actual_height = image.size
        dimensions_ok = image.size == (width, height)
        grayscale = image.convert("L")
        entropy = float(grayscale.entropy())
        extrema = grayscale.getextrema()
        low, high = extrema
        if isinstance(low, tuple) or isinstance(high, tuple):
            raise ValueError("preview grayscale extrema are invalid")
        spread = int(high) - int(low)
        variance = float(ImageStat.Stat(grayscale).var[0])
    passed = dimensions_ok and entropy >= 0.35 and spread >= 8 and variance >= 4.0
    return {
        "passed": passed,
        "width": width,
        "height": height,
        "actual_width": actual_width,
        "actual_height": actual_height,
        "entropy": round(entropy, 4),
        "luminance_spread": spread,
        "luminance_variance": round(variance, 4),
    }


def _require_png_dimensions(path: Path, *, width: int, height: int, label: str) -> None:
    with Image.open(path) as image:
        if image.size != (width, height):
            raise ValueError(
                f"{label} must be {width}x{height}, got {image.width}x{image.height}"
            )


def _rgb(value: str, fallback: str = "#111827") -> Any:
    from pptx.dml.color import RGBColor

    candidate = value.strip().lstrip("#")
    if not re.fullmatch(r"[0-9A-Fa-f]{6}", candidate):
        candidate = fallback.lstrip("#")
    return RGBColor.from_string(candidate.upper())


def _set_paragraph_bullet(paragraph: Any, bullet: NativeTextBullet, level: int) -> None:
    from pptx.oxml.ns import qn
    from pptx.oxml.xmlchemy import OxmlElement

    properties = paragraph._p.get_or_add_pPr()  # noqa: SLF001
    properties.set("lvl", str(max(0, min(8, level))))
    properties.set("marL", str(int((18 + level * 18) * 12_700)))
    properties.set("indent", str(-int(10 * 12_700)))
    for child in list(properties):
        if child.tag in {qn("a:buNone"), qn("a:buChar"), qn("a:buAutoNum")}:
            properties.remove(child)
    if bullet.kind == "number":
        marker = OxmlElement("a:buAutoNum")
        marker.set("type", "arabicPeriod")
        if bullet.start is not None:
            marker.set("startAt", str(bullet.start))
    else:
        marker = OxmlElement("a:buChar")
        marker.set("char", bullet.marker or "•")
    properties.append(marker)


def _set_run_style(run: Any, value: NativeTextRun) -> None:
    from pptx.util import Pt

    family = (value.font_family or "Arial").split(",")[0].strip(" '\"")
    run.font.name = family or "Arial"
    run.font.size = Pt(float(value.font_size or 24) * 72 / 96)
    run.font.bold = value.bold
    run.font.italic = value.italic
    run.font.underline = value.underline
    run.font.color.rgb = _rgb(value.color or "#111827")
    spacing = round(value.letter_spacing * 0.75 * 100)
    properties = run._r.get_or_add_rPr()  # noqa: SLF001
    properties.set("spc", str(spacing))


def _add_styled_run(paragraph: Any, value: NativeTextRun) -> None:
    parts = value.text.split("\n")
    for index, part in enumerate(parts):
        if index:
            paragraph.add_line_break()
        if part or len(parts) == 1:
            run = paragraph.add_run()
            run.text = part
            _set_run_style(run, value)


def _element_paragraphs(element: NativeTextElement) -> list[NativeTextParagraph]:
    if element.paragraphs:
        return element.paragraphs
    return [
        NativeTextParagraph(
            runs=[
                NativeTextRun(
                    text=element.text or "",
                    font_family=element.font_family,
                    font_size=element.font_size,
                    bold=element.bold,
                    italic=element.italic,
                    underline=element.underline,
                    color=element.color,
                )
            ]
        )
    ]


def _add_text_overlay(slide: Any, element: NativeTextElement, *, scale: float) -> None:
    from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
    from pptx.util import Inches, Pt

    shape = slide.shapes.add_textbox(
        Inches(element.x / scale),
        Inches(element.y / scale),
        Inches(element.width / scale),
        Inches(element.height / scale),
    )
    frame = shape.text_frame
    frame.clear()
    frame.auto_size = MSO_AUTO_SIZE.NONE
    frame.word_wrap = True
    frame.margin_left = Inches(element.padding.left / scale)
    frame.margin_right = Inches(element.padding.right / scale)
    frame.margin_top = Inches(element.padding.top / scale)
    frame.margin_bottom = Inches(element.padding.bottom / scale)
    frame.vertical_anchor = {
        "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }.get(element.vertical_align, MSO_ANCHOR.TOP)
    alignments = {
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }
    for index, value in enumerate(_element_paragraphs(element)):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = alignments.get(
            value.text_align or element.text_align, PP_ALIGN.LEFT
        )
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)
        paragraph.line_spacing = element.line_height_ratio
        level = value.bullet.level if value.bullet else value.level
        paragraph.level = max(0, min(8, level))
        if value.bullet:
            _set_paragraph_bullet(paragraph, value.bullet, level)
        for run in value.runs:
            _add_styled_run(paragraph, run)
    shape.rotation = element.rotation
    shape.name = element.name[:160]


def _add_image_overlay(
    slide: Any,
    element: NativeImageElement,
    assets: dict[str, Path],
    *,
    scale: float,
) -> None:
    from pptx.util import Inches

    asset_id = element.asset_id
    source = assets.get(asset_id)
    if source is None or source.suffix.lower() not in {
        ".png",
        ".jpg",
        ".jpeg",
        ".gif",
        ".webp",
    }:
        raise ValueError(
            f"editable image asset is not a supported raster image: {asset_id}"
        )
    picture = slide.shapes.add_picture(
        str(source),
        Inches(element.x / scale),
        Inches(element.y / scale),
        Inches(element.width / scale),
        Inches(element.height / scale),
    )
    picture.name = element.name[:160]
    descr = element.alt or source.stem
    picture._element.nvPicPr.cNvPr.set("descr", descr)


def _compose_pptx(
    project: HtmlPptxProject,
    rendered: list[dict[str, Any]],
    slide_assets: list[dict[str, Path]],
    output: Path,
) -> dict[str, int]:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    presentation.slide_width = Inches(project.width / 96)
    presentation.slide_height = Inches(project.height / 96)
    for plan, item, assets in zip(project.slides, rendered, slide_assets, strict=True):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        slide.shapes.add_picture(
            str(item["shell_path"]),
            0,
            0,
            presentation.slide_width,
            presentation.slide_height,
        )
        for element in item["editable_elements"]:
            if isinstance(element, NativeTextElement):
                _add_text_overlay(slide, element, scale=96)
            elif isinstance(element, NativeImageElement):
                _add_image_overlay(slide, element, assets, scale=96)
        if plan.speaker_notes:
            slide.notes_slide.notes_text_frame.text = plan.speaker_notes
    output.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(output))
    # Structural check only; the caller separately reopens and rasterizes output.
    check = Presentation(str(output))
    if len(check.slides) != len(project.slides):
        raise ValueError("PPTX round-trip slide count changed")
    expected_text = sum(
        isinstance(element, NativeTextElement)
        for item in rendered
        for element in item["editable_elements"]
    )
    actual_text = sum(
        bool(getattr(shape, "has_text_frame", False))
        for slide in check.slides
        for shape in slide.shapes
    )
    if actual_text != expected_text:
        raise ValueError(
            "PPTX round-trip editable text count changed: "
            f"expected {expected_text}, got {actual_text}"
        )
    return {
        "editable_text_objects": actual_text,
        "editable_image_objects": sum(
            isinstance(element, NativeImageElement)
            for item in rendered
            for element in item["editable_elements"]
        ),
    }


def _raster_parity(source: Path, reopened: Path) -> dict[str, Any]:
    """Return deterministic pixel parity without an optional numeric dependency."""

    with Image.open(source).convert("RGB") as source_image:
        with Image.open(reopened).convert("RGB") as reopened_image:
            source_size = source_image.size
            reopened_size = reopened_image.size
            if source_size != reopened_size:
                return {
                    "metric": "normalized-rgb-rmse-similarity-v1",
                    "score": 0.0,
                    "source_size": list(source_size),
                    "reopened_size": list(reopened_size),
                    "dimensions_match": False,
                }
            difference = ImageChops.difference(source_image, reopened_image)
            rms = ImageStat.Stat(difference).rms
            normalized_rmse = (sum(value * value for value in rms) / 3) ** 0.5 / 255
    return {
        "metric": "normalized-rgb-rmse-similarity-v1",
        "score": round(max(0.0, 1.0 - normalized_rmse), 6),
        "source_size": list(source_size),
        "reopened_size": list(reopened_size),
        "dimensions_match": True,
    }


def _render_surface_ledger(
    *,
    previews: list[Path],
    shells: list[Path],
    reopened: list[Path],
    source_passed: bool,
    reopened_status: str,
) -> list[dict[str, Any]]:
    return [
        {
            "id": "source-preview",
            "status": "verified" if source_passed else "failed",
            "evidence": [str(path) for path in previews],
        },
        {
            "id": "html-visual-shell",
            "status": "verified" if len(shells) == len(previews) else "failed",
            "evidence": [str(path) for path in shells],
        },
        {
            "id": "reopened-plugin-preview",
            "status": reopened_status,
            "evidence": [str(path) for path in reopened],
        },
        {
            "id": "powerpoint-reference",
            "status": "unverified",
            "evidence": [],
        },
    ]


async def compose_html_pptx_project(
    project_path: Path,
    output: Path,
    *,
    workspace_root: Path,
    work_dir: Path,
    session_id: str | None,
) -> HtmlPptxPipelineResult:
    del workspace_root
    project = load_html_pptx_project(project_path)
    validate_html_pptx_project(project, project_path)
    project_dir = project_path.parent.resolve()
    work_dir.mkdir(parents=True, exist_ok=True)
    broker = get_html_slide_render_broker()
    issues: list[dict[str, Any]] = []
    rendered: list[dict[str, Any]] = []
    previews: list[Path] = []
    shells: list[Path] = []
    reopened_previews: list[Path] = []
    layouts: list[Path] = []
    all_assets: list[dict[str, Path]] = []
    quality: list[dict[str, Any]] = []
    text_coverage: list[dict[str, Any]] = []
    editable_text_count = 0
    editable_image_count = 0
    for slide_number, slide in enumerate(project.slides, start=1):
        html, css, assets = _slide_sources(slide, project_dir)
        inlined_html, inlined_css, asset_metadata = _inline_assets(html, css, assets)
        response = HtmlRenderResponse.model_validate(
            await broker.request(
                session_id,
                {
                    "slide_id": slide.id,
                    "width": project.width,
                    "height": project.height,
                    "html": inlined_html,
                    "css": inlined_css,
                    "assets": asset_metadata,
                },
            )
        )
        preview = _decode_png(
            response.preview_png_base64,
            work_dir / "previews" / f"slide-{slide_number:03d}.png",
            label=f"slide {slide_number} preview",
        )
        shell = _decode_png(
            response.shell_png_base64,
            work_dir / "shells" / f"slide-{slide_number:03d}.png",
            label=f"slide {slide_number} shell",
        )
        _require_png_dimensions(
            shell,
            width=project.width * 2,
            height=project.height * 2,
            label=f"slide {slide_number} shell",
        )
        metric = _preview_quality(preview, width=project.width, height=project.height)
        metric.update({"slide": slide_number, "slide_id": slide.id})
        quality.append(metric)
        if not metric["passed"]:
            issues.append(
                {
                    "severity": "error",
                    "code": "invalid-html-slide-preview",
                    "message": f"Slide {slide_number} preview is blank or malformed",
                    "slide": slide_number,
                    "quality": metric,
                }
            )
        for value in response.issues:
            if isinstance(value, dict):
                issues.append({**value, "slide": slide_number, "slide_id": slide.id})
        elements = response.editable_elements
        slide_text_count = sum(
            isinstance(element, NativeTextElement) for element in elements
        )
        editable_text_count += slide_text_count
        editable_image_count += sum(
            isinstance(element, NativeImageElement) for element in elements
        )
        coverage = response.text_coverage.model_dump()
        flattened_blocks = len(response.text_coverage.flattened)
        flattened_characters = sum(
            record.characters for record in response.text_coverage.flattened
        )
        accounted_blocks = response.text_coverage.native_blocks + flattened_blocks
        accounted_characters = (
            response.text_coverage.native_characters + flattened_characters
        )
        coverage.update(
            {
                "slide": slide_number,
                "slide_id": slide.id,
                "accounted_blocks": accounted_blocks,
                "accounted_characters": accounted_characters,
                "native_character_ratio": (
                    round(
                        response.text_coverage.native_characters
                        / response.text_coverage.visible_characters,
                        6,
                    )
                    if response.text_coverage.visible_characters
                    else 1.0
                ),
            }
        )
        text_coverage.append(coverage)
        if slide_text_count != response.text_coverage.native_blocks:
            issues.append(
                {
                    "severity": "error",
                    "code": "editable-text-manifest-mismatch",
                    "message": (
                        f"Slide {slide_number} declares "
                        f"{response.text_coverage.native_blocks} native text blocks "
                        f"but returned {slide_text_count}"
                    ),
                    "slide": slide_number,
                    "slide_id": slide.id,
                }
            )
        if (
            accounted_blocks != response.text_coverage.visible_blocks
            or accounted_characters != response.text_coverage.visible_characters
        ):
            issues.append(
                {
                    "severity": "error",
                    "code": "unaccounted-html-text",
                    "message": (
                        f"Slide {slide_number} did not account for every visible text "
                        "block and character as native or intentionally flattened"
                    ),
                    "slide": slide_number,
                    "slide_id": slide.id,
                    "coverage": coverage,
                }
            )
        layout_path = work_dir / "layouts" / f"slide-{slide_number:03d}.json"
        layout_path.parent.mkdir(parents=True, exist_ok=True)
        layout_path.write_text(
            json.dumps(
                {
                    "slide": slide_number,
                    "slide_id": slide.id,
                    "editable_elements": [
                        element.model_dump(mode="json") for element in elements
                    ],
                    "text_coverage": coverage,
                },
                ensure_ascii=False,
                indent=2,
            )
            + "\n",
            encoding="utf-8",
        )
        rendered.append(
            {
                "preview_path": preview,
                "shell_path": shell,
                "editable_elements": elements,
            }
        )
        previews.append(preview)
        shells.append(shell)
        layouts.append(layout_path)
        all_assets.append(assets)
    source_passed = all(item["passed"] for item in quality) and not any(
        issue.get("severity") == "error" for issue in issues
    )
    candidate: Path | None = output
    reopened_parity: list[dict[str, Any]] = []
    reopened_status = "not-run"
    structural_counts = {
        "editable_text_objects": 0,
        "editable_image_objects": 0,
    }
    if not any(issue.get("severity") == "error" for issue in issues):
        structural_counts = await asyncio.to_thread(
            _compose_pptx, project, rendered, all_assets, output
        )
        try:
            reopened_previews = await asyncio.to_thread(
                render_pptx_pages,
                output,
                work_dir / "reopened-previews",
                width=project.width,
            )
        except Exception as exc:
            issues.append(
                {
                    "severity": "error",
                    "code": "unrendered-reopened-slide",
                    "message": f"Generated PPTX could not be rendered after reopen: {exc}",
                }
            )
            reopened_status = "failed"
        if len(reopened_previews) != len(previews):
            issues.append(
                {
                    "severity": "error",
                    "code": "unrendered-reopened-slide",
                    "message": (
                        "Generated PPTX reopened with "
                        f"{len(reopened_previews)} rendered slides; expected {len(previews)}"
                    ),
                }
            )
            reopened_status = "failed"
        else:
            for slide_number, (source_preview, reopened_preview) in enumerate(
                zip(previews, reopened_previews, strict=True), start=1
            ):
                metric = _raster_parity(source_preview, reopened_preview)
                metric.update(
                    {
                        "slide": slide_number,
                        "slide_id": project.slides[slide_number - 1].id,
                        "threshold": PER_SLIDE_SIMILARITY_MIN,
                    }
                )
                reopened_parity.append(metric)
                if metric["score"] < PER_SLIDE_SIMILARITY_MIN:
                    issues.append(
                        {
                            "severity": "error",
                            "code": "render-parity-below-target",
                            "message": (
                                f"Slide {slide_number} reopened parity "
                                f"{metric['score']:.3f} is below "
                                f"{PER_SLIDE_SIMILARITY_MIN:.3f}"
                            ),
                            "slide": slide_number,
                            "slide_id": project.slides[slide_number - 1].id,
                            "parity": metric,
                        }
                    )
            deck_median = median(item["score"] for item in reopened_parity)
            if deck_median < DECK_MEDIAN_SIMILARITY_MIN:
                issues.append(
                    {
                        "severity": "error",
                        "code": "render-parity-below-target",
                        "message": (
                            f"Deck reopened median parity {deck_median:.3f} is below "
                            f"{DECK_MEDIAN_SIMILARITY_MIN:.3f}"
                        ),
                        "median": round(deck_median, 6),
                        "threshold": DECK_MEDIAN_SIMILARITY_MIN,
                    }
                )
            reopened_status = (
                "failed"
                if any(
                    issue.get("code")
                    in {
                        "render-parity-below-target",
                        "unrendered-reopened-slide",
                    }
                    for issue in issues
                )
                else "verified"
            )
    else:
        output.unlink(missing_ok=True)
        candidate = None
    surfaces = _render_surface_ledger(
        previews=previews,
        shells=shells,
        reopened=reopened_previews,
        source_passed=source_passed,
        reopened_status=reopened_status,
    )
    median_parity = (
        round(median(item["score"] for item in reopened_parity), 6)
        if reopened_parity
        else None
    )
    accepted = not any(issue.get("severity") == "error" for issue in issues)
    if not accepted:
        output.unlink(missing_ok=True)
        candidate = None
    manifest_path = work_dir / "html-pptx-manifest.json"
    manifest = {
        "schemaVersion": 7,
        "engine": "evoflux-html-shell-editable-text",
        "slideCount": len(project.slides),
        "editableTextObjectCount": editable_text_count,
        "editableImageObjectCount": editable_image_count,
        "previewQuality": quality,
        "visualSource": "desktop-webview",
        "roundTrip": "structural-openxml",
        "qualityPolicy": {
            "evidence": "runtime-render-only",
            "perSlideSimilarityMin": PER_SLIDE_SIMILARITY_MIN,
            "deckMedianSimilarityMin": DECK_MEDIAN_SIMILARITY_MIN,
        },
        "textConversion": text_coverage,
        "structuralCounts": structural_counts,
        "renderSurfaces": surfaces,
        "reopenedParity": {
            "metric": "normalized-rgb-rmse-similarity-v1",
            "perSlide": reopened_parity,
            "median": median_parity,
        },
        "accepted": accepted,
    }
    manifest_path.write_text(json.dumps(manifest, indent=2) + "\n", encoding="utf-8")
    return HtmlPptxPipelineResult(
        action="compose",
        work_dir=work_dir,
        output=candidate,
        previews=reopened_previews or previews,
        layout_paths=layouts,
        manifest_path=manifest_path,
        issues=issues,
        metadata={
            "engine": "evoflux-html-shell-editable-text",
            "slide_count": len(project.slides),
            "editable_object_count": editable_text_count + editable_image_count,
            "editable_text_object_count": editable_text_count,
            "semantic_editable_object_count": editable_text_count,
            "text_conversion": text_coverage,
            "preview_quality": quality,
            "visual_verification": (
                "reopened-plugin-preview"
                if reopened_status == "verified"
                else "failed-reopened-plugin-preview"
            ),
            "round_trip_verification": "structural-openxml",
            "quality_policy": {
                "evidence": "runtime-render-only",
                "per_slide_similarity_min": PER_SLIDE_SIMILARITY_MIN,
                "deck_median_similarity_min": DECK_MEDIAN_SIMILARITY_MIN,
            },
            "structural_counts": structural_counts,
            "accepted": accepted,
            "render_surfaces": surfaces,
            "reopened_parity": {
                "metric": "normalized-rgb-rmse-similarity-v1",
                "per_slide": reopened_parity,
                "median": median_parity,
            },
        },
    )


__all__ = [
    "HtmlPptxPipelineResult",
    "HtmlPptxProject",
    "HtmlSlidePlan",
    "compose_html_pptx_project",
    "html_pptx_catalog",
    "load_html_pptx_project",
    "validate_html_pptx_project",
]
