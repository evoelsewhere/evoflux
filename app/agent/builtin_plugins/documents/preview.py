"""Plugin-owned cached HTML preview engine for the shared document reader.

The renderer is intentionally self-contained: DOCX, XLSX, and PPTX are parsed
with the Python OpenXML libraries already bundled in the EvoFlux sidecar.  It
does not execute document macros, launch an office application, or depend on an
external renderer binary.
"""

from __future__ import annotations

import base64
import hashlib
import html
import math
import os
import shutil
import tempfile
import threading
from pathlib import Path
from typing import Any

from loguru import logger

from app.agent.builtin_plugins.documents.rendering.xlsx_formula import (
    FormulaEvaluation,
    evaluate_workbook_formulas,
    format_computed_value,
)
from app.agent.builtin_plugins.documents.preview_security import (
    cached_preview_is_valid,
    maintain_preview_cache,
    mark_cached_preview_used,
    preflight_ooxml_package,
    prepare_preview_cache_directory,
)
from app.core.config import settings
from app.plugin_platform.previews import (
    DOCUMENT_PREVIEW_CSP,
    DocumentPreviewError,
    DocumentPreviewProvider,
    DocumentPreviewUnsupportedError,
)

SUPPORTED_DOCUMENT_PREVIEW_EXTENSIONS = frozenset({".docx", ".xlsx", ".pptx", ".pdf"})
MAX_DOCUMENT_PREVIEW_BYTES = 100 * 1024 * 1024
MAX_DOCUMENT_PREVIEW_HTML_BYTES = 36 * 1024 * 1024
MAX_PDF_PREVIEW_PAGES = 80
MAX_PDF_PREVIEW_RASTER_BYTES = 24 * 1024 * 1024
MAX_PDF_PREVIEW_PIXELS_PER_PAGE = 10_000_000
_CACHE_SCHEMA_VERSION = "evoflux-document-html-v14"
_render_locks = tuple(threading.Lock() for _ in range(32))


def _render_lock_for(output: Path) -> threading.Lock:
    """Serialize one cache key without blocking unrelated document previews."""

    return _render_locks[int(output.stem[:8], 16) % len(_render_locks)]


def _cache_path(source: Path) -> Path:
    fingerprint = hashlib.sha256()
    fingerprint.update(_CACHE_SCHEMA_VERSION.encode())
    fingerprint.update(b"\0")
    fingerprint.update(source.suffix.casefold().encode())
    fingerprint.update(b"\0")
    fingerprint.update(os.fsencode(source.name))
    fingerprint.update(b"\0")
    with source.open("rb") as handle:
        while chunk := handle.read(1024 * 1024):
            fingerprint.update(chunk)
    digest = fingerprint.hexdigest()
    return Path(settings.EVOFLUX_CACHE_DIR) / "document-previews" / f"{digest}.html"


def _page(*, title: str, body: str, css: str) -> str:
    policy = html.escape(DOCUMENT_PREVIEW_CSP, quote=True)
    safe_title = html.escape(title)
    return (
        '<!doctype html><html><head><meta charset="utf-8">'
        f'<meta http-equiv="Content-Security-Policy" content="{policy}">'
        '<meta name="referrer" content="no-referrer">'
        f"<title>{safe_title}</title><style>{css}</style></head>"
        f"<body>{body}</body></html>"
    )


def _css_color(value: Any, default: str = "transparent") -> str:
    """Convert an OpenXML color object or raw ARGB string to CSS."""
    if value is None:
        return default
    try:
        raw = getattr(value, "rgb", value)
    except (AttributeError, TypeError, ValueError):
        return default
    if raw is None:
        return default
    text = str(raw)
    if len(text) == 8:
        text = text[2:]
    if len(text) == 6 and all(ch in "0123456789abcdefABCDEF" for ch in text):
        return f"#{text}"
    return default


def _css_font_family(name: Any) -> str:
    family = str(name).replace("\\", "\\\\").replace('"', '\\"')
    fallback = (
        '"Helvetica Neue",Arial,sans-serif'
        if family.casefold() == "aptos display"
        else "Arial,sans-serif"
    )
    return html.escape(f'"{family}",{fallback}', quote=True)


def _render_docx(source: Path) -> str:
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    document = Document(str(source))
    section = document.sections[0]

    def pixels(length: Any, default_inches: float) -> float:
        return float(length.inches) * 96 if length is not None else default_inches * 96

    page_width = pixels(section.page_width, 8.5)
    page_height = pixels(section.page_height, 11)
    margin_top = pixels(section.top_margin, 1)
    margin_right = pixels(section.right_margin, 1)
    margin_bottom = pixels(section.bottom_margin, 1)
    margin_left = pixels(section.left_margin, 1)
    alignment = {
        WD_ALIGN_PARAGRAPH.CENTER: "center",
        WD_ALIGN_PARAGRAPH.RIGHT: "right",
        WD_ALIGN_PARAGRAPH.JUSTIFY: "justify",
        WD_ALIGN_PARAGRAPH.DISTRIBUTE: "justify",
    }
    list_number = 0

    def render_paragraph(paragraph: Paragraph) -> str:
        nonlocal list_number
        style_name = (paragraph.style.name if paragraph.style else "").lower()
        tag = "p"
        if style_name == "title":
            tag = "h1"
        for level in range(1, 4):
            if style_name.startswith(f"heading {level}"):
                tag = f"h{level + 1}"
                break
        is_numbered = style_name.startswith("list number")
        is_bullet = style_name.startswith("list bullet")
        if is_numbered:
            list_number += 1
        elif not is_bullet:
            list_number = 0

        runs: list[str] = []
        for run in paragraph.runs:
            styles: list[str] = []
            if run.bold:
                styles.append("font-weight:700")
            if run.italic:
                styles.append("font-style:italic")
            if run.underline:
                styles.append("text-decoration:underline")
            if run.font.size:
                styles.append(f"font-size:{run.font.size.pt:.2f}pt")
            if run.font.name:
                styles.append(f"font-family:{_css_font_family(run.font.name)}")
            if run.font.color and run.font.color.rgb:
                styles.append(f"color:#{run.font.color.rgb}")
            text = html.escape(run.text).replace("\n", "<br>")
            style_attr = f' style="{";".join(styles)}"' if styles else ""
            runs.append(f"<span{style_attr}>{text}</span>")
            for blip in run._element.xpath(".//a:blip"):
                relation_id = blip.get(qn("r:embed"))
                related = paragraph.part.related_parts.get(relation_id)
                if related is None or not hasattr(related, "blob"):
                    continue
                content_type = getattr(related, "content_type", "image/png")
                encoded = base64.b64encode(related.blob).decode("ascii")
                extent = run._element.xpath(".//wp:extent")
                image_styles = ["max-width:100%", "height:auto"]
                if extent:
                    width_px = int(extent[0].get("cx", "0")) / 914400 * 96
                    if width_px > 0:
                        image_styles.append(f"width:{width_px:.2f}px")
                runs.append(
                    f'<img class="inline-image" style="{";".join(image_styles)}" '
                    f'src="data:{content_type};base64,{encoded}" alt="">'
                )
        content = "".join(runs) or html.escape(paragraph.text)
        if not content:
            content = "&nbsp;"
        paragraph_styles: list[str] = []
        style = paragraph.style
        if style is not None:
            if style.font.name:
                paragraph_styles.append(
                    f"font-family:{_css_font_family(style.font.name)}"
                )
            if style.font.size:
                paragraph_styles.append(f"font-size:{style.font.size.pt:.2f}pt")
            if style.font.bold:
                paragraph_styles.append("font-weight:700")
            if style.font.color and style.font.color.rgb:
                paragraph_styles.append(f"color:#{style.font.color.rgb}")
        if paragraph.alignment in alignment:
            paragraph_styles.append(f"text-align:{alignment[paragraph.alignment]}")
        formatting = paragraph.paragraph_format
        style_formatting = style.paragraph_format if style is not None else None
        space_before = formatting.space_before or (
            style_formatting.space_before if style_formatting is not None else None
        )
        space_after = formatting.space_after or (
            style_formatting.space_after if style_formatting is not None else None
        )
        line_spacing = formatting.line_spacing or (
            style_formatting.line_spacing if style_formatting is not None else None
        )
        if space_before:
            paragraph_styles.append(f"margin-top:{space_before.pt:.2f}pt")
        if space_after:
            paragraph_styles.append(f"margin-bottom:{space_after.pt:.2f}pt")
        if isinstance(line_spacing, (int, float)):
            paragraph_styles.append(f"line-height:{float(line_spacing):.3f}")
        elif line_spacing is not None:
            paragraph_styles.append(f"line-height:{line_spacing.pt:.2f}pt")
        if formatting.left_indent:
            paragraph_styles.append(f"margin-left:{formatting.left_indent.pt:.2f}pt")
        if formatting.right_indent:
            paragraph_styles.append(f"margin-right:{formatting.right_indent.pt:.2f}pt")
        if formatting.first_line_indent:
            paragraph_styles.append(
                f"text-indent:{formatting.first_line_indent.pt:.2f}pt"
            )
        style_attr = (
            f' style="{";".join(paragraph_styles)}"' if paragraph_styles else ""
        )
        page_break = (
            "true"
            if formatting.page_break_before
            or bool(paragraph._element.xpath(".//w:br[@w:type='page']"))
            else "false"
        )
        label = html.escape(paragraph.text[:80], quote=True)
        if is_numbered or is_bullet:
            marker = f"{list_number}." if is_numbered else "•"
            return (
                f'<div class="list-item"{style_attr} '
                f'data-page-break-before="{page_break}" data-qa-label="{label}">'
                f'<span class="list-marker">{marker}</span>'
                f'<span class="list-content">{content}</span></div>'
            )
        return (
            f'<{tag}{style_attr} data-page-break-before="{page_break}" '
            f'data-qa-label="{label}">{content}</{tag}>'
        )

    def render_table(table: Table) -> str:
        rows: list[str] = []
        for row_index, row in enumerate(table.rows):
            cells: list[str] = []
            cell_tag = "th" if row_index == 0 else "td"
            for cell in row.cells:
                cell_body = "".join(render_paragraph(p) for p in cell.paragraphs)
                cells.append(f"<{cell_tag}>{cell_body}</{cell_tag}>")
            rows.append(f"<tr>{''.join(cells)}</tr>")
        return f"<table>{''.join(rows)}</table>"

    blocks: list[str] = []
    for child in document.element.body.iterchildren():
        if isinstance(child, CT_P):
            blocks.append(render_paragraph(Paragraph(child, document)))
        elif isinstance(child, CT_Tbl):
            blocks.append(render_table(Table(child, document)))

    header = "".join(render_paragraph(p) for p in section.header.paragraphs)
    footer = "".join(render_paragraph(p) for p in section.footer.paragraphs)
    geometry = (
        f"--page-width:{page_width:.2f}px;--page-height:{page_height:.2f}px;"
        f"--margin-top:{margin_top:.2f}px;--margin-right:{margin_right:.2f}px;"
        f"--margin-bottom:{margin_bottom:.2f}px;--margin-left:{margin_left:.2f}px"
    )
    css = """
    *{box-sizing:border-box}body{margin:0;padding:32px;background:#e8eaed;
    color:#202124;font-family:Arial,sans-serif}.document{margin:0 auto}
    .document-flow,.document-page{width:var(--page-width);background:#fff;
    box-shadow:0 2px 12px #0002}.document-flow{min-height:var(--page-height);
    padding:var(--margin-top) var(--margin-right) var(--margin-bottom)
    var(--margin-left)}.document-pages{display:flex;flex-direction:column;gap:24px;
    align-items:center}.document-page{height:var(--page-height);position:relative;
    padding:var(--margin-top) var(--margin-right) var(--margin-bottom)
    var(--margin-left);overflow:hidden}.document-page-body{height:100%;overflow:hidden}
    .document-page-header,.document-page-footer{position:absolute;left:var(--margin-left);
    right:var(--margin-right);font-size:9pt;color:#5f6368;overflow:hidden}
    .document-page-header{top:16px;height:calc(var(--margin-top) - 20px)}
    .document-page-footer{bottom:12px;height:calc(var(--margin-bottom) - 16px)}
    .document-header-template,.document-footer-template{display:none}
    h1{font-size:26pt;line-height:1.15;margin:0 0
    18pt}h2{font-size:18pt;margin:20pt 0 8pt}h3{font-size:14pt;margin:16pt 0
    6pt}h4{font-size:12pt;margin:12pt 0 5pt}p{font-size:11pt;line-height:1.5;
    margin:0 0 8pt;white-space:pre-wrap}.list-item{display:grid;
    grid-template-columns:1.4em 1fr;column-gap:.25em;align-items:start;
    margin:0 0 6pt;white-space:pre-wrap}.list-marker{text-align:right;
    font-weight:600}.list-content{min-width:0}
    table{width:100%;border-collapse:collapse;margin:12pt 0 18pt;table-layout:fixed}
    th,td{border:1px solid #c7c7c7;padding:7px 9px;vertical-align:top;overflow-wrap:anywhere}
    th{background:#f1f3f4;text-align:left}td p,th p{margin:0;font-size:10pt}
    .inline-image{display:inline-block;vertical-align:middle}
    """
    return _page(
        title=source.name,
        body=(
            f'<main class="document" style="{geometry}">'
            f'<div class="document-header-template">{header}</div>'
            f'<div class="document-footer-template">{footer}</div>'
            '<div class="document-flow" data-preview-item '
            'data-preview-label="Document">'
            f"{''.join(blocks)}</div></main>"
        ),
        css=css,
    )


def _worksheet_bounds(sheet: Any) -> tuple[int, int]:
    """Return bounded preview dimensions for a worksheet."""
    max_row = min(max(sheet.max_row, 1), 500)
    max_column = min(max(sheet.max_column, 1), 100)
    return max_row, max_column


def _display_cell(
    cell: Any,
    *,
    evaluation: FormulaEvaluation | None = None,
    sheet_name: str | None = None,
) -> str:
    value = cell.value
    if value is None:
        return ""
    if cell.data_type == "f":
        value = (
            evaluation.display_value(sheet_name, cell.coordinate, "#N/A")
            if evaluation is not None and sheet_name is not None
            else "#N/A"
        )
    return format_computed_value(value, str(getattr(cell, "number_format", "") or ""))


def _border_css(side: Any, name: str) -> str | None:
    if not side or not side.style:
        return None
    color = _css_color(side.color, "#80868b")
    width = "2px" if side.style in {"medium", "thick", "double"} else "1px"
    line = "dashed" if "dash" in side.style else "solid"
    return f"border-{name}:{width} {line} {color}"


def _xlsx_series_values(
    workbook: Any,
    series: Any,
    evaluation: FormulaEvaluation | None = None,
) -> list[float]:
    """Resolve a chart series from its cached values or workbook cell range."""
    from openpyxl.utils.cell import range_boundaries

    data_source = getattr(series, "val", None) or getattr(series, "yVal", None)
    if data_source is None:
        return []
    literal = getattr(data_source, "numLit", None)
    if literal is not None:
        values = []
        for point in getattr(literal, "pt", ()):
            try:
                values.append(float(point.v))
            except (TypeError, ValueError):
                values.append(0.0)
        return values

    reference = getattr(data_source, "numRef", None)
    formula = str(getattr(reference, "f", "") or "")
    if "!" not in formula:
        return []
    sheet_name, cell_range = formula.rsplit("!", 1)
    sheet_name = sheet_name.strip("'").replace("''", "'")
    if sheet_name not in workbook.sheetnames:
        return []
    min_column, min_row, max_column, max_row = range_boundaries(cell_range)
    values: list[float] = []
    sheet = workbook[sheet_name]
    for row in sheet.iter_rows(
        min_row=min_row,
        max_row=max_row,
        min_col=min_column,
        max_col=max_column,
    ):
        for cell in row:
            try:
                value = (
                    evaluation.display_value(sheet_name, cell.coordinate, 0)
                    if evaluation is not None
                    else cell.value
                )
                values.append(float(value or 0))
            except (TypeError, ValueError):
                values.append(0.0)
    return values


def _xlsx_category_values(
    workbook: Any,
    series: Any,
    evaluation: FormulaEvaluation | None = None,
) -> list[str]:
    """Resolve chart category labels from a worksheet reference."""
    from openpyxl.utils.cell import range_boundaries

    data_source = getattr(series, "cat", None) or getattr(series, "xVal", None)
    if data_source is None:
        return []
    reference = getattr(data_source, "strRef", None) or getattr(
        data_source, "numRef", None
    )
    formula = str(getattr(reference, "f", "") or "")
    if "!" not in formula:
        return []
    sheet_name, cell_range = formula.rsplit("!", 1)
    sheet_name = sheet_name.strip("'").replace("''", "'")
    if sheet_name not in workbook.sheetnames:
        return []
    min_column, min_row, max_column, max_row = range_boundaries(cell_range)
    sheet = workbook[sheet_name]
    values: list[str] = []
    for row in sheet.iter_rows(
        min_row=min_row,
        max_row=max_row,
        min_col=min_column,
        max_col=max_column,
    ):
        for cell in row:
            value = (
                evaluation.display_value(sheet_name, cell.coordinate, "")
                if evaluation is not None
                else cell.value
            )
            values.append(str(value or ""))
    return values


def _xlsx_chart_svg(
    workbook: Any,
    chart: Any,
    evaluation: FormulaEvaluation | None = None,
) -> str:
    """Render common workbook charts as a deterministic SVG preview."""
    palette = ("#2563eb", "#0f766e", "#f59e0b", "#dc2626", "#7c3aed")
    series_values = [
        values
        for series in chart.ser
        if (values := _xlsx_series_values(workbook, series, evaluation))
    ]
    if not series_values:
        return '<div class="chart-empty">Chart data unavailable</div>'
    maximum = max(
        (abs(value) for values in series_values for value in values),
        default=1,
    )
    maximum = maximum or 1
    count = max(len(values) for values in series_values)
    chart_type = type(chart).__name__.upper()
    categories = (
        _xlsx_category_values(workbook, chart.ser[0], evaluation) if chart.ser else []
    )

    def category_labels() -> str:
        if not categories:
            return ""
        return "".join(
            f'<text x="{52 + index * 520 / max(len(categories) - 1, 1):.2f}" '
            'y="278" text-anchor="middle" font-size="13" fill="#6b7280">'
            f"{html.escape(label)}</text>"
            for index, label in enumerate(categories)
        )

    if "PIE" in chart_type or "DOUGHNUT" in chart_type:
        values = [abs(value) for value in series_values[0]]
        total = sum(values) or 1
        start_angle = -math.pi / 2
        slices: list[str] = []
        for index, value in enumerate(values):
            end_angle = start_angle + 2 * math.pi * value / total
            start_x = 300 + 108 * math.cos(start_angle)
            start_y = 150 + 108 * math.sin(start_angle)
            end_x = 300 + 108 * math.cos(end_angle)
            end_y = 150 + 108 * math.sin(end_angle)
            large_arc = 1 if end_angle - start_angle > math.pi else 0
            slices.append(
                f'<path d="M300 150 L{start_x:.2f} {start_y:.2f} '
                f'A108 108 0 {large_arc} 1 {end_x:.2f} {end_y:.2f} Z" '
                f'fill="{palette[index % len(palette)]}"/>'
            )
            start_angle = end_angle
        hole = (
            '<circle cx="300" cy="150" r="65" fill="white"/>'
            if "DOUGHNUT" in chart_type
            else ""
        )
        return (
            '<svg class="workbook-chart-svg" viewBox="0 0 600 300">'
            f"{''.join(slices)}{hole}</svg>"
        )

    if "AREA" in chart_type or "LINE" in chart_type:
        paths: list[str] = []
        for series_index, values in enumerate(series_values):
            points = [
                (
                    52 + value_index * 520 / max(len(values) - 1, 1),
                    252 - (value / maximum) * 205,
                )
                for value_index, value in enumerate(values)
            ]
            point_text = " ".join(f"{x:.2f},{y:.2f}" for x, y in points)
            color = palette[series_index % len(palette)]
            area = ""
            if "AREA" in chart_type and points:
                polygon = [
                    f"{points[0][0]:.2f},252",
                    *[f"{x:.2f},{y:.2f}" for x, y in points],
                    f"{points[-1][0]:.2f},252",
                ]
                area = (
                    f'<polygon class="workbook-chart-area" '
                    f'points="{" ".join(polygon)}" fill="{color}" '
                    'fill-opacity=".24"/>'
                )
            markers = "".join(
                f'<circle cx="{x:.2f}" cy="{y:.2f}" r="4" fill="{color}"/>'
                for x, y in points
            )
            paths.append(
                area + f'<polyline points="{point_text}" fill="none" '
                f'stroke="{color}" stroke-width="4"/>{markers}'
            )
        return (
            '<svg class="workbook-chart-svg" viewBox="0 0 600 300" '
            'preserveAspectRatio="none">'
            '<line x1="45" y1="252" x2="575" y2="252" stroke="#9ca3af"/>'
            f"{''.join(paths)}{category_labels()}</svg>"
        )

    group_width = 520 / max(count, 1)
    bar_width = max(group_width / max(len(series_values), 1) * 0.72, 2)
    bars: list[str] = []
    for series_index, values in enumerate(series_values):
        for value_index, value in enumerate(values):
            height = abs(value) / maximum * 210
            x = 52 + value_index * group_width + series_index * bar_width
            y = 252 - height if value >= 0 else 252
            bars.append(
                f'<rect x="{x:.2f}" y="{y:.2f}" width="{bar_width:.2f}" '
                f'height="{height:.2f}" fill="{palette[series_index % len(palette)]}"/>'
            )
    return (
        '<svg class="workbook-chart-svg" viewBox="0 0 600 300" '
        'preserveAspectRatio="none">'
        '<line x1="45" y1="252" x2="575" y2="252" stroke="#9ca3af"/>'
        f"{''.join(bars)}</svg>"
    )


def _xlsx_image_data_uri(image: Any) -> str:
    extension = str(getattr(image, "format", "png") or "png").lower()
    mime = {
        "jpg": "image/jpeg",
        "jpeg": "image/jpeg",
        "gif": "image/gif",
        "webp": "image/webp",
        "bmp": "image/bmp",
    }.get(extension, "image/png")
    encoded = base64.b64encode(image._data()).decode("ascii")
    return f"data:{mime};base64,{encoded}"


def _xlsx_column_width_px(sheet: Any, column: int) -> float:
    from openpyxl.utils import get_column_letter

    dimension = sheet.column_dimensions[get_column_letter(column)]
    if dimension.hidden:
        return 1
    return max(1, min((dimension.width or 10) * 7, 280))


def _xlsx_row_height_px(sheet: Any, row: int) -> float:
    dimension = sheet.row_dimensions[row]
    if dimension.hidden:
        return 1
    return (dimension.height or 18.75) * 96 / 72


def _xlsx_anchor_box(sheet: Any, drawing: Any) -> tuple[float, float, float, float]:
    """Resolve an openpyxl drawing anchor to CSS pixels."""
    anchor = drawing.anchor
    marker = anchor._from
    left = 44 + sum(
        _xlsx_column_width_px(sheet, column) for column in range(1, int(marker.col) + 1)
    )
    top = 25 + sum(
        _xlsx_row_height_px(sheet, row) for row in range(1, int(marker.row) + 1)
    )
    left += float(marker.colOff or 0) / 914400 * 96
    top += float(marker.rowOff or 0) / 914400 * 96
    extent = getattr(anchor, "ext", None)
    if extent is not None:
        width = float(extent.cx) / 914400 * 96
        height = float(extent.cy) / 914400 * 96
    else:
        width = float(getattr(drawing, "width", 320) or 320)
        height = float(getattr(drawing, "height", 180) or 180)
    return left, top, max(width, 1), max(height, 1)


def _render_xlsx(source: Path) -> str:
    from openpyxl import load_workbook
    from openpyxl.utils import get_column_letter

    workbook = load_workbook(source, data_only=False, read_only=False)
    evaluation = evaluate_workbook_formulas(workbook)
    sections: list[str] = []
    for sheet in workbook.worksheets:
        max_row, max_column = _worksheet_bounds(sheet)
        merged_starts = {
            (cell_range.min_row, cell_range.min_col): cell_range
            for cell_range in sheet.merged_cells.ranges
        }
        merged_covered: set[tuple[int, int]] = set()
        for cell_range in sheet.merged_cells.ranges:
            for row in range(cell_range.min_row, cell_range.max_row + 1):
                for column in range(cell_range.min_col, cell_range.max_col + 1):
                    if (row, column) != (cell_range.min_row, cell_range.min_col):
                        merged_covered.add((row, column))

        column_widths = [
            _xlsx_column_width_px(sheet, column) for column in range(1, max_column + 1)
        ]
        row_heights = [_xlsx_row_height_px(sheet, row) for row in range(1, max_row + 1)]
        columns = ['<col style="width:44px">']
        column_headers = ['<th class="corner"></th>']
        for column in range(1, max_column + 1):
            width_px = column_widths[column - 1]
            columns.append(f'<col style="width:{width_px:.0f}px">')
            column_headers.append(
                f'<th style="min-width:{width_px:.0f}px">'
                f"{get_column_letter(column)}</th>"
            )

        rows = [f"<tr>{''.join(column_headers)}</tr>"]
        for row_index in range(1, max_row + 1):
            row_height = row_heights[row_index - 1]
            cells = [f'<th class="row-number">{row_index}</th>']
            for column_index in range(1, max_column + 1):
                if (row_index, column_index) in merged_covered:
                    continue
                cell = sheet.cell(row_index, column_index)
                merged = merged_starts.get((row_index, column_index))
                spans = ""
                if merged:
                    spans = (
                        f' rowspan="{merged.max_row - merged.min_row + 1}"'
                        f' colspan="{merged.max_col - merged.min_col + 1}"'
                    )
                styles: list[str] = []
                fill = _css_color(cell.fill.fgColor)
                if fill != "transparent" and cell.fill.fill_type:
                    styles.append(f"background:{fill}")
                if cell.font.bold:
                    styles.append("font-weight:700")
                if cell.font.italic:
                    styles.append("font-style:italic")
                if cell.font.name:
                    styles.append(f"font-family:{_css_font_family(cell.font.name)}")
                if cell.font.sz:
                    styles.append(f"font-size:{float(cell.font.sz):.2f}pt")
                color = _css_color(cell.font.color)
                if color != "transparent":
                    styles.append(f"color:{color}")
                if cell.alignment.horizontal:
                    styles.append(f"text-align:{cell.alignment.horizontal}")
                if cell.alignment.vertical:
                    styles.append(f"vertical-align:{cell.alignment.vertical}")
                if cell.alignment.wrap_text:
                    styles.append("white-space:normal")
                for side_name in ("top", "right", "bottom", "left"):
                    border = _border_css(getattr(cell.border, side_name), side_name)
                    if border:
                        styles.append(border)
                style_attr = f' style="{";".join(styles)}"' if styles else ""
                value = html.escape(
                    _display_cell(
                        cell,
                        evaluation=evaluation,
                        sheet_name=sheet.title,
                    )
                )
                formula_class = " formula" if cell.data_type == "f" else ""
                formula_attr = ""
                if cell.data_type == "f":
                    formula = str(cell.value or "")
                    if not formula.startswith("="):
                        formula = f"={formula}"
                    formula_attr = f' data-formula="{html.escape(formula, quote=True)}"'
                cells.append(
                    f'<td class="cell{formula_class}" data-cell="{cell.coordinate}" '
                    f'data-qa-label="{html.escape(cell.coordinate, quote=True)}"'
                    f"{formula_attr}{spans}{style_attr}>{value}</td>"
                )
            rows.append(f'<tr style="height:{row_height:.2f}px">{"".join(cells)}</tr>')
        truncated = sheet.max_row > max_row or sheet.max_column > max_column
        notice = (
            '<p class="notice">Preview limited to the first 500 rows and 100 columns.</p>'
            if truncated
            else ""
        )
        objects: list[str] = []
        stage_width = 44 + sum(column_widths)
        # CSS cells have a 25px minimum height.  Compact workbook profiles may
        # request slightly shorter rows, so use the rendered minimum here or
        # the final rows can extend beyond the stage and be falsely clipped.
        stage_height = 32 + sum(max(height, 25) + 1 for height in row_heights)
        for chart_index, chart in enumerate(sheet._charts, start=1):
            title = f"Chart {chart_index}"
            try:
                text = chart.title.tx.rich.p[0].r[0].t
                if text:
                    title = str(text)
            except (AttributeError, IndexError, TypeError):
                pass
            left, top, width, height = _xlsx_anchor_box(sheet, chart)
            stage_width = max(stage_width, left + width)
            stage_height = max(stage_height, top + height)
            objects.append(
                '<figure class="workbook-chart" '
                f'data-qa-label="{html.escape(title, quote=True)}" '
                f'style="left:{left:.2f}px;top:{top:.2f}px;'
                f'width:{width:.2f}px;height:{height:.2f}px">'
                f"<figcaption>{html.escape(title)}</figcaption>"
                f"{_xlsx_chart_svg(workbook, chart, evaluation)}</figure>"
            )
        for image_index, image in enumerate(sheet._images, start=1):
            left, top, width, height = _xlsx_anchor_box(sheet, image)
            stage_width = max(stage_width, left + width)
            stage_height = max(stage_height, top + height)
            objects.append(
                f'<img class="workbook-image" data-qa-label="Image {image_index}" '
                f'style="left:{left:.2f}px;top:{top:.2f}px;'
                f'width:{width:.2f}px;height:{height:.2f}px" '
                f'src="{_xlsx_image_data_uri(image)}" alt="Image {image_index}">'
            )
        sections.append(
            f'<section class="sheet" data-preview-item '
            f'data-preview-label="{html.escape(sheet.title, quote=True)}">'
            f"<h2>{html.escape(sheet.title)}</h2>{notice}"
            f'<div class="grid-wrap"><div class="grid-stage" '
            f'style="width:{stage_width:.2f}px;height:{stage_height:.2f}px">'
            f"<table><colgroup>{''.join(columns)}</colgroup>"
            f"{''.join(rows)}</table>{''.join(objects)}</div></div></section>"
        )

    css = """
    *{box-sizing:border-box}body{margin:0;padding:24px;background:#f4f5f7;
    color:#202124;font-family:Arial,sans-serif}.sheet{margin:0 auto 28px;background:#fff;
    border:1px solid #d9dce1;border-radius:8px;box-shadow:0 2px 8px #0001;overflow:hidden}
    h2{position:sticky;left:0;margin:0;padding:14px 18px;background:#fff;border-bottom:
    1px solid #e3e6ea;font-size:15px}.notice{margin:0;padding:8px 18px;background:#fff8dc;
    font-size:12px}.grid-wrap{overflow:auto;max-height:76vh}.grid-stage{position:relative}
    table{position:absolute;left:0;top:0;border-collapse:separate;
    border-spacing:0;font-size:12px}th,td{height:25px;padding:3px 6px;border-right:1px
    solid #e1e4e8;border-bottom:1px solid #e1e4e8;white-space:pre;vertical-align:middle}
    th{position:sticky;top:0;z-index:2;background:#f1f3f4;color:#5f6368;text-align:center;
    font-weight:500}.row-number{left:0;z-index:1;min-width:44px}.corner{left:0;z-index:3}
    td{overflow:hidden;text-overflow:ellipsis}
    .formula{color:#174ea6}.cell:empty{background:#fff}
    .workbook-chart{position:absolute;margin:0;border:1px solid #d9dce1;
    background:#fff;padding:10px;overflow:hidden}.workbook-chart figcaption{
    font-weight:600;margin-bottom:5px}.workbook-chart-svg{display:block;width:100%;
    height:calc(100% - 22px)}.workbook-image{position:absolute;object-fit:contain}
    .chart-empty{color:#6b7280}
    """
    return _page(title=source.name, body="".join(sections), css=css)


def _picture_data_uri(shape: Any) -> str:
    relation_id = shape._pic.blip_rId
    related_part = shape.part.related_part(relation_id)
    extension = (related_part.partname.ext or "png").lower()
    mime = {
        "jpg": "image/jpeg",
        "jpeg": "image/jpeg",
        "gif": "image/gif",
        "webp": "image/webp",
        "svg": "image/svg+xml",
    }.get(extension, getattr(related_part, "content_type", "image/png"))
    encoded = base64.b64encode(related_part.blob).decode("ascii")
    return f"data:{mime};base64,{encoded}"


def _local_name(element: Any) -> str:
    return str(getattr(element, "tag", "")).rsplit("}", 1)[-1]


def _xml_child(element: Any | None, name: str) -> Any | None:
    if element is None:
        return None
    for child in element:
        if _local_name(child) == name:
            return child
    return None


def _shape_transform_css(shape: Any) -> str | None:
    transforms = shape._element.xpath(".//*[local-name()='xfrm']")
    transform = transforms[0] if transforms else None
    values: list[str] = []
    rotation = float(getattr(shape, "rotation", 0) or 0)
    if rotation:
        values.append(f"rotate({rotation:.3f}deg)")
    if transform is not None and transform.get("flipH") == "1":
        values.append("scaleX(-1)")
    if transform is not None and transform.get("flipV") == "1":
        values.append("scaleY(-1)")
    return " ".join(values) or None


def _shape_box(
    shape: Any,
    *,
    slide_width: int,
    slide_height: int,
    coordinate_left: int = 0,
    coordinate_top: int = 0,
    coordinate_width: int | None = None,
    coordinate_height: int | None = None,
) -> str:
    width_basis = coordinate_width or slide_width
    height_basis = coordinate_height or slide_height
    left = 100 * (int(shape.left) - coordinate_left) / max(width_basis, 1)
    top = 100 * (int(shape.top) - coordinate_top) / max(height_basis, 1)
    width = 100 * int(shape.width) / max(width_basis, 1)
    height = 100 * int(shape.height) / max(height_basis, 1)
    styles = f"left:{left:.3f}%;top:{top:.3f}%;width:{width:.3f}%;height:{height:.3f}%"
    transform = _shape_transform_css(shape)
    if transform:
        styles += f";transform:{transform}"
    return styles


def _theme_part_root(slide: Any) -> Any | None:
    from xml.etree import ElementTree

    for relation in slide.slide_layout.slide_master.part.rels.values():
        if relation.reltype.endswith("/theme"):
            return ElementTree.fromstring(relation.target_part.blob)
    return None


def _pptx_theme_palette(slide: Any) -> dict[str, str]:
    """Extract the color scheme attached to the slide's master."""
    root = _theme_part_root(slide)
    if root is None:
        return {}
    namespace = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    scheme = root.find(".//a:themeElements/a:clrScheme", namespace)
    if scheme is None:
        return {}
    palette: dict[str, str] = {}
    for entry in scheme:
        if not list(entry):
            continue
        color = list(entry)[0]
        raw = color.get("val") or color.get("lastClr")
        css = _css_color(raw)
        if css != "transparent":
            palette[entry.tag.rsplit("}", 1)[-1]] = css
    palette.update(
        {
            "bg1": palette.get("lt1", "#ffffff"),
            "tx1": palette.get("dk1", "#000000"),
            "bg2": palette.get("lt2", "#e7e6e6"),
            "tx2": palette.get("dk2", "#44546a"),
        }
    )
    return palette


def _pptx_theme_fonts(slide: Any) -> dict[str, str]:
    root = _theme_part_root(slide)
    if root is None:
        return {}
    namespace = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    fonts: dict[str, str] = {}
    paths = {
        "major-latin": ".//a:fontScheme/a:majorFont/a:latin",
        "major-east-asian": ".//a:fontScheme/a:majorFont/a:ea",
        "major-complex": ".//a:fontScheme/a:majorFont/a:cs",
        "minor-latin": ".//a:fontScheme/a:minorFont/a:latin",
        "minor-east-asian": ".//a:fontScheme/a:minorFont/a:ea",
        "minor-complex": ".//a:fontScheme/a:minorFont/a:cs",
    }
    for key, path in paths.items():
        node = root.find(path, namespace)
        if node is not None and node.get("typeface"):
            fonts[key] = str(node.get("typeface"))
    return fonts


def _apply_brightness(color: str, brightness: float) -> str:
    if not color.startswith("#") or len(color) != 7 or not brightness:
        return color
    channels = [int(color[index : index + 2], 16) for index in (1, 3, 5)]
    if brightness > 0:
        channels = [
            round(channel + (255 - channel) * brightness) for channel in channels
        ]
    else:
        channels = [round(channel * (1 + brightness)) for channel in channels]
    return "#" + "".join(f"{max(0, min(channel, 255)):02x}" for channel in channels)


def _pptx_color(
    color_format: Any,
    palette: dict[str, str],
    default: str = "transparent",
) -> str:
    direct = _css_color(color_format)
    if direct != "transparent":
        color = direct
    else:
        theme_color = getattr(color_format, "theme_color", None)
        theme_name = getattr(theme_color, "name", "")
        theme_key = {
            "DARK_1": "dk1",
            "LIGHT_1": "lt1",
            "DARK_2": "dk2",
            "LIGHT_2": "lt2",
            "ACCENT_1": "accent1",
            "ACCENT_2": "accent2",
            "ACCENT_3": "accent3",
            "ACCENT_4": "accent4",
            "ACCENT_5": "accent5",
            "ACCENT_6": "accent6",
            "HYPERLINK": "hlink",
            "FOLLOWED_HYPERLINK": "folHlink",
            "BACKGROUND_1": "bg1",
            "TEXT_1": "tx1",
            "BACKGROUND_2": "bg2",
            "TEXT_2": "tx2",
        }.get(theme_name)
        color = palette.get(theme_key or "", default)
    try:
        brightness = float(color_format.brightness or 0)
    except (AttributeError, TypeError, ValueError):
        brightness = 0
    return _apply_brightness(color, brightness)


def _slide_background(slide: Any, palette: dict[str, str]) -> str:
    """Resolve an explicit slide/layout/master background color."""
    drawing = "http://schemas.openxmlformats.org/drawingml/2006/main"
    for layer in (slide, slide.slide_layout, slide.slide_layout.slide_master):
        background = layer._element.cSld.bg
        if background is None:
            continue
        for tag, attribute in (
            (f".//{{{drawing}}}srgbClr", "val"),
            (f".//{{{drawing}}}sysClr", "lastClr"),
        ):
            node = background.find(tag)
            if node is not None:
                color = _css_color(node.get(attribute))
                if color != "transparent":
                    return color
        scheme = background.find(f".//{{{drawing}}}schemeClr")
        if scheme is not None:
            return palette.get(scheme.get("val", ""), "#ffffff")
    return "#ffffff"


def _xml_drawing_color_css(
    color_node: Any | None, palette: dict[str, str]
) -> str | None:
    """Resolve one DrawingML color node, including common luminance transforms."""

    if color_node is None:
        return None
    kind = _local_name(color_node)
    if kind == "srgbClr":
        color = _css_color(color_node.get("val"))
    elif kind == "sysClr":
        color = _css_color(color_node.get("lastClr"))
    elif kind == "schemeClr":
        color = palette.get(str(color_node.get("val") or ""), "transparent")
    else:
        color = "transparent"
    if color == "transparent":
        return None

    brightness = 0.0
    tint = _xml_child(color_node, "tint")
    shade = _xml_child(color_node, "shade")
    luminance_mod = _xml_child(color_node, "lumMod")
    luminance_offset = _xml_child(color_node, "lumOff")
    if tint is not None:
        brightness = max(0.0, min(int(tint.get("val", "0")) / 100000, 1.0))
    elif shade is not None:
        brightness = max(-1.0, min(int(shade.get("val", "100000")) / 100000 - 1, 0.0))
    elif luminance_mod is not None or luminance_offset is not None:
        mod = (
            int(luminance_mod.get("val", "100000")) / 100000
            if luminance_mod is not None
            else 1
        )
        offset = (
            int(luminance_offset.get("val", "0")) / 100000
            if luminance_offset is not None
            else 0
        )
        brightness = offset if offset else mod - 1
    return _apply_brightness(color, brightness)


def _xml_color_css(element: Any | None, palette: dict[str, str]) -> str | None:
    if element is None:
        return None
    fill = _xml_child(element, "solidFill")
    if fill is None:
        return None
    return _xml_drawing_color_css(next(iter(fill), None), palette)


def _matching_placeholder(container: Any, shape: Any) -> Any | None:
    if not getattr(shape, "is_placeholder", False):
        return None
    try:
        index = int(shape.placeholder_format.idx)
        placeholder_type = getattr(shape.placeholder_format.type, "name", "")
    except (AttributeError, TypeError, ValueError):
        return None
    same_type = None
    for candidate in container.placeholders:
        try:
            if int(candidate.placeholder_format.idx) == index:
                return candidate
            if (
                same_type is None
                and getattr(candidate.placeholder_format.type, "name", "")
                == placeholder_type
            ):
                same_type = candidate
        except (AttributeError, TypeError, ValueError):
            continue
    return same_type


def _placeholder_style_name(shape: Any) -> str:
    try:
        name = str(shape.placeholder_format.type.name)
    except (AttributeError, TypeError, ValueError):
        return "otherStyle"
    if "TITLE" in name:
        return "titleStyle"
    if name in {"BODY", "OBJECT", "TEXT", "VERTICAL_BODY"}:
        return "bodyStyle"
    return "otherStyle"


def _level_properties(list_style: Any | None, level: int) -> Any | None:
    if list_style is None:
        return None
    nodes = list_style.xpath(f"./*[local-name()='lvl{min(max(level, 0), 8) + 1}pPr']")
    if nodes:
        return nodes[0]
    defaults = list_style.xpath("./*[local-name()='defPPr']")
    return defaults[0] if defaults else None


def _placeholder_paragraph_properties(placeholder: Any, level: int) -> Any | None:
    try:
        paragraphs = placeholder.text_frame.paragraphs
    except (AttributeError, TypeError, ValueError):
        return None
    fallback = None
    for paragraph in paragraphs:
        properties = getattr(paragraph._p, "pPr", None)
        if properties is None:
            continue
        if fallback is None:
            fallback = properties
        if int(getattr(paragraph, "level", 0) or 0) == level:
            return properties
    return fallback


def _effective_paragraph_nodes(shape: Any, paragraph: Any, slide: Any) -> list[Any]:
    """Return paragraph properties from most specific to least specific."""

    level = int(getattr(paragraph, "level", 0) or 0)
    result: list[Any] = []
    direct = getattr(paragraph._p, "pPr", None)
    if direct is not None:
        result.append(direct)

    shape_list = getattr(shape.text_frame._txBody, "lstStyle", None)
    local = _level_properties(shape_list, level)
    if local is not None:
        result.append(local)

    layout_placeholder = _matching_placeholder(slide.slide_layout, shape)
    master_placeholder = None
    if layout_placeholder is not None:
        layout_paragraph = _placeholder_paragraph_properties(layout_placeholder, level)
        if layout_paragraph is not None:
            result.append(layout_paragraph)
        layout_list = getattr(layout_placeholder.text_frame._txBody, "lstStyle", None)
        layout_properties = _level_properties(layout_list, level)
        if layout_properties is not None:
            result.append(layout_properties)
        master_placeholder = _matching_placeholder(
            slide.slide_layout.slide_master, layout_placeholder
        )
    elif getattr(shape, "is_placeholder", False):
        master_placeholder = _matching_placeholder(
            slide.slide_layout.slide_master, shape
        )

    if master_placeholder is not None:
        master_paragraph = _placeholder_paragraph_properties(master_placeholder, level)
        if master_paragraph is not None:
            result.append(master_paragraph)
        master_list = getattr(master_placeholder.text_frame._txBody, "lstStyle", None)
        master_properties = _level_properties(master_list, level)
        if master_properties is not None:
            result.append(master_properties)

    master = slide.slide_layout.slide_master
    tx_styles = master._element.xpath("./*[local-name()='txStyles']")
    if tx_styles:
        style_name = _placeholder_style_name(shape)
        style = tx_styles[0].xpath(f"./*[local-name()='{style_name}']")
        if style:
            master_properties = _level_properties(style[0], level)
            if master_properties is not None:
                result.append(master_properties)
    return result


def _effective_body_nodes(shape: Any, slide: Any) -> list[Any]:
    result: list[Any] = []
    direct = getattr(shape.text_frame._txBody, "bodyPr", None)
    if direct is not None:
        result.append(direct)
    layout_placeholder = _matching_placeholder(slide.slide_layout, shape)
    master_placeholder = None
    if layout_placeholder is not None:
        body = getattr(layout_placeholder.text_frame._txBody, "bodyPr", None)
        if body is not None:
            result.append(body)
        master_placeholder = _matching_placeholder(
            slide.slide_layout.slide_master, layout_placeholder
        )
    elif getattr(shape, "is_placeholder", False):
        master_placeholder = _matching_placeholder(
            slide.slide_layout.slide_master, shape
        )
    if master_placeholder is not None:
        body = getattr(master_placeholder.text_frame._txBody, "bodyPr", None)
        if body is not None:
            result.append(body)
    return result


def _first_attribute(
    nodes: list[Any], name: str, default: str | None = None
) -> str | None:
    for node in nodes:
        value = node.get(name)
        if value is not None:
            return str(value)
    return default


def _run_property_nodes(
    run_properties: Any | None, paragraph_nodes: list[Any]
) -> list[Any]:
    result = [run_properties] if run_properties is not None else []
    for paragraph_node in paragraph_nodes:
        default_run = _xml_child(paragraph_node, "defRPr")
        if default_run is not None:
            result.append(default_run)
    return result


def _theme_font_name(typeface: str, fonts: dict[str, str]) -> str:
    mapping = {
        "+mj-lt": "major-latin",
        "+mj-ea": "major-east-asian",
        "+mj-cs": "major-complex",
        "+mn-lt": "minor-latin",
        "+mn-ea": "minor-east-asian",
        "+mn-cs": "minor-complex",
    }
    return fonts.get(mapping.get(typeface, ""), "") or typeface


def _text_length_cqw(value: float, slide_width: int) -> str:
    return f"{value / max(slide_width, 1) * 100:.4f}cqw"


def _point_cqw(value: float, slide_width: int) -> str:
    return _text_length_cqw(value * 12700, slide_width)


def _xml_bool(value: str | None) -> bool | None:
    if value is None:
        return None
    return value.casefold() in {"1", "true", "on", "yes"}


def _effective_run_styles(
    run_properties: Any | None,
    paragraph_nodes: list[Any],
    *,
    palette: dict[str, str],
    fonts: dict[str, str],
    slide_width: int,
    default_size_pt: float = 18.0,
) -> tuple[list[str], float]:
    nodes = _run_property_nodes(run_properties, paragraph_nodes)
    raw_size = _first_attribute(nodes, "sz")
    try:
        size_pt = int(raw_size) / 100 if raw_size is not None else default_size_pt
    except ValueError:
        size_pt = default_size_pt
    styles = [f"font-size:{_point_cqw(size_pt, slide_width)}"]

    typeface = ""
    for node in nodes:
        latin = _xml_child(node, "latin")
        if latin is not None and latin.get("typeface"):
            typeface = _theme_font_name(str(latin.get("typeface")), fonts)
            break
    if not typeface:
        typeface = fonts.get("minor-latin", "Arial")
    styles.append(f"font-family:{_css_font_family(typeface)}")

    raw_bold = _first_attribute(nodes, "b")
    raw_italic = _first_attribute(nodes, "i")
    bold = _xml_bool(raw_bold)
    italic = _xml_bool(raw_italic)
    underline = _first_attribute(nodes, "u")
    strike = _first_attribute(nodes, "strike")
    if bold is not None:
        styles.append("font-weight:700" if bold else "font-weight:400")
    if italic is not None:
        styles.append("font-style:italic" if italic else "font-style:normal")
    decorations: list[str] = []
    if underline and underline not in {"none", "0", "false"}:
        decorations.append("underline")
    if strike and strike not in {"noStrike", "none", "0", "false"}:
        decorations.append("line-through")
    if decorations:
        styles.append(f"text-decoration-line:{' '.join(decorations)}")
        if strike == "dblStrike":
            styles.append("text-decoration-style:double")
    elif underline is not None or strike is not None:
        styles.append("text-decoration:none")
    capitalization = _first_attribute(nodes, "cap")
    if capitalization == "all":
        styles.append("text-transform:uppercase")
    elif capitalization == "small":
        styles.append("font-variant:small-caps")
    raw_baseline = _first_attribute(nodes, "baseline")
    if raw_baseline is not None:
        try:
            baseline = int(raw_baseline)
        except ValueError:
            baseline = 0
        if baseline:
            styles.append(
                "vertical-align:super" if baseline > 0 else "vertical-align:sub"
            )
    raw_spacing = _first_attribute(nodes, "spc")
    if raw_spacing is not None:
        try:
            styles.append(
                f"letter-spacing:{_point_cqw(int(raw_spacing) / 100, slide_width)}"
            )
        except ValueError:
            pass
    for node in nodes:
        color = _xml_color_css(node, palette)
        if color:
            styles.append(f"color:{color}")
            break
    for node in nodes:
        highlight = _xml_child(node, "highlight")
        if highlight is None:
            continue
        color = _xml_drawing_color_css(next(iter(highlight), None), palette)
        if color:
            styles.append(f"background-color:{color}")
            break
    return styles, size_pt


def _spacing_css(
    paragraph_nodes: list[Any],
    element_name: str,
    *,
    slide_width: int,
) -> str | None:
    for properties in paragraph_nodes:
        container = _xml_child(properties, element_name)
        if container is None:
            continue
        points = _xml_child(container, "spcPts")
        percentage = _xml_child(container, "spcPct")
        if points is not None:
            try:
                return _point_cqw(int(points.get("val", "0")) / 100, slide_width)
            except ValueError:
                return None
        if percentage is not None:
            try:
                return f"{int(percentage.get('val', '0')) / 100000:.4f}em"
            except ValueError:
                return None
    return None


def _number_marker(number_format: str, value: int) -> str:
    def alpha(index: int) -> str:
        result = ""
        while index > 0:
            index -= 1
            result = chr(65 + index % 26) + result
            index //= 26
        return result

    def roman(index: int) -> str:
        values = (
            (1000, "M"),
            (900, "CM"),
            (500, "D"),
            (400, "CD"),
            (100, "C"),
            (90, "XC"),
            (50, "L"),
            (40, "XL"),
            (10, "X"),
            (9, "IX"),
            (5, "V"),
            (4, "IV"),
            (1, "I"),
        )
        result = ""
        for amount, token in values:
            while index >= amount:
                result += token
                index -= amount
        return result

    if number_format.startswith("alpha"):
        marker = alpha(value)
        if "Lc" in number_format:
            marker = marker.lower()
    elif number_format.startswith("roman"):
        marker = roman(value)
        if "Lc" in number_format:
            marker = marker.lower()
    else:
        marker = str(value)
    if number_format.endswith("ParenBoth"):
        return f"({marker})"
    if number_format.endswith("ParenR"):
        return f"{marker})"
    return f"{marker}." if number_format.endswith("Period") else marker


def _effective_bullet(
    paragraph_nodes: list[Any], counters: dict[int, int], level: int
) -> str:
    for properties in paragraph_nodes:
        for child in properties:
            name = _local_name(child)
            if name == "buNone":
                return ""
            if name == "buChar":
                return str(child.get("char") or "•")
            if name == "buAutoNum":
                start_at = child.get("startAt")
                if start_at is not None:
                    try:
                        counters[level] = int(start_at)
                    except ValueError:
                        pass
                value = counters.get(level, 1)
                counters[level] = value + 1
                return _number_marker(str(child.get("type") or "arabicPeriod"), value)
    return ""


def _first_run_properties(paragraph: Any) -> Any | None:
    for child in paragraph._p:
        if _local_name(child) in {"r", "fld"}:
            return getattr(child, "rPr", None)
    return None


def _bullet_marker_styles(
    paragraph: Any,
    paragraph_nodes: list[Any],
    *,
    palette: dict[str, str],
    fonts: dict[str, str],
    slide_width: int,
) -> tuple[list[str], float]:
    """Resolve bullet typography; DrawingML bullets follow the first run by default."""

    styles, size_pt = _effective_run_styles(
        _first_run_properties(paragraph),
        paragraph_nodes,
        palette=palette,
        fonts=fonts,
        slide_width=slide_width,
    )

    for properties in paragraph_nodes:
        bullet_font = next(
            (
                child
                for child in properties
                if _local_name(child) in {"buFont", "buFontTx"}
            ),
            None,
        )
        if bullet_font is None:
            continue
        if _local_name(bullet_font) == "buFont" and bullet_font.get("typeface"):
            typeface = _theme_font_name(str(bullet_font.get("typeface")), fonts)
            styles.append(f"font-family:{_css_font_family(typeface)}")
        break

    for properties in paragraph_nodes:
        bullet_color = next(
            (
                child
                for child in properties
                if _local_name(child) in {"buClr", "buClrTx"}
            ),
            None,
        )
        if bullet_color is None:
            continue
        if _local_name(bullet_color) == "buClr":
            color = _xml_drawing_color_css(next(iter(bullet_color), None), palette)
            if color:
                styles.append(f"color:{color}")
        break

    for properties in paragraph_nodes:
        bullet_size = next(
            (
                child
                for child in properties
                if _local_name(child) in {"buSzPct", "buSzPts", "buSzTx"}
            ),
            None,
        )
        if bullet_size is None:
            continue
        name = _local_name(bullet_size)
        try:
            if name == "buSzPct":
                size_pt *= int(bullet_size.get("val", "100000")) / 100000
                styles.append(f"font-size:{_point_cqw(size_pt, slide_width)}")
            elif name == "buSzPts":
                size_pt = int(bullet_size.get("val", "1800")) / 100
                styles.append(f"font-size:{_point_cqw(size_pt, slide_width)}")
        except ValueError:
            pass
        break

    styles.extend(("line-height:inherit", "white-space:nowrap"))
    return styles, size_pt


def _escaped_run_text(value: str) -> str:
    return html.escape(value).replace("\v", "<br>").replace("\n", "<br>")


def _paragraph_runs_html(
    paragraph: Any,
    paragraph_nodes: list[Any],
    *,
    palette: dict[str, str],
    fonts: dict[str, str],
    slide_width: int,
    default_size_pt: float = 18.0,
) -> tuple[str, float]:
    run_lookup = {id(run._r): run for run in paragraph.runs}
    parts: list[str] = []
    first_size = default_size_pt
    for child in paragraph._p:
        name = _local_name(child)
        if name == "br":
            parts.append("<br>")
            continue
        if name not in {"r", "fld"}:
            continue
        run = run_lookup.get(id(child))
        run_properties = getattr(child, "rPr", None)
        styles, size_pt = _effective_run_styles(
            run_properties,
            paragraph_nodes,
            palette=palette,
            fonts=fonts,
            slide_width=slide_width,
            default_size_pt=default_size_pt,
        )
        if not parts:
            first_size = size_pt
        text_value = run.text if run is not None else "".join(child.itertext())
        parts.append(
            f'<span style="{";".join(styles)}">{_escaped_run_text(text_value)}</span>'
        )
    if not parts and paragraph.text:
        styles, first_size = _effective_run_styles(
            None,
            paragraph_nodes,
            palette=palette,
            fonts=fonts,
            slide_width=slide_width,
            default_size_pt=default_size_pt,
        )
        parts.append(
            f'<span style="{";".join(styles)}">{_escaped_run_text(paragraph.text)}</span>'
        )
    return "".join(parts), first_size


def _shape_text(
    shape: Any,
    palette: dict[str, str],
    *,
    fonts: dict[str, str],
    slide: Any,
    slide_width: int,
) -> str:
    shape_name = str(getattr(shape, "name", "")).casefold()
    is_slide_title = "[role:title]" in shape_name
    body_nodes = _effective_body_nodes(shape, slide)
    wrap_mode = str(_first_attribute(body_nodes, "wrap", "square"))
    paragraphs: list[str] = []
    counters: dict[int, int] = {}
    for paragraph in shape.text_frame.paragraphs:
        paragraph_nodes = _effective_paragraph_nodes(shape, paragraph, slide)
        runs, _ = _paragraph_runs_html(
            paragraph,
            paragraph_nodes,
            palette=palette,
            fonts=fonts,
            slide_width=slide_width,
        )
        paragraph_styles: list[str] = []
        if is_slide_title:
            paragraph_styles.extend(
                (
                    "white-space:nowrap",
                    "overflow-wrap:normal",
                    "line-height:1.05",
                    "letter-spacing:-.015em",
                )
            )
        if wrap_mode == "none":
            paragraph_styles.append("white-space:pre")
        alignment = _first_attribute(paragraph_nodes, "algn", "l")
        paragraph_styles.append(
            "text-align:"
            + {
                "ctr": "center",
                "r": "right",
                "just": "justify",
                "dist": "justify",
                "thaiDist": "justify",
            }.get(str(alignment), "left")
        )
        before = _spacing_css(paragraph_nodes, "spcBef", slide_width=slide_width)
        after = _spacing_css(paragraph_nodes, "spcAft", slide_width=slide_width)
        line = _spacing_css(paragraph_nodes, "lnSpc", slide_width=slide_width)
        if before:
            paragraph_styles.append(f"margin-top:{before}")
        if after:
            paragraph_styles.append(f"margin-bottom:{after}")
        if line:
            paragraph_styles.append(f"line-height:{line}")

        level = int(getattr(paragraph, "level", 0) or 0)
        for nested_level in tuple(counters):
            if nested_level > level:
                counters.pop(nested_level, None)
        marker = _effective_bullet(paragraph_nodes, counters, level)
        raw_margin = _first_attribute(paragraph_nodes, "marL", "0")
        raw_right = _first_attribute(paragraph_nodes, "marR", "0")
        raw_indent = _first_attribute(paragraph_nodes, "indent", "0")
        try:
            margin = int(raw_margin or 0)
            right = int(raw_right or 0)
            indent = int(raw_indent or 0)
        except ValueError:
            margin = right = indent = 0
        if right:
            paragraph_styles.append(
                f"padding-right:{_text_length_cqw(right, slide_width)}"
            )
        marker_html = ""
        if marker:
            marker_position = max(margin + indent, 0)
            marker_styles, marker_size = _bullet_marker_styles(
                paragraph,
                paragraph_nodes,
                palette=palette,
                fonts=fonts,
                slide_width=slide_width,
            )
            marker_width = max(
                margin - marker_position,
                round(marker_size * 12700 * 0.8),
                1,
            )
            paragraph_styles.extend(
                (
                    "display:grid",
                    f"padding-left:{_text_length_cqw(marker_position, slide_width)}",
                    f"grid-template-columns:{_text_length_cqw(marker_width, slide_width)} 1fr",
                    "column-gap:0",
                )
            )
            marker_html = (
                f'<span class="bullet-marker" style="{";".join(marker_styles)}">'
                f"{html.escape(marker)}</span>"
            )
        elif margin or indent:
            paragraph_styles.append(
                f"padding-left:{_text_length_cqw(margin, slide_width)}"
            )
            if indent:
                paragraph_styles.append(
                    f"text-indent:{_text_length_cqw(indent, slide_width)}"
                )
        style_attr = (
            f' style="{";".join(paragraph_styles)}"' if paragraph_styles else ""
        )
        paragraphs.append(
            f'<p{style_attr}>{marker_html}<span class="paragraph-content">'
            f"{runs}</span></p>"
        )
    try:
        columns = max(int(_first_attribute(body_nodes, "numCol", "1") or 1), 1)
    except ValueError:
        columns = 1
    vertical = {
        "t": "flex-start",
        "ctr": "center",
        "b": "flex-end",
        "just": "space-between",
        "dist": "space-around",
    }.get(str(_first_attribute(body_nodes, "anchor", "t")), "flex-start")
    defaults = {"lIns": 91440, "rIns": 91440, "tIns": 45720, "bIns": 45720}
    margins: dict[str, int] = {}
    for name, default in defaults.items():
        try:
            margins[name] = int(
                _first_attribute(body_nodes, name, str(default)) or default
            )
        except ValueError:
            margins[name] = default
    padding = " ".join(
        _text_length_cqw(margins[name], slide_width)
        for name in ("tIns", "rIns", "bIns", "lIns")
    )
    common = f"padding:{padding};justify-content:{vertical};"
    if columns > 1:
        try:
            spacing = int(_first_attribute(body_nodes, "spcCol", "0") or 0)
        except ValueError:
            spacing = 0
        frame_style = (
            f"{common}column-count:{columns};"
            f"column-gap:{_text_length_cqw(spacing, slide_width)};"
            "column-fill:auto;display:block;overflow:hidden"
        )
    else:
        frame_style = f"{common}display:flex;flex-direction:column;overflow:hidden"
    return f'<div class="text-frame" style="{frame_style}">{"".join(paragraphs)}</div>'


def _shape_identity(shape: Any, layer: str) -> str:
    shape_id = html.escape(str(getattr(shape, "shape_id", "")), quote=True)
    name = html.escape(str(getattr(shape, "name", "")), quote=True)
    return (
        f'data-shape-id="{shape_id}" data-shape-name="{name}" '
        f'data-qa-label="{name or shape_id}" data-source-layer="{layer}"'
    )


def _shape_xml_color_css(
    color_node: Any | None,
    palette: dict[str, str],
    *,
    placeholder: str | None = None,
    default_alpha: float = 1.0,
) -> str | None:
    """Resolve a DrawingML color node, including theme transforms and alpha."""
    if color_node is None:
        return None
    kind = _local_name(color_node)
    if kind == "srgbClr":
        color = _css_color(color_node.get("val"))
    elif kind == "sysClr":
        color = _css_color(color_node.get("lastClr"))
    elif kind == "schemeClr":
        key = str(color_node.get("val") or "")
        aliases = {"lt1": "bg1", "dk1": "tx1", "lt2": "bg2", "dk2": "tx2"}
        color = (
            placeholder
            if key == "phClr"
            else palette.get(key) or palette.get(aliases.get(key, ""))
        )
    else:
        color = None
    if not color or color == "transparent":
        return None

    tint = _xml_child(color_node, "tint")
    shade = _xml_child(color_node, "shade")
    luminance_mod = _xml_child(color_node, "lumMod")
    luminance_offset = _xml_child(color_node, "lumOff")
    brightness = 0.0
    try:
        if tint is not None:
            brightness = max(0.0, min(int(tint.get("val", "0")) / 100000, 1.0))
        elif shade is not None:
            brightness = max(
                -1.0,
                min(int(shade.get("val", "100000")) / 100000 - 1, 0.0),
            )
        elif luminance_mod is not None or luminance_offset is not None:
            mod = (
                int(luminance_mod.get("val", "100000")) / 100000
                if luminance_mod is not None
                else 1
            )
            offset = (
                int(luminance_offset.get("val", "0")) / 100000
                if luminance_offset is not None
                else 0
            )
            brightness = offset if offset else mod - 1
    except (TypeError, ValueError):
        brightness = 0.0
    color = _apply_brightness(color, brightness)

    alpha = default_alpha
    alpha_node = _xml_child(color_node, "alpha")
    if alpha_node is not None:
        try:
            alpha = max(0.0, min(int(alpha_node.get("val", "100000")) / 100000, 1))
        except (TypeError, ValueError):
            alpha = default_alpha
    if alpha >= 0.999 or not color.startswith("#"):
        return color
    red, green, blue = (int(color[index : index + 2], 16) for index in (1, 3, 5))
    return f"rgba({red},{green},{blue},{alpha:.3f})"


def _shape_style_reference(shape: Any, name: str) -> Any | None:
    try:
        references = shape._element.xpath(
            f"./*[local-name()='style']/*[local-name()='{name}']"
        )
    except (AttributeError, TypeError, ValueError):
        return None
    return references[0] if references else None


def _shape_reference_color(
    shape: Any, name: str, palette: dict[str, str]
) -> str | None:
    reference = _shape_style_reference(shape, name)
    if reference is None:
        return None
    return _shape_xml_color_css(next(iter(reference), None), palette)


def _shape_theme_style_node(
    shape: Any,
    slide: Any,
    *,
    reference_name: str,
    list_name: str,
) -> Any | None:
    reference = _shape_style_reference(shape, reference_name)
    if reference is None:
        return None
    try:
        index = int(reference.get("idx", "0")) - 1
    except (TypeError, ValueError):
        return None
    if index < 0:
        return None
    theme = _theme_part_root(slide)
    if theme is None:
        return None
    drawing = "http://schemas.openxmlformats.org/drawingml/2006/main"
    styles = theme.findall(f".//{{{drawing}}}fmtScheme/{{{drawing}}}{list_name}/*")
    return styles[index] if index < len(styles) else None


def _shape_line_properties(
    shape: Any,
    *,
    slide: Any,
    palette: dict[str, str],
) -> tuple[str | None, float, str | None]:
    """Return effective line color, CSS-pixel width, and preset dash."""
    try:
        properties = shape._element.xpath(
            "./*[local-name()='spPr']/*[local-name()='ln']"
        )
    except (AttributeError, TypeError, ValueError):
        properties = []
    line = properties[0] if properties else None
    placeholder = _shape_reference_color(shape, "lnRef", palette)
    theme_line = _shape_theme_style_node(
        shape,
        slide,
        reference_name="lnRef",
        list_name="lnStyleLst",
    )
    effective = line if line is not None else theme_line
    if effective is None or _xml_child(effective, "noFill") is not None:
        return None, 1.0, None
    solid_fill = _xml_child(effective, "solidFill")
    color = _shape_xml_color_css(
        next(iter(solid_fill), None) if solid_fill is not None else None,
        palette,
        placeholder=placeholder,
    )
    if color is None and line is not None and theme_line is not None:
        solid_fill = _xml_child(theme_line, "solidFill")
        color = _shape_xml_color_css(
            next(iter(solid_fill), None) if solid_fill is not None else None,
            palette,
            placeholder=placeholder,
        )
    try:
        raw_width = int(
            line.get("w")
            if line is not None and line.get("w") is not None
            else theme_line.get("w", "9525")
            if theme_line is not None
            else "9525"
        )
    except (TypeError, ValueError):
        raw_width = 9525
    dash_node = _xml_child(line, "prstDash") if line is not None else None
    if dash_node is None and theme_line is not None:
        dash_node = _xml_child(theme_line, "prstDash")
    dash = str(dash_node.get("val")) if dash_node is not None else None
    return color, max(raw_width / 9525, 0.5), dash


def _line_dash_array(preset: str | None, width: float) -> str | None:
    units = {
        "dot": (1, 2.4),
        "sysDot": (1, 2),
        "dash": (4, 3),
        "sysDash": (3, 2),
        "lgDash": (8, 3),
        "dashDot": (4, 2, 1, 2),
        "sysDashDot": (3, 2, 1, 2),
        "lgDashDot": (8, 3, 1, 3),
        "lgDashDotDot": (8, 3, 1, 3, 1, 3),
        "sysDashDotDot": (3, 2, 1, 2, 1, 2),
    }.get(str(preset))
    if not units:
        return None
    return " ".join(f"{max(value * width, 1):.2f}" for value in units)


def _shape_shadow_node(shape: Any, slide: Any) -> Any | None:
    try:
        shadows = shape._element.xpath(
            "./*[local-name()='spPr']/*[local-name()='effectLst']"
            "/*[local-name()='outerShdw']"
        )
    except (AttributeError, TypeError, ValueError):
        shadows = []
    if shadows:
        return shadows[0]
    theme_effect = _shape_theme_style_node(
        shape,
        slide,
        reference_name="effectRef",
        list_name="effectStyleLst",
    )
    if theme_effect is None:
        return None
    return next(
        (node for node in theme_effect.iter() if _local_name(node) == "outerShdw"),
        None,
    )


def _shape_shadow_css(shape: Any, slide: Any, palette: dict[str, str]) -> str | None:
    shadow = _shape_shadow_node(shape, slide)
    if shadow is None:
        return None
    color_node = next(
        (
            child
            for child in shadow
            if _local_name(child) in {"srgbClr", "sysClr", "schemeClr"}
        ),
        None,
    )
    placeholder = _shape_reference_color(shape, "effectRef", palette)
    color = (
        _shape_xml_color_css(
            color_node,
            palette,
            placeholder=placeholder,
            default_alpha=0.16,
        )
        or "rgba(0,0,0,.16)"
    )
    try:
        distance_pt = int(shadow.get("dist", "0")) / 12700
        angle = math.radians(int(shadow.get("dir", "0")) / 60000)
        blur = int(shadow.get("blurRad", "0")) / 12700 * 96 / 72
    except (TypeError, ValueError):
        distance_pt = blur = 0
        angle = 0
    offset_x = math.cos(angle) * distance_pt * 96 / 72
    offset_y = math.sin(angle) * distance_pt * 96 / 72
    return f"{offset_x:.2f}px {offset_y:.2f}px {blur:.2f}px {color}"


def _connector_adjustment(shape: Any, default: float = 0.5) -> float:
    try:
        guides = shape._element.xpath(
            "./*[local-name()='spPr']/*[local-name()='prstGeom']"
            "/*[local-name()='avLst']/*[local-name()='gd']"
        )
    except (AttributeError, TypeError, ValueError):
        return default
    for guide in guides:
        formula = str(guide.get("fmla") or "")
        parts = formula.split()
        if len(parts) == 2 and parts[0] == "val":
            try:
                return max(0.0, min(float(parts[1]) / 100000, 1.0))
            except ValueError:
                continue
    return default


def _connector_marker(
    marker: Any | None,
    *,
    marker_id: str,
    color: str,
) -> tuple[str, str]:
    if marker is None:
        return "", ""
    marker_type = str(marker.get("type") or "none")
    if marker_type in {"", "none"}:
        return "", ""
    size = {"sm": 2.8, "med": 3.6, "lg": 4.4}.get(
        str(marker.get("w") or marker.get("len") or "med"), 3.6
    )
    if marker_type == "diamond":
        path = f'<path d="M0 5L5 0L10 5L5 10Z" fill="{color}"/>'
        reference_x = "9"
    elif marker_type == "oval":
        path = f'<circle cx="5" cy="5" r="4" fill="{color}"/>'
        reference_x = "9"
    elif marker_type in {"arrow", "open"}:
        path = (
            f'<path d="M1 1L9 5L1 9" fill="none" stroke="{color}" '
            'stroke-width="1.7" stroke-linecap="round" stroke-linejoin="round"/>'
        )
        reference_x = "9"
    elif marker_type == "stealth":
        path = f'<path d="M0 0L10 5L0 10L3.5 5Z" fill="{color}"/>'
        reference_x = "9"
    else:
        path = f'<path d="M0 0L10 5L0 10Z" fill="{color}"/>'
        reference_x = "9"
    definition = (
        f'<marker id="{marker_id}" viewBox="0 0 10 10" refX="{reference_x}" '
        f'refY="5" markerWidth="{size:.1f}" markerHeight="{size:.1f}" '
        f'orient="auto-start-reverse">{path}</marker>'
    )
    return definition, f"url(#{marker_id})"


def _connector_svg(
    shape: Any,
    *,
    layer: str,
    slide: Any,
    slide_width: int,
    slide_height: int,
    palette: dict[str, str],
    coordinate_left: int = 0,
    coordinate_top: int = 0,
    coordinate_width: int | None = None,
    coordinate_height: int | None = None,
) -> str:
    width_basis = coordinate_width or slide_width
    height_basis = coordinate_height or slide_height
    left = 100 * (int(shape.left) - coordinate_left) / max(width_basis, 1)
    top = 100 * (int(shape.top) - coordinate_top) / max(height_basis, 1)
    width = 100 * int(shape.width) / max(width_basis, 1)
    height = 100 * int(shape.height) / max(height_basis, 1)
    transform = shape._element.xpath(".//*[local-name()='xfrm']")
    flip_horizontal = bool(transform and transform[0].get("flipH") == "1")
    flip_vertical = bool(transform and transform[0].get("flipV") == "1")
    x1, x2 = (1000, 0) if flip_horizontal else (0, 1000)
    y1, y2 = (1000, 0) if flip_vertical else (0, 1000)
    if shape.width == 0:
        x1 = x2 = 500
    if shape.height == 0:
        y1 = y2 = 500
    color, line_width, dash = _shape_line_properties(
        shape, slide=slide, palette=palette
    )
    color = color or "#66717c"
    try:
        line_nodes = shape._element.xpath(
            "./*[local-name()='spPr']/*[local-name()='ln']"
        )
    except (AttributeError, TypeError, ValueError):
        line_nodes = []
    line_node = line_nodes[0] if line_nodes else None
    head = _xml_child(line_node, "headEnd")
    tail = _xml_child(line_node, "tailEnd")
    marker_prefix = hashlib.sha1(
        f"{getattr(shape.part, 'partname', '')}:{getattr(shape, 'shape_id', 0)}".encode()
    ).hexdigest()[:12]
    head_definition, marker_start = _connector_marker(
        head,
        marker_id=f"connector-head-{marker_prefix}",
        color=color,
    )
    tail_definition, marker_end = _connector_marker(
        tail,
        marker_id=f"connector-tail-{marker_prefix}",
        color=color,
    )
    definitions = (
        f"<defs>{head_definition}{tail_definition}</defs>"
        if head_definition or tail_definition
        else ""
    )
    marker_attributes = ""
    if marker_start:
        marker_attributes += f' marker-start="{marker_start}"'
    if marker_end:
        marker_attributes += f' marker-end="{marker_end}"'
    dash_array = _line_dash_array(dash, line_width)
    dash_attribute = (
        f' stroke-dasharray="{dash_array}"' if dash_array is not None else ""
    )
    preset = _preset_geometry_name(shape)
    if preset.startswith("bentConnector"):
        adjustment = _connector_adjustment(shape)
        middle_x = x1 + (x2 - x1) * adjustment
        path = f"M{x1} {y1}H{middle_x:.2f}V{y2}H{x2}"
    elif preset.startswith("curvedConnector"):
        adjustment = _connector_adjustment(shape)
        middle_x = x1 + (x2 - x1) * adjustment
        path = f"M{x1} {y1}C{middle_x:.2f} {y1} {middle_x:.2f} {y2} {x2} {y2}"
    else:
        path = f"M{x1} {y1}L{x2} {y2}"
    shadow = _shape_shadow_css(shape, slide, palette)
    shadow_style = f' style="filter:drop-shadow({shadow})"' if shadow else ""
    identity = _shape_identity(shape, layer)
    style = (
        f"left:{left:.3f}%;top:{top:.3f}%;"
        f"width:max({width:.3f}%,2px);height:max({height:.3f}%,2px)"
    )
    rotation = float(getattr(shape, "rotation", 0) or 0)
    if rotation:
        style += f";transform:rotate({rotation:.3f}deg)"
    return (
        f'<svg class="shape connector" {identity} style="{style}" '
        f'data-connector-geometry="{html.escape(preset, quote=True)}" '
        'viewBox="0 0 1000 1000" preserveAspectRatio="none">'
        f'{definitions}<path d="{path}" fill="none" stroke="{color}" '
        f'stroke-width="{line_width:.2f}" stroke-linejoin="miter" '
        f'vector-effect="non-scaling-stroke"{dash_attribute}{marker_attributes}'
        f"{shadow_style}/></svg>"
    )


def _shape_svg_key(shape: Any, suffix: str) -> str:
    value = f"{getattr(shape.part, 'partname', '')}:{getattr(shape, 'shape_id', 0)}:{suffix}"
    return hashlib.sha1(value.encode()).hexdigest()[:12]


def _shape_gradient_svg(
    shape: Any,
    gradient: Any,
    *,
    palette: dict[str, str],
    placeholder: str | None,
) -> tuple[str, str] | None:
    stops: list[str] = []
    for stop in gradient:
        if _local_name(stop) != "gsLst":
            continue
        for gradient_stop in stop:
            if _local_name(gradient_stop) != "gs":
                continue
            color_node = next(iter(gradient_stop), None)
            color = _shape_xml_color_css(color_node, palette, placeholder=placeholder)
            if color is None:
                continue
            try:
                offset = int(gradient_stop.get("pos", "0")) / 1000
            except (TypeError, ValueError):
                offset = 0
            stops.append(
                f'<stop offset="{offset:.2f}%" stop-color="{html.escape(color, quote=True)}"/>'
            )
    if len(stops) < 2:
        return None
    gradient_id = f"shape-gradient-{_shape_svg_key(shape, 'gradient')}"
    linear = _xml_child(gradient, "lin")
    path = _xml_child(gradient, "path")
    if path is not None:
        definition = (
            f'<radialGradient id="{gradient_id}" cx="50%" cy="50%" r="70%">'
            f"{''.join(stops)}</radialGradient>"
        )
    else:
        try:
            office_angle = (
                int(linear.get("ang", "0")) / 60000 if linear is not None else 0
            )
        except (TypeError, ValueError):
            office_angle = 0
        definition = (
            f'<linearGradient id="{gradient_id}" x1="0" y1="0.5" x2="1" y2="0.5" '
            f'gradientTransform="rotate({office_angle:.3f} .5 .5)">'
            f"{''.join(stops)}</linearGradient>"
        )
    return definition, f"url(#{gradient_id})"


def _shape_pattern_svg(
    shape: Any,
    pattern: Any,
    *,
    palette: dict[str, str],
    placeholder: str | None,
) -> tuple[str, str, str]:
    foreground_node = _xml_child(_xml_child(pattern, "fgClr"), "srgbClr")
    background_node = _xml_child(_xml_child(pattern, "bgClr"), "srgbClr")
    if foreground_node is None:
        foreground = _xml_child(pattern, "fgClr")
        foreground_node = (
            next(iter(foreground), None) if foreground is not None else None
        )
    if background_node is None:
        background = _xml_child(pattern, "bgClr")
        background_node = (
            next(iter(background), None) if background is not None else None
        )
    foreground_color = (
        _shape_xml_color_css(foreground_node, palette, placeholder=placeholder)
        or "#000000"
    )
    background_color = (
        _shape_xml_color_css(background_node, palette, placeholder=placeholder)
        or "#ffffff"
    )
    preset = str(pattern.get("prst") or "diagCross")
    pattern_id = f"shape-pattern-{_shape_svg_key(shape, preset)}"
    if preset in {"cross", "smGrid", "lgGrid"}:
        marks = '<path d="M50 0V100M0 50H100"/>'
    elif preset in {"horz", "ltHorz", "dkHorz", "narHorz", "dashHorz"}:
        marks = '<path d="M0 50H100"/>'
    elif preset in {"vert", "ltVert", "dkVert", "narVert", "dashVert"}:
        marks = '<path d="M50 0V100"/>'
    elif preset in {"upDiag", "ltUpDiag", "dkUpDiag", "wdUpDiag", "dashUpDiag"}:
        marks = '<path d="M-25 75L75-25M25 125L125 25"/>'
    elif preset in {"dnDiag", "ltDnDiag", "dkDnDiag", "wdDnDiag", "dashDnDiag"}:
        marks = '<path d="M-25 25L75 125M25-25L125 75"/>'
    elif preset in {"pct5", "pct10", "pct20", "pct25", "pct30", "pct40"}:
        marks = '<circle cx="25" cy="25" r="4"/><circle cx="75" cy="75" r="4"/>'
    else:
        marks = '<path d="M-25 25L75 125M25-25L125 75M-25 75L75-25M25 125L125 25"/>'
    definition = (
        f'<pattern id="{pattern_id}" width="100" height="100" '
        'patternUnits="userSpaceOnUse">'
        f'<rect width="100" height="100" fill="{html.escape(background_color, quote=True)}"/>'
        f'<g fill="{html.escape(foreground_color, quote=True)}" '
        f'stroke="{html.escape(foreground_color, quote=True)}" stroke-width="3">'
        f"{marks}</g></pattern>"
    )
    return definition, f"url(#{pattern_id})", preset


def _shape_fill_svg(
    shape: Any,
    *,
    slide: Any,
    palette: dict[str, str],
) -> tuple[str, str, str]:
    try:
        properties = shape._element.xpath("./*[local-name()='spPr']")
    except (AttributeError, TypeError, ValueError):
        properties = []
    shape_properties = properties[0] if properties else None
    fill = (
        next(
            (
                child
                for child in shape_properties
                if _local_name(child) in {"solidFill", "gradFill", "pattFill", "noFill"}
            ),
            None,
        )
        if shape_properties is not None
        else None
    )
    placeholder = _shape_reference_color(shape, "fillRef", palette)
    if fill is None:
        fill = _shape_theme_style_node(
            shape,
            slide,
            reference_name="fillRef",
            list_name="fillStyleLst",
        )
    if fill is None or _local_name(fill) == "noFill":
        return "", "none", 'data-shape-fill="none"'
    name = _local_name(fill)
    if name == "solidFill":
        color = _shape_xml_color_css(
            next(iter(fill), None), palette, placeholder=placeholder
        )
        return "", color or "none", 'data-shape-fill="solid"'
    if name == "gradFill":
        gradient = _shape_gradient_svg(
            shape,
            fill,
            palette=palette,
            placeholder=placeholder,
        )
        if gradient:
            return gradient[0], gradient[1], 'data-shape-fill="gradient"'
    if name == "pattFill":
        definition, pattern_fill, preset = _shape_pattern_svg(
            shape,
            fill,
            palette=palette,
            placeholder=placeholder,
        )
        return (
            definition,
            pattern_fill,
            f'data-shape-fill="pattern" data-pattern="{html.escape(preset, quote=True)}"',
        )
    return "", placeholder or "none", 'data-shape-fill="fallback"'


def _preset_svg_geometry(preset: str) -> tuple[str, bool]:
    polygons = {
        "triangle": "500,0 1000,1000 0,1000",
        "rtTriangle": "0,0 1000,1000 0,1000",
        "diamond": "500,0 1000,500 500,1000 0,500",
        "parallelogram": "150,0 1000,0 850,1000 0,1000",
        "trapezoid": "200,0 800,0 1000,1000 0,1000",
        "pentagon": "500,0 1000,380 820,1000 180,1000 0,380",
        "hexagon": "250,0 750,0 1000,500 750,1000 250,1000 0,500",
        "octagon": "290,0 710,0 1000,290 1000,710 710,1000 290,1000 0,710 0,290",
        "chevron": "0,0 750,0 1000,500 750,1000 0,1000 250,500",
        "homePlate": "0,0 750,0 1000,500 750,1000 0,1000",
        "rightArrow": "0,200 650,200 650,0 1000,500 650,1000 650,800 0,800",
        "leftArrow": "350,0 350,200 1000,200 1000,800 350,800 350,1000 0,500",
        "upArrow": "500,0 1000,350 800,350 800,1000 200,1000 200,350 0,350",
        "downArrow": "200,0 800,0 800,650 1000,650 500,1000 0,650 200,650",
        "leftRightArrow": "200,0 200,250 800,250 800,0 1000,500 800,1000 800,750 200,750 200,1000 0,500",
        "star5": "500,0 610,350 980,350 680,570 790,920 500,700 210,920 320,570 20,350 390,350",
        "star4": "500,0 610,390 1000,500 610,610 500,1000 390,610 0,500 390,390",
        "plus": "350,0 650,0 650,350 1000,350 1000,650 650,650 650,1000 350,1000 350,650 0,650 0,350 350,350",
        "heart": "500,1000 70,550 0,320 80,120 250,20 420,80 500,220 580,80 750,20 920,120 1000,320 930,550",
    }
    if preset in {"roundRect", "round1Rect", "round2SameRect", "round2DiagRect"}:
        return '<rect x="0" y="0" width="1000" height="1000" rx="80" ry="80"/>', True
    if preset == "ellipse":
        return '<ellipse cx="500" cy="500" rx="500" ry="500"/>', True
    if preset == "arc":
        return '<path d="M500 0A500 500 0 0 1 1000 500L500 500Z"/>', True
    if preset == "cloud":
        return (
            '<path d="M110 665C20 630 15 505 95 450C55 330 160 235 275 265'
            "C320 120 475 75 585 175C690 70 860 130 880 275"
            "C1010 295 1045 450 950 535C980 665 845 770 725 710"
            'C650 845 455 850 365 720C255 795 135 760 110 665Z"/>',
            True,
        )
    if preset in polygons:
        return f'<polygon points="{polygons[preset]}"/>', True
    return '<rect x="0" y="0" width="1000" height="1000"/>', preset == "rect"


def _shape_surface_svg(
    shape: Any,
    *,
    slide: Any,
    palette: dict[str, str],
) -> str:
    preset = _preset_geometry_name(shape)
    geometry, supported = _preset_svg_geometry(preset)
    definitions, fill, fill_metadata = _shape_fill_svg(
        shape, slide=slide, palette=palette
    )
    line, line_width, dash = _shape_line_properties(shape, slide=slide, palette=palette)
    if fill == "none" and line is None:
        return ""
    dash_array = _line_dash_array(dash, line_width)
    geometry_attributes = (
        f'fill="{html.escape(fill, quote=True)}" '
        f'stroke="{html.escape(line, quote=True) if line else "none"}" '
        f'stroke-width="{line_width:.2f}" stroke-linejoin="round" '
        'vector-effect="non-scaling-stroke"'
    )
    if dash_array:
        geometry_attributes += f' stroke-dasharray="{dash_array}"'
    shadow = _shape_shadow_css(shape, slide, palette)
    surface_style = f' style="filter:drop-shadow({shadow})"' if shadow else ""
    fallback = "" if supported else ' data-geometry-fallback="rect"'
    return (
        f'<svg class="shape-surface" aria-hidden="true" '
        f'data-preset-geometry="{html.escape(preset, quote=True)}" '
        f'{fill_metadata}{fallback} viewBox="0 0 1000 1000" '
        f'preserveAspectRatio="none"{surface_style}>'
        f"<defs>{definitions}</defs>{geometry[:-2]} {geometry_attributes}/></svg>"
    )


def _shape_text_container_attributes(
    shape: Any, preset: str, palette: dict[str, str]
) -> tuple[list[str], str]:
    styles = ["overflow:visible"]
    font_color = _shape_reference_color(shape, "fontRef", palette)
    has_explicit_run_color = bool(
        shape._element.xpath(
            "./*[local-name()='txBody']//*[local-name()='rPr']"
            "/*[local-name()='solidFill' or local-name()='gradFill' or "
            "local-name()='pattFill']"
        )
    )
    color_attribute = ""
    if font_color and not has_explicit_run_color:
        styles.extend((f"color:{font_color}", f"--shape-text-color:{font_color}"))
        color_attribute = ' data-shape-font-color="theme"'
    explicit_alignment = any(
        getattr(paragraph._p, "pPr", None) is not None
        and paragraph._p.pPr.get("algn") is not None
        for paragraph in shape.text_frame.paragraphs
    )
    centered_presets = {
        "rightArrow",
        "leftArrow",
        "upArrow",
        "downArrow",
        "leftRightArrow",
    }
    attributes = (
        ' data-default-text-align="center"'
        if not explicit_alignment and preset in centered_presets
        else ""
    )
    attributes += color_attribute
    transforms = shape._element.xpath(".//*[local-name()='xfrm']")
    transform = transforms[0] if transforms else None
    flips: list[str] = []
    if transform is not None and transform.get("flipH") == "1":
        flips.append("horizontal")
    if transform is not None and transform.get("flipV") == "1":
        flips.append("vertical")
    if flips:
        attributes += f' data-text-unflip="{" ".join(flips)}"'
    return styles, attributes


def _picture_crop(shape: Any) -> tuple[float, float, float, float]:
    values: list[float] = []
    for name in ("crop_left", "crop_right", "crop_top", "crop_bottom"):
        try:
            values.append(float(getattr(shape, name) or 0))
        except (AttributeError, TypeError, ValueError):
            values.append(0.0)
    return values[0], values[1], values[2], values[3]


def _preset_geometry_name(shape: Any) -> str:
    try:
        geometries = shape._element.xpath(
            "./*[local-name()='spPr']/*[local-name()='prstGeom']"
        )
    except (AttributeError, TypeError, ValueError):
        return "rect"
    if not geometries:
        return "rect"
    return str(geometries[0].get("prst") or "rect")


def _preset_clip_css(preset: str) -> str | None:
    clip_paths = {
        "ellipse": "ellipse(50% 50% at 50% 50%)",
        "triangle": "polygon(50% 0,100% 100%,0 100%)",
        "rtTriangle": "polygon(0 0,100% 100%,0 100%)",
        "diamond": "polygon(50% 0,100% 50%,50% 100%,0 50%)",
        "parallelogram": "polygon(15% 0,100% 0,85% 100%,0 100%)",
        "trapezoid": "polygon(20% 0,80% 0,100% 100%,0 100%)",
        "pentagon": "polygon(50% 0,100% 38%,82% 100%,18% 100%,0 38%)",
        "hexagon": "polygon(25% 0,75% 0,100% 50%,75% 100%,25% 100%,0 50%)",
        "octagon": "polygon(29% 0,71% 0,100% 29%,100% 71%,71% 100%,29% 100%,0 71%,0 29%)",
        "chevron": "polygon(0 0,75% 0,100% 50%,75% 100%,0 100%,25% 50%)",
        "homePlate": "polygon(0 0,75% 0,100% 50%,75% 100%,0 100%)",
        "rightArrow": "polygon(0 20%,65% 20%,65% 0,100% 50%,65% 100%,65% 80%,0 80%)",
        "leftArrow": "polygon(35% 0,35% 20%,100% 20%,100% 80%,35% 80%,35% 100%,0 50%)",
        "upArrow": "polygon(50% 0,100% 35%,80% 35%,80% 100%,20% 100%,20% 35%,0 35%)",
        "downArrow": "polygon(20% 0,80% 0,80% 65%,100% 65%,50% 100%,0 65%,20% 65%)",
        "leftRightArrow": "polygon(20% 0,20% 25%,80% 25%,80% 0,100% 50%,80% 100%,80% 75%,20% 75%,20% 100%,0 50%)",
        "star5": "polygon(50% 0,61% 35%,98% 35%,68% 57%,79% 92%,50% 70%,21% 92%,32% 57%,2% 35%,39% 35%)",
        "star4": "polygon(50% 0,61% 39%,100% 50%,61% 61%,50% 100%,39% 61%,0 50%,39% 39%)",
        "plus": "polygon(35% 0,65% 0,65% 35%,100% 35%,100% 65%,65% 65%,65% 100%,35% 100%,35% 65%,0 65%,0 35%,35% 35%)",
        "heart": "polygon(50% 100%,7% 55%,0 32%,8% 12%,25% 2%,42% 8%,50% 22%,58% 8%,75% 2%,92% 12%,100% 32%,93% 55%)",
    }
    if preset == "roundRect":
        return "inset(0 round 8%)"
    return clip_paths.get(preset)


def _picture_alt_text(shape: Any) -> str:
    try:
        properties = shape._element.xpath(
            "./*[local-name()='nvPicPr']/*[local-name()='cNvPr']"
        )
        if properties:
            return str(properties[0].get("descr") or properties[0].get("title") or "")
    except (AttributeError, TypeError, ValueError):
        pass
    return ""


def _render_picture(
    shape: Any,
    *,
    layer: str,
    slide_width: int,
    slide_height: int,
    coordinate_left: int = 0,
    coordinate_top: int = 0,
    coordinate_width: int | None = None,
    coordinate_height: int | None = None,
) -> str:
    left_crop, right_crop, top_crop, bottom_crop = _picture_crop(shape)
    visible_width = max(1 - left_crop - right_crop, 0.0001)
    visible_height = max(1 - top_crop - bottom_crop, 0.0001)
    image_left = -left_crop / visible_width * 100
    image_top = -top_crop / visible_height * 100
    image_width = 100 / visible_width
    image_height = 100 / visible_height
    preset = _preset_geometry_name(shape)
    styles = [
        _shape_box(
            shape,
            slide_width=slide_width,
            slide_height=slide_height,
            coordinate_left=coordinate_left,
            coordinate_top=coordinate_top,
            coordinate_width=coordinate_width,
            coordinate_height=coordinate_height,
        )
    ]
    clip = _preset_clip_css(preset)
    if clip:
        styles.append(f"clip-path:{clip}")
    identity = _shape_identity(shape, layer)
    crop_attributes = (
        f'data-crop-left="{left_crop:.5f}" '
        f'data-crop-right="{right_crop:.5f}" '
        f'data-crop-top="{top_crop:.5f}" '
        f'data-crop-bottom="{bottom_crop:.5f}" '
        f'data-preset-geometry="{html.escape(preset, quote=True)}"'
    )
    image_style = (
        f"left:{image_left:.4f}%;top:{image_top:.4f}%;"
        f"width:{image_width:.4f}%;height:{image_height:.4f}%;object-fit:fill"
    )
    alt_text = html.escape(_picture_alt_text(shape), quote=True)
    return (
        f'<div class="shape picture-frame" {identity} {crop_attributes} '
        f'style="{";".join(styles)}">'
        f'<img class="picture" src="{_picture_data_uri(shape)}" alt="{alt_text}" '
        f'style="{image_style}" draggable="false"></div>'
    )


_CHART_FALLBACK_PALETTE = (
    "#4472c4",
    "#ed7d31",
    "#a5a5a5",
    "#ffc000",
    "#5b9bd5",
    "#70ad47",
)


def _chart_type_name(chart: Any) -> str:
    try:
        return str(chart.chart_type.name).upper()
    except (AttributeError, TypeError, ValueError):
        return "UNKNOWN"


def _chart_cache_values(series: Any, branch: str) -> list[float]:
    try:
        points = series._ser.xpath(
            f"./*[local-name()='{branch}']//*[local-name()='pt']"
        )
    except (AttributeError, TypeError, ValueError):
        return []
    indexed: list[tuple[int, float]] = []
    for point in points:
        value_nodes = point.xpath("./*[local-name()='v']")
        if not value_nodes:
            continue
        try:
            index = int(point.get("idx", len(indexed)))
            value = float(value_nodes[0].text or 0)
        except (TypeError, ValueError):
            continue
        indexed.append((index, value))
    return [value for _, value in sorted(indexed)]


def _chart_series_color(series: Any, index: int, theme_palette: dict[str, str]) -> str:
    for source in ("fill", "line"):
        try:
            color_format = (
                series.format.fill.fore_color
                if source == "fill"
                else series.format.line.color
            )
            color = _pptx_color(color_format, theme_palette)
        except (AttributeError, TypeError, ValueError):
            continue
        if color != "transparent":
            return color
    return _CHART_FALLBACK_PALETTE[index % len(_CHART_FALLBACK_PALETTE)]


def _chart_data_labels_enabled(chart: Any) -> bool:
    try:
        return any(bool(plot.has_data_labels) for plot in chart.plots)
    except (AttributeError, TypeError, ValueError):
        return False


def _chart_has_legend(chart: Any) -> bool:
    try:
        return bool(chart.has_legend)
    except (AttributeError, TypeError, ValueError):
        return False


def _chart_major_gridlines(chart: Any) -> bool:
    try:
        return _chart_axis_major_gridlines(chart.value_axis)
    except (AttributeError, TypeError, ValueError):
        return False


def _chart_axis_major_gridlines(axis: Any) -> bool:
    try:
        return bool(axis.has_major_gridlines)
    except (AttributeError, TypeError, ValueError):
        return False


def _chart_axis_major_unit(axis: Any, lower: float, upper: float) -> float:
    try:
        major_unit = float(axis.major_unit or 0)
    except (AttributeError, TypeError, ValueError):
        major_unit = 0
    if major_unit > 0:
        return major_unit

    span = max(upper - lower, 1e-12)
    if (
        span <= 10
        and math.isclose(lower, round(lower), abs_tol=1e-9)
        and math.isclose(upper, round(upper), abs_tol=1e-9)
    ):
        return 1.0

    rough_unit = span / 7
    magnitude = 10 ** math.floor(math.log10(rough_unit))
    fraction = rough_unit / magnitude
    nice_fraction = min(
        (1.0, 2.0, 2.5, 5.0, 10.0),
        key=lambda candidate: abs(math.log(candidate / fraction)),
    )
    return nice_fraction * magnitude


def _chart_scale(
    chart: Any,
    values: list[float],
    *,
    axis: Any | None = None,
) -> tuple[float, float]:
    data_lower = min([0.0, *values])
    data_upper = max([0.0, *values])
    if axis is None:
        try:
            axis = chart.value_axis
        except (AttributeError, TypeError, ValueError):
            axis = None
    explicit_lower = explicit_upper = None
    if axis is not None:
        try:
            explicit_lower = axis.minimum_scale
            explicit_upper = axis.maximum_scale
        except (AttributeError, TypeError, ValueError):
            pass

    lower = float(explicit_lower) if explicit_lower is not None else data_lower
    upper = float(explicit_upper) if explicit_upper is not None else data_upper
    if math.isclose(lower, upper):
        upper = lower + 1

    major_unit = _chart_axis_major_unit(axis, lower, upper)
    if explicit_lower is None:
        lower = math.floor((data_lower + major_unit * 1e-9) / major_unit) * major_unit
        if data_lower < 0 and math.isclose(
            lower, data_lower, abs_tol=major_unit * 1e-9
        ):
            lower -= major_unit
    if explicit_upper is None:
        upper = math.ceil((data_upper - major_unit * 1e-9) / major_unit) * major_unit
        if data_upper > 0 and math.isclose(
            upper, data_upper, abs_tol=major_unit * 1e-9
        ):
            upper += major_unit
    if math.isclose(lower, upper):
        upper = lower + major_unit
    return lower, upper


def _chart_tick_values(
    chart: Any,
    lower: float,
    upper: float,
    *,
    axis: Any | None = None,
) -> list[float]:
    if axis is None:
        try:
            axis = chart.value_axis
        except (AttributeError, TypeError, ValueError):
            axis = None
    major_unit = _chart_axis_major_unit(axis, lower, upper)
    first = math.ceil((lower - major_unit * 1e-9) / major_unit) * major_unit
    ticks: list[float] = []
    value = first
    while value <= upper + major_unit * 1e-9 and len(ticks) < 40:
        ticks.append(
            0.0 if math.isclose(value, 0, abs_tol=major_unit * 1e-9) else value
        )
        value += major_unit
    return ticks or [lower, upper]


def _chart_axes(
    chart: Any,
    *,
    lower: float,
    upper: float,
    baseline: float,
    gridlines: bool,
    axis: Any | None = None,
    plot_left: float = 62,
    plot_right: float = 574,
    plot_top: float = 28,
    plot_bottom: float = 248,
) -> str:
    parts: list[str] = []
    for tick in reversed(_chart_tick_values(chart, lower, upper, axis=axis)):
        y = plot_bottom - (tick - lower) / (upper - lower) * (plot_bottom - plot_top)
        if gridlines:
            parts.append(
                f'<line class="chart-gridline" x1="{plot_left:.2f}" '
                f'y1="{y:.2f}" x2="{plot_right:.2f}" y2="{y:.2f}" '
                'stroke="#d9d9d9" stroke-width="1"/>'
            )
        parts.append(
            f'<text class="chart-axis-label chart-axis-label-y" '
            f'x="{plot_left - 7:.2f}" y="{y + 4:.2f}" '
            'text-anchor="end" font-size="11" fill="#666">'
            f"{tick:g}</text>"
        )
    parts.extend(
        (
            f'<line class="chart-axis chart-axis-y" x1="{plot_left:.2f}" '
            f'y1="{plot_top:.2f}" x2="{plot_left:.2f}" '
            f'y2="{plot_bottom:.2f}" stroke="#7f7f7f"/>',
            f'<line class="chart-axis chart-axis-x" x1="{plot_left:.2f}" '
            f'y1="{baseline:.2f}" x2="{plot_right:.2f}" '
            f'y2="{baseline:.2f}" stroke="#7f7f7f"/>',
        )
    )
    return "".join(parts)


def _chart_legend_position(chart: Any) -> str:
    try:
        value = str(chart.legend.position.name).casefold()
    except (AttributeError, TypeError, ValueError):
        return "right"
    return value if value in {"top", "bottom", "left", "right", "corner"} else "right"


def _chart_plot_horizontal_bounds(chart: Any) -> tuple[float, float]:
    if not _chart_has_legend(chart):
        return 62.0, 574.0
    position = _chart_legend_position(chart)
    if position in {"right", "corner"}:
        return 62.0, 430.0
    if position == "left":
        return 170.0, 574.0
    return 62.0, 574.0


def _chart_marker_style(series: Any, *, default: str = "CIRCLE") -> str:
    try:
        style = series.marker.style
        if style is not None:
            return str(style.name).upper()
    except (AttributeError, TypeError, ValueError):
        pass
    return default


def _chart_marker_svg(
    series: Any,
    *,
    x: float,
    y: float,
    color: str,
    x_compensation: float = 1.0,
    default_style: str = "CIRCLE",
) -> str:
    try:
        size = float(series.marker.size or 0)
    except (AttributeError, TypeError, ValueError):
        size = 0
    radius = min(max(size / 2 if size else 4.5, 3.0), 12.0)
    radius_x = radius * max(x_compensation, 0.25)
    style = _chart_marker_style(series, default=default_style)
    common = f'class="chart-marker" fill="{color}"'
    if style in {"DIAMOND", "AUTO"}:
        return (
            f'<polygon {common} data-marker-style="diamond" points="'
            f"{x:.2f},{y - radius:.2f} {x + radius_x:.2f},{y:.2f} "
            f'{x:.2f},{y + radius:.2f} {x - radius_x:.2f},{y:.2f}"/>'
        )
    if style in {"SQUARE", "DASH"}:
        return (
            f'<rect {common} data-marker-style="square" '
            f'x="{x - radius_x:.2f}" y="{y - radius:.2f}" '
            f'width="{2 * radius_x:.2f}" height="{2 * radius:.2f}"/>'
        )
    if style in {"TRIANGLE"}:
        return (
            f'<polygon {common} data-marker-style="triangle" points="'
            f"{x:.2f},{y - radius:.2f} {x + radius_x:.2f},{y + radius:.2f} "
            f'{x - radius_x:.2f},{y + radius:.2f}"/>'
        )
    return (
        f'<ellipse {common} data-marker-style="circle" cx="{x:.2f}" cy="{y:.2f}" '
        f'rx="{radius_x:.2f}" ry="{radius:.2f}"/>'
    )


def _chart_legend(
    chart: Any,
    names: list[str],
    colors: list[str],
    *,
    marker_styles: list[str] | None = None,
) -> str:
    position = _chart_legend_position(chart)
    items: list[str] = []
    if position in {"top", "bottom"}:
        y = 5 if position == "top" else 283
        item_widths = [27 + min(len(name), 28) * 7 for name in names]
        total_width = sum(item_widths) + max(len(names) - 1, 0) * 16
        x = max((600 - total_width) / 2, 8)
        for index, (name, item_width) in enumerate(
            zip(names, item_widths, strict=True)
        ):
            items.append(
                f'<rect x="{x:.2f}" y="{y}" width="11" height="11" '
                f'fill="{colors[index % len(colors)]}"/>'
                f'<text x="{x + 16:.2f}" y="{y + 10}" font-size="11" '
                f'fill="#595959">{html.escape(name)}</text>'
            )
            x += item_width + 16
    else:
        x = 7 if position == "left" else 455
        start_y = 150 - max(len(names), 1) * 25 / 2
        if position == "corner":
            start_y = 35
        for index, name in enumerate(names):
            y = start_y + index * 25
            marker_style = marker_styles[index] if marker_styles else "SQUARE"
            if marker_style == "DIAMOND":
                symbol = (
                    f'<polygon class="chart-legend-marker" '
                    f'data-marker-style="diamond" fill="{colors[index % len(colors)]}" '
                    f'points="{x + 5.5:.2f},{y:.2f} {x + 11:.2f},{y + 5.5:.2f} '
                    f'{x + 5.5:.2f},{y + 11:.2f} {x:.2f},{y + 5.5:.2f}"/>'
                )
            else:
                symbol = (
                    f'<rect class="chart-legend-marker" x="{x}" y="{y}" '
                    f'width="11" height="11" fill="{colors[index % len(colors)]}"/>'
                )
            items.append(
                f'{symbol}<text x="{x + 16}" y="{y + 10}" font-size="11" '
                f'fill="#595959">{html.escape(name)}</text>'
            )
    return (
        f'<g class="chart-legend chart-legend-{position}" '
        f'data-legend-position="{position}">{"".join(items)}</g>'
    )


def _chart_categories(chart: Any) -> list[str]:
    try:
        return [str(category) for category in chart.plots[0].categories]
    except (AttributeError, IndexError, TypeError, ValueError):
        return []


def _chart_series_name(series: Any, index: int) -> str:
    fallback = f"Series {index + 1}"
    try:
        name = str(series.name or "").strip()
    except (AttributeError, TypeError, ValueError):
        name = ""
    if name:
        return name

    # python-pptx currently resolves c:tx/c:strRef but returns an empty name for
    # the equally valid literal form emitted by some OOXML producers:
    # <c:tx><c:v>Series name</c:v></c:tx>.
    try:
        value_nodes = series._ser.xpath("./*[local-name()='tx']//*[local-name()='v']")
    except (AttributeError, TypeError, ValueError):
        return fallback
    for value_node in value_nodes:
        name = str(value_node.text or "").strip()
        if name:
            return name
    return fallback


def _chart_series_records(
    chart: Any, theme_palette: dict[str, str]
) -> list[tuple[str, list[float], str, Any]]:
    records: list[tuple[str, list[float], str, Any]] = []
    for index, series in enumerate(chart.series):
        values: list[float] = []
        for raw_value in getattr(series, "values", ()):
            try:
                values.append(float(raw_value or 0))
            except (TypeError, ValueError):
                values.append(0.0)
        if not values:
            values = _chart_cache_values(series, "yVal")
        if not values:
            continue
        records.append(
            (
                _chart_series_name(series, index),
                values,
                _chart_series_color(series, index, theme_palette),
                series,
            )
        )
    return records


def _chart_pie_svg(
    chart: Any,
    chart_type: str,
    records: list[tuple[str, list[float], str, Any]],
    categories: list[str],
) -> str:
    values = [abs(value) for value in records[0][1]]
    total = sum(values) or 1
    data_labels = _chart_data_labels_enabled(chart)
    has_legend = _chart_has_legend(chart)
    legend_position = _chart_legend_position(chart) if has_legend else "right"
    if legend_position == "left":
        center_x, center_y, radius = 405, 154, 103
    elif legend_position in {"top", "bottom"}:
        center_x = 300
        center_y = 166 if legend_position == "top" else 132
        radius = 92
    else:
        center_x, center_y, radius = 190, 154, 108
    parts: list[str] = []
    start_angle = -math.pi / 2
    colors = [
        _CHART_FALLBACK_PALETTE[index % len(_CHART_FALLBACK_PALETTE)]
        for index in range(len(values))
    ]
    for index, value in enumerate(values):
        end_angle = start_angle + 2 * math.pi * value / total
        if math.isclose(value, total):
            parts.append(
                f'<circle class="chart-slice" cx="{center_x}" cy="{center_y}" '
                f'r="{radius}" '
                f'fill="{colors[index]}"/>'
            )
        else:
            start_x = center_x + radius * math.cos(start_angle)
            start_y = center_y + radius * math.sin(start_angle)
            end_x = center_x + radius * math.cos(end_angle)
            end_y = center_y + radius * math.sin(end_angle)
            large_arc = 1 if end_angle - start_angle > math.pi else 0
            parts.append(
                f'<path class="chart-slice" d="M{center_x} {center_y} '
                f"L{start_x:.2f} {start_y:.2f} A{radius} {radius} 0 "
                f"{large_arc} 1 {end_x:.2f} "
                f'{end_y:.2f} Z" fill="{colors[index]}"/>'
            )
        if data_labels and value:
            middle = (start_angle + end_angle) / 2
            label_x = center_x + radius * 0.68 * math.cos(middle)
            label_y = center_y + radius * 0.68 * math.sin(middle)
            parts.append(
                f'<text class="chart-data-label" x="{label_x:.2f}" '
                f'y="{label_y:.2f}" text-anchor="middle" font-size="12" '
                f'fill="#fff">{value / total:.0%}</text>'
            )
        start_angle = end_angle
    if "DOUGHNUT" in chart_type:
        parts.append(
            f'<circle cx="{center_x}" cy="{center_y}" '
            f'r="{radius * 0.57:.2f}" fill="#fff"/>'
        )
    if has_legend:
        legend_items: list[str] = []
        for index, value in enumerate(values):
            label = categories[index] if index < len(categories) else ""
            suffix = f" ({value / total:.0%})" if data_labels else ""
            if legend_position in {"top", "bottom"}:
                item_x = 45 + index * (510 / max(len(values), 1))
                item_y = 5 if legend_position == "top" else 283
            else:
                item_x = 8 if legend_position == "left" else 355
                item_y = 54 + index * 31
            legend_items.append(
                f'<rect x="{item_x:.2f}" y="{item_y:.2f}" width="12" height="12" '
                f'fill="{colors[index]}"/>'
                f'<text x="{item_x + 21:.2f}" y="{item_y + 11:.2f}" font-size="12" '
                f'fill="#595959">{html.escape(label)}{suffix}</text>'
            )
        parts.append(
            f'<g class="chart-legend chart-legend-{legend_position}" '
            f'data-legend-position="{legend_position}">{"".join(legend_items)}</g>'
        )
    semantic_type = "doughnut" if "DOUGHNUT" in chart_type else "pie"
    return (
        f'<svg class="chart-svg chart-{semantic_type}" viewBox="0 0 600 300">'
        f"{''.join(parts)}</svg>"
    )


def _chart_category_labels(
    categories: list[str],
    *,
    bar: bool = False,
    plot_left: float = 62,
    plot_right: float = 574,
    plot_top: float = 28,
    plot_bottom: float = 248,
) -> str:
    if bar:
        slot = (plot_bottom - plot_top) / max(len(categories), 1)
        return "".join(
            f'<text class="chart-category-label" x="{plot_left - 11:.2f}" '
            f'y="{plot_top + (index + 0.5) * slot + 4:.2f}" '
            'text-anchor="end" font-size="12" fill="#666">'
            f"{html.escape(label)}</text>"
            for index, label in enumerate(categories)
        )
    slot = (plot_right - plot_left) / max(len(categories), 1)
    return "".join(
        f'<text class="chart-category-label" '
        f'x="{plot_left + (index + 0.5) * slot:.2f}" '
        f'y="{plot_bottom + 24:.2f}" text-anchor="middle" '
        'font-size="12" fill="#666">'
        f"{html.escape(label)}</text>"
        for index, label in enumerate(categories)
    )


def _chart_line_or_area_svg(
    chart: Any,
    chart_type: str,
    records: list[tuple[str, list[float], str, Any]],
    categories: list[str],
) -> str:
    stacked = "STACKED" in chart_type
    hundred_percent = "100" in chart_type
    count = max(len(record[1]) for record in records)
    display_values: list[list[float]] = []
    cumulative = [0.0] * count
    totals = [
        sum(abs(record[1][index]) for record in records if index < len(record[1]))
        for index in range(count)
    ]
    for _, values, _, _ in records:
        current: list[float] = []
        for index in range(count):
            value = values[index] if index < len(values) else 0.0
            if hundred_percent:
                value = value / (totals[index] or 1)
            if stacked:
                cumulative[index] += value
                value = cumulative[index]
            current.append(value)
        display_values.append(current)
    flat_values = [value for values in display_values for value in values]
    lower, upper = _chart_scale(chart, flat_values)
    plot_left, plot_right = _chart_plot_horizontal_bounds(chart)
    plot_top, plot_bottom = 28.0, 248.0

    def map_y(value: float) -> float:
        return plot_bottom - (value - lower) / (upper - lower) * (
            plot_bottom - plot_top
        )

    baseline = map_y(0)
    parts = [
        _chart_axes(
            chart,
            lower=lower,
            upper=upper,
            baseline=baseline,
            gridlines=_chart_major_gridlines(chart),
            plot_left=plot_left,
            plot_right=plot_right,
            plot_top=plot_top,
            plot_bottom=plot_bottom,
        )
    ]
    data_labels = _chart_data_labels_enabled(chart)
    show_markers = "MARKERS" in chart_type and "NO_MARKERS" not in chart_type
    for series_index, (_, _, color, _) in enumerate(records):
        values = display_values[series_index]
        points = [
            (
                plot_left + index * (plot_right - plot_left) / max(len(values) - 1, 1),
                map_y(value),
            )
            for index, value in enumerate(values)
        ]
        point_text = " ".join(f"{x:.2f},{y:.2f}" for x, y in points)
        if "AREA" in chart_type and points:
            previous = (
                display_values[series_index - 1]
                if stacked and series_index > 0
                else [0.0] * len(points)
            )
            polygon_points = [*[f"{x:.2f},{y:.2f}" for x, y in points]]
            polygon_points.extend(
                f"{points[index][0]:.2f},{map_y(previous[index]):.2f}"
                for index in range(len(points) - 1, -1, -1)
            )
            parts.append(
                f'<polygon class="chart-area" points="{" ".join(polygon_points)}" '
                f'fill="{color}" fill-opacity=".24"/>'
            )
        parts.append(
            f'<polyline class="chart-series-line" points="{point_text}" fill="none" '
            f'stroke="{color}" stroke-width="3"/>'
        )
        if show_markers:
            parts.extend(
                f'<circle class="chart-marker" cx="{x:.2f}" cy="{y:.2f}" '
                f'r="4" fill="{color}"/>'
                for x, y in points
            )
        if data_labels:
            parts.extend(
                f'<text class="chart-data-label" x="{x:.2f}" y="{y - 7:.2f}" '
                f'text-anchor="middle" font-size="11" fill="#404040">'
                f"{values[index]:g}</text>"
                for index, (x, y) in enumerate(points)
            )
    parts.append(
        _chart_category_labels(
            categories,
            plot_left=plot_left,
            plot_right=plot_right,
            plot_top=plot_top,
            plot_bottom=plot_bottom,
        )
    )
    if _chart_has_legend(chart):
        parts.append(
            _chart_legend(
                chart,
                [record[0] for record in records],
                [record[2] for record in records],
            )
        )
    semantic = "area" if "AREA" in chart_type else "line"
    stacked_class = " chart-stacked" if stacked else ""
    return (
        f'<svg class="chart-svg chart-{semantic}{stacked_class}" '
        'viewBox="0 0 600 300" preserveAspectRatio="none">'
        f"{''.join(parts)}</svg>"
    )


def _chart_scatter_svg(
    chart: Any,
    chart_type: str,
    records: list[tuple[str, list[float], str, Any]],
    *,
    chart_aspect: float,
) -> str:
    series_points: list[list[tuple[float, float]]] = []
    for _, y_values, _, series in records:
        x_values = _chart_cache_values(series, "xVal")
        if not x_values:
            x_values = [float(index + 1) for index in range(len(y_values))]
        series_points.append(list(zip(x_values, y_values, strict=False)))
    all_x = [x for points in series_points for x, _ in points] or [0.0, 1.0]
    all_y = [y for points in series_points for _, y in points] or [0.0, 1.0]
    try:
        x_axis = chart.category_axis
    except (AttributeError, TypeError, ValueError):
        x_axis = None
    try:
        y_axis = chart.value_axis
    except (AttributeError, TypeError, ValueError):
        y_axis = None
    x_lower, x_upper = _chart_scale(chart, all_x, axis=x_axis)
    y_lower, y_upper = _chart_scale(chart, all_y, axis=y_axis)
    plot_left, plot_right = _chart_plot_horizontal_bounds(chart)
    plot_top, plot_bottom = 28.0, 248.0

    def map_x(value: float) -> float:
        return plot_left + (value - x_lower) / (x_upper - x_lower) * (
            plot_right - plot_left
        )

    def map_y(value: float) -> float:
        return plot_bottom - (value - y_lower) / (y_upper - y_lower) * (
            plot_bottom - plot_top
        )

    baseline = map_y(0)
    parts = [
        _chart_axes(
            chart,
            lower=y_lower,
            upper=y_upper,
            baseline=baseline,
            gridlines=_chart_axis_major_gridlines(y_axis),
            axis=y_axis,
            plot_left=plot_left,
            plot_right=plot_right,
            plot_top=plot_top,
            plot_bottom=plot_bottom,
        )
    ]
    for tick in _chart_tick_values(chart, x_lower, x_upper, axis=x_axis):
        x = map_x(tick)
        if _chart_axis_major_gridlines(x_axis):
            parts.append(
                f'<line class="chart-gridline chart-gridline-x" '
                f'x1="{x:.2f}" y1="{plot_top:.2f}" x2="{x:.2f}" '
                f'y2="{plot_bottom:.2f}" stroke="#d9d9d9" stroke-width="1"/>'
            )
        parts.append(
            f'<line class="chart-axis-tick chart-axis-tick-x" '
            f'x1="{x:.2f}" y1="{plot_bottom:.2f}" x2="{x:.2f}" '
            f'y2="{plot_bottom + 4:.2f}" stroke="#7f7f7f"/>'
            f'<text class="chart-axis-label chart-axis-label-x" '
            f'x="{x:.2f}" y="{plot_bottom + 20:.2f}" text-anchor="middle" '
            f'font-size="11" fill="#666">{tick:g}</text>'
        )
    has_lines = "LINES" in chart_type or "SMOOTH" in chart_type
    has_markers = "NO_MARKERS" not in chart_type
    data_labels = _chart_data_labels_enabled(chart)
    x_compensation = min(max(2.0 / max(chart_aspect, 0.1), 0.5), 4.0)
    for series_index, points in enumerate(series_points):
        color = records[series_index][2]
        series = records[series_index][3]
        mapped = [(map_x(x), map_y(y), x, y) for x, y in points]
        if has_lines:
            coordinates = " ".join(f"{x:.2f},{y:.2f}" for x, y, _, _ in mapped)
            parts.append(
                f'<polyline class="chart-series-line" points="{coordinates}" '
                f'fill="none" stroke="{color}" stroke-width="3"/>'
            )
        if has_markers:
            parts.extend(
                _chart_marker_svg(
                    series,
                    x=x,
                    y=y,
                    color=color,
                    x_compensation=x_compensation,
                    default_style="DIAMOND",
                )
                for x, y, _, _ in mapped
            )
        if data_labels:
            parts.extend(
                f'<text class="chart-data-label" x="{x:.2f}" '
                f'y="{max(plot_top + 10, min(y - 7, plot_bottom - 4)):.2f}" '
                f'text-anchor="middle" font-size="11" fill="#404040">'
                f"{raw_y:g}</text>"
                for x, y, _, raw_y in mapped
            )
    if _chart_has_legend(chart):
        parts.append(
            _chart_legend(
                chart,
                [record[0] for record in records],
                [record[2] for record in records],
                marker_styles=[
                    _chart_marker_style(record[3], default="DIAMOND")
                    for record in records
                ],
            )
        )
    return (
        '<svg class="chart-svg chart-scatter" viewBox="0 0 600 300" '
        f'preserveAspectRatio="none">{"".join(parts)}</svg>'
    )


def _chart_bar_or_column_svg(
    chart: Any,
    chart_type: str,
    records: list[tuple[str, list[float], str, Any]],
    categories: list[str],
) -> str:
    horizontal = chart_type.startswith("BAR_")
    stacked = "STACKED" in chart_type
    hundred_percent = "100" in chart_type
    category_count = max(len(record[1]) for record in records)
    values = [record[1] for record in records]
    totals = [
        sum(abs(series[index]) for series in values if index < len(series))
        for index in range(category_count)
    ]
    if stacked:
        positive = [0.0] * category_count
        negative = [0.0] * category_count
        for series in values:
            for index, raw_value in enumerate(series):
                value = (
                    raw_value / (totals[index] or 1) if hundred_percent else raw_value
                )
                if value >= 0:
                    positive[index] += value
                else:
                    negative[index] += value
        scale_values = [*positive, *negative]
    else:
        scale_values = [value for series in values for value in series]
    lower, upper = _chart_scale(chart, scale_values)
    data_labels = _chart_data_labels_enabled(chart)
    parts: list[str] = []
    plot_top, plot_bottom = 28.0, 248.0

    if horizontal:
        plot_left = 150.0
        plot_right = (
            430.0
            if _chart_has_legend(chart)
            and _chart_legend_position(chart) in {"right", "corner"}
            else 564.0
        )

        def map_x(value: float) -> float:
            return plot_left + (value - lower) / (upper - lower) * (
                plot_right - plot_left
            )

        baseline = map_x(0)
        for tick in _chart_tick_values(chart, lower, upper):
            x = map_x(tick)
            if _chart_major_gridlines(chart):
                parts.append(
                    f'<line class="chart-gridline" x1="{x:.2f}" '
                    f'y1="{plot_top:.2f}" x2="{x:.2f}" '
                    f'y2="{plot_bottom:.2f}" stroke="#d9d9d9"/>'
                )
            parts.append(
                f'<text class="chart-axis-label chart-axis-label-x" '
                f'x="{x:.2f}" y="{plot_bottom + 20:.2f}" '
                f'text-anchor="middle" font-size="11" fill="#666">{tick:g}</text>'
            )
        parts.append(
            f'<line class="chart-axis chart-axis-y" x1="{baseline:.2f}" '
            f'y1="{plot_top:.2f}" x2="{baseline:.2f}" '
            f'y2="{plot_bottom:.2f}" stroke="#7f7f7f"/>'
        )
        category_slot = (plot_bottom - plot_top) / max(category_count, 1)
        clustered_height = category_slot * 0.72 / max(len(records), 1)
        stacked_height = category_slot * 0.72
        positive = [0.0] * category_count
        negative = [0.0] * category_count
        for series_index, (_, series, color, _) in enumerate(records):
            for category_index in range(category_count):
                raw_value = (
                    series[category_index] if category_index < len(series) else 0.0
                )
                value = (
                    raw_value / (totals[category_index] or 1)
                    if hundred_percent
                    else raw_value
                )
                if stacked:
                    start = (
                        positive[category_index]
                        if value >= 0
                        else negative[category_index]
                    )
                    end = start + value
                    if value >= 0:
                        positive[category_index] = end
                    else:
                        negative[category_index] = end
                    height = stacked_height
                    y = plot_top + category_index * category_slot + category_slot * 0.14
                else:
                    start, end = 0.0, value
                    height = clustered_height
                    y = (
                        plot_top
                        + category_index * category_slot
                        + category_slot * 0.14
                        + series_index * height
                    )
                x1, x2 = map_x(start), map_x(end)
                visible_x1 = min(max(x1, plot_left), plot_right)
                visible_x2 = min(max(x2, plot_left), plot_right)
                parts.append(
                    f'<rect class="chart-bar" '
                    f'x="{min(visible_x1, visible_x2):.2f}" y="{y:.2f}" '
                    f'width="{abs(visible_x2 - visible_x1):.2f}" '
                    f'height="{height:.2f}" '
                    f'fill="{color}"/>'
                )
                if data_labels:
                    if stacked:
                        label_x = (visible_x1 + visible_x2) / 2
                        anchor = "middle"
                    else:
                        label_x = min(
                            max(x2 + (4 if value >= 0 else -4), plot_left + 4),
                            plot_right - 4,
                        )
                        anchor = "start" if value >= 0 else "end"
                    parts.append(
                        f'<text class="chart-data-label" x="{label_x:.2f}" '
                        f'y="{y + height / 2 + 4:.2f}" text-anchor="{anchor}" '
                        f'font-size="11" fill="#404040">{raw_value:g}</text>'
                    )
        parts.append(
            _chart_category_labels(
                categories,
                bar=True,
                plot_left=plot_left,
                plot_right=plot_right,
                plot_top=plot_top,
                plot_bottom=plot_bottom,
            )
        )
    else:
        plot_left, plot_right = _chart_plot_horizontal_bounds(chart)

        def map_y(value: float) -> float:
            return plot_bottom - (value - lower) / (upper - lower) * (
                plot_bottom - plot_top
            )

        baseline = map_y(0)
        parts.append(
            _chart_axes(
                chart,
                lower=lower,
                upper=upper,
                baseline=baseline,
                gridlines=_chart_major_gridlines(chart),
                plot_left=plot_left,
                plot_right=plot_right,
                plot_top=plot_top,
                plot_bottom=plot_bottom,
            )
        )
        category_slot = (plot_right - plot_left) / max(category_count, 1)
        clustered_width = category_slot * 0.72 / max(len(records), 1)
        stacked_width = category_slot * 0.72
        positive = [0.0] * category_count
        negative = [0.0] * category_count
        for series_index, (_, series, color, _) in enumerate(records):
            for category_index in range(category_count):
                raw_value = (
                    series[category_index] if category_index < len(series) else 0.0
                )
                value = (
                    raw_value / (totals[category_index] or 1)
                    if hundred_percent
                    else raw_value
                )
                if stacked:
                    start = (
                        positive[category_index]
                        if value >= 0
                        else negative[category_index]
                    )
                    end = start + value
                    if value >= 0:
                        positive[category_index] = end
                    else:
                        negative[category_index] = end
                    width = stacked_width
                    x = (
                        plot_left
                        + category_index * category_slot
                        + category_slot * 0.14
                    )
                else:
                    start, end = 0.0, value
                    width = clustered_width
                    x = (
                        plot_left
                        + category_index * category_slot
                        + category_slot * 0.14
                        + series_index * width
                    )
                y1, y2 = map_y(start), map_y(end)
                visible_y1 = min(max(y1, plot_top), plot_bottom)
                visible_y2 = min(max(y2, plot_top), plot_bottom)
                parts.append(
                    f'<rect class="chart-column" x="{x:.2f}" '
                    f'y="{min(visible_y1, visible_y2):.2f}" width="{width:.2f}" '
                    f'height="{abs(visible_y2 - visible_y1):.2f}" fill="{color}"/>'
                )
                if data_labels:
                    if stacked:
                        label_y = (visible_y1 + visible_y2) / 2 + 4
                    else:
                        label_y = min(
                            max(
                                y2 + (-4 if value >= 0 else 12),
                                plot_top + 10,
                            ),
                            plot_bottom - 4,
                        )
                    parts.append(
                        f'<text class="chart-data-label" x="{x + width / 2:.2f}" '
                        f'y="{label_y:.2f}" text-anchor="middle" font-size="11" '
                        f'fill="#404040">{raw_value:g}</text>'
                    )
        parts.append(
            _chart_category_labels(
                categories,
                plot_left=plot_left,
                plot_right=plot_right,
                plot_top=plot_top,
                plot_bottom=plot_bottom,
            )
        )
    if _chart_has_legend(chart):
        parts.append(
            _chart_legend(
                chart,
                [record[0] for record in records],
                [record[2] for record in records],
            )
        )
    semantic = "bar" if horizontal else "column"
    grouping = "stacked" if stacked else "clustered"
    return (
        f'<svg class="chart-svg chart-{semantic} chart-{grouping}" '
        f'data-chart-grouping="{grouping}" viewBox="0 0 600 300" '
        f'preserveAspectRatio="none">{"".join(parts)}</svg>'
    )


def _chart_svg(
    chart: Any,
    theme_palette: dict[str, str],
    *,
    chart_aspect: float = 2.0,
) -> str:
    chart_type = _chart_type_name(chart)
    semantic_type = chart_type.casefold().replace("_", "-")
    unsupported = (
        chart_type == "UNKNOWN"
        or "_3D" in chart_type
        or not chart_type.startswith(
            ("COLUMN_", "BAR_", "LINE", "AREA", "PIE", "DOUGHNUT", "XY_SCATTER")
        )
    )
    if unsupported:
        safe_type = html.escape(semantic_type, quote=True)
        return (
            f'<div class="chart-unsupported" data-chart-type="{safe_type}">'
            f"Unsupported chart type: {html.escape(chart_type)}</div>"
        )
    records = _chart_series_records(chart, theme_palette)
    if not records:
        return '<div class="chart-empty">Chart data unavailable</div>'
    categories = _chart_categories(chart)
    if chart_type.startswith(("PIE", "DOUGHNUT")):
        return _chart_pie_svg(chart, chart_type, records, categories)
    if chart_type.startswith("XY_SCATTER"):
        return _chart_scatter_svg(
            chart,
            chart_type,
            records,
            chart_aspect=chart_aspect,
        )
    if chart_type.startswith(("LINE", "AREA")):
        return _chart_line_or_area_svg(chart, chart_type, records, categories)
    return _chart_bar_or_column_svg(chart, chart_type, records, categories)


def _placeholder_key(shape: Any) -> int | None:
    if not getattr(shape, "is_placeholder", False):
        return None
    try:
        return int(shape.placeholder_format.idx)
    except (AttributeError, TypeError, ValueError):
        return None


_MEDIUM_STYLE_2_ACCENT_1 = "5C22544A-7EE6-4342-B048-85BDC9FD1C3A"


def _table_style_context(table: Any, palette: dict[str, str]) -> dict[str, Any]:
    properties = getattr(table._tbl, "tblPr", None)
    style_id = ""
    if properties is not None:
        style_node = _xml_child(properties, "tableStyleId")
        if style_node is not None and style_node.text:
            style_id = str(style_node.text).strip().strip("{}").upper()
    return {
        "style_id": style_id,
        "medium_style_2_accent_1": style_id == _MEDIUM_STYLE_2_ACCENT_1,
        "accent": palette.get("accent1", "#4f81bd"),
        "first_row": bool(_xml_bool(properties.get("firstRow")))
        if properties is not None
        else False,
        "last_row": bool(_xml_bool(properties.get("lastRow")))
        if properties is not None
        else False,
        "first_column": bool(_xml_bool(properties.get("firstCol")))
        if properties is not None
        else False,
        "last_column": bool(_xml_bool(properties.get("lastCol")))
        if properties is not None
        else False,
        "band_rows": bool(_xml_bool(properties.get("bandRow")))
        if properties is not None
        else False,
        "band_columns": bool(_xml_bool(properties.get("bandCol")))
        if properties is not None
        else False,
    }


def _table_style_cell_defaults(
    context: dict[str, Any],
    *,
    row_index: int,
    column_index: int,
    row_count: int,
    column_count: int,
    slide_width: int,
) -> tuple[list[str], list[str]]:
    styles: list[str] = []
    classes = ["table-cell"]
    if not context["medium_style_2_accent_1"]:
        return styles, classes

    accent = str(context["accent"])
    border_width = _text_length_cqw(12700, slide_width)
    styles.append(f"border:{border_width} solid #ffffff")
    first_row = bool(context["first_row"] and row_index == 0)
    last_row = bool(context["last_row"] and row_index == row_count - 1)
    if first_row:
        classes.append("table-first-row")
        styles.extend((f"background:{accent}", "color:#ffffff", "font-weight:700"))
    elif last_row:
        classes.append("table-last-row")
        styles.extend((f"background:{accent}", "color:#ffffff", "font-weight:700"))
    elif context["band_rows"]:
        body_row = row_index - (1 if context["first_row"] else 0)
        band_brightness = 0.78 if body_row % 2 == 0 else 0.9
        classes.append(
            "table-band-row-primary"
            if body_row % 2 == 0
            else "table-band-row-secondary"
        )
        styles.append(f"background:{_apply_brightness(accent, band_brightness)}")

    if context["band_columns"]:
        body_column = column_index - (1 if context["first_column"] else 0)
        if body_column >= 0 and body_column % 2 == 0:
            classes.append("table-band-column")
            styles.append(f"background:{_apply_brightness(accent, 0.82)}")
    if context["first_column"] and column_index == 0:
        classes.append("table-first-column")
        styles.append("font-weight:700")
    if context["last_column"] and column_index == column_count - 1:
        classes.append("table-last-column")
        styles.append("font-weight:700")
    return styles, classes


def _table_cell_fill(cell: Any, palette: dict[str, str]) -> tuple[bool, str | None]:
    properties = getattr(cell._tc, "tcPr", None)
    if properties is None:
        return False, None
    color = _xml_color_css(properties, palette)
    if color:
        return True, color
    if _xml_child(properties, "noFill") is not None:
        return True, None
    return False, None


def _table_cell_border_styles(
    cell: Any, palette: dict[str, str], *, slide_width: int
) -> list[str]:
    properties = getattr(cell._tc, "tcPr", None)
    if properties is None:
        return []
    result: list[str] = []
    dash_styles = {
        "dash": "dashed",
        "dashDot": "dashed",
        "dot": "dotted",
        "lgDash": "dashed",
        "lgDashDot": "dashed",
        "lgDashDotDot": "dashed",
        "sysDash": "dashed",
        "sysDashDot": "dashed",
        "sysDashDotDot": "dashed",
        "sysDot": "dotted",
    }
    for xml_name, css_name in (
        ("lnL", "left"),
        ("lnR", "right"),
        ("lnT", "top"),
        ("lnB", "bottom"),
    ):
        line = _xml_child(properties, xml_name)
        if line is None:
            continue
        if _xml_child(line, "noFill") is not None:
            result.append(f"border-{css_name}:none")
            continue
        color = _xml_color_css(line, palette) or "#808080"
        try:
            width = max(int(line.get("w", "12700")), 1)
        except ValueError:
            width = 12700
        dash = _xml_child(line, "prstDash")
        dash_name = str(dash.get("val") or "solid") if dash is not None else "solid"
        line_style = dash_styles.get(dash_name, "solid")
        result.append(
            f"border-{css_name}:{_text_length_cqw(width, slide_width)} "
            f"{line_style} {color}"
        )
    return result


def _table_cell_paragraph_html(
    paragraph: Any,
    *,
    palette: dict[str, str],
    fonts: dict[str, str],
    slide_width: int,
    list_style: Any | None,
    counters: dict[int, int],
) -> str:
    paragraph_properties = getattr(paragraph._p, "pPr", None)
    paragraph_nodes = [paragraph_properties] if paragraph_properties is not None else []
    level = int(getattr(paragraph, "level", 0) or 0)
    local_properties = _level_properties(list_style, level)
    if local_properties is not None:
        paragraph_nodes.append(local_properties)
    runs, _ = _paragraph_runs_html(
        paragraph,
        paragraph_nodes,
        palette=palette,
        fonts=fonts,
        slide_width=slide_width,
        default_size_pt=18.0,
    )
    styles = ["margin:0", "white-space:pre-wrap", "line-height:1.12"]
    alignment = _first_attribute(paragraph_nodes, "algn")
    if alignment:
        styles.append(
            "text-align:"
            + {
                "ctr": "center",
                "r": "right",
                "just": "justify",
                "dist": "justify",
            }.get(alignment, "left")
        )
    before = _spacing_css(paragraph_nodes, "spcBef", slide_width=slide_width)
    after = _spacing_css(paragraph_nodes, "spcAft", slide_width=slide_width)
    line = _spacing_css(paragraph_nodes, "lnSpc", slide_width=slide_width)
    if before:
        styles.append(f"margin-top:{before}")
    if after:
        styles.append(f"margin-bottom:{after}")
    if line:
        styles.append(f"line-height:{line}")

    marker = _effective_bullet(paragraph_nodes, counters, level)
    marker_html = ""
    if marker:
        try:
            margin = int(_first_attribute(paragraph_nodes, "marL", "0") or 0)
            indent = int(_first_attribute(paragraph_nodes, "indent", "0") or 0)
        except ValueError:
            margin = indent = 0
        marker_position = max(margin + indent, 0)
        marker_styles, marker_size = _bullet_marker_styles(
            paragraph,
            paragraph_nodes,
            palette=palette,
            fonts=fonts,
            slide_width=slide_width,
        )
        marker_width = max(
            margin - marker_position,
            round(marker_size * 12700 * 0.8),
            1,
        )
        styles.extend(
            (
                "display:grid",
                f"padding-left:{_text_length_cqw(marker_position, slide_width)}",
                f"grid-template-columns:{_text_length_cqw(marker_width, slide_width)} 1fr",
                "column-gap:0",
            )
        )
        marker_html = (
            f'<span class="bullet-marker" style="{";".join(marker_styles)}">'
            f"{html.escape(marker)}</span>"
        )
    return (
        f'<p style="{";".join(styles)}">{marker_html}'
        f'<span class="paragraph-content">{runs}</span></p>'
    )


def _table_cell_html(
    cell: Any,
    palette: dict[str, str],
    *,
    fonts: dict[str, str],
    slide_width: int,
    colspan: int = 1,
    rowspan: int = 1,
    row_index: int,
    column_index: int,
    row_count: int,
    column_count: int,
    table_style: dict[str, Any],
) -> str:
    styles, classes = _table_style_cell_defaults(
        table_style,
        row_index=row_index,
        column_index=column_index,
        row_count=row_count,
        column_count=column_count,
        slide_width=slide_width,
    )
    has_fill, fill = _table_cell_fill(cell, palette)
    if has_fill:
        styles.append(f"background:{fill or 'transparent'}")
    styles.extend(_table_cell_border_styles(cell, palette, slide_width=slide_width))
    try:
        vertical_name = str(cell.vertical_anchor.name)
    except (AttributeError, TypeError, ValueError):
        vertical_name = "TOP"
    styles.append(
        "vertical-align:"
        + {"MIDDLE": "middle", "BOTTOM": "bottom"}.get(vertical_name, "top")
    )
    margin_defaults = {
        "margin_top": 45720,
        "margin_right": 91440,
        "margin_bottom": 45720,
        "margin_left": 91440,
    }
    margins: list[str] = []
    for name, default in margin_defaults.items():
        try:
            raw_value = getattr(cell, name)
            value = default if raw_value is None else int(raw_value)
        except (AttributeError, TypeError, ValueError):
            value = default
        margins.append(_text_length_cqw(value, slide_width))
    styles.append(f"padding:{' '.join(margins)}")
    list_style = getattr(cell.text_frame._txBody, "lstStyle", None)
    counters: dict[int, int] = {}
    paragraphs = "".join(
        _table_cell_paragraph_html(
            paragraph,
            palette=palette,
            fonts=fonts,
            slide_width=slide_width,
            list_style=list_style,
            counters=counters,
        )
        for paragraph in cell.text_frame.paragraphs
    )
    spans = ""
    if colspan > 1:
        spans += f' colspan="{colspan}"'
    if rowspan > 1:
        spans += f' rowspan="{rowspan}"'
    spans += f' data-row-index="{row_index}" data-column-index="{column_index}"'
    class_attr = f' class="{" ".join(classes)}"'
    style_attr = f' style="{";".join(styles)}"' if styles else ""
    return f"<td{class_attr}{spans}{style_attr}>{paragraphs}</td>"


def _render_table(
    shape: Any,
    *,
    layer: str,
    palette: dict[str, str],
    fonts: dict[str, str],
    slide_width: int,
    slide_height: int,
    coordinate_left: int = 0,
    coordinate_top: int = 0,
    coordinate_width: int | None = None,
    coordinate_height: int | None = None,
) -> str:
    table = shape.table
    table_style = _table_style_context(table, palette)
    row_count = len(table.rows)
    column_count = len(table.columns)
    total_width = sum(int(column.width) for column in table.columns) or 1
    total_height = sum(int(row.height) for row in table.rows) or 1
    columns = "".join(
        f'<col style="width:{100 * int(column.width) / total_width:.4f}%">'
        for column in table.columns
    )
    rows: list[str] = []
    for row_index, row in enumerate(table.rows):
        cells: list[str] = []
        for column_index, cell in enumerate(row.cells):
            if cell.is_spanned:
                continue
            colspan = int(cell.span_width) if cell.is_merge_origin else 1
            rowspan = int(cell.span_height) if cell.is_merge_origin else 1
            cells.append(
                _table_cell_html(
                    cell,
                    palette,
                    fonts=fonts,
                    slide_width=slide_width,
                    colspan=colspan,
                    rowspan=rowspan,
                    row_index=row_index,
                    column_index=column_index,
                    row_count=row_count,
                    column_count=column_count,
                    table_style=table_style,
                )
            )
        row_height = 100 * int(row.height) / total_height
        rows.append(f'<tr style="height:{row_height:.4f}%">{"".join(cells)}</tr>')
    box = _shape_box(
        shape,
        slide_width=slide_width,
        slide_height=slide_height,
        coordinate_left=coordinate_left,
        coordinate_top=coordinate_top,
        coordinate_width=coordinate_width,
        coordinate_height=coordinate_height,
    )
    identity = _shape_identity(shape, layer)
    style_id = html.escape(str(table_style["style_id"]), quote=True)
    return (
        f'<div class="shape table-shape" {identity} style="{box}">'
        f'<table data-table-style-id="{style_id}"><colgroup>{columns}</colgroup>'
        f"<tbody>{''.join(rows)}</tbody>"
        "</table></div>"
    )


def _composite_slide_shapes(slide: Any) -> list[tuple[str, Any]]:
    layout = slide.slide_layout
    result: list[tuple[str, Any]] = []
    for shape in layout.slide_master.shapes:
        if _placeholder_key(shape) is None:
            result.append(("master", shape))
    for shape in layout.shapes:
        if _placeholder_key(shape) is None:
            result.append(("layout", shape))
    for shape in slide.shapes:
        result.append(("slide", shape))
    return result


def _group_coordinate_system(shape: Any) -> tuple[int, int, int, int]:
    try:
        transforms = shape._element.xpath(
            "./*[local-name()='grpSpPr']/*[local-name()='xfrm']"
        )
    except (AttributeError, TypeError, ValueError):
        transforms = []
    if transforms:
        child_offset = _xml_child(transforms[0], "chOff")
        child_extent = _xml_child(transforms[0], "chExt")
        if child_offset is not None and child_extent is not None:
            try:
                return (
                    int(child_offset.get("x", "0")),
                    int(child_offset.get("y", "0")),
                    max(int(child_extent.get("cx", "1")), 1),
                    max(int(child_extent.get("cy", "1")), 1),
                )
            except (TypeError, ValueError):
                pass
    return (0, 0, max(int(shape.width), 1), max(int(shape.height), 1))


def _render_pptx_shape(
    shape: Any,
    *,
    layer: str,
    slide: Any,
    palette: dict[str, str],
    fonts: dict[str, str],
    slide_width: int,
    slide_height: int,
    coordinate_left: int = 0,
    coordinate_top: int = 0,
    coordinate_width: int | None = None,
    coordinate_height: int | None = None,
) -> str:
    shape_type = getattr(getattr(shape, "shape_type", None), "name", "")
    if shape_type == "GROUP":
        box = _shape_box(
            shape,
            slide_width=slide_width,
            slide_height=slide_height,
            coordinate_left=coordinate_left,
            coordinate_top=coordinate_top,
            coordinate_width=coordinate_width,
            coordinate_height=coordinate_height,
        )
        child_left, child_top, child_width, child_height = _group_coordinate_system(
            shape
        )
        children = "".join(
            _render_pptx_shape(
                child,
                layer=layer,
                slide=slide,
                palette=palette,
                fonts=fonts,
                slide_width=slide_width,
                slide_height=slide_height,
                coordinate_left=child_left,
                coordinate_top=child_top,
                coordinate_width=child_width,
                coordinate_height=child_height,
            )
            for child in shape.shapes
        )
        identity = _shape_identity(shape, layer)
        return (
            f'<div class="shape shape-group" {identity} '
            f'data-coordinate-left="{child_left}" data-coordinate-top="{child_top}" '
            f'data-coordinate-width="{child_width}" '
            f'data-coordinate-height="{child_height}" style="{box}">{children}</div>'
        )

    if shape_type == "LINE":
        return _connector_svg(
            shape,
            layer=layer,
            slide=slide,
            slide_width=slide_width,
            slide_height=slide_height,
            palette=palette,
            coordinate_left=coordinate_left,
            coordinate_top=coordinate_top,
            coordinate_width=coordinate_width,
            coordinate_height=coordinate_height,
        )
    if shape_type == "PICTURE":
        return _render_picture(
            shape,
            layer=layer,
            slide_width=slide_width,
            slide_height=slide_height,
            coordinate_left=coordinate_left,
            coordinate_top=coordinate_top,
            coordinate_width=coordinate_width,
            coordinate_height=coordinate_height,
        )
    if getattr(shape, "has_table", False):
        return _render_table(
            shape,
            layer=layer,
            palette=palette,
            fonts=fonts,
            slide_width=slide_width,
            slide_height=slide_height,
            coordinate_left=coordinate_left,
            coordinate_top=coordinate_top,
            coordinate_width=coordinate_width,
            coordinate_height=coordinate_height,
        )

    box = _shape_box(
        shape,
        slide_width=slide_width,
        slide_height=slide_height,
        coordinate_left=coordinate_left,
        coordinate_top=coordinate_top,
        coordinate_width=coordinate_width,
        coordinate_height=coordinate_height,
    )
    identity = _shape_identity(shape, layer)
    preset = _preset_geometry_name(shape)
    geometry = f'data-preset-geometry="{html.escape(preset, quote=True)}"'
    if getattr(shape, "has_chart", False):
        chart_title = ""
        try:
            if shape.chart.has_title:
                chart_title = shape.chart.chart_title.text_frame.text.strip()
        except (AttributeError, TypeError, ValueError):
            chart_title = ""
        title_html = (
            f'<strong class="chart-title">{html.escape(chart_title)}</strong>'
            if chart_title
            else ""
        )
        try:
            chart_aspect = float(shape.width) / max(float(shape.height), 1.0)
        except (AttributeError, TypeError, ValueError):
            chart_aspect = 2.0
        return (
            f'<div class="shape chart" {identity} style="{box}">'
            f"{title_html}{_chart_svg(shape.chart, palette, chart_aspect=chart_aspect)}"
            "</div>"
        )

    if getattr(shape, "has_text_frame", False) and shape.text.strip():
        surface = _shape_surface_svg(shape, slide=slide, palette=palette)
        container_styles, alignment_attribute = _shape_text_container_attributes(
            shape, preset, palette
        )
        styles = [box, *container_styles]
        native_class = " html-native" if "[evoflux-html]" in shape.name else ""
        return (
            f'<div class="shape text-shape{native_class}" {identity} {geometry}'
            f"{alignment_attribute} "
            f'style="{";".join(styles)}">'
            f"{surface}"
            f"{_shape_text(shape, palette, fonts=fonts, slide=slide, slide_width=slide_width)}"
            "</div>"
        )

    surface = _shape_surface_svg(shape, slide=slide, palette=palette)
    if surface:
        return (
            f'<div class="shape vector-shape" {identity} {geometry} '
            f'style="{box};overflow:visible">{surface}</div>'
        )
    return ""


def _slide_preview_title(slide: Any) -> str:
    candidates: list[Any] = []
    try:
        if slide.shapes.title is not None:
            candidates.append(slide.shapes.title)
    except (AttributeError, TypeError, ValueError):
        pass
    candidates.extend(
        shape
        for shape in slide.shapes
        if shape not in candidates
        and "[role:title]" in str(getattr(shape, "name", "")).casefold()
    )
    candidates.extend(shape for shape in slide.shapes if shape not in candidates)
    for shape in candidates:
        try:
            title = " ".join(shape.text.split())
        except (AttributeError, TypeError, ValueError):
            continue
        if title:
            return title[:160]
    return ""


def _slide_notes_text(slide: Any) -> str:
    try:
        if not slide.has_notes_slide:
            return ""
        return slide.notes_slide.notes_text_frame.text.strip()
    except (AttributeError, TypeError, ValueError):
        return ""


def _render_pptx(source: Path) -> str:
    from pptx import Presentation

    presentation = Presentation(str(source))
    slide_width = int(presentation.slide_width or 0)
    slide_height = int(presentation.slide_height or 0)
    if slide_width <= 0 or slide_height <= 0:
        raise ValueError("Presentation has invalid slide dimensions")
    rendered_slides: list[str] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        palette = _pptx_theme_palette(slide)
        fonts = _pptx_theme_fonts(slide)
        shapes = [
            _render_pptx_shape(
                shape,
                layer=layer,
                slide=slide,
                palette=palette,
                fonts=fonts,
                slide_width=slide_width,
                slide_height=slide_height,
            )
            for layer, shape in _composite_slide_shapes(slide)
        ]

        background = _slide_background(slide, palette)
        ratio = f"{slide_width}/{slide_height}"
        title = _slide_preview_title(slide)
        label = f"Slide {slide_number}" + (f" — {title}" if title else "")
        notes = _slide_notes_text(slide)
        notes_metadata = ""
        if notes:
            escaped_notes = html.escape(notes)
            notes_metadata = (
                '<span class="slide-notes-metadata" hidden aria-hidden="true" '
                f'data-preview-notes="{html.escape(notes, quote=True)}">'
                f"{escaped_notes}</span>"
            )
        rendered_slides.append(
            f'<article class="slide-wrap" data-preview-item '
            f'data-preview-label="{html.escape(label, quote=True)}" '
            f'data-preview-title="{html.escape(title, quote=True)}">'
            f'<div class="slide-number">{slide_number}</div>'
            f'<section class="slide" style="aspect-ratio:{ratio};background:{background}">'
            f"{''.join(shapes)}</section>{notes_metadata}</article>"
        )

    css = """
    *{box-sizing:border-box}body{margin:0;padding:28px;background:#e8eaed;color:#202124;
    font-family:Arial,sans-serif}.slide-wrap{position:relative;width:min(1120px,94vw);
    margin:0 auto 30px}.slide{position:relative;width:100%;background:#fff;
    container-type:inline-size;
    overflow:hidden;box-shadow:0 3px 14px #0003}.slide-number{position:absolute;right:calc(100%
    + 8px);top:0;color:#6b7280;font-size:12px}.shape{position:absolute;overflow:hidden}
    .shape-group,.text-shape,.vector-shape{overflow:visible}.text-shape{padding:0}
    .shape-surface{position:absolute;inset:0;display:block;width:100%;height:100%;
    overflow:visible;pointer-events:none}.text-frame{position:relative;z-index:1;width:100%;height:100%;display:flex;
    flex-direction:column;overflow:hidden}.text-shape[data-default-text-align="center"] p{
    text-align:center!important}.text-shape[data-shape-font-color="theme"] .text-frame span{
    color:var(--shape-text-color)!important}.text-shape[data-text-unflip~="horizontal"] .text-frame{
    transform:scaleX(-1)}.text-shape[data-text-unflip~="vertical"] .text-frame{
    transform:scaleY(-1)}.text-shape[data-text-unflip~="horizontal"][data-text-unflip~="vertical"] .text-frame{
    transform:scale(-1,-1)}
    .text-shape p{margin:0 0 .2em;line-height:1.12;white-space:pre-wrap;
    overflow-wrap:normal;word-break:normal;break-inside:avoid}.bullet-marker{text-align:center}
    .text-shape.html-native{padding:0}.text-shape.html-native p{margin:0}
    .picture-frame{overflow:hidden}.picture{position:absolute;display:block;max-width:none;
    max-height:none}.table-shape table{width:100%;height:100%;border-collapse:collapse;
    table-layout:fixed}.table-shape td{border:1px solid #b8bdc6;overflow:hidden}
    .table-shape p{overflow:hidden}.chart{display:flex;flex-direction:column;align-items:center;
    justify-content:flex-start;border:0;background:transparent;color:#374151;
    gap:.3cqw}.chart-title{font-size:1.875cqw;line-height:1.1}.chart-svg{width:100%;
    flex:1;min-height:0}.chart-empty,.chart-unsupported{font-size:1.1cqw;color:#6b7280}
    .vector-shape{pointer-events:none}.connector{overflow:visible;pointer-events:none}
    .slide-notes-metadata{display:none!important}
    """
    return _page(title=source.name, body="".join(rendered_slides), css=css)


def _render_pdf(source: Path) -> str:
    from app.agent.builtin_plugins.documents.rendering.internal import (
        count_pdf_pages,
        render_pdf_pages,
    )

    cache_path = _cache_path(source)
    render_root = cache_path.parent / f"{cache_path.stem}-pages"
    total_pages = count_pdf_pages(source)
    pages: list[Path] = []
    try:
        pages = render_pdf_pages(
            source,
            render_root,
            dpi=120,
            max_pages=MAX_PDF_PREVIEW_PAGES,
            max_total_bytes=MAX_PDF_PREVIEW_RASTER_BYTES,
            max_pixels_per_page=MAX_PDF_PREVIEW_PIXELS_PER_PAGE,
        )
        rendered: list[str] = []
        for page_number, page in enumerate(pages, start=1):
            encoded = base64.b64encode(page.read_bytes()).decode("ascii")
            rendered.append(
                '<article class="pdf-page-wrap" data-preview-item '
                f'data-preview-label="Page {page_number}">'
                f'<span class="pdf-page-number">{page_number}</span>'
                f'<img class="pdf-page" src="data:image/png;base64,{encoded}" '
                f'alt="Page {page_number}"></article>'
            )
        if len(pages) < total_pages:
            rendered.append(
                '<p class="pdf-preview-limit" role="status">'
                f"Preview shows {len(pages)} of {total_pages} pages. "
                "Download the document to inspect the remaining pages.</p>"
            )
    finally:
        shutil.rmtree(render_root, ignore_errors=True)
    css = """
    *{box-sizing:border-box}body{margin:0;padding:28px;background:#e5e7eb;
    color:#242424;font-family:Arial,sans-serif}.pdf-page-wrap{position:relative;
    width:min(960px,94vw);margin:0 auto 28px}.pdf-page{display:block;width:100%;
    height:auto;background:#fff;box-shadow:0 2px 14px #0003}.pdf-page-number{
    position:absolute;right:calc(100% + 9px);top:0;color:#616161;font-size:12px}
    .pdf-preview-limit{max-width:720px;margin:20px auto;padding:14px 18px;border-radius:8px;
    background:#fff4ce;color:#5c2e00;box-shadow:0 1px 4px #0002}
    """
    return _page(title=source.name, body="".join(rendered), css=css)


def _render_source(source: Path) -> str:
    renderers = {
        ".docx": _render_docx,
        ".xlsx": _render_xlsx,
        ".pptx": _render_pptx,
        ".pdf": _render_pdf,
    }
    try:
        return renderers[source.suffix.lower()](source)
    except DocumentPreviewError:
        raise
    except Exception as exc:
        logger.warning(
            "document_preview_render_failed file={} error={}",
            source.name,
            str(exc)[:500],
        )
        raise DocumentPreviewError(f"Could not render this document: {exc}") from exc


def _render_document_preview(source: Path) -> Path:
    """Render ``source`` to a cached, self-contained HTML document."""
    suffix = source.suffix.lower()
    if suffix not in SUPPORTED_DOCUMENT_PREVIEW_EXTENSIONS:
        raise DocumentPreviewUnsupportedError(
            f"{suffix or 'This file type'} is not supported for document preview."
        )
    if source.stat().st_size > MAX_DOCUMENT_PREVIEW_BYTES:
        raise DocumentPreviewUnsupportedError(
            "Document preview is limited to "
            f"{MAX_DOCUMENT_PREVIEW_BYTES // (1024 * 1024)} MB."
        )
    preflight_ooxml_package(source, suffix)

    output = _cache_path(source)
    with _render_lock_for(output):
        prepare_preview_cache_directory(output.parent)
        if cached_preview_is_valid(
            output,
            max_bytes=MAX_DOCUMENT_PREVIEW_HTML_BYTES,
        ):
            mark_cached_preview_used(output)
            maintain_preview_cache(output.parent, preserve=output)
            if cached_preview_is_valid(
                output,
                max_bytes=MAX_DOCUMENT_PREVIEW_HTML_BYTES,
            ):
                return output
        try:
            output.unlink()
        except FileNotFoundError:
            pass
        except IsADirectoryError as exc:
            raise DocumentPreviewError(
                "The document preview cache entry is unsafe."
            ) from exc

        maintain_preview_cache(output.parent)
        descriptor, temporary_name = tempfile.mkstemp(
            prefix=f".{output.stem}-",
            suffix=".tmp",
            dir=output.parent,
        )
        temporary = Path(temporary_name)
        try:
            rendered = _render_source(source).encode("utf-8")
            if len(rendered) > MAX_DOCUMENT_PREVIEW_HTML_BYTES:
                raise DocumentPreviewUnsupportedError(
                    "The generated preview is too large for the in-app viewer. "
                    "Download the document to inspect it externally."
                )
            with os.fdopen(descriptor, "wb") as handle:
                descriptor = -1
                handle.write(rendered)
            os.replace(temporary, output)
        finally:
            if descriptor >= 0:
                os.close(descriptor)
            temporary.unlink(missing_ok=True)
        maintain_preview_cache(output.parent, preserve=output)
    return output


document_preview_provider = DocumentPreviewProvider(
    name="evoflux.documents",
    extensions=SUPPORTED_DOCUMENT_PREVIEW_EXTENSIONS,
    render=_render_document_preview,
)
render_document_preview = _render_document_preview


__all__ = [
    "DOCUMENT_PREVIEW_CSP",
    "DocumentPreviewError",
    "DocumentPreviewUnsupportedError",
    "MAX_DOCUMENT_PREVIEW_BYTES",
    "MAX_DOCUMENT_PREVIEW_HTML_BYTES",
    "MAX_PDF_PREVIEW_PAGES",
    "MAX_PDF_PREVIEW_PIXELS_PER_PAGE",
    "MAX_PDF_PREVIEW_RASTER_BYTES",
    "SUPPORTED_DOCUMENT_PREVIEW_EXTENSIONS",
    "document_preview_provider",
    "render_document_preview",
]
