"""Native and SVG-fidelity PPTX authoring without browser or Node runtimes."""

from __future__ import annotations

import asyncio
from dataclasses import dataclass, field
import json
from pathlib import Path
import re
from typing import Annotated, Any, Literal
from xml.etree import ElementTree as ET

from PIL import Image, ImageChops
from pydantic import (
    BaseModel,
    ConfigDict,
    Field,
    TypeAdapter,
    field_validator,
    model_validator,
)

from app.services.office.internal_rendering import render_pptx_pages, render_svg
from app.services.office.runtime import file_sha256

MAX_SLIDES = 80
MAX_ELEMENTS_PER_SLIDE = 160
_PLACEHOLDER = re.compile(r"\b(?:lorem ipsum|todo|tbd|click to add)\b", re.I)


class SlidePosition(BaseModel):
    model_config = ConfigDict(extra="forbid")
    left: float = Field(ge=0)
    top: float = Field(ge=0)
    width: float = Field(gt=0)
    height: float = Field(gt=0)


class TextElement(BaseModel):
    model_config = ConfigDict(extra="forbid")
    type: Literal["text"] = "text"
    name: str | None = Field(default=None, max_length=160)
    position: SlidePosition
    text: str = Field(min_length=1, max_length=20000)
    font_size: float = Field(default=24, ge=6, le=160)
    typeface: str | None = Field(default=None, max_length=120)
    color: str = "#111827"
    bold: bool = False
    italic: bool = False
    alignment: Literal["left", "center", "right", "justify"] = "left"
    vertical_alignment: Literal["top", "middle", "bottom"] = "top"
    auto_fit: Literal["none", "shrinkText", "resizeShapeToFitText"] = "shrinkText"
    fill: str = "none"
    line_fill: str = "none"
    line_width: float = Field(default=0, ge=0, le=20)


class ShapeElement(BaseModel):
    model_config = ConfigDict(extra="forbid")
    type: Literal["shape"] = "shape"
    name: str | None = Field(default=None, max_length=160)
    position: SlidePosition
    geometry: str = Field(default="rect", min_length=1, max_length=80)
    fill: str = "#ffffff"
    line_fill: str = "none"
    line_width: float = Field(default=0, ge=0, le=20)
    text: str | None = Field(default=None, max_length=12000)
    font_size: float = Field(default=20, ge=6, le=160)
    text_color: str = "#111827"
    bold: bool = False
    alignment: Literal["left", "center", "right", "justify"] = "left"
    vertical_alignment: Literal["top", "middle", "bottom"] = "middle"


class ImageElement(BaseModel):
    model_config = ConfigDict(extra="forbid")
    type: Literal["image"] = "image"
    name: str | None = Field(default=None, max_length=160)
    position: SlidePosition
    asset_path: str = Field(min_length=1, max_length=2000)
    alt: str = Field(min_length=1, max_length=1000)
    fit: Literal["cover", "contain"] = "cover"


class TableElement(BaseModel):
    model_config = ConfigDict(extra="forbid")
    type: Literal["table"] = "table"
    name: str | None = Field(default=None, max_length=160)
    position: SlidePosition
    values: list[list[str | int | float | None]] = Field(min_length=1, max_length=100)
    header_row: bool = True
    header_fill: str = "#e2e8f0"
    header_text_color: str = "#0f172a"
    body_fill: str = "#ffffff"
    body_text_color: str = "#334155"
    font_size: float = Field(default=16, ge=6, le=72)

    @field_validator("values")
    @classmethod
    def rectangular(cls, rows: list[list[Any]]) -> list[list[Any]]:
        width = len(rows[0])
        if width == 0 or width > 40 or any(len(row) != width for row in rows):
            raise ValueError("table values must be a non-empty rectangular matrix")
        return rows


class ChartSeries(BaseModel):
    model_config = ConfigDict(extra="forbid")
    name: str = Field(min_length=1, max_length=240)
    values: list[float] = Field(min_length=1, max_length=500)
    fill: str | None = None


class ChartElement(BaseModel):
    model_config = ConfigDict(extra="forbid")
    type: Literal["chart"] = "chart"
    name: str | None = Field(default=None, max_length=160)
    position: SlidePosition
    chart_type: Literal["bar", "line", "pie", "doughnut", "area"] = "bar"
    title: str | None = Field(default=None, max_length=500)
    categories: list[str] = Field(min_length=1, max_length=500)
    series: list[ChartSeries] = Field(min_length=1, max_length=40)
    has_legend: bool = True
    show_values: bool = False

    @model_validator(mode="after")
    def matching_series(self) -> ChartElement:
        if any(len(series.values) != len(self.categories) for series in self.series):
            raise ValueError("every chart series must match categories length")
        return self


SlideElement = Annotated[
    TextElement | ShapeElement | ImageElement | TableElement | ChartElement,
    Field(discriminator="type"),
]


class VisualShell(BaseModel):
    """Static SVG composition rendered by the bundled Rust renderer."""

    model_config = ConfigDict(extra="forbid")
    svg_path: str = Field(min_length=1, max_length=2000)
    alt: str = Field(min_length=1, max_length=1000)
    reference_svg_path: str | None = Field(default=None, max_length=2000)
    max_changed_pixel_ratio: float = Field(default=0.02, ge=0, le=0.25)
    max_mean_absolute_error: float = Field(default=0.01, ge=0, le=0.25)


class NativeSlide(BaseModel):
    model_config = ConfigDict(extra="forbid")
    id: str = Field(pattern=r"^[A-Za-z0-9][A-Za-z0-9_-]{0,79}$")
    background: str = "#ffffff"
    visual_shell: VisualShell | None = None
    elements: list[SlideElement] = Field(
        default_factory=list, max_length=MAX_ELEMENTS_PER_SLIDE
    )
    speaker_notes: str | None = Field(default=None, max_length=30000)


class NativePptxProject(BaseModel):
    model_config = ConfigDict(extra="forbid")
    schema_version: Literal[3] = 3
    mode: Literal["new"] = "new"
    quality_profile: Literal["fidelity", "hybrid", "native"] = "fidelity"
    title: str = Field(min_length=1, max_length=500)
    width: int = Field(default=1280, ge=640, le=3840)
    height: int = Field(default=720, ge=360, le=2160)
    slides: list[NativeSlide] = Field(min_length=1, max_length=MAX_SLIDES)

    @model_validator(mode="after")
    def validate_geometry(self) -> NativePptxProject:
        ids = [slide.id for slide in self.slides]
        if len(ids) != len(set(ids)):
            raise ValueError("slide ids must be unique")
        for slide in self.slides:
            if self.quality_profile == "fidelity":
                if slide.visual_shell is None or slide.elements:
                    raise ValueError(
                        f"slide {slide.id} requires only visual_shell in fidelity mode"
                    )
            elif self.quality_profile == "hybrid":
                if slide.visual_shell is None or not slide.elements:
                    raise ValueError(
                        f"slide {slide.id} requires visual_shell and native elements in hybrid mode"
                    )
                if slide.visual_shell.reference_svg_path is None:
                    raise ValueError(
                        f"slide {slide.id} requires reference_svg_path in hybrid mode"
                    )
            elif slide.visual_shell is not None or not slide.elements:
                raise ValueError(
                    f"slide {slide.id} requires native elements and no visual_shell in native mode"
                )
            for element in slide.elements:
                position = element.position
                if position.left + position.width > self.width + 0.01:
                    raise ValueError(
                        f"slide {slide.id} element exceeds the right slide boundary"
                    )
                if position.top + position.height > self.height + 0.01:
                    raise ValueError(
                        f"slide {slide.id} element exceeds the bottom slide boundary"
                    )
        return self


_PROJECT_ADAPTER = TypeAdapter(NativePptxProject)


@dataclass
class NativePptxPipelineResult:
    action: str
    work_dir: Path
    output: Path | None = None
    previews: list[Path] = field(default_factory=list)
    layout_paths: list[Path] = field(default_factory=list)
    manifest_path: Path | None = None
    issues: list[dict[str, Any]] = field(default_factory=list)
    metadata: dict[str, Any] = field(default_factory=dict)

    @property
    def passed(self) -> bool:
        return not any(issue.get("severity") == "error" for issue in self.issues)


def load_native_pptx_project(path: Path) -> NativePptxProject:
    return _PROJECT_ADAPTER.validate_json(path.read_text(encoding="utf-8"))


def _project_file(project_dir: Path, value: str, *, label: str) -> Path:
    relative = Path(value)
    if relative.is_absolute():
        raise ValueError(f"{label} must be relative to the project directory")
    root = project_dir.resolve()
    resolved = (root / relative).resolve()
    if not resolved.is_relative_to(root):
        raise ValueError(f"{label} escapes the project directory")
    if not resolved.is_file():
        raise FileNotFoundError(f"{label} does not exist: {resolved}")
    return resolved


def _validate_svg(source: Path) -> None:
    """Reject active or externally loaded SVG content before rasterization."""
    try:
        root = ET.parse(source).getroot()
    except ET.ParseError as exc:
        raise ValueError(f"visual shell is not valid SVG: {source.name}") from exc
    forbidden = {"script", "foreignObject", "iframe", "audio", "video"}
    href_names = {"href", "{http://www.w3.org/1999/xlink}href"}
    for element in root.iter():
        tag = element.tag.rsplit("}", 1)[-1]
        if tag in forbidden:
            raise ValueError(f"visual shell contains forbidden <{tag}> content")
        for name, value in element.attrib.items():
            if name.lower().startswith("on"):
                raise ValueError("visual shell contains an event handler")
            if name in href_names and value and not value.startswith(("#", "data:")):
                raise ValueError("visual shell references an external resource")


def native_pptx_catalog() -> dict[str, Any]:
    return {
        "workflow": "evoflux-openxml-svg-pptx",
        "invariants": [
            "Render visual shells from local SVG with the bundled Rust renderer.",
            "Keep native text, shapes, images, tables, and charts editable.",
            "Render every generated slide with the internal OOXML renderer.",
            "Reject placeholders, missing assets, out-of-slide geometry, and parity drift.",
        ],
        "quality_profiles": {
            "fidelity": "Full-slide SVG composition embedded as a static visual.",
            "hybrid": "SVG shell plus editable native overlays and reference parity gate.",
            "native": "Fully editable OpenXML primitives.",
        },
        "supported_elements": ["text", "shape", "image", "table", "chart"],
        "project_json_schema": NativePptxProject.model_json_schema(),
    }


def validate_native_pptx_project(
    project: NativePptxProject, project_path: Path
) -> dict[str, Any]:
    project_dir = project_path.parent.resolve()
    for slide in project.slides:
        if slide.visual_shell is not None:
            shell = _project_file(
                project_dir,
                slide.visual_shell.svg_path,
                label=f"slide {slide.id} visual_shell.svg_path",
            )
            if shell.suffix.lower() != ".svg":
                raise ValueError("visual_shell.svg_path must reference an SVG file")
            _validate_svg(shell)
            if slide.visual_shell.reference_svg_path:
                reference = _project_file(
                    project_dir,
                    slide.visual_shell.reference_svg_path,
                    label=f"slide {slide.id} visual_shell.reference_svg_path",
                )
                if reference.suffix.lower() != ".svg":
                    raise ValueError(
                        "visual_shell.reference_svg_path must reference an SVG file"
                    )
                _validate_svg(reference)
        for element in slide.elements:
            if isinstance(element, ImageElement):
                path = Path(element.asset_path)
                if not path.is_absolute():
                    path = project_dir / path
                if not path.is_file():
                    raise FileNotFoundError(
                        f"presentation image does not exist: {path}"
                    )
    return {
        "valid": True,
        "mode": "new",
        "quality_profile": project.quality_profile,
        "slide_count": len(project.slides),
        "element_count": sum(len(slide.elements) for slide in project.slides),
        "project_sha256": file_sha256(project_path),
    }


def _rgb(value: str) -> Any:
    from pptx.dml.color import RGBColor

    normalized = value.strip().lstrip("#")
    if len(normalized) != 6:
        raise ValueError(f"expected #RRGGBB color, got {value!r}")
    return RGBColor.from_string(normalized.upper())


def _coordinates(position: SlidePosition) -> tuple[Any, Any, Any, Any]:
    from pptx.util import Inches

    return (
        Inches(position.left / 96),
        Inches(position.top / 96),
        Inches(position.width / 96),
        Inches(position.height / 96),
    )


def _set_fill(fill: Any, color: str) -> None:
    if color == "none":
        fill.background()
    else:
        fill.solid()
        fill.fore_color.rgb = _rgb(color)


def _set_line(line: Any, color: str, width: float) -> None:
    from pptx.util import Pt

    if color == "none" or width == 0:
        line.fill.background()
        return
    line.color.rgb = _rgb(color)
    line.width = Pt(width)


def _style_text_frame(frame: Any, element: TextElement | ShapeElement) -> None:
    from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
    from pptx.util import Pt

    alignments = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }
    anchors = {
        "top": MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }
    frame.clear()
    frame.vertical_anchor = anchors[element.vertical_alignment]
    if isinstance(element, TextElement):
        frame.auto_size = {
            "none": MSO_AUTO_SIZE.NONE,
            "shrinkText": MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE,
            "resizeShapeToFitText": MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT,
        }[element.auto_fit]
    paragraph = frame.paragraphs[0]
    paragraph.alignment = alignments[element.alignment]
    run = paragraph.add_run()
    run.text = element.text or ""
    run.font.size = Pt(element.font_size)
    run.font.bold = element.bold
    run.font.color.rgb = _rgb(
        element.color if isinstance(element, TextElement) else element.text_color
    )
    if isinstance(element, TextElement):
        run.font.italic = element.italic
        if element.typeface:
            run.font.name = element.typeface


def _shape_geometry(value: str) -> Any:
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

    mapping = {
        "rect": MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        "roundrect": MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        "ellipse": MSO_AUTO_SHAPE_TYPE.OVAL,
        "triangle": MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE,
        "diamond": MSO_AUTO_SHAPE_TYPE.DIAMOND,
        "hexagon": MSO_AUTO_SHAPE_TYPE.HEXAGON,
        "chevron": MSO_AUTO_SHAPE_TYPE.CHEVRON,
    }
    return mapping.get(value.casefold(), MSO_AUTO_SHAPE_TYPE.RECTANGLE)


def _add_element(slide: Any, element: SlideElement, project_dir: Path) -> None:
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION
    from pptx.util import Pt

    left, top, width, height = _coordinates(element.position)
    if isinstance(element, TextElement):
        shape = slide.shapes.add_textbox(left, top, width, height)
        shape.name = element.name or shape.name
        _set_fill(shape.fill, element.fill)
        _set_line(shape.line, element.line_fill, element.line_width)
        _style_text_frame(shape.text_frame, element)
        return
    if isinstance(element, ShapeElement):
        shape = slide.shapes.add_shape(
            _shape_geometry(element.geometry), left, top, width, height
        )
        shape.name = element.name or shape.name
        _set_fill(shape.fill, element.fill)
        _set_line(shape.line, element.line_fill, element.line_width)
        if element.text:
            _style_text_frame(shape.text_frame, element)
        return
    if isinstance(element, ImageElement):
        path = Path(element.asset_path)
        if not path.is_absolute():
            path = project_dir / path
        with Image.open(path) as image:
            image_ratio = image.width / image.height
        frame_ratio = float(width) / float(height)
        if element.fit == "contain":
            if image_ratio > frame_ratio:
                picture_width = width
                picture_height = int(width / image_ratio)
                picture_left = left
                picture_top = top + int((height - picture_height) / 2)
            else:
                picture_height = height
                picture_width = int(height * image_ratio)
                picture_left = left + int((width - picture_width) / 2)
                picture_top = top
            shape = slide.shapes.add_picture(
                str(path), picture_left, picture_top, picture_width, picture_height
            )
        else:
            shape = slide.shapes.add_picture(str(path), left, top, width, height)
            if image_ratio > frame_ratio:
                crop = (1 - frame_ratio / image_ratio) / 2
                shape.crop_left = crop
                shape.crop_right = crop
            elif image_ratio < frame_ratio:
                crop = (1 - image_ratio / frame_ratio) / 2
                shape.crop_top = crop
                shape.crop_bottom = crop
        shape.name = element.name or shape.name
        return
    if isinstance(element, TableElement):
        shape = slide.shapes.add_table(
            len(element.values), len(element.values[0]), left, top, width, height
        )
        shape.name = element.name or shape.name
        table = shape.table
        for row_index, row in enumerate(element.values):
            for column_index, value in enumerate(row):
                cell = table.cell(row_index, column_index)
                cell.text = "" if value is None else str(value)
                _set_fill(
                    cell.fill,
                    element.header_fill
                    if element.header_row and row_index == 0
                    else element.body_fill,
                )
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.runs[0].font.size = Pt(element.font_size)
                paragraph.runs[0].font.color.rgb = _rgb(
                    element.header_text_color
                    if element.header_row and row_index == 0
                    else element.body_text_color
                )
                paragraph.runs[0].font.bold = element.header_row and row_index == 0
        return
    data = CategoryChartData()
    data.categories = element.categories
    for series in element.series:
        data.add_series(series.name, series.values)
    chart_types = {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE_MARKERS,
        "pie": XL_CHART_TYPE.PIE,
        "doughnut": XL_CHART_TYPE.DOUGHNUT,
        "area": XL_CHART_TYPE.AREA,
    }
    chart = slide.shapes.add_chart(
        chart_types[element.chart_type], left, top, width, height, data
    ).chart
    chart.has_legend = element.has_legend
    if element.title:
        chart.has_title = True
        chart.chart_title.text_frame.text = element.title
    if element.show_values:
        for plot in chart.plots:
            plot.has_data_labels = True
            plot.data_labels.show_value = True
            plot.data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
    for source_series, native_series in zip(element.series, chart.series, strict=True):
        if source_series.fill:
            native_series.format.fill.solid()
            native_series.format.fill.fore_color.rgb = _rgb(source_series.fill)


def _visual_parity(
    reference_path: Path,
    preview_path: Path,
    *,
    max_changed_pixel_ratio: float,
    max_mean_absolute_error: float,
) -> dict[str, Any]:
    reference = Image.open(reference_path).convert("RGB")
    preview = (
        Image.open(preview_path)
        .convert("RGB")
        .resize(reference.size, Image.Resampling.LANCZOS)
    )
    detail_size = (min(960, reference.width), min(540, reference.height))
    structural_size = (min(240, reference.width), min(135, reference.height))
    structural_reference = reference.resize(structural_size, Image.Resampling.LANCZOS)
    structural_preview = preview.resize(structural_size, Image.Resampling.LANCZOS)
    difference = ImageChops.difference(structural_reference, structural_preview)
    histogram = difference.histogram()
    changed = sum(count for index, count in enumerate(histogram) if index % 256 != 0)
    pixel_count = structural_size[0] * structural_size[1] * 3
    changed_pixel_ratio = changed / max(1, pixel_count)
    mean_absolute_error = sum(
        (index % 256) * count for index, count in enumerate(histogram)
    ) / max(1, pixel_count * 255)
    return {
        "reference_path": str(reference_path),
        "preview_path": str(preview_path),
        "comparison_size": list(structural_size),
        "detail_comparison_size": list(detail_size),
        "changed_pixel_ratio": changed_pixel_ratio,
        "mean_absolute_error": mean_absolute_error,
        "max_changed_pixel_ratio": max_changed_pixel_ratio,
        "max_mean_absolute_error": max_mean_absolute_error,
        "passed": changed_pixel_ratio <= max_changed_pixel_ratio
        and mean_absolute_error <= max_mean_absolute_error,
    }


def _compose(
    project: NativePptxProject, project_path: Path, output: Path, work_dir: Path
) -> NativePptxPipelineResult:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    presentation.slide_width = Inches(project.width / 96)
    presentation.slide_height = Inches(project.height / 96)
    project_dir = project_path.parent.resolve()
    reference_paths: list[Path | None] = []
    layout_paths: list[Path] = []
    issues: list[dict[str, Any]] = []
    for index, plan in enumerate(project.slides, start=1):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = _rgb(plan.background)
        reference_path: Path | None = None
        if plan.visual_shell:
            shell_source = _project_file(
                project_dir,
                plan.visual_shell.svg_path,
                label=f"slide {plan.id} visual_shell.svg_path",
            )
            shell_png = render_svg(
                shell_source,
                work_dir / "shells" / f"slide-{index:03d}.png",
                width=project.width,
                height=project.height,
            )
            slide.shapes.add_picture(
                str(shell_png),
                0,
                0,
                presentation.slide_width,
                presentation.slide_height,
            )
            if plan.visual_shell.reference_svg_path:
                reference_source = _project_file(
                    project_dir,
                    plan.visual_shell.reference_svg_path,
                    label=f"slide {plan.id} visual_shell.reference_svg_path",
                )
                reference_path = render_svg(
                    reference_source,
                    work_dir / "references" / f"slide-{index:03d}.png",
                    width=project.width,
                    height=project.height,
                )
            else:
                reference_path = shell_png
        for element in plan.elements:
            text = getattr(element, "text", None)
            if text and _PLACEHOLDER.search(text):
                issues.append(
                    {
                        "severity": "error",
                        "code": "unresolved-placeholder",
                        "message": f"Slide {index} contains placeholder text: {text}",
                        "slide": index,
                    }
                )
            _add_element(slide, element, project_dir)
        if plan.speaker_notes:
            slide.notes_slide.notes_text_frame.text = plan.speaker_notes
        layout_path = work_dir / "layouts" / f"slide-{index:03d}.layout.json"
        layout_path.parent.mkdir(parents=True, exist_ok=True)
        layout_path.write_text(
            json.dumps(
                {
                    "slide": index,
                    "id": plan.id,
                    "elements": [item.model_dump() for item in plan.elements],
                },
                ensure_ascii=False,
                indent=2,
            )
            + "\n",
            encoding="utf-8",
        )
        layout_paths.append(layout_path)
        reference_paths.append(reference_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(output))
    previews = render_pptx_pages(output, work_dir / "previews", width=project.width)
    parity: list[dict[str, Any]] = []
    for index, (slide, reference, preview) in enumerate(
        zip(project.slides, reference_paths, previews, strict=True), start=1
    ):
        if reference is None or slide.visual_shell is None:
            continue
        metric = _visual_parity(
            reference,
            preview,
            max_changed_pixel_ratio=slide.visual_shell.max_changed_pixel_ratio,
            max_mean_absolute_error=slide.visual_shell.max_mean_absolute_error,
        )
        metric.update({"slide": index, "slide_id": slide.id})
        parity.append(metric)
        if not metric["passed"]:
            issues.append(
                {
                    "severity": "error",
                    "code": "visual-parity-drift",
                    "message": f"Slide {index} differs from its SVG reference",
                    "slide": index,
                }
            )
    manifest_path = work_dir / "native-pptx-manifest.json"
    manifest_path.write_text(
        json.dumps(
            {
                "schemaVersion": 3,
                "engine": "evoflux-openxml-svg",
                "qualityProfile": project.quality_profile,
                "slideCount": len(project.slides),
                "layoutPaths": [str(path) for path in layout_paths],
            },
            indent=2,
        )
        + "\n",
        encoding="utf-8",
    )
    if any(issue.get("severity") == "error" for issue in issues):
        output.unlink(missing_ok=True)
        candidate = None
    else:
        candidate = output
    return NativePptxPipelineResult(
        action="compose",
        work_dir=work_dir,
        output=candidate,
        previews=previews,
        layout_paths=layout_paths,
        manifest_path=manifest_path,
        issues=issues,
        metadata={
            "engine": "evoflux-openxml-svg",
            "slide_count": len(project.slides),
            "quality_profile": project.quality_profile,
            "visual_parity": parity,
            "editable_object_count": sum(
                len(slide.elements) for slide in project.slides
            ),
            "semantic_editable_object_count": sum(
                len(slide.elements) for slide in project.slides
            ),
        },
    )


async def compose_native_pptx_project(
    project_path: Path, output: Path, *, workspace_root: Path, work_dir: Path
) -> NativePptxPipelineResult:
    del workspace_root
    project = load_native_pptx_project(project_path)
    validate_native_pptx_project(project, project_path)
    work_dir.mkdir(parents=True, exist_ok=True)
    return await asyncio.to_thread(_compose, project, project_path, output, work_dir)


__all__ = [
    "ChartElement",
    "ImageElement",
    "NativePptxPipelineResult",
    "NativePptxProject",
    "ShapeElement",
    "TableElement",
    "TextElement",
    "VisualShell",
    "_visual_parity",
    "compose_native_pptx_project",
    "load_native_pptx_project",
    "native_pptx_catalog",
    "validate_native_pptx_project",
]
