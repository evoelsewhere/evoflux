"""Declarative PowerPoint artifact compiler for EvoFlux.

The engine keeps presentation authoring behind a validated JSON contract so an
agent chooses narrative content, layouts, and native Office objects without
writing low-level coordinate scripts.  Existing decks continue to use the
package-preserving template editor; this module owns deterministic creation,
rendering, and validation of new decks.
"""

from __future__ import annotations

import os
import re
import tempfile
import zipfile
from collections.abc import Mapping
from dataclasses import dataclass
from pathlib import Path
from typing import Annotated, Any, Literal

from pydantic import BaseModel, ConfigDict, Field, field_validator, model_validator
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from app.agent.builtin_skills.pptx.scripts.icons import add_icon
from app.agent.builtin_skills.pptx.scripts.office_features import (
    RichParagraph,
    RichTextRun,
    add_grouped_process,
    add_native_chart,
    add_native_table,
    add_rich_text,
    set_accessibility,
    set_slide_transition,
)
from app.agent.builtin_skills.pptx.scripts.qa import inspect_pptx
from app.agent.builtin_skills.pptx.scripts.stylekit import (
    LayoutGuard,
    LayoutName,
    LayoutProfileName,
    PresentationTheme,
    QualityLedger,
    add_footer,
    add_image_cover,
    add_speaker_source_note,
    add_text,
    add_title,
    apply_layout_profile,
    estimate_text_fit,
    layout_plan,
    layout_profile,
    new_wide_presentation,
    set_background,
)
from app.services.office_visual_qa_service import render_office_images

ChartKind = Literal["column", "bar", "line", "area", "pie", "doughnut"]
TextAlign = Literal["left", "center", "right"]
TransitionKind = Literal["none", "fade", "push", "wipe", "cut", "morph"]


LAYOUT_SLOTS: dict[LayoutName, tuple[str, ...]] = {
    "cover": ("primary", "visual"),
    "hero": ("canvas",),
    "split": ("text", "visual"),
    "visual-left": ("visual", "text"),
    "comparison": ("left", "right"),
    "statement": ("statement", "evidence"),
    "workstreams": ("column-1", "column-2", "column-3", "column-4", "summary"),
    "section": ("statement", "marker"),
    "agenda": ("intro", "agenda"),
    "three-column": ("column-1", "column-2", "column-3"),
    "four-quadrant": ("top-left", "top-right", "bottom-left", "bottom-right"),
    "chart-focus": ("chart", "insight"),
    "table-focus": ("table", "note"),
    "timeline": ("timeline",),
    "metrics": ("metric-1", "metric-2", "metric-3", "detail"),
    "quote": ("quote", "attribution"),
    "process": ("process", "note"),
    "image-full": ("canvas",),
}

LAYOUT_DEFAULT_PROFILES: dict[LayoutName, LayoutProfileName] = {
    "workstreams": "operational",
    "four-quadrant": "executive-dense",
    "table-focus": "executive-dense",
}

PPTX_CAPABILITIES: dict[str, dict[str, str]] = {
    "text": {"create": "full", "edit": "full", "preview": "medium"},
    "images_and_crop": {"create": "full", "edit": "full", "preview": "medium"},
    "charts": {"create": "category charts", "edit": "data/style", "preview": "medium"},
    "tables": {"create": "full", "edit": "cells/style", "preview": "medium"},
    "groups_and_connectors": {
        "create": "simple process",
        "edit": "partial",
        "preview": "medium",
    },
    "speaker_notes": {"create": "full", "edit": "full", "preview": "structural"},
    "transitions_and_morph": {
        "create": "basic",
        "edit": "basic",
        "preview": "structural",
    },
    "masters_layouts_themes": {
        "create": "tokenized",
        "edit": "template-first",
        "preview": "medium",
    },
    "smartart": {
        "create": "unsupported",
        "edit": "preserve-only",
        "preview": "structural",
    },
    "audio_video_ole": {
        "create": "template-only",
        "edit": "preserve-only",
        "preview": "structural",
    },
    "complex_animations": {
        "create": "template-only",
        "edit": "preserve-only",
        "preview": "structural",
    },
}

_SIGNED_AXIS_ID = re.compile(rb'(<c:(?:axId|crossAx)\s+val=")(-\d+)(")')


def _hex(value: str) -> str:
    normalized = value.strip().removeprefix("#").upper()
    if len(normalized) != 6 or any(
        character not in "0123456789ABCDEF" for character in normalized
    ):
        raise ValueError(f"Expected a six-digit RGB color, received {value!r}")
    return normalized


class PptxThemeSpec(BaseModel):
    model_config = ConfigDict(extra="forbid")

    title_font: str = "Aptos Display"
    body_font: str = "Aptos"
    background: str = "F7F5F0"
    ink: str = "20303C"
    muted: str = "66717C"
    accent: str = "2F6D68"
    highlight: str = "D6A756"
    title_pt: int = Field(default=40, ge=24, le=64)
    body_pt: int = Field(default=20, ge=8, le=32)
    margin_inches: float = Field(default=0.72, ge=0.3, le=1.4)

    @field_validator("background", "ink", "muted", "accent", "highlight")
    @classmethod
    def validate_color(cls, value: str) -> str:
        return _hex(value)

    def to_native(self) -> PresentationTheme:
        return PresentationTheme(**self.model_dump())


class TextBlock(BaseModel):
    model_config = ConfigDict(extra="forbid")

    type: Literal["text"] = "text"
    text: str = Field(min_length=1, max_length=1800)
    size: float | None = Field(default=None, ge=7, le=54)
    bold: bool = False
    color: str | None = None
    align: TextAlign = "left"
    role: Literal[
        "body", "subheading", "section-heading", "label", "caption", "metadata"
    ] = "body"
    max_lines: int | None = Field(default=None, ge=1, le=30)

    @field_validator("color")
    @classmethod
    def validate_optional_color(cls, value: str | None) -> str | None:
        return _hex(value) if value else None


class BulletsBlock(BaseModel):
    model_config = ConfigDict(extra="forbid")

    type: Literal["bullets"] = "bullets"
    items: list[str] = Field(min_length=1, max_length=12)
    size: float | None = Field(default=None, ge=7, le=32)
    columns: int = Field(default=1, ge=1, le=4)
    color: str | None = None

    @field_validator("items")
    @classmethod
    def validate_items(cls, value: list[str]) -> list[str]:
        if any(not item.strip() for item in value):
            raise ValueError("Bullet items cannot be blank")
        return value

    @field_validator("color")
    @classmethod
    def validate_optional_color(cls, value: str | None) -> str | None:
        return _hex(value) if value else None


class ImageBlock(BaseModel):
    model_config = ConfigDict(extra="forbid")

    type: Literal["image"] = "image"
    path: str = Field(min_length=1)
    alt_text: str = Field(min_length=1, max_length=500)
    focal_x: float = Field(default=0.5, ge=0, le=1)
    focal_y: float = Field(default=0.5, ge=0, le=1)


class ChartBlock(BaseModel):
    model_config = ConfigDict(extra="forbid")

    type: Literal["chart"] = "chart"
    kind: ChartKind = "column"
    categories: list[str] = Field(min_length=1, max_length=40)
    series: dict[str, list[float]] = Field(min_length=1, max_length=8)
    title: str | None = Field(default=None, max_length=160)
    number_format: str = "0"
    show_legend: bool | None = None
    show_data_labels: bool = False
    alt_text: str | None = Field(default=None, max_length=500)

    @model_validator(mode="after")
    def validate_series_lengths(self) -> ChartBlock:
        expected = len(self.categories)
        invalid = [
            name for name, values in self.series.items() if len(values) != expected
        ]
        if invalid:
            raise ValueError(
                "Chart series must match category count: " + ", ".join(invalid)
            )
        return self


class TableBlock(BaseModel):
    model_config = ConfigDict(extra="forbid")

    type: Literal["table"] = "table"
    headers: list[str] = Field(min_length=1, max_length=12)
    rows: list[list[str | int | float]] = Field(min_length=1, max_length=30)
    column_weights: list[float] | None = None
    alt_text: str | None = Field(default=None, max_length=500)

    @model_validator(mode="after")
    def validate_table_shape(self) -> TableBlock:
        width = len(self.headers)
        if any(len(row) != width for row in self.rows):
            raise ValueError("Every table row must match the header width")
        if self.column_weights is not None and len(self.column_weights) != width:
            raise ValueError("column_weights must match the header width")
        return self


class ProcessBlock(BaseModel):
    model_config = ConfigDict(extra="forbid")

    type: Literal["process"] = "process"
    steps: list[str] = Field(min_length=2, max_length=6)
    alt_text: str | None = Field(default=None, max_length=500)


class IconBlock(BaseModel):
    model_config = ConfigDict(extra="forbid")

    type: Literal["icon"] = "icon"
    name: str = Field(min_length=1, max_length=80)
    color: str | None = None
    size_inches: float | None = Field(default=None, ge=0.18, le=2.5)
    alt_text: str | None = Field(default=None, max_length=500)

    @field_validator("color")
    @classmethod
    def validate_optional_color(cls, value: str | None) -> str | None:
        return _hex(value) if value else None


class MetricBlock(BaseModel):
    model_config = ConfigDict(extra="forbid")

    type: Literal["metric"] = "metric"
    value: str = Field(min_length=1, max_length=32)
    label: str = Field(min_length=1, max_length=100)
    context: str | None = Field(default=None, max_length=180)
    accent: str | None = None

    @field_validator("accent")
    @classmethod
    def validate_optional_color(cls, value: str | None) -> str | None:
        return _hex(value) if value else None


class QuoteBlock(BaseModel):
    model_config = ConfigDict(extra="forbid")

    type: Literal["quote"] = "quote"
    quote: str = Field(min_length=1, max_length=700)
    attribution: str | None = Field(default=None, max_length=180)


SlideBlock = Annotated[
    TextBlock
    | BulletsBlock
    | ImageBlock
    | ChartBlock
    | TableBlock
    | ProcessBlock
    | IconBlock
    | MetricBlock
    | QuoteBlock,
    Field(discriminator="type"),
]


class PptxSlideSpec(BaseModel):
    model_config = ConfigDict(extra="forbid")

    title: str = Field(min_length=1, max_length=180)
    layout: LayoutName
    profile: LayoutProfileName | None = None
    kicker: str | None = Field(default=None, max_length=80)
    slots: dict[str, SlideBlock] = Field(min_length=1)
    footer: str | None = Field(default=None, max_length=180)
    sources: list[str] = Field(default_factory=list, max_length=30)
    transition: TransitionKind = "none"

    @model_validator(mode="after")
    def validate_slots(self) -> PptxSlideSpec:
        allowed = set(LAYOUT_SLOTS[self.layout])
        unknown = sorted(set(self.slots) - allowed)
        if unknown:
            raise ValueError(
                f"Layout {self.layout!r} does not contain slots: {', '.join(unknown)}; "
                f"allowed slots are {', '.join(LAYOUT_SLOTS[self.layout])}"
            )
        return self

    @property
    def active_profile(self) -> LayoutProfileName:
        return self.profile or LAYOUT_DEFAULT_PROFILES.get(self.layout, "editorial")


class PresentationSpec(BaseModel):
    model_config = ConfigDict(extra="forbid")

    title: str = Field(min_length=1, max_length=180)
    theme: PptxThemeSpec = Field(default_factory=PptxThemeSpec)
    slides: list[PptxSlideSpec] = Field(min_length=1, max_length=80)
    allow_shape_only: bool = False


@dataclass(frozen=True)
class PptxBuildResult:
    output: Path
    report: dict[str, Any]
    render: dict[str, Any] | None

    @property
    def passed(self) -> bool:
        render_errors = list((self.render or {}).get("errors", []))
        return not self.report.get("errors") and not render_errors

    def to_dict(self) -> dict[str, Any]:
        return {
            "output": str(self.output),
            "passed": self.passed,
            "report": self.report,
            "render": self.render,
        }


def layout_catalog() -> list[dict[str, Any]]:
    """Return the stable layout/slot contract exposed to agents and tests."""

    return [
        {
            "name": name,
            "slots": list(slots),
            "default_profile": LAYOUT_DEFAULT_PROFILES.get(name, "editorial"),
        }
        for name, slots in LAYOUT_SLOTS.items()
    ]


def _alignment(value: TextAlign) -> PP_ALIGN:
    return {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }[value]


def _profile_body_size(profile: LayoutProfileName) -> float:
    return {
        "editorial": 20,
        "executive-dense": 14,
        "operational": 10,
    }[profile]


def _validated_size(
    size: float | None, profile: LayoutProfileName, role: str = "body"
) -> float:
    policy = layout_profile(profile)
    minimums = {
        "body": policy.body_min_pt,
        "subheading": policy.subheading_min_pt,
        "section-heading": policy.subheading_min_pt,
        "label": policy.body_min_pt,
        "caption": policy.caption_min_pt,
        "metadata": policy.metadata_min_pt,
    }
    actual = float(size if size is not None else _profile_body_size(profile))
    minimum = minimums.get(role, policy.body_min_pt)
    if actual < minimum:
        raise ValueError(
            f"{role} text uses {actual:g}pt in {profile}; minimum is {minimum:g}pt"
        )
    return actual


def _resolve_asset(path: str, asset_root: Path) -> Path:
    candidate = Path(path).expanduser()
    if not candidate.is_absolute():
        candidate = asset_root / candidate
    candidate = candidate.resolve()
    if not candidate.is_file():
        raise FileNotFoundError(f"PPTX image asset does not exist: {candidate}")
    return candidate


def _normalize_chart_axis_ids(source: Path) -> None:
    """Rewrite signed chart-axis IDs to their equivalent unsigned values.

    ``python-pptx`` may serialize randomly generated 32-bit axis IDs as signed
    integers. PowerPoint tolerates them, but strict OpenXML readers model the
    value as ``UInt32`` and reject the package. Keeping the same bit pattern in
    unsigned decimal form makes the deck portable without changing references.
    """

    def replace_axis_id(match: re.Match[bytes]) -> bytes:
        signed = int(match.group(2))
        unsigned = signed & 0xFFFFFFFF
        return match.group(1) + str(unsigned).encode("ascii") + match.group(3)

    with tempfile.NamedTemporaryFile(
        prefix=f".{source.stem}-",
        suffix=".pptx",
        dir=source.parent,
        delete=False,
    ) as temporary:
        temporary_path = Path(temporary.name)
    try:
        with (
            zipfile.ZipFile(source) as package,
            zipfile.ZipFile(temporary_path, "w") as rewritten,
        ):
            for info in package.infolist():
                data = package.read(info.filename)
                if info.filename.startswith("ppt/charts/") and info.filename.endswith(
                    ".xml"
                ):
                    data = _SIGNED_AXIS_ID.sub(replace_axis_id, data)
                rewritten.writestr(info, data)
        os.replace(temporary_path, source)
    finally:
        temporary_path.unlink(missing_ok=True)


def _fit_rich_text(
    text: str,
    *,
    width: int,
    height: int,
    size: float,
    ledger: QualityLedger,
    label: str,
) -> None:
    fit = estimate_text_fit(text, width=width, height=height, size=size)
    if not fit.fits:
        ledger.error(
            f"{label} needs about {fit.estimated_lines} lines but allows "
            f"{fit.maximum_lines}; shorten the copy or choose another layout"
        )


def _render_cover(
    slide: Any,
    specification: PptxSlideSpec,
    *,
    plan: Any,
    guard: LayoutGuard,
    theme: PresentationTheme,
    asset_root: Path,
) -> None:
    primary = plan.region("primary")
    visual = plan.region("visual")
    add_text(
        slide,
        specification.title,
        left=primary.left,
        top=primary.top,
        width=primary.width,
        height=int(primary.height * 0.48),
        font=theme.title_font,
        size=max(theme.title_pt, 50),
        color=theme.ink,
        bold=True,
        role="deck-title",
        max_lines=2,
        guard=guard,
    )
    primary_block = specification.slots.get("primary")
    if primary_block is not None:
        body_top = int(primary.top + primary.height * 0.58)
        body_height = int(primary.bottom - body_top)
        _render_block(
            slide,
            primary_block,
            left=primary.left,
            top=body_top,
            width=primary.width,
            height=body_height,
            profile=specification.active_profile,
            theme=theme,
            guard=guard,
            asset_root=asset_root,
        )
    visual_block = specification.slots.get("visual")
    if visual_block is not None:
        _render_block(
            slide,
            visual_block,
            left=visual.left,
            top=visual.top,
            width=visual.width,
            height=visual.height,
            profile=specification.active_profile,
            theme=theme,
            guard=guard,
            asset_root=asset_root,
        )


def _render_block(
    slide: Any,
    block: SlideBlock,
    *,
    left: int,
    top: int,
    width: int,
    height: int,
    profile: LayoutProfileName,
    theme: PresentationTheme,
    guard: LayoutGuard,
    asset_root: Path,
) -> Any:
    if isinstance(block, TextBlock):
        size = _validated_size(block.size, profile, block.role)
        return add_text(
            slide,
            block.text,
            left=left,
            top=top,
            width=width,
            height=height,
            font=theme.body_font,
            size=size,
            color=block.color or theme.ink,
            bold=block.bold,
            align=_alignment(block.align),
            role=block.role,
            max_lines=block.max_lines,
            profile=profile,
            guard=guard,
        )
    if isinstance(block, BulletsBlock):
        size = _validated_size(block.size, profile)
        joined = "\n".join(block.items)
        _fit_rich_text(
            joined,
            width=width,
            height=height,
            size=size,
            ledger=guard.ledger or QualityLedger(),
            label="Bullet block",
        )
        paragraphs = tuple(
            RichParagraph(
                runs=(
                    RichTextRun(
                        item,
                        font=theme.body_font,
                        size=int(size),
                        color=block.color or theme.ink,
                    ),
                ),
                bullet=True,
                space_after_pt=max(size * 0.42, 5),
            )
            for item in block.items
        )
        return add_rich_text(
            slide,
            paragraphs,
            left=left,
            top=top,
            width=width,
            height=height,
            columns=block.columns,
            theme=theme,
            guard=guard,
        )
    if isinstance(block, ImageBlock):
        return add_image_cover(
            slide,
            _resolve_asset(block.path, asset_root),
            left=left,
            top=top,
            width=width,
            height=height,
            focal_x=block.focal_x,
            focal_y=block.focal_y,
            alt_text=block.alt_text,
            guard=guard,
        )
    if isinstance(block, ChartBlock):
        shape = add_native_chart(
            slide,
            block.categories,
            block.series,
            left=left,
            top=top,
            width=width,
            height=height,
            kind=block.kind,
            title=block.title,
            number_format=block.number_format,
            show_legend=block.show_legend,
            show_data_labels=block.show_data_labels,
            theme=theme,
            guard=guard,
        )
        if block.alt_text:
            set_accessibility(
                shape, title=block.title or "Chart", description=block.alt_text
            )
        return shape
    if isinstance(block, TableBlock):
        shape = add_native_table(
            slide,
            block.headers,
            block.rows,
            left=left,
            top=top,
            width=width,
            height=height,
            column_weights=block.column_weights,
            theme=theme,
            guard=guard,
        )
        if block.alt_text:
            set_accessibility(shape, title="Table", description=block.alt_text)
        return shape
    if isinstance(block, ProcessBlock):
        shape = add_grouped_process(
            slide,
            block.steps,
            left=left,
            top=top,
            width=width,
            height=height,
            theme=theme,
            guard=guard,
        )
        if block.alt_text:
            set_accessibility(shape, title="Process", description=block.alt_text)
        return shape
    if isinstance(block, IconBlock):
        side = int(
            min(
                Inches(block.size_inches or (1.05 if profile == "editorial" else 0.65)),
                width,
                height,
            )
        )
        shape = add_icon(
            slide,
            block.name,
            left=int(left + (width - side) / 2),
            top=int(top + (height - side) / 2),
            size=side,
            color=block.color or theme.accent,
            guard=guard,
        )
        if block.alt_text:
            set_accessibility(shape, title=block.name, description=block.alt_text)
        return shape
    if isinstance(block, MetricBlock):
        accent = block.accent or theme.accent
        paragraphs = [
            RichParagraph(
                runs=(
                    RichTextRun(
                        block.value,
                        size=38 if profile == "editorial" else 28,
                        color=accent,
                        bold=True,
                    ),
                ),
                space_after_pt=4,
            ),
            RichParagraph(
                runs=(
                    RichTextRun(
                        block.label,
                        size=int(_validated_size(None, profile)),
                        color=theme.ink,
                        bold=True,
                    ),
                ),
                space_after_pt=4,
            ),
        ]
        if block.context:
            paragraphs.append(
                RichParagraph(
                    runs=(
                        RichTextRun(
                            block.context,
                            size=max(int(_validated_size(None, profile) - 2), 8),
                            color=theme.muted,
                        ),
                    ),
                    space_after_pt=0,
                )
            )
        return add_rich_text(
            slide,
            tuple(paragraphs),
            left=left,
            top=top,
            width=width,
            height=height,
            theme=theme,
            guard=guard,
        )
    if isinstance(block, QuoteBlock):
        quote_size = 30 if profile == "editorial" else 22
        paragraphs = [
            RichParagraph(
                runs=(
                    RichTextRun(
                        f"“{block.quote}”", size=quote_size, color=theme.ink, bold=True
                    ),
                ),
                space_after_pt=14,
            )
        ]
        if block.attribution:
            paragraphs.append(
                RichParagraph(
                    runs=(
                        RichTextRun(
                            block.attribution,
                            size=int(_validated_size(None, profile)),
                            color=theme.accent,
                            bold=True,
                        ),
                    ),
                    space_after_pt=0,
                )
            )
        return add_rich_text(
            slide,
            tuple(paragraphs),
            left=left,
            top=top,
            width=width,
            height=height,
            theme=theme,
            guard=guard,
        )
    raise TypeError(f"Unsupported slide block: {type(block).__name__}")


def build_presentation(
    specification: PresentationSpec | Mapping[str, Any],
    output: Path,
    *,
    asset_root: Path | None = None,
    render_dir: Path | None = None,
) -> PptxBuildResult:
    """Compile, save, structurally inspect, and optionally render a PPTX."""

    spec = (
        specification
        if isinstance(specification, PresentationSpec)
        else PresentationSpec.model_validate(specification)
    )
    output = output.expanduser().resolve()
    if output.suffix.lower() != ".pptx":
        raise ValueError("PowerPoint output must use the .pptx extension")
    output.parent.mkdir(parents=True, exist_ok=True)
    assets = (asset_root or output.parent).expanduser().resolve()
    theme = spec.theme.to_native()
    presentation = new_wide_presentation()
    blank_layout = presentation.slide_layouts[6]
    ledger = QualityLedger()

    for slide_number, slide_spec in enumerate(spec.slides, start=1):
        error_start = len(ledger.errors)
        warning_start = len(ledger.warnings)
        slide = presentation.slides.add_slide(blank_layout)
        profile = slide_spec.active_profile
        apply_layout_profile(slide, profile)
        set_background(slide, theme.background)
        plan = layout_plan(
            presentation, slide_spec.layout, theme=theme, profile=profile
        )
        guard = LayoutGuard(plan, ledger=ledger)

        if slide_spec.layout == "cover":
            _render_cover(
                slide,
                slide_spec,
                plan=plan,
                guard=guard,
                theme=theme,
                asset_root=assets,
            )
        else:
            add_title(
                slide,
                slide_spec.title,
                theme=theme,
                kicker=slide_spec.kicker,
                profile=profile,
                guard=guard,
            )
            for slot_name, block in slide_spec.slots.items():
                region = plan.region(slot_name)
                _render_block(
                    slide,
                    block,
                    left=region.left,
                    top=region.top,
                    width=region.width,
                    height=region.height,
                    profile=profile,
                    theme=theme,
                    guard=guard,
                    asset_root=assets,
                )

        add_footer(
            slide,
            slide_number=slide_number,
            source=slide_spec.footer,
            theme=theme,
        )
        if slide_spec.sources:
            add_speaker_source_note(slide, slide_spec.sources)
        if slide_spec.transition != "none":
            set_slide_transition(slide, slide_spec.transition)
        ledger.errors[error_start:] = [
            f"Slide {slide_number}: {message}"
            for message in ledger.errors[error_start:]
        ]
        ledger.warnings[warning_start:] = [
            f"Slide {slide_number}: {message}"
            for message in ledger.warnings[warning_start:]
        ]

    ledger.raise_if_errors()
    presentation.core_properties.title = spec.title
    presentation.core_properties.subject = "Generated by EvoOffice PPTX Engine"
    presentation.core_properties.keywords = "EvoFlux, EvoOffice, editable PowerPoint"
    presentation.save(str(output))
    _normalize_chart_axis_ids(output)

    report = inspect_pptx(output, allow_shape_only=spec.allow_shape_only)
    render = render_office_images(output, render_dir.resolve()) if render_dir else None
    if render and render.get("status") == "rendered":
        report["errors"].extend(render.get("errors", []))
        report["warnings"].extend(render.get("warnings", []))
    return PptxBuildResult(output=output, report=report, render=render)


def validate_presentation(
    source: Path,
    *,
    render_dir: Path | None = None,
    reference: Path | None = None,
    allow_shape_only: bool = False,
) -> dict[str, Any]:
    """Validate an existing PPTX through the same structural and visual gates."""

    source = source.expanduser().resolve()
    report = inspect_pptx(
        source,
        reference=reference.expanduser().resolve() if reference else None,
        allow_shape_only=allow_shape_only,
    )
    if render_dir:
        report["render"] = render_office_images(source, render_dir.resolve())
        rendered = report["render"]
        if rendered.get("status") == "rendered":
            report["errors"].extend(rendered.get("errors", []))
            report["warnings"].extend(rendered.get("warnings", []))
    return report


__all__ = [
    "LAYOUT_SLOTS",
    "PPTX_CAPABILITIES",
    "PptxBuildResult",
    "PresentationSpec",
    "build_presentation",
    "layout_catalog",
    "validate_presentation",
]
