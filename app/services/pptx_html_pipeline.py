"""HTML/Tailwind-first PowerPoint authoring with selective native overlays."""

from __future__ import annotations

import asyncio
import base64
import binascii
from dataclasses import dataclass, field
import json
from pathlib import Path
import re
from typing import Any, Literal

from PIL import Image, ImageStat
from pydantic import BaseModel, ConfigDict, Field, field_validator, model_validator

from app.services.html_slide_render_service import get_html_slide_render_broker

MAX_SLIDES = 80
MAX_ASSETS_PER_SLIDE = 80
MAX_HTML_BYTES = 2_000_000
MAX_CSS_BYTES = 2_000_000
MAX_ASSET_BYTES = 20_000_000
MAX_TOTAL_ASSET_BYTES = 60_000_000
_PLACEHOLDER = re.compile(r"\b(?:lorem ipsum|todo|tbd|click to add)\b", re.I)
_UNSAFE_HTML = re.compile(
    r"<\s*/?\s*(?:script|iframe|object|embed|video|audio|canvas|form|input|button|textarea|select|html|head|body|meta|link|base|style|template)\b",
    re.I,
)
_EVENT_HANDLER = re.compile(r"(?:\s|/)on[a-z0-9_-]+\s*=", re.I)
_UNSAFE_URL = re.compile(
    r"(?:javascript:|vbscript:|data:text/html|https?://|(?<!asset:)//)", re.I
)
_UNSAFE_CSS = re.compile(
    r"(?:</?style\b|@import|expression\s*\(|-moz-binding|behavior\s*:)", re.I
)
_ASSET_TOKEN = re.compile(r"asset://([A-Za-z][A-Za-z0-9_-]{0,63})")
_ASSET_SUFFIXES = {
    ".png",
    ".jpg",
    ".jpeg",
    ".gif",
    ".webp",
    ".svg",
    ".woff",
    ".woff2",
    ".ttf",
    ".otf",
}
_MIME_TYPES = {
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".gif": "image/gif",
    ".webp": "image/webp",
    ".svg": "image/svg+xml",
    ".woff": "font/woff",
    ".woff2": "font/woff2",
    ".ttf": "font/ttf",
    ".otf": "font/otf",
}


class HtmlSlidePlan(BaseModel):
    model_config = ConfigDict(extra="forbid")

    id: str = Field(pattern=r"^[A-Za-z][A-Za-z0-9_-]{0,79}$")
    html_path: str = Field(min_length=1, max_length=2000)
    style_paths: list[str] = Field(default_factory=list, max_length=8)
    assets: dict[str, str] = Field(default_factory=dict)
    speaker_notes: str | None = Field(default=None, max_length=40_000)

    @field_validator("assets")
    @classmethod
    def validate_assets(cls, value: dict[str, str]) -> dict[str, str]:
        if len(value) > MAX_ASSETS_PER_SLIDE:
            raise ValueError(f"a slide may reference at most {MAX_ASSETS_PER_SLIDE} assets")
        for key, path in value.items():
            if not re.fullmatch(r"[A-Za-z][A-Za-z0-9_-]{0,63}", key):
                raise ValueError(f"invalid asset key: {key}")
            if not path or len(path) > 2000:
                raise ValueError(f"invalid asset path for {key}")
        return value


class HtmlPptxProject(BaseModel):
    model_config = ConfigDict(extra="forbid")

    schema_version: Literal[4] = 4
    title: str = Field(min_length=1, max_length=240)
    width: int = Field(default=1280, ge=640, le=3840)
    height: int = Field(default=720, ge=360, le=2160)
    slides: list[HtmlSlidePlan] = Field(min_length=1, max_length=MAX_SLIDES)

    @model_validator(mode="after")
    def unique_slides(self) -> HtmlPptxProject:
        ids = [slide.id for slide in self.slides]
        if len(ids) != len(set(ids)):
            raise ValueError("slide ids must be unique")
        return self


@dataclass
class HtmlPptxPipelineResult:
    action: str
    work_dir: Path
    output: Path | None = None
    previews: list[Path] = field(default_factory=list)
    layout_paths: list[Path] = field(default_factory=list)
    manifest_path: Path | None = None
    issues: list[dict[str, Any]] = field(default_factory=list)
    metadata: dict[str, Any] = field(default_factory=dict)

    @property
    def passed(self) -> bool:
        return not any(issue.get("severity") == "error" for issue in self.issues)


def load_html_pptx_project(path: Path) -> HtmlPptxProject:
    return HtmlPptxProject.model_validate_json(path.read_text(encoding="utf-8"))


def _project_file(project_dir: Path, value: str, *, label: str) -> Path:
    candidate = (project_dir / value).resolve(strict=False)
    try:
        candidate.relative_to(project_dir)
    except (ValueError, binascii.Error) as exc:
        raise ValueError(f"{label} must stay inside the project directory") from exc
    if not candidate.is_file():
        raise FileNotFoundError(f"{label} does not exist: {candidate}")
    return candidate


def _read_limited(path: Path, limit: int, *, label: str) -> str:
    if path.stat().st_size > limit:
        raise ValueError(f"{label} exceeds {limit} bytes")
    return path.read_text(encoding="utf-8")


def _validate_markup(html: str, css: str, *, slide_id: str) -> None:
    if _UNSAFE_HTML.search(html):
        raise ValueError(f"slide {slide_id} contains an unsafe HTML element")
    if _EVENT_HANDLER.search(html):
        raise ValueError(f"slide {slide_id} contains an inline event handler")
    if _UNSAFE_URL.search(html) or _UNSAFE_URL.search(css):
        raise ValueError(f"slide {slide_id} contains a network or executable URL")
    if _UNSAFE_CSS.search(css):
        raise ValueError(f"slide {slide_id} contains unsafe CSS")
    if "data-slide-root" not in html:
        raise ValueError(f"slide {slide_id} must contain one data-slide-root element")


def _slide_sources(
    slide: HtmlSlidePlan, project_dir: Path
) -> tuple[str, str, dict[str, Path]]:
    html_path = _project_file(
        project_dir, slide.html_path, label=f"slide {slide.id} html_path"
    )
    html = _read_limited(html_path, MAX_HTML_BYTES, label=f"slide {slide.id} HTML")
    css_parts: list[str] = []
    for value in slide.style_paths:
        path = _project_file(
            project_dir, value, label=f"slide {slide.id} style path"
        )
        css_parts.append(
            _read_limited(path, MAX_CSS_BYTES, label=f"slide {slide.id} CSS")
        )
    css = "\n".join(css_parts)
    if len(css.encode("utf-8")) > MAX_CSS_BYTES:
        raise ValueError(f"slide {slide.id} CSS exceeds {MAX_CSS_BYTES} bytes")
    assets: dict[str, Path] = {}
    total_asset_bytes = 0
    for key, value in slide.assets.items():
        path = _project_file(
            project_dir, value, label=f"slide {slide.id} asset {key}"
        )
        if path.suffix.lower() not in _ASSET_SUFFIXES:
            raise ValueError(f"unsupported slide asset type: {path.suffix}")
        if path.stat().st_size > MAX_ASSET_BYTES:
            raise ValueError(f"slide asset exceeds {MAX_ASSET_BYTES} bytes: {path.name}")
        total_asset_bytes += path.stat().st_size
        if total_asset_bytes > MAX_TOTAL_ASSET_BYTES:
            raise ValueError(
                f"slide {slide.id} assets exceed {MAX_TOTAL_ASSET_BYTES} bytes"
            )
        assets[key] = path
    _validate_markup(html, css, slide_id=slide.id)
    referenced = set(_ASSET_TOKEN.findall(html + "\n" + css))
    missing = sorted(referenced - set(assets))
    if missing:
        raise ValueError(
            f"slide {slide.id} references undeclared assets: {', '.join(missing)}"
        )
    return html, css, assets


def validate_html_pptx_project(
    project: HtmlPptxProject, project_path: Path
) -> dict[str, Any]:
    project_dir = project_path.parent.resolve()
    placeholder_slides: list[str] = []
    for slide in project.slides:
        html, _css, _assets = _slide_sources(slide, project_dir)
        if _PLACEHOLDER.search(html):
            placeholder_slides.append(slide.id)
    if placeholder_slides:
        raise ValueError(
            "unresolved placeholder text in slides: " + ", ".join(placeholder_slides)
        )
    return {
        "valid": True,
        "schema_version": project.schema_version,
        "slide_count": len(project.slides),
        "canvas": {"width": project.width, "height": project.height, "unit": "px"},
        "rendering": "desktop-webview-html-tailwind",
        "editable_kinds": ["text", "image"],
    }


def html_pptx_catalog() -> dict[str, Any]:
    return {
        "workflow": "html-tailwind-hybrid-pptx",
        "schema_version": 4,
        "canvas": {"width": 1280, "height": 720, "unit": "CSS px"},
        "project": {
            "required": ["schema_version", "title", "slides"],
            "slide_required": ["id", "html_path"],
            "slide_optional": ["style_paths", "assets", "speaker_notes"],
        },
        "html_contract": {
            "root": "exactly one element with data-slide-root",
            "asset_urls": "asset://<declared-key>",
            "editable": {
                "text": 'data-pptx-editable="text"',
                "image": (
                    'data-pptx-editable="image" data-pptx-asset="<declared-key>"'
                ),
            },
            "scripts": False,
            "network": False,
        },
        "runtime": {
            "renderer": "connected EvoFlux desktop WebView",
            "tailwind": "build-time curated utility runtime plus project-local CSS",
            "headless_fallback": False,
        },
        "editability": (
            "HTML is the visual source of truth. Only explicitly marked simple text "
            "and raster images become native PowerPoint objects; all other styling "
            "is flattened into the slide background."
        ),
    }


def _inline_assets(
    html: str, css: str, assets: dict[str, Path]
) -> tuple[str, str, dict[str, dict[str, str]]]:
    metadata: dict[str, dict[str, str]] = {}
    for key, path in assets.items():
        mime = _MIME_TYPES[path.suffix.lower()]
        encoded = base64.b64encode(path.read_bytes()).decode("ascii")
        data_url = f"data:{mime};base64,{encoded}"
        token = f"asset://{key}"
        html = html.replace(token, data_url)
        css = css.replace(token, data_url)
        metadata[key] = {
            "mime_type": mime,
            "suffix": path.suffix.lower(),
        }
    return html, css, metadata


def _decode_png(value: str, destination: Path, *, label: str) -> Path:
    try:
        data = base64.b64decode(value, validate=True)
    except ValueError as exc:
        raise ValueError(f"{label} is not valid base64") from exc
    if len(data) > 40_000_000:
        raise ValueError(f"{label} exceeds 40 MB")
    destination.parent.mkdir(parents=True, exist_ok=True)
    destination.write_bytes(data)
    try:
        with Image.open(destination) as image:
            if image.format != "PNG":
                raise ValueError(f"{label} must be PNG")
            image.verify()
    except Exception:
        destination.unlink(missing_ok=True)
        raise
    return destination


def _preview_quality(path: Path, *, width: int, height: int) -> dict[str, Any]:
    with Image.open(path).convert("RGB") as image:
        actual_width, actual_height = image.size
        dimensions_ok = image.size == (width, height)
        grayscale = image.convert("L")
        entropy = float(grayscale.entropy())
        extrema = grayscale.getextrema()
        low, high = extrema
        if isinstance(low, tuple) or isinstance(high, tuple):
            raise ValueError("preview grayscale extrema are invalid")
        spread = int(high) - int(low)
        variance = float(ImageStat.Stat(grayscale).var[0])
    passed = dimensions_ok and entropy >= 0.35 and spread >= 8 and variance >= 4.0
    return {
        "passed": passed,
        "width": width,
        "height": height,
        "actual_width": actual_width,
        "actual_height": actual_height,
        "entropy": round(entropy, 4),
        "luminance_spread": spread,
        "luminance_variance": round(variance, 4),
    }


def _require_png_dimensions(
    path: Path, *, width: int, height: int, label: str
) -> None:
    with Image.open(path) as image:
        if image.size != (width, height):
            raise ValueError(
                f"{label} must be {width}x{height}, got {image.width}x{image.height}"
            )


def _rgb(value: str, fallback: str = "#111827") -> Any:
    from pptx.dml.color import RGBColor

    candidate = value.strip().lstrip("#")
    if not re.fullmatch(r"[0-9A-Fa-f]{6}", candidate):
        candidate = fallback.lstrip("#")
    return RGBColor.from_string(candidate.upper())


def _add_text_overlay(slide: Any, element: dict[str, Any], *, scale: float) -> None:
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    x = float(element["x"])
    y = float(element["y"])
    width = float(element["width"])
    height = float(element["height"])
    shape = slide.shapes.add_textbox(
        Inches(x / scale),
        Inches(y / scale),
        Inches(width / scale),
        Inches(height / scale),
    )
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = {
        "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }.get(str(element.get("vertical_align")), MSO_ANCHOR.TOP)
    paragraph = frame.paragraphs[0]
    paragraph.text = str(element.get("text") or "")
    paragraph.alignment = {
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }.get(str(element.get("text_align")), PP_ALIGN.LEFT)
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)
    paragraph.line_spacing = float(element.get("line_height_ratio") or 1.0)
    for run in paragraph.runs:
        run.font.name = str(element.get("font_family") or "Arial").split(",")[0].strip(
            " '\""
        )
        run.font.size = Pt(float(element.get("font_size") or 24) * 72 / 96)
        run.font.bold = bool(element.get("bold"))
        run.font.italic = bool(element.get("italic"))
        run.font.underline = bool(element.get("underline"))
        run.font.color.rgb = _rgb(str(element.get("color") or "#111827"))
    shape.rotation = float(element.get("rotation") or 0)
    shape.name = str(element.get("name") or "Editable text")[:160]


def _add_image_overlay(
    slide: Any,
    element: dict[str, Any],
    assets: dict[str, Path],
    *,
    scale: float,
) -> None:
    from pptx.util import Inches

    asset_id = str(element.get("asset_id") or "")
    source = assets.get(asset_id)
    if source is None or source.suffix.lower() not in {
        ".png",
        ".jpg",
        ".jpeg",
        ".gif",
        ".webp",
    }:
        raise ValueError(f"editable image asset is not a supported raster image: {asset_id}")
    x = float(element["x"])
    y = float(element["y"])
    width = float(element["width"])
    height = float(element["height"])
    picture = slide.shapes.add_picture(
        str(source),
        Inches(x / scale),
        Inches(y / scale),
        Inches(width / scale),
        Inches(height / scale),
    )
    picture.name = str(element.get("name") or source.stem)[:160]
    descr = str(element.get("alt") or source.stem)
    picture._element.nvPicPr.cNvPr.set("descr", descr)


def _compose_pptx(
    project: HtmlPptxProject,
    rendered: list[dict[str, Any]],
    slide_assets: list[dict[str, Path]],
    output: Path,
) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    presentation.slide_width = Inches(project.width / 96)
    presentation.slide_height = Inches(project.height / 96)
    for plan, item, assets in zip(
        project.slides, rendered, slide_assets, strict=True
    ):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        slide.shapes.add_picture(
            str(item["shell_path"]),
            0,
            0,
            presentation.slide_width,
            presentation.slide_height,
        )
        for element in item["editable_elements"]:
            kind = str(element.get("kind"))
            if kind == "text":
                _add_text_overlay(slide, element, scale=96)
            elif kind == "image":
                _add_image_overlay(slide, element, assets, scale=96)
        if plan.speaker_notes:
            slide.notes_slide.notes_text_frame.text = plan.speaker_notes
    output.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(output))
    # Structural round-trip check. Visual QA is the immutable HTML preview.
    check = Presentation(str(output))
    if len(check.slides) != len(project.slides):
        raise ValueError("PPTX round-trip slide count changed")


async def compose_html_pptx_project(
    project_path: Path,
    output: Path,
    *,
    workspace_root: Path,
    work_dir: Path,
    session_id: str | None,
) -> HtmlPptxPipelineResult:
    del workspace_root
    project = load_html_pptx_project(project_path)
    validate_html_pptx_project(project, project_path)
    project_dir = project_path.parent.resolve()
    work_dir.mkdir(parents=True, exist_ok=True)
    broker = get_html_slide_render_broker()
    issues: list[dict[str, Any]] = []
    rendered: list[dict[str, Any]] = []
    previews: list[Path] = []
    layouts: list[Path] = []
    all_assets: list[dict[str, Path]] = []
    quality: list[dict[str, Any]] = []
    editable_count = 0
    for slide_number, slide in enumerate(project.slides, start=1):
        html, css, assets = _slide_sources(slide, project_dir)
        inlined_html, inlined_css, asset_metadata = _inline_assets(html, css, assets)
        response = await broker.request(
            session_id,
            {
                "slide_id": slide.id,
                "width": project.width,
                "height": project.height,
                "html": inlined_html,
                "css": inlined_css,
                "assets": asset_metadata,
            },
        )
        preview = _decode_png(
            str(response.get("preview_png_base64") or ""),
            work_dir / "previews" / f"slide-{slide_number:03d}.png",
            label=f"slide {slide_number} preview",
        )
        shell = _decode_png(
            str(response.get("shell_png_base64") or ""),
            work_dir / "shells" / f"slide-{slide_number:03d}.png",
            label=f"slide {slide_number} shell",
        )
        _require_png_dimensions(
            shell,
            width=project.width * 2,
            height=project.height * 2,
            label=f"slide {slide_number} shell",
        )
        metric = _preview_quality(preview, width=project.width, height=project.height)
        metric.update({"slide": slide_number, "slide_id": slide.id})
        quality.append(metric)
        if not metric["passed"]:
            issues.append(
                {
                    "severity": "error",
                    "code": "invalid-html-slide-preview",
                    "message": f"Slide {slide_number} preview is blank or malformed",
                    "slide": slide_number,
                    "quality": metric,
                }
            )
        for value in response.get("issues", []):
            if isinstance(value, dict):
                issues.append({**value, "slide": slide_number, "slide_id": slide.id})
        elements = response.get("editable_elements", [])
        if not isinstance(elements, list):
            raise ValueError("renderer editable_elements must be a list")
        editable_count += len(elements)
        layout_path = work_dir / "layouts" / f"slide-{slide_number:03d}.json"
        layout_path.parent.mkdir(parents=True, exist_ok=True)
        layout_path.write_text(
            json.dumps(
                {
                    "slide": slide_number,
                    "slide_id": slide.id,
                    "editable_elements": elements,
                },
                ensure_ascii=False,
                indent=2,
            )
            + "\n",
            encoding="utf-8",
        )
        rendered.append(
            {
                "preview_path": preview,
                "shell_path": shell,
                "editable_elements": elements,
            }
        )
        previews.append(preview)
        layouts.append(layout_path)
        all_assets.append(assets)
    candidate: Path | None = output
    if not any(issue.get("severity") == "error" for issue in issues):
        await asyncio.to_thread(
            _compose_pptx, project, rendered, all_assets, output
        )
    else:
        output.unlink(missing_ok=True)
        candidate = None
    manifest_path = work_dir / "html-pptx-manifest.json"
    manifest = {
        "schemaVersion": 4,
        "engine": "evoflux-html-tailwind-hybrid",
        "slideCount": len(project.slides),
        "editableObjectCount": editable_count,
        "previewQuality": quality,
        "visualSource": "desktop-webview",
        "roundTrip": "structural-openxml",
    }
    manifest_path.write_text(json.dumps(manifest, indent=2) + "\n", encoding="utf-8")
    return HtmlPptxPipelineResult(
        action="compose",
        work_dir=work_dir,
        output=candidate,
        previews=previews,
        layout_paths=layouts,
        manifest_path=manifest_path,
        issues=issues,
        metadata={
            "engine": "evoflux-html-tailwind-hybrid",
            "slide_count": len(project.slides),
            "editable_object_count": editable_count,
            "semantic_editable_object_count": editable_count,
            "preview_quality": quality,
            "visual_verification": "html-source-preview",
            "round_trip_verification": "structural-openxml",
        },
    )


__all__ = [
    "HtmlPptxPipelineResult",
    "HtmlPptxProject",
    "HtmlSlidePlan",
    "compose_html_pptx_project",
    "html_pptx_catalog",
    "load_html_pptx_project",
    "validate_html_pptx_project",
]
