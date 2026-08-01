"""Cached, sandbox-safe previews for OpenXML workspace documents.

The renderer is intentionally self-contained: DOCX, XLSX, and PPTX are parsed
with the Python OpenXML libraries already bundled in the EvoFlux sidecar.  It
does not execute document macros, launch an office application, or depend on an
external renderer binary.
"""

from __future__ import annotations

import base64
import hashlib
import html
import math
import os
import threading
from pathlib import Path
from typing import Any

from loguru import logger

from app.core.config import settings

SUPPORTED_OFFICE_PREVIEW_EXTENSIONS = frozenset({".docx", ".xlsx", ".pptx"})
MAX_OFFICE_PREVIEW_BYTES = 100 * 1024 * 1024
OFFICE_PREVIEW_CSP = (
    "default-src 'none'; "
    "base-uri 'none'; "
    "connect-src 'none'; "
    "font-src data:; "
    "form-action 'none'; "
    "frame-src 'none'; "
    "img-src data: blob:; "
    "media-src data: blob:; "
    "object-src 'none'; "
    "script-src 'none'; "
    "style-src 'unsafe-inline'"
)

_CACHE_SCHEMA_VERSION = "python-openxml-html-v7"
_render_lock = threading.Lock()


class OfficePreviewError(RuntimeError):
    """Base class for office preview failures."""


class OfficePreviewUnsupportedError(OfficePreviewError):
    """Raised when the requested file cannot be rendered safely."""


def _cache_path(source: Path) -> Path:
    stat = source.stat()
    fingerprint = "\0".join(
        (
            _CACHE_SCHEMA_VERSION,
            str(source.resolve()),
            str(stat.st_size),
            str(stat.st_mtime_ns),
        )
    )
    digest = hashlib.sha256(fingerprint.encode()).hexdigest()
    return Path(settings.EVOFLUX_CACHE_DIR) / "office-previews" / f"{digest}.html"


def _page(*, title: str, body: str, css: str) -> str:
    policy = html.escape(OFFICE_PREVIEW_CSP, quote=True)
    safe_title = html.escape(title)
    return (
        '<!doctype html><html><head><meta charset="utf-8">'
        f'<meta http-equiv="Content-Security-Policy" content="{policy}">'
        '<meta name="referrer" content="no-referrer">'
        f"<title>{safe_title}</title><style>{css}</style></head>"
        f"<body>{body}</body></html>"
    )


def _css_color(value: Any, default: str = "transparent") -> str:
    """Convert an OpenXML color object or raw ARGB string to CSS."""
    if value is None:
        return default
    try:
        raw = getattr(value, "rgb", value)
    except (AttributeError, TypeError, ValueError):
        return default
    if raw is None:
        return default
    text = str(raw)
    if len(text) == 8:
        text = text[2:]
    if len(text) == 6 and all(ch in "0123456789abcdefABCDEF" for ch in text):
        return f"#{text}"
    return default


def _css_font_family(name: Any) -> str:
    family = str(name).replace("\\", "\\\\").replace('"', '\\"')
    fallback = (
        '"Helvetica Neue",Arial,sans-serif'
        if family.casefold() == "aptos display"
        else "Arial,sans-serif"
    )
    return html.escape(f'"{family}",{fallback}', quote=True)


def _render_docx(source: Path) -> str:
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    document = Document(str(source))
    section = document.sections[0]

    def pixels(length: Any, default_inches: float) -> float:
        return float(length.inches) * 96 if length is not None else default_inches * 96

    page_width = pixels(section.page_width, 8.5)
    page_height = pixels(section.page_height, 11)
    margin_top = pixels(section.top_margin, 1)
    margin_right = pixels(section.right_margin, 1)
    margin_bottom = pixels(section.bottom_margin, 1)
    margin_left = pixels(section.left_margin, 1)
    alignment = {
        WD_ALIGN_PARAGRAPH.CENTER: "center",
        WD_ALIGN_PARAGRAPH.RIGHT: "right",
        WD_ALIGN_PARAGRAPH.JUSTIFY: "justify",
        WD_ALIGN_PARAGRAPH.DISTRIBUTE: "justify",
    }
    list_number = 0

    def render_paragraph(paragraph: Paragraph) -> str:
        nonlocal list_number
        style_name = (paragraph.style.name if paragraph.style else "").lower()
        tag = "p"
        if style_name == "title":
            tag = "h1"
        for level in range(1, 4):
            if style_name.startswith(f"heading {level}"):
                tag = f"h{level + 1}"
                break
        is_numbered = style_name.startswith("list number")
        is_bullet = style_name.startswith("list bullet")
        if is_numbered:
            list_number += 1
        elif not is_bullet:
            list_number = 0

        runs: list[str] = []
        for run in paragraph.runs:
            styles: list[str] = []
            if run.bold:
                styles.append("font-weight:700")
            if run.italic:
                styles.append("font-style:italic")
            if run.underline:
                styles.append("text-decoration:underline")
            if run.font.size:
                styles.append(f"font-size:{run.font.size.pt:.2f}pt")
            if run.font.name:
                styles.append(f"font-family:{_css_font_family(run.font.name)}")
            if run.font.color and run.font.color.rgb:
                styles.append(f"color:#{run.font.color.rgb}")
            text = html.escape(run.text).replace("\n", "<br>")
            style_attr = f' style="{";".join(styles)}"' if styles else ""
            runs.append(f"<span{style_attr}>{text}</span>")
            for blip in run._element.xpath(".//a:blip"):
                relation_id = blip.get(qn("r:embed"))
                related = paragraph.part.related_parts.get(relation_id)
                if related is None or not hasattr(related, "blob"):
                    continue
                content_type = getattr(related, "content_type", "image/png")
                encoded = base64.b64encode(related.blob).decode("ascii")
                extent = run._element.xpath(".//wp:extent")
                image_styles = ["max-width:100%", "height:auto"]
                if extent:
                    width_px = int(extent[0].get("cx", "0")) / 914400 * 96
                    if width_px > 0:
                        image_styles.append(f"width:{width_px:.2f}px")
                runs.append(
                    f'<img class="inline-image" style="{";".join(image_styles)}" '
                    f'src="data:{content_type};base64,{encoded}" alt="">'
                )
        content = "".join(runs) or html.escape(paragraph.text)
        if not content:
            content = "&nbsp;"
        paragraph_styles: list[str] = []
        style = paragraph.style
        if style is not None:
            if style.font.name:
                paragraph_styles.append(
                    f"font-family:{_css_font_family(style.font.name)}"
                )
            if style.font.size:
                paragraph_styles.append(f"font-size:{style.font.size.pt:.2f}pt")
            if style.font.bold:
                paragraph_styles.append("font-weight:700")
            if style.font.color and style.font.color.rgb:
                paragraph_styles.append(f"color:#{style.font.color.rgb}")
        if paragraph.alignment in alignment:
            paragraph_styles.append(f"text-align:{alignment[paragraph.alignment]}")
        formatting = paragraph.paragraph_format
        style_formatting = style.paragraph_format if style is not None else None
        space_before = formatting.space_before or (
            style_formatting.space_before if style_formatting is not None else None
        )
        space_after = formatting.space_after or (
            style_formatting.space_after if style_formatting is not None else None
        )
        line_spacing = formatting.line_spacing or (
            style_formatting.line_spacing if style_formatting is not None else None
        )
        if space_before:
            paragraph_styles.append(f"margin-top:{space_before.pt:.2f}pt")
        if space_after:
            paragraph_styles.append(f"margin-bottom:{space_after.pt:.2f}pt")
        if isinstance(line_spacing, (int, float)):
            paragraph_styles.append(f"line-height:{float(line_spacing):.3f}")
        elif line_spacing is not None:
            paragraph_styles.append(f"line-height:{line_spacing.pt:.2f}pt")
        if formatting.left_indent:
            paragraph_styles.append(f"margin-left:{formatting.left_indent.pt:.2f}pt")
        if formatting.right_indent:
            paragraph_styles.append(f"margin-right:{formatting.right_indent.pt:.2f}pt")
        if formatting.first_line_indent:
            paragraph_styles.append(
                f"text-indent:{formatting.first_line_indent.pt:.2f}pt"
            )
        style_attr = (
            f' style="{";".join(paragraph_styles)}"' if paragraph_styles else ""
        )
        page_break = (
            "true"
            if formatting.page_break_before
            or bool(paragraph._element.xpath(".//w:br[@w:type='page']"))
            else "false"
        )
        label = html.escape(paragraph.text[:80], quote=True)
        if is_numbered or is_bullet:
            marker = f"{list_number}." if is_numbered else "•"
            return (
                f'<div class="list-item"{style_attr} '
                f'data-page-break-before="{page_break}" data-qa-label="{label}">'
                f'<span class="list-marker">{marker}</span>'
                f'<span class="list-content">{content}</span></div>'
            )
        return (
            f'<{tag}{style_attr} data-page-break-before="{page_break}" '
            f'data-qa-label="{label}">{content}</{tag}>'
        )

    def render_table(table: Table) -> str:
        rows: list[str] = []
        for row_index, row in enumerate(table.rows):
            cells: list[str] = []
            cell_tag = "th" if row_index == 0 else "td"
            for cell in row.cells:
                cell_body = "".join(render_paragraph(p) for p in cell.paragraphs)
                cells.append(f"<{cell_tag}>{cell_body}</{cell_tag}>")
            rows.append(f"<tr>{''.join(cells)}</tr>")
        return f"<table>{''.join(rows)}</table>"

    blocks: list[str] = []
    for child in document.element.body.iterchildren():
        if isinstance(child, CT_P):
            blocks.append(render_paragraph(Paragraph(child, document)))
        elif isinstance(child, CT_Tbl):
            blocks.append(render_table(Table(child, document)))

    header = "".join(render_paragraph(p) for p in section.header.paragraphs)
    footer = "".join(render_paragraph(p) for p in section.footer.paragraphs)
    geometry = (
        f"--page-width:{page_width:.2f}px;--page-height:{page_height:.2f}px;"
        f"--margin-top:{margin_top:.2f}px;--margin-right:{margin_right:.2f}px;"
        f"--margin-bottom:{margin_bottom:.2f}px;--margin-left:{margin_left:.2f}px"
    )
    css = """
    *{box-sizing:border-box}body{margin:0;padding:32px;background:#e8eaed;
    color:#202124;font-family:Arial,sans-serif}.document{margin:0 auto}
    .document-flow,.document-page{width:var(--page-width);background:#fff;
    box-shadow:0 2px 12px #0002}.document-flow{min-height:var(--page-height);
    padding:var(--margin-top) var(--margin-right) var(--margin-bottom)
    var(--margin-left)}.document-pages{display:flex;flex-direction:column;gap:24px;
    align-items:center}.document-page{height:var(--page-height);position:relative;
    padding:var(--margin-top) var(--margin-right) var(--margin-bottom)
    var(--margin-left);overflow:hidden}.document-page-body{height:100%;overflow:hidden}
    .document-page-header,.document-page-footer{position:absolute;left:var(--margin-left);
    right:var(--margin-right);font-size:9pt;color:#5f6368;overflow:hidden}
    .document-page-header{top:16px;height:calc(var(--margin-top) - 20px)}
    .document-page-footer{bottom:12px;height:calc(var(--margin-bottom) - 16px)}
    .document-header-template,.document-footer-template{display:none}
    h1{font-size:26pt;line-height:1.15;margin:0 0
    18pt}h2{font-size:18pt;margin:20pt 0 8pt}h3{font-size:14pt;margin:16pt 0
    6pt}h4{font-size:12pt;margin:12pt 0 5pt}p{font-size:11pt;line-height:1.5;
    margin:0 0 8pt;white-space:pre-wrap}.list-item{display:grid;
    grid-template-columns:1.4em 1fr;column-gap:.25em;align-items:start;
    margin:0 0 6pt;white-space:pre-wrap}.list-marker{text-align:right;
    font-weight:600}.list-content{min-width:0}
    table{width:100%;border-collapse:collapse;margin:12pt 0 18pt;table-layout:fixed}
    th,td{border:1px solid #c7c7c7;padding:7px 9px;vertical-align:top;overflow-wrap:anywhere}
    th{background:#f1f3f4;text-align:left}td p,th p{margin:0;font-size:10pt}
    .inline-image{display:inline-block;vertical-align:middle}
    """
    return _page(
        title=source.name,
        body=(
            f'<main class="document" style="{geometry}">'
            f'<div class="document-header-template">{header}</div>'
            f'<div class="document-footer-template">{footer}</div>'
            f'<div class="document-flow">{"".join(blocks)}</div></main>'
        ),
        css=css,
    )


def _worksheet_bounds(sheet: Any) -> tuple[int, int]:
    """Return bounded preview dimensions for a worksheet."""
    max_row = min(max(sheet.max_row, 1), 500)
    max_column = min(max(sheet.max_column, 1), 100)
    return max_row, max_column


def _display_cell(cell: Any) -> str:
    value = cell.value
    if value is None:
        return ""
    if cell.data_type == "f":
        formula = str(value)
        return formula if formula.startswith("=") else f"={formula}"
    if getattr(cell, "is_date", False) and hasattr(value, "strftime"):
        return value.strftime("%Y-%m-%d")
    number_format = str(getattr(cell, "number_format", "") or "")
    if isinstance(value, (int, float)):
        if "%" in number_format:
            decimals = 0
            if "." in number_format:
                decimals = len(number_format.split(".", 1)[1].split("%", 1)[0])
            return f"{value * 100:.{decimals}f}%"
        currency = next(
            (symbol for symbol in ("$", "€", "£", "¥") if symbol in number_format),
            None,
        )
        if currency:
            decimals = 2 if ".00" in number_format else 0
            return f"{currency}{value:,.{decimals}f}"
        if "," in number_format:
            decimals = 2 if ".00" in number_format else 0
            return f"{value:,.{decimals}f}"
    return str(value)


def _border_css(side: Any, name: str) -> str | None:
    if not side or not side.style:
        return None
    color = _css_color(side.color, "#80868b")
    width = "2px" if side.style in {"medium", "thick", "double"} else "1px"
    line = "dashed" if "dash" in side.style else "solid"
    return f"border-{name}:{width} {line} {color}"


def _xlsx_series_values(workbook: Any, series: Any) -> list[float]:
    """Resolve a chart series from its cached values or workbook cell range."""
    from openpyxl.utils.cell import range_boundaries

    data_source = getattr(series, "val", None) or getattr(series, "yVal", None)
    if data_source is None:
        return []
    literal = getattr(data_source, "numLit", None)
    if literal is not None:
        values = []
        for point in getattr(literal, "pt", ()):
            try:
                values.append(float(point.v))
            except (TypeError, ValueError):
                values.append(0.0)
        return values

    reference = getattr(data_source, "numRef", None)
    formula = str(getattr(reference, "f", "") or "")
    if "!" not in formula:
        return []
    sheet_name, cell_range = formula.rsplit("!", 1)
    sheet_name = sheet_name.strip("'").replace("''", "'")
    if sheet_name not in workbook.sheetnames:
        return []
    min_column, min_row, max_column, max_row = range_boundaries(cell_range)
    values: list[float] = []
    sheet = workbook[sheet_name]
    for row in sheet.iter_rows(
        min_row=min_row,
        max_row=max_row,
        min_col=min_column,
        max_col=max_column,
    ):
        for cell in row:
            try:
                values.append(float(cell.value or 0))
            except (TypeError, ValueError):
                values.append(0.0)
    return values


def _xlsx_category_values(workbook: Any, series: Any) -> list[str]:
    """Resolve chart category labels from a worksheet reference."""
    from openpyxl.utils.cell import range_boundaries

    data_source = getattr(series, "cat", None) or getattr(series, "xVal", None)
    if data_source is None:
        return []
    reference = getattr(data_source, "strRef", None) or getattr(
        data_source, "numRef", None
    )
    formula = str(getattr(reference, "f", "") or "")
    if "!" not in formula:
        return []
    sheet_name, cell_range = formula.rsplit("!", 1)
    sheet_name = sheet_name.strip("'").replace("''", "'")
    if sheet_name not in workbook.sheetnames:
        return []
    min_column, min_row, max_column, max_row = range_boundaries(cell_range)
    sheet = workbook[sheet_name]
    return [
        str(cell.value or "")
        for row in sheet.iter_rows(
            min_row=min_row,
            max_row=max_row,
            min_col=min_column,
            max_col=max_column,
        )
        for cell in row
    ]


def _xlsx_chart_svg(workbook: Any, chart: Any) -> str:
    """Render common workbook charts as a deterministic SVG preview."""
    palette = ("#2563eb", "#0f766e", "#f59e0b", "#dc2626", "#7c3aed")
    series_values = [
        values
        for series in chart.ser
        if (values := _xlsx_series_values(workbook, series))
    ]
    if not series_values:
        return '<div class="chart-empty">Chart data unavailable</div>'
    maximum = max(
        (abs(value) for values in series_values for value in values),
        default=1,
    )
    maximum = maximum or 1
    count = max(len(values) for values in series_values)
    chart_type = type(chart).__name__.upper()
    categories = _xlsx_category_values(workbook, chart.ser[0]) if chart.ser else []

    def category_labels() -> str:
        if not categories:
            return ""
        return "".join(
            f'<text x="{52 + index * 520 / max(len(categories) - 1, 1):.2f}" '
            'y="278" text-anchor="middle" font-size="13" fill="#6b7280">'
            f"{html.escape(label)}</text>"
            for index, label in enumerate(categories)
        )

    if "PIE" in chart_type or "DOUGHNUT" in chart_type:
        values = [abs(value) for value in series_values[0]]
        total = sum(values) or 1
        start_angle = -math.pi / 2
        slices: list[str] = []
        for index, value in enumerate(values):
            end_angle = start_angle + 2 * math.pi * value / total
            start_x = 300 + 108 * math.cos(start_angle)
            start_y = 150 + 108 * math.sin(start_angle)
            end_x = 300 + 108 * math.cos(end_angle)
            end_y = 150 + 108 * math.sin(end_angle)
            large_arc = 1 if end_angle - start_angle > math.pi else 0
            slices.append(
                f'<path d="M300 150 L{start_x:.2f} {start_y:.2f} '
                f'A108 108 0 {large_arc} 1 {end_x:.2f} {end_y:.2f} Z" '
                f'fill="{palette[index % len(palette)]}"/>'
            )
            start_angle = end_angle
        hole = (
            '<circle cx="300" cy="150" r="65" fill="white"/>'
            if "DOUGHNUT" in chart_type
            else ""
        )
        return (
            '<svg class="workbook-chart-svg" viewBox="0 0 600 300">'
            f"{''.join(slices)}{hole}</svg>"
        )

    if "AREA" in chart_type or "LINE" in chart_type:
        paths: list[str] = []
        for series_index, values in enumerate(series_values):
            points = [
                (
                    52 + value_index * 520 / max(len(values) - 1, 1),
                    252 - (value / maximum) * 205,
                )
                for value_index, value in enumerate(values)
            ]
            point_text = " ".join(f"{x:.2f},{y:.2f}" for x, y in points)
            color = palette[series_index % len(palette)]
            area = ""
            if "AREA" in chart_type and points:
                polygon = [
                    f"{points[0][0]:.2f},252",
                    *[f"{x:.2f},{y:.2f}" for x, y in points],
                    f"{points[-1][0]:.2f},252",
                ]
                area = (
                    f'<polygon class="workbook-chart-area" '
                    f'points="{" ".join(polygon)}" fill="{color}" '
                    'fill-opacity=".24"/>'
                )
            markers = "".join(
                f'<circle cx="{x:.2f}" cy="{y:.2f}" r="4" fill="{color}"/>'
                for x, y in points
            )
            paths.append(
                area + f'<polyline points="{point_text}" fill="none" '
                f'stroke="{color}" stroke-width="4"/>{markers}'
            )
        return (
            '<svg class="workbook-chart-svg" viewBox="0 0 600 300" '
            'preserveAspectRatio="none">'
            '<line x1="45" y1="252" x2="575" y2="252" stroke="#9ca3af"/>'
            f"{''.join(paths)}{category_labels()}</svg>"
        )

    group_width = 520 / max(count, 1)
    bar_width = max(group_width / max(len(series_values), 1) * 0.72, 2)
    bars: list[str] = []
    for series_index, values in enumerate(series_values):
        for value_index, value in enumerate(values):
            height = abs(value) / maximum * 210
            x = 52 + value_index * group_width + series_index * bar_width
            y = 252 - height if value >= 0 else 252
            bars.append(
                f'<rect x="{x:.2f}" y="{y:.2f}" width="{bar_width:.2f}" '
                f'height="{height:.2f}" fill="{palette[series_index % len(palette)]}"/>'
            )
    return (
        '<svg class="workbook-chart-svg" viewBox="0 0 600 300" '
        'preserveAspectRatio="none">'
        '<line x1="45" y1="252" x2="575" y2="252" stroke="#9ca3af"/>'
        f"{''.join(bars)}</svg>"
    )


def _xlsx_image_data_uri(image: Any) -> str:
    extension = str(getattr(image, "format", "png") or "png").lower()
    mime = {
        "jpg": "image/jpeg",
        "jpeg": "image/jpeg",
        "gif": "image/gif",
        "webp": "image/webp",
        "bmp": "image/bmp",
    }.get(extension, "image/png")
    encoded = base64.b64encode(image._data()).decode("ascii")
    return f"data:{mime};base64,{encoded}"


def _xlsx_column_width_px(sheet: Any, column: int) -> float:
    from openpyxl.utils import get_column_letter

    dimension = sheet.column_dimensions[get_column_letter(column)]
    if dimension.hidden:
        return 1
    return max(1, min((dimension.width or 10) * 7, 280))


def _xlsx_row_height_px(sheet: Any, row: int) -> float:
    dimension = sheet.row_dimensions[row]
    if dimension.hidden:
        return 1
    return (dimension.height or 18.75) * 96 / 72


def _xlsx_anchor_box(sheet: Any, drawing: Any) -> tuple[float, float, float, float]:
    """Resolve an openpyxl drawing anchor to CSS pixels."""
    anchor = drawing.anchor
    marker = anchor._from
    left = 44 + sum(
        _xlsx_column_width_px(sheet, column) for column in range(1, int(marker.col) + 1)
    )
    top = 25 + sum(
        _xlsx_row_height_px(sheet, row) for row in range(1, int(marker.row) + 1)
    )
    left += float(marker.colOff or 0) / 914400 * 96
    top += float(marker.rowOff or 0) / 914400 * 96
    extent = getattr(anchor, "ext", None)
    if extent is not None:
        width = float(extent.cx) / 914400 * 96
        height = float(extent.cy) / 914400 * 96
    else:
        width = float(getattr(drawing, "width", 320) or 320)
        height = float(getattr(drawing, "height", 180) or 180)
    return left, top, max(width, 1), max(height, 1)


def _render_xlsx(source: Path) -> str:
    from openpyxl import load_workbook
    from openpyxl.utils import get_column_letter

    workbook = load_workbook(source, data_only=False, read_only=False)
    sections: list[str] = []
    for sheet in workbook.worksheets:
        max_row, max_column = _worksheet_bounds(sheet)
        merged_starts = {
            (cell_range.min_row, cell_range.min_col): cell_range
            for cell_range in sheet.merged_cells.ranges
        }
        merged_covered: set[tuple[int, int]] = set()
        for cell_range in sheet.merged_cells.ranges:
            for row in range(cell_range.min_row, cell_range.max_row + 1):
                for column in range(cell_range.min_col, cell_range.max_col + 1):
                    if (row, column) != (cell_range.min_row, cell_range.min_col):
                        merged_covered.add((row, column))

        column_widths = [
            _xlsx_column_width_px(sheet, column) for column in range(1, max_column + 1)
        ]
        row_heights = [_xlsx_row_height_px(sheet, row) for row in range(1, max_row + 1)]
        columns = ['<col style="width:44px">']
        column_headers = ['<th class="corner"></th>']
        for column in range(1, max_column + 1):
            width_px = column_widths[column - 1]
            columns.append(f'<col style="width:{width_px:.0f}px">')
            column_headers.append(
                f'<th style="min-width:{width_px:.0f}px">'
                f"{get_column_letter(column)}</th>"
            )

        rows = [f"<tr>{''.join(column_headers)}</tr>"]
        for row_index in range(1, max_row + 1):
            row_height = row_heights[row_index - 1]
            cells = [f'<th class="row-number">{row_index}</th>']
            for column_index in range(1, max_column + 1):
                if (row_index, column_index) in merged_covered:
                    continue
                cell = sheet.cell(row_index, column_index)
                merged = merged_starts.get((row_index, column_index))
                spans = ""
                if merged:
                    spans = (
                        f' rowspan="{merged.max_row - merged.min_row + 1}"'
                        f' colspan="{merged.max_col - merged.min_col + 1}"'
                    )
                styles: list[str] = []
                fill = _css_color(cell.fill.fgColor)
                if fill != "transparent" and cell.fill.fill_type:
                    styles.append(f"background:{fill}")
                if cell.font.bold:
                    styles.append("font-weight:700")
                if cell.font.italic:
                    styles.append("font-style:italic")
                if cell.font.name:
                    styles.append(f"font-family:{_css_font_family(cell.font.name)}")
                if cell.font.sz:
                    styles.append(f"font-size:{float(cell.font.sz):.2f}pt")
                color = _css_color(cell.font.color)
                if color != "transparent":
                    styles.append(f"color:{color}")
                if cell.alignment.horizontal:
                    styles.append(f"text-align:{cell.alignment.horizontal}")
                if cell.alignment.vertical:
                    styles.append(f"vertical-align:{cell.alignment.vertical}")
                if cell.alignment.wrap_text:
                    styles.append("white-space:normal")
                for side_name in ("top", "right", "bottom", "left"):
                    border = _border_css(getattr(cell.border, side_name), side_name)
                    if border:
                        styles.append(border)
                style_attr = f' style="{";".join(styles)}"' if styles else ""
                value = html.escape(_display_cell(cell))
                formula_class = " formula" if cell.data_type == "f" else ""
                cells.append(
                    f'<td class="cell{formula_class}" data-cell="{cell.coordinate}" '
                    f'data-qa-label="{html.escape(cell.coordinate, quote=True)}"'
                    f"{spans}{style_attr}>{value}</td>"
                )
            rows.append(f'<tr style="height:{row_height:.2f}px">{"".join(cells)}</tr>')
        truncated = sheet.max_row > max_row or sheet.max_column > max_column
        notice = (
            '<p class="notice">Preview limited to the first 500 rows and 100 columns.</p>'
            if truncated
            else ""
        )
        objects: list[str] = []
        stage_width = 44 + sum(column_widths)
        # CSS cells have a 25px minimum height.  Compact workbook profiles may
        # request slightly shorter rows, so use the rendered minimum here or
        # the final rows can extend beyond the stage and be falsely clipped.
        stage_height = 32 + sum(max(height, 25) + 1 for height in row_heights)
        for chart_index, chart in enumerate(sheet._charts, start=1):
            title = f"Chart {chart_index}"
            try:
                text = chart.title.tx.rich.p[0].r[0].t
                if text:
                    title = str(text)
            except (AttributeError, IndexError, TypeError):
                pass
            left, top, width, height = _xlsx_anchor_box(sheet, chart)
            stage_width = max(stage_width, left + width)
            stage_height = max(stage_height, top + height)
            objects.append(
                '<figure class="workbook-chart" '
                f'data-qa-label="{html.escape(title, quote=True)}" '
                f'style="left:{left:.2f}px;top:{top:.2f}px;'
                f'width:{width:.2f}px;height:{height:.2f}px">'
                f"<figcaption>{html.escape(title)}</figcaption>"
                f"{_xlsx_chart_svg(workbook, chart)}</figure>"
            )
        for image_index, image in enumerate(sheet._images, start=1):
            left, top, width, height = _xlsx_anchor_box(sheet, image)
            stage_width = max(stage_width, left + width)
            stage_height = max(stage_height, top + height)
            objects.append(
                f'<img class="workbook-image" data-qa-label="Image {image_index}" '
                f'style="left:{left:.2f}px;top:{top:.2f}px;'
                f'width:{width:.2f}px;height:{height:.2f}px" '
                f'src="{_xlsx_image_data_uri(image)}" alt="Image {image_index}">'
            )
        sections.append(
            f'<section class="sheet"><h2>{html.escape(sheet.title)}</h2>{notice}'
            f'<div class="grid-wrap"><div class="grid-stage" '
            f'style="width:{stage_width:.2f}px;height:{stage_height:.2f}px">'
            f"<table><colgroup>{''.join(columns)}</colgroup>"
            f"{''.join(rows)}</table>{''.join(objects)}</div></div></section>"
        )

    css = """
    *{box-sizing:border-box}body{margin:0;padding:24px;background:#f4f5f7;
    color:#202124;font-family:Arial,sans-serif}.sheet{margin:0 auto 28px;background:#fff;
    border:1px solid #d9dce1;border-radius:8px;box-shadow:0 2px 8px #0001;overflow:hidden}
    h2{position:sticky;left:0;margin:0;padding:14px 18px;background:#fff;border-bottom:
    1px solid #e3e6ea;font-size:15px}.notice{margin:0;padding:8px 18px;background:#fff8dc;
    font-size:12px}.grid-wrap{overflow:auto;max-height:76vh}.grid-stage{position:relative}
    table{position:absolute;left:0;top:0;border-collapse:separate;
    border-spacing:0;font-size:12px}th,td{height:25px;padding:3px 6px;border-right:1px
    solid #e1e4e8;border-bottom:1px solid #e1e4e8;white-space:pre;vertical-align:middle}
    th{position:sticky;top:0;z-index:2;background:#f1f3f4;color:#5f6368;text-align:center;
    font-weight:500}.row-number{left:0;z-index:1;min-width:44px}.corner{left:0;z-index:3}
    td{overflow:hidden;text-overflow:ellipsis}
    .formula{color:#174ea6}.cell:empty{background:#fff}
    .workbook-chart{position:absolute;margin:0;border:1px solid #d9dce1;
    background:#fff;padding:10px;overflow:hidden}.workbook-chart figcaption{
    font-weight:600;margin-bottom:5px}.workbook-chart-svg{display:block;width:100%;
    height:calc(100% - 22px)}.workbook-image{position:absolute;object-fit:contain}
    .chart-empty{color:#6b7280}
    """
    return _page(title=source.name, body="".join(sections), css=css)


def _picture_data_uri(shape: Any) -> str:
    relation_id = shape._pic.blip_rId
    related_part = shape.part.related_part(relation_id)
    extension = (related_part.partname.ext or "png").lower()
    mime = {
        "jpg": "image/jpeg",
        "jpeg": "image/jpeg",
        "gif": "image/gif",
        "webp": "image/webp",
        "svg": "image/svg+xml",
    }.get(extension, getattr(related_part, "content_type", "image/png"))
    encoded = base64.b64encode(related_part.blob).decode("ascii")
    return f"data:{mime};base64,{encoded}"


def _shape_box(shape: Any, *, slide_width: int, slide_height: int) -> str:
    left = 100 * shape.left / slide_width
    top = 100 * shape.top / slide_height
    width = 100 * shape.width / slide_width
    height = 100 * shape.height / slide_height
    styles = f"left:{left:.3f}%;top:{top:.3f}%;width:{width:.3f}%;height:{height:.3f}%"
    rotation = float(getattr(shape, "rotation", 0) or 0)
    if rotation:
        styles += f";transform:rotate({rotation:.3f}deg)"
    return styles


def _pptx_theme_palette(slide: Any) -> dict[str, str]:
    """Extract the color scheme attached to the slide's master."""
    from xml.etree import ElementTree

    theme_part = None
    for relation in slide.slide_layout.slide_master.part.rels.values():
        if relation.reltype.endswith("/theme"):
            theme_part = relation.target_part
            break
    if theme_part is None:
        return {}
    root = ElementTree.fromstring(theme_part.blob)
    namespace = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    scheme = root.find(".//a:themeElements/a:clrScheme", namespace)
    if scheme is None:
        return {}
    palette: dict[str, str] = {}
    for entry in scheme:
        if not list(entry):
            continue
        color = list(entry)[0]
        raw = color.get("val") or color.get("lastClr")
        css = _css_color(raw)
        if css != "transparent":
            palette[entry.tag.rsplit("}", 1)[-1]] = css
    palette.update(
        {
            "bg1": palette.get("lt1", "#ffffff"),
            "tx1": palette.get("dk1", "#000000"),
            "bg2": palette.get("lt2", "#e7e6e6"),
            "tx2": palette.get("dk2", "#44546a"),
        }
    )
    return palette


def _apply_brightness(color: str, brightness: float) -> str:
    if not color.startswith("#") or len(color) != 7 or not brightness:
        return color
    channels = [int(color[index : index + 2], 16) for index in (1, 3, 5)]
    if brightness > 0:
        channels = [
            round(channel + (255 - channel) * brightness) for channel in channels
        ]
    else:
        channels = [round(channel * (1 + brightness)) for channel in channels]
    return "#" + "".join(f"{max(0, min(channel, 255)):02x}" for channel in channels)


def _pptx_color(
    color_format: Any,
    palette: dict[str, str],
    default: str = "transparent",
) -> str:
    direct = _css_color(color_format)
    if direct != "transparent":
        color = direct
    else:
        theme_color = getattr(color_format, "theme_color", None)
        theme_name = getattr(theme_color, "name", "")
        theme_key = {
            "DARK_1": "dk1",
            "LIGHT_1": "lt1",
            "DARK_2": "dk2",
            "LIGHT_2": "lt2",
            "ACCENT_1": "accent1",
            "ACCENT_2": "accent2",
            "ACCENT_3": "accent3",
            "ACCENT_4": "accent4",
            "ACCENT_5": "accent5",
            "ACCENT_6": "accent6",
            "HYPERLINK": "hlink",
            "FOLLOWED_HYPERLINK": "folHlink",
            "BACKGROUND_1": "bg1",
            "TEXT_1": "tx1",
            "BACKGROUND_2": "bg2",
            "TEXT_2": "tx2",
        }.get(theme_name)
        color = palette.get(theme_key or "", default)
    try:
        brightness = float(color_format.brightness or 0)
    except (AttributeError, TypeError, ValueError):
        brightness = 0
    return _apply_brightness(color, brightness)


def _slide_background(slide: Any, palette: dict[str, str]) -> str:
    """Resolve an explicit slide/layout/master background color."""
    drawing = "http://schemas.openxmlformats.org/drawingml/2006/main"
    for layer in (slide, slide.slide_layout, slide.slide_layout.slide_master):
        background = layer._element.cSld.bg
        if background is None:
            continue
        for tag, attribute in (
            (f".//{{{drawing}}}srgbClr", "val"),
            (f".//{{{drawing}}}sysClr", "lastClr"),
        ):
            node = background.find(tag)
            if node is not None:
                color = _css_color(node.get(attribute))
                if color != "transparent":
                    return color
        scheme = background.find(f".//{{{drawing}}}schemeClr")
        if scheme is not None:
            return palette.get(scheme.get("val", ""), "#ffffff")
    return "#ffffff"


def _shape_text(
    shape: Any,
    palette: dict[str, str],
    *,
    slide_width: int,
) -> str:
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

    shape_name = str(getattr(shape, "name", "")).casefold()
    is_slide_title = "[role:title]" in shape_name
    is_html_native = "[evoflux-html]" in shape_name
    slide_width_pt = slide_width / 12700
    alignment = {
        PP_ALIGN.CENTER: "center",
        PP_ALIGN.RIGHT: "right",
        PP_ALIGN.JUSTIFY: "justify",
        PP_ALIGN.DISTRIBUTE: "justify",
    }
    paragraphs: list[str] = []
    auto_number = 1
    for paragraph in shape.text_frame.paragraphs:
        runs: list[str] = []
        for run in paragraph.runs:
            styles: list[str] = []
            if run.font.bold:
                styles.append("font-weight:700")
            if run.font.italic:
                styles.append("font-style:italic")
            if run.font.underline:
                styles.append("text-decoration:underline")
            if run.font.size:
                size_pt = run.font.size.pt
                if is_html_native:
                    styles.append(
                        f"font-size:{size_pt / slide_width_pt * 100:.4f}cqw"
                    )
                else:
                    size_pt *= 0.9 if is_slide_title else 1
                    styles.append(f"font-size:{size_pt:.2f}pt")
            if run.font.name:
                styles.append(f"font-family:{_css_font_family(run.font.name)}")
            run_properties = run._r.rPr
            raw_spacing = (
                run_properties.get("spc") if run_properties is not None else None
            )
            if raw_spacing is not None:
                try:
                    spacing_pt = int(raw_spacing) / 100
                    if is_html_native:
                        styles.append(
                            "letter-spacing:"
                            f"{spacing_pt / slide_width_pt * 100:.4f}cqw"
                        )
                    else:
                        styles.append(f"letter-spacing:{spacing_pt:.2f}pt")
                except ValueError:
                    pass
            if run.font.color and run.font.color.type is not None:
                color = _pptx_color(run.font.color, palette)
                if color != "transparent":
                    styles.append(f"color:{color}")
            style_attr = f' style="{";".join(styles)}"' if styles else ""
            runs.append(f"<span{style_attr}>{html.escape(run.text)}</span>")
        paragraph_styles: list[str] = []
        if is_slide_title and not is_html_native:
            paragraph_styles.extend(
                (
                    "white-space:nowrap",
                    "overflow-wrap:normal",
                    "line-height:1.05",
                    "letter-spacing:-.015em",
                )
            )
        if paragraph.alignment in alignment:
            paragraph_styles.append(f"text-align:{alignment[paragraph.alignment]}")
        line_spacing = paragraph.line_spacing
        if is_html_native and isinstance(line_spacing, (int, float)):
            paragraph_styles.append(f"line-height:{float(line_spacing):.3f}")
        if paragraph.level:
            paragraph_styles.append(f"padding-left:{paragraph.level * 1.1:.2f}em")
        properties = paragraph._p.pPr
        marker = ""
        if properties is not None:
            bullets = properties.xpath("./*[local-name()='buChar']")
            auto_numbers = properties.xpath("./*[local-name()='buAutoNum']")
            if bullets:
                marker = html.escape(bullets[0].get("char", "•"))
            elif auto_numbers:
                start_at = auto_numbers[0].get("startAt")
                if start_at:
                    auto_number = int(start_at)
                marker = f"{auto_number}."
                auto_number += 1
        marker_html = f'<span class="bullet-marker">{marker}</span>' if marker else ""
        if marker:
            paragraph_styles.extend(
                ("display:grid", "grid-template-columns:1em 1fr", "column-gap:.2em")
            )
        style_attr = (
            f' style="{";".join(paragraph_styles)}"' if paragraph_styles else ""
        )
        paragraphs.append(
            f'<p{style_attr}>{marker_html}<span class="paragraph-content">'
            f"{''.join(runs) or html.escape(paragraph.text)}</span></p>"
        )
    body_properties = shape.text_frame._txBody.bodyPr
    columns = max(int(body_properties.get("numCol", "1")), 1)
    vertical = {
        MSO_ANCHOR.TOP: "flex-start",
        MSO_ANCHOR.MIDDLE: "center",
        MSO_ANCHOR.BOTTOM: "flex-end",
    }.get(shape.text_frame.vertical_anchor, "center")
    if columns > 1:
        spacing = int(body_properties.get("spcCol", "0")) / 12700
        frame_style = (
            f"column-count:{columns};column-gap:{spacing:.2f}pt;"
            "display:block;overflow:hidden"
        )
    else:
        frame_style = (
            f"justify-content:{vertical};display:flex;"
            "flex-direction:column;overflow:hidden"
        )
    return f'<div class="text-frame" style="{frame_style}">{"".join(paragraphs)}</div>'


def _shape_identity(shape: Any, layer: str) -> str:
    shape_id = html.escape(str(getattr(shape, "shape_id", "")), quote=True)
    name = html.escape(str(getattr(shape, "name", "")), quote=True)
    return (
        f'data-shape-id="{shape_id}" data-shape-name="{name}" '
        f'data-qa-label="{name or shape_id}" data-source-layer="{layer}"'
    )


def _connector_svg(
    shape: Any,
    *,
    layer: str,
    slide_width: int,
    slide_height: int,
    palette: dict[str, str],
) -> str:
    left = 100 * shape.left / slide_width
    top = 100 * shape.top / slide_height
    width = 100 * shape.width / slide_width
    height = 100 * shape.height / slide_height
    transform = shape._element.xpath(".//*[local-name()='xfrm']")
    flip_horizontal = bool(transform and transform[0].get("flipH") == "1")
    flip_vertical = bool(transform and transform[0].get("flipV") == "1")
    x1, x2 = (1000, 0) if flip_horizontal else (0, 1000)
    y1, y2 = (1000, 0) if flip_vertical else (0, 1000)
    if shape.width == 0:
        x1 = x2 = 500
    if shape.height == 0:
        y1 = y2 = 500
    try:
        color = _pptx_color(shape.line.color, palette, "#66717c")
        line_width = shape.line.width.pt if shape.line.width else 1.2
    except (AttributeError, TypeError, ValueError):
        color = "#66717c"
        line_width = 1.2
    endings = shape._element.xpath(".//*[local-name()='tailEnd']")
    has_arrow = bool(endings and endings[0].get("type", "none") not in {"none", ""})
    marker_id = f"arrow-{getattr(shape, 'shape_id', 0)}"
    definitions = ""
    marker_end = ""
    if has_arrow:
        definitions = (
            f'<defs><marker id="{marker_id}" viewBox="0 0 10 10" refX="9" '
            'refY="5" markerWidth="6" markerHeight="6" orient="auto-start-reverse">'
            f'<path d="M 0 0 L 10 5 L 0 10 z" fill="{color}"/></marker></defs>'
        )
        marker_end = f' marker-end="url(#{marker_id})"'
    identity = _shape_identity(shape, layer)
    style = (
        f"left:{left:.3f}%;top:{top:.3f}%;"
        f"width:max({width:.3f}%,2px);height:max({height:.3f}%,2px)"
    )
    return (
        f'<svg class="shape connector" {identity} style="{style}" '
        'viewBox="0 0 1000 1000" preserveAspectRatio="none">'
        f'{definitions}<line x1="{x1}" y1="{y1}" x2="{x2}" y2="{y2}" '
        f'stroke="{color}" stroke-width="{line_width * 1.333:.2f}" '
        f'vector-effect="non-scaling-stroke"{marker_end}/></svg>'
    )


def _shape_gradient_css(shape: Any) -> str | None:
    try:
        gradients = shape._element.xpath(
            "./*[local-name()='spPr']/*[local-name()='gradFill']"
        )
    except (AttributeError, TypeError, ValueError):
        return None
    if not gradients:
        return None
    gradient = gradients[0]
    stops: list[str] = []
    for stop in gradient.xpath("./*[local-name()='gsLst']/*[local-name()='gs']"):
        colors = stop.xpath("./*[local-name()='srgbClr']")
        if not colors:
            continue
        color = colors[0].get("val", "000000")
        alpha_nodes = colors[0].xpath("./*[local-name()='alpha']")
        opacity = (
            int(alpha_nodes[0].get("val", "100000")) / 100000 if alpha_nodes else 1
        )
        red, green, blue = (int(color[index : index + 2], 16) for index in (0, 2, 4))
        position = int(stop.get("pos", "0")) / 1000
        stops.append(f"rgba({red},{green},{blue},{opacity:.3f}) {position:.1f}%")
    if len(stops) < 2:
        return None
    linear = gradient.xpath("./*[local-name()='lin']")
    office_angle = int(linear[0].get("ang", "0")) / 60000 if linear else 0
    return f"linear-gradient({office_angle + 90:.1f}deg,{','.join(stops)})"


def _shape_shadow_css(shape: Any) -> str | None:
    try:
        shadows = shape._element.xpath(
            "./*[local-name()='spPr']/*[local-name()='effectLst']"
            "/*[local-name()='outerShdw']"
        )
    except (AttributeError, TypeError, ValueError):
        return None
    if not shadows:
        return None
    shadow = shadows[0]
    colors = shadow.xpath("./*[local-name()='srgbClr']")
    if not colors:
        return None
    color = colors[0].get("val", "000000")
    alpha_nodes = colors[0].xpath("./*[local-name()='alpha']")
    opacity = int(alpha_nodes[0].get("val", "16000")) / 100000 if alpha_nodes else 0.16
    red, green, blue = (int(color[index : index + 2], 16) for index in (0, 2, 4))
    distance_pt = int(shadow.get("dist", "0")) / 12700
    angle = math.radians(int(shadow.get("dir", "0")) / 60000)
    offset_x = math.cos(angle) * distance_pt * 96 / 72
    offset_y = math.sin(angle) * distance_pt * 96 / 72
    blur = int(shadow.get("blurRad", "0")) / 12700 * 96 / 72
    return (
        f"{offset_x:.2f}px {offset_y:.2f}px {blur:.2f}px "
        f"rgba({red},{green},{blue},{opacity:.3f})"
    )


def _shape_styles(shape: Any, palette: dict[str, str]) -> list[str]:
    styles: list[str] = []
    gradient = _shape_gradient_css(shape)
    if gradient:
        styles.append(f"background:{gradient}")
    try:
        fill = _pptx_color(shape.fill.fore_color, palette)
        if fill != "transparent" and not gradient:
            styles.append(f"background:{fill}")
    except (AttributeError, TypeError, ValueError):
        pass
    try:
        line_type = shape.line.color.type
        line = _pptx_color(shape.line.color, palette)
        width = shape.line.width.pt if shape.line.width else 1
        if line_type is not None and line != "transparent":
            styles.append(f"border:{max(width, 0.5):.2f}px solid {line}")
    except (AttributeError, TypeError, ValueError):
        pass
    try:
        auto_shape = str(shape.auto_shape_type or "").lower()
    except (AttributeError, ValueError):
        auto_shape = ""
    if "oval" in auto_shape or "ellipse" in auto_shape:
        styles.append("border-radius:50%")
    shadow = _shape_shadow_css(shape)
    if shadow:
        styles.append(f"box-shadow:{shadow}")
    return styles


def _chart_svg(chart: Any, theme_palette: dict[str, str]) -> str:
    fallback_palette = ("#2563eb", "#0f766e", "#f59e0b", "#dc2626", "#7c3aed")
    series_values: list[list[float]] = []
    colors: list[str] = []
    for series_index, series in enumerate(chart.series):
        values = []
        for value in getattr(series, "values", ()):
            try:
                values.append(float(value or 0))
            except (TypeError, ValueError):
                values.append(0.0)
        if values:
            series_values.append(values)
            color = "transparent"
            for source in ("fill", "line"):
                try:
                    candidate = (
                        series.format.fill.fore_color
                        if source == "fill"
                        else series.format.line.color
                    )
                    color = _pptx_color(candidate, theme_palette)
                except (AttributeError, TypeError, ValueError):
                    continue
                if color != "transparent":
                    break
            colors.append(
                color
                if color != "transparent"
                else fallback_palette[series_index % len(fallback_palette)]
            )
    if not series_values:
        return '<div class="chart-empty">Chart data unavailable</div>'
    chart_type = getattr(getattr(chart, "chart_type", None), "name", "")
    chart_type = chart_type.upper()
    maximum = max(
        (abs(value) for values in series_values for value in values), default=1
    )
    maximum = maximum or 1
    count = max(len(values) for values in series_values)
    try:
        categories = [str(category) for category in chart.plots[0].categories]
    except (AttributeError, IndexError, TypeError):
        categories = []

    def category_labels() -> str:
        if not categories:
            return ""
        return "".join(
            f'<text x="{52 + index * 520 / max(len(categories) - 1, 1):.2f}" '
            'y="278" text-anchor="middle" font-size="13" fill="#6b7280">'
            f"{html.escape(label)}</text>"
            for index, label in enumerate(categories)
        )

    if "DOUGHNUT" in chart_type or "PIE" in chart_type:
        values = [abs(value) for value in series_values[0]]
        total = sum(values) or 1
        start_angle = -math.pi / 2
        slices: list[str] = []
        for index, value in enumerate(values):
            end_angle = start_angle + 2 * math.pi * value / total
            start_x = 180 + 108 * math.cos(start_angle)
            start_y = 150 + 108 * math.sin(start_angle)
            end_x = 180 + 108 * math.cos(end_angle)
            end_y = 150 + 108 * math.sin(end_angle)
            large_arc = 1 if end_angle - start_angle > math.pi else 0
            slices.append(
                f'<path d="M180 150 L{start_x:.2f} {start_y:.2f} '
                f'A108 108 0 {large_arc} 1 {end_x:.2f} {end_y:.2f} Z" '
                f'fill="{fallback_palette[index % len(fallback_palette)]}"/>'
            )
            start_angle = end_angle
        legend = "".join(
            (
                f'<rect x="340" y="{54 + index * 34}" width="14" height="14" '
                f'fill="{fallback_palette[index % len(fallback_palette)]}"/>'
                f'<text x="364" y="{66 + index * 34}" font-size="14" fill="#374151">'
                f"{html.escape(categories[index] if index < len(categories) else f'Item {index + 1}')} "
                f"({value / total:.0%})</text>"
            )
            for index, value in enumerate(values)
        )
        hole = (
            '<circle cx="180" cy="150" r="65" fill="white"/>'
            if "DOUGHNUT" in chart_type
            else ""
        )
        return (
            '<svg class="chart-svg" viewBox="0 0 600 300">'
            f"{''.join(slices)}{hole}{legend}</svg>"
        )

    if "AREA" in chart_type or "LINE" in chart_type:
        paths: list[str] = []
        for series_index, values in enumerate(series_values):
            points = []
            for value_index, value in enumerate(values):
                x = 52 + value_index * 520 / max(len(values) - 1, 1)
                y = 252 - (value / maximum) * 205
                points.append((x, y))
            point_text = " ".join(f"{x:.2f},{y:.2f}" for x, y in points)
            markers = "".join(
                f'<circle cx="{x:.2f}" cy="{y:.2f}" r="4" fill="{colors[series_index]}"/>'
                for x, y in points
            )
            area = ""
            if "AREA" in chart_type and points:
                polygon_points = [
                    f"{points[0][0]:.2f},252",
                    *[f"{x:.2f},{y:.2f}" for x, y in points],
                    f"{points[-1][0]:.2f},252",
                ]
                area = (
                    f'<polygon class="chart-area" points="{" ".join(polygon_points)}" '
                    f'fill="{colors[series_index]}" fill-opacity=".24"/>'
                )
            paths.append(
                area + f'<polyline points="{point_text}" fill="none" '
                f'stroke="{colors[series_index]}" stroke-width="4"/>'
                f"{markers}"
            )
        return (
            '<svg class="chart-svg" viewBox="0 0 600 300" preserveAspectRatio="none">'
            '<line x1="45" y1="252" x2="575" y2="252" stroke="#d1d5db"/>'
            f"{''.join(paths)}{category_labels()}</svg>"
        )

    group_width = 520 / max(count, 1)
    bar_width = max(group_width / max(len(series_values), 1) * 0.72, 2)
    bars: list[str] = []
    for series_index, values in enumerate(series_values):
        for value_index, value in enumerate(values):
            if "BAR" in chart_type:
                height = max(205 / max(count, 1) * 0.7, 4)
                y = 34 + value_index * 205 / max(count, 1) + series_index * height
                bar_length = abs(value) / maximum * 400
                bars.append(
                    f'<rect x="150" y="{y:.2f}" width="{bar_length:.2f}" '
                    f'height="{height:.2f}" fill="{colors[series_index]}"/>'
                )
            else:
                bar_height = abs(value) / maximum * 210
                x = 52 + value_index * group_width + series_index * bar_width
                y = 252 - bar_height if value >= 0 else 252
                bars.append(
                    f'<rect x="{x:.2f}" y="{y:.2f}" width="{bar_width:.2f}" '
                    f'height="{bar_height:.2f}" fill="{colors[series_index]}"/>'
                )
    if "BAR" in chart_type:
        labels = "".join(
            f'<text x="140" y="{47 + index * 205 / max(count, 1):.2f}" '
            'text-anchor="end" font-size="13" fill="#6b7280">'
            f"{html.escape(label)}</text>"
            for index, label in enumerate(categories)
        )
        axis = '<line x1="150" y1="25" x2="150" y2="260" stroke="#9ca3af"/>'
    else:
        labels = category_labels()
        axis = '<line x1="45" y1="252" x2="575" y2="252" stroke="#9ca3af"/>'
    return (
        '<svg class="chart-svg" viewBox="0 0 600 300" preserveAspectRatio="none">'
        f"{axis}{''.join(bars)}{labels}</svg>"
    )


def _placeholder_key(shape: Any) -> int | None:
    if not getattr(shape, "is_placeholder", False):
        return None
    try:
        return int(shape.placeholder_format.idx)
    except (AttributeError, TypeError, ValueError):
        return None


def _table_cell_html(cell: Any, palette: dict[str, str]) -> str:
    styles: list[str] = []
    try:
        fill = _pptx_color(cell.fill.fore_color, palette)
        if fill != "transparent":
            styles.append(f"background:{fill}")
    except (AttributeError, TypeError, ValueError):
        pass
    paragraph = cell.text_frame.paragraphs[0] if cell.text_frame.paragraphs else None
    run = paragraph.runs[0] if paragraph and paragraph.runs else None
    if run is not None:
        if run.font.color is not None:
            color = _pptx_color(run.font.color, palette)
            if color != "transparent":
                styles.append(f"color:{color}")
        if run.font.size:
            styles.append(f"font-size:{run.font.size.pt:.2f}pt")
        if run.font.name:
            styles.append(f"font-family:{_css_font_family(run.font.name)}")
        if run.font.bold:
            styles.append("font-weight:700")
    alignment = {
        2: "center",
        3: "right",
        4: "justify",
    }.get(int(paragraph.alignment or 0) if paragraph else 0)
    if alignment:
        styles.append(f"text-align:{alignment}")
    style_attr = f' style="{";".join(styles)}"' if styles else ""
    return f"<td{style_attr}>{html.escape(cell.text)}</td>"


def _composite_slide_shapes(slide: Any) -> list[tuple[str, Any]]:
    def append_shape(layer: str, shape: Any) -> None:
        shape_type = getattr(getattr(shape, "shape_type", None), "name", "")
        if shape_type == "GROUP":
            for child in shape.shapes:
                append_shape(layer, child)
            return
        result.append((layer, shape))

    layout = slide.slide_layout
    result: list[tuple[str, Any]] = []
    for shape in layout.slide_master.shapes:
        if _placeholder_key(shape) is None:
            append_shape("master", shape)
    for shape in layout.shapes:
        if _placeholder_key(shape) is None:
            append_shape("layout", shape)
    for shape in slide.shapes:
        append_shape("slide", shape)
    return result


def _render_pptx(source: Path) -> str:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    presentation = Presentation(str(source))
    slide_width = int(presentation.slide_width or 0)
    slide_height = int(presentation.slide_height or 0)
    if slide_width <= 0 or slide_height <= 0:
        raise ValueError("Presentation has invalid slide dimensions")
    rendered_slides: list[str] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        palette = _pptx_theme_palette(slide)
        shapes: list[str] = []
        for layer, shape in _composite_slide_shapes(slide):
            if shape.shape_type == MSO_SHAPE_TYPE.LINE:
                shapes.append(
                    _connector_svg(
                        shape,
                        layer=layer,
                        slide_width=slide_width,
                        slide_height=slide_height,
                        palette=palette,
                    )
                )
                continue
            box = _shape_box(
                shape,
                slide_width=slide_width,
                slide_height=slide_height,
            )
            identity = _shape_identity(shape, layer)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                shapes.append(
                    f'<img class="shape picture" {identity} style="{box}" '
                    f'src="{_picture_data_uri(shape)}" alt="">'
                )
                continue

            if getattr(shape, "has_table", False):
                table_rows: list[str] = []
                for row in shape.table.rows:
                    cells = "".join(
                        _table_cell_html(cell, palette) for cell in row.cells
                    )
                    table_rows.append(f"<tr>{cells}</tr>")
                shapes.append(
                    f'<div class="shape table-shape" {identity} style="{box}"><table>'
                    f"{''.join(table_rows)}</table></div>"
                )
                continue

            if getattr(shape, "has_chart", False):
                chart_title = ""
                if shape.chart.has_title:
                    chart_title = shape.chart.chart_title.text_frame.text
                shapes.append(
                    f'<div class="shape chart" {identity} style="{box}">'
                    f"<strong>{html.escape(chart_title or 'Chart')}</strong>"
                    f"{_chart_svg(shape.chart, palette)}</div>"
                )
                continue

            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                styles = [box, *_shape_styles(shape, palette)]
                native_class = " html-native" if "[evoflux-html]" in shape.name else ""
                shapes.append(
                    f'<div class="shape text-shape{native_class}" {identity} '
                    f'style="{";".join(styles)}">'
                    f"{_shape_text(shape, palette, slide_width=slide_width)}</div>"
                )
                continue

            generic_styles = [box, *_shape_styles(shape, palette)]
            if len(generic_styles) > 1:
                shapes.append(
                    f'<div class="shape vector-shape" {identity} '
                    f'style="{";".join(generic_styles)}"></div>'
                )

        background = _slide_background(slide, palette)
        ratio = f"{slide_width}/{slide_height}"
        rendered_slides.append(
            f'<article class="slide-wrap"><div class="slide-number">{slide_number}</div>'
            f'<section class="slide" style="aspect-ratio:{ratio};background:{background}">'
            f"{''.join(shapes)}</section></article>"
        )

    css = """
    *{box-sizing:border-box}body{margin:0;padding:28px;background:#e8eaed;color:#202124;
    font-family:Arial,sans-serif}.slide-wrap{position:relative;width:min(1120px,94vw);
    margin:0 auto 30px}.slide{position:relative;width:100%;background:#fff;
    container-type:inline-size;
    overflow:hidden;box-shadow:0 3px 14px #0003}.slide-number{position:absolute;right:calc(100%
    + 8px);top:0;color:#6b7280;font-size:12px}.shape{position:absolute;overflow:hidden}
    .text-shape{padding:4px}.text-frame{width:100%;height:100%;display:flex;
    flex-direction:column;overflow:hidden}
    .text-shape p{margin:0 0 .2em;line-height:1.12;white-space:pre-wrap;
    overflow-wrap:normal;word-break:normal;break-inside:avoid}.bullet-marker{text-align:center}
    .text-shape.html-native{padding:0}.text-shape.html-native p{margin:0}
    .picture{object-fit:cover}.table-shape table{width:100%;height:100%;border-collapse:collapse;
    table-layout:fixed}.table-shape td{border:1px solid #b8bdc6;padding:4px;font-size:12px;
    overflow:hidden}.chart{display:flex;flex-direction:column;align-items:center;justify-content:
    center;border:1px solid #ccd1d9;background:#fff;color:#374151;gap:5px}.chart strong{
    font-size:12px}.chart-svg{width:100%;height:calc(100% - 20px)}.chart-empty{
    font-size:11px;color:#6b7280}.vector-shape{pointer-events:none}
    .connector{overflow:visible;pointer-events:none}
    """
    return _page(title=source.name, body="".join(rendered_slides), css=css)


def _render_source(source: Path) -> str:
    renderers = {
        ".docx": _render_docx,
        ".xlsx": _render_xlsx,
        ".pptx": _render_pptx,
    }
    try:
        return renderers[source.suffix.lower()](source)
    except Exception as exc:
        logger.warning(
            "office_preview_render_failed file={} error={}",
            source.name,
            str(exc)[:500],
        )
        raise OfficePreviewError(f"Could not render this document: {exc}") from exc


def render_office_preview(source: Path) -> Path:
    """Render ``source`` to a cached, self-contained HTML document."""
    suffix = source.suffix.lower()
    if suffix not in SUPPORTED_OFFICE_PREVIEW_EXTENSIONS:
        raise OfficePreviewUnsupportedError(
            f"{suffix or 'This file type'} is not supported for Office preview."
        )
    if source.stat().st_size > MAX_OFFICE_PREVIEW_BYTES:
        raise OfficePreviewUnsupportedError(
            f"Office preview is limited to {MAX_OFFICE_PREVIEW_BYTES // (1024 * 1024)} MB."
        )

    output = _cache_path(source)
    if output.is_file():
        return output

    with _render_lock:
        if output.is_file():
            return output
        output.parent.mkdir(parents=True, exist_ok=True)
        temporary = output.with_suffix(".tmp")
        try:
            temporary.write_text(_render_source(source), encoding="utf-8")
            os.replace(temporary, output)
        finally:
            temporary.unlink(missing_ok=True)
    return output
