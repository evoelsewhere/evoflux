"""High-preservation PPTX template editing through direct OOXML cloning."""

from __future__ import annotations

import asyncio
from copy import deepcopy
from dataclasses import dataclass, field
import json
from pathlib import Path
from typing import Any, Literal

from pydantic import BaseModel, ConfigDict, Field, model_validator

from app.services.office.internal_rendering import render_pptx_pages
from app.services.office.runtime import file_sha256

pptx_sha256 = file_sha256
MAX_TEMPLATE_SLIDES = 80
MAX_EDITS_PER_SLIDE = 160


class TemplateObjectEdit(BaseModel):
    model_config = ConfigDict(extra="forbid")
    operation: Literal[
        "set_text",
        "replace_text",
        "replace_image",
        "set_table_cell",
        "set_chart_series",
    ]
    target_id: str = Field(pattern=r"^(?:sh|im|tb|ch)/[1-9][0-9]*/[1-9][0-9]*$")
    text: str | None = Field(default=None, max_length=20_000)
    find: str | None = Field(default=None, min_length=1, max_length=4_000)
    replace: str | None = Field(default=None, max_length=20_000)
    asset_path: str | None = Field(default=None, max_length=2_000)
    alt: str | None = Field(default=None, max_length=1_000)
    row: int | None = Field(default=None, ge=0, le=500)
    column: int | None = Field(default=None, ge=0, le=500)
    series_index: int | None = Field(default=None, ge=0, le=100)
    values: list[float] | None = Field(default=None, min_length=1, max_length=2_000)

    @model_validator(mode="after")
    def validate_operation_fields(self) -> TemplateObjectEdit:
        required = {
            "set_text": self.text is not None,
            "replace_text": self.find is not None and self.replace is not None,
            "replace_image": bool(self.asset_path),
            "set_table_cell": self.row is not None
            and self.column is not None
            and self.text is not None,
            "set_chart_series": self.series_index is not None
            and self.values is not None,
        }
        if not required[self.operation]:
            raise ValueError(f"{self.operation} is missing required fields")
        return self


class TemplateSlidePlan(BaseModel):
    model_config = ConfigDict(extra="forbid")
    output_slide: int = Field(ge=1, le=MAX_TEMPLATE_SLIDES)
    source_slide: int = Field(ge=1, le=MAX_TEMPLATE_SLIDES)
    narrative_role: str = Field(min_length=1, max_length=240)
    reuse_mode: Literal["duplicate-slide"] = "duplicate-slide"
    edits: list[TemplateObjectEdit] = Field(
        default_factory=list, max_length=MAX_EDITS_PER_SLIDE
    )
    speaker_notes: str | None = Field(default=None, max_length=40_000)


class TemplateDeckProject(BaseModel):
    model_config = ConfigDict(extra="forbid")
    schema_version: Literal[2] = 2
    title: str = Field(min_length=1, max_length=240)
    source_sha256: str = Field(pattern=r"^[0-9a-f]{64}$")
    template_confirmed: Literal[True]
    output_slides: list[TemplateSlidePlan] = Field(
        min_length=1, max_length=MAX_TEMPLATE_SLIDES
    )
    omitted_source_slides: list[int] = Field(default_factory=list)

    @model_validator(mode="after")
    def validate_slide_map(self) -> TemplateDeckProject:
        expected = list(range(1, len(self.output_slides) + 1))
        actual = [slide.output_slide for slide in self.output_slides]
        if actual != expected:
            raise ValueError("output_slides must be ordered sequentially from 1")
        if len(self.omitted_source_slides) != len(set(self.omitted_source_slides)):
            raise ValueError("omitted_source_slides must not contain duplicates")
        return self


@dataclass
class TemplatePipelineResult:
    action: str
    source_pptx: Path
    work_dir: Path
    manifest_path: Path | None = None
    output: Path | None = None
    previews: list[Path] = field(default_factory=list)
    layout_paths: list[Path] = field(default_factory=list)
    slide_count: int = 0
    issues: list[dict[str, Any]] = field(default_factory=list)
    metadata: dict[str, Any] = field(default_factory=dict)

    @property
    def passed(self) -> bool:
        return not any(issue.get("severity") == "error" for issue in self.issues)


def load_template_project(path: Path) -> TemplateDeckProject:
    return TemplateDeckProject.model_validate_json(path.read_text(encoding="utf-8"))


def load_template_manifest(path: Path) -> dict[str, Any]:
    value = json.loads(path.read_text(encoding="utf-8"))
    if not isinstance(value, dict):
        raise ValueError("template manifest must be an object")
    return value


def _record_kind(shape: Any) -> tuple[str, str]:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    if getattr(shape, "has_table", False):
        return "tb", "table"
    if getattr(shape, "has_chart", False):
        return "ch", "chart"
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return "im", "image"
    return "sh", "textbox" if getattr(shape, "has_text_frame", False) else "shape"


def _shape_record(shape: Any, slide_number: int) -> dict[str, Any]:
    prefix, kind = _record_kind(shape)
    record: dict[str, Any] = {
        "id": f"{prefix}/{slide_number}/{shape.shape_id}",
        "slide": slide_number,
        "shapeId": shape.shape_id,
        "kind": kind,
        "name": shape.name,
        "left": int(shape.left),
        "top": int(shape.top),
        "width": int(shape.width),
        "height": int(shape.height),
    }
    if getattr(shape, "has_text_frame", False):
        record["text"] = shape.text
    if getattr(shape, "has_table", False):
        record.update(
            {"rows": len(shape.table.rows), "columns": len(shape.table.columns)}
        )
    if getattr(shape, "has_chart", False):
        record["seriesCount"] = len(shape.chart.series)
    return record


def _inspect_records(presentation: Any) -> list[dict[str, Any]]:
    return [
        _shape_record(shape, slide_number)
        for slide_number, slide in enumerate(presentation.slides, start=1)
        for shape in slide.shapes
    ]


def validate_template_project(
    project: TemplateDeckProject, manifest: dict[str, Any], *, source_pptx: Path
) -> dict[str, Any]:
    actual_digest = file_sha256(source_pptx)
    expected_digest = str(manifest.get("sourceSha256") or "")
    if actual_digest != project.source_sha256 or actual_digest != expected_digest:
        raise ValueError("source PPTX changed after inspection; inspect it again")
    source_slide_count = int(manifest.get("slideCount", 0))
    used = {slide.source_slide for slide in project.output_slides}
    if any(slide < 1 or slide > source_slide_count for slide in used):
        raise ValueError("source_slide is outside the inspected presentation")
    expected_omissions = set(range(1, source_slide_count + 1)) - used
    if set(project.omitted_source_slides) != expected_omissions:
        raise ValueError(
            "omitted_source_slides must explicitly list every unused source slide"
        )
    records = {str(record.get("id")): record for record in manifest.get("records", [])}
    operation_kind = {
        "set_text": {"textbox", "shape"},
        "replace_text": {"textbox", "shape"},
        "replace_image": {"image"},
        "set_table_cell": {"table"},
        "set_chart_series": {"chart"},
    }
    for slide in project.output_slides:
        for edit in slide.edits:
            record = records.get(edit.target_id)
            if record is None:
                raise ValueError(
                    f"target_id was not found in inspect manifest: {edit.target_id}"
                )
            if int(record.get("slide", 0)) != slide.source_slide:
                raise ValueError(
                    f"target {edit.target_id} does not belong to source slide {slide.source_slide}"
                )
            kind = str(record.get("kind"))
            if kind not in operation_kind[edit.operation]:
                expected = next(iter(operation_kind[edit.operation]))
                raise ValueError(f"{edit.operation} requires a {expected} target")
    return {
        "valid": True,
        "source_slide_count": source_slide_count,
        "output_slide_count": len(project.output_slides),
        "edit_count": sum(len(slide.edits) for slide in project.output_slides),
        "preserve_only_slide_count": sum(
            1 for slide in project.output_slides if not slide.edits
        ),
    }


def template_catalog() -> dict[str, Any]:
    return {
        "workflow": "direct-openxml-pptx-template",
        "style_behavior": {
            "uploaded_template_is_style_confirmation": True,
            "ask_style_question": False,
            "ambiguous_upload": "Ask whether the PPTX is a visual template or only a content source.",
        },
        "invariants": [
            "Clone source slide XML and relationships without rebuilding its layout.",
            "Preserve masters, layouts, themes, transitions, timing, and untouched objects.",
            "Edit only stable shape IDs declared by inspect.",
            "Render every result with the bundled internal OOXML renderer.",
        ],
        "supported_edits": {
            "set_text": "Replace text while retaining the first run style.",
            "replace_text": "Replace a substring while retaining surrounding style.",
            "replace_image": "Swap image bytes while preserving frame geometry.",
            "set_table_cell": "Update one native table cell.",
            "set_chart_series": "Update one native chart series.",
            "speaker_notes": "Set notes on the cloned output slide.",
        },
        "project_json_schema": TemplateDeckProject.model_json_schema(),
    }


def _write_inspect_artifacts(source: Path, work_dir: Path) -> TemplatePipelineResult:
    from pptx import Presentation

    presentation = Presentation(str(source))
    records = _inspect_records(presentation)
    previews = render_pptx_pages(source, work_dir / "previews")
    layout_paths: list[Path] = []
    for slide_number in range(1, len(presentation.slides) + 1):
        layout_path = (
            work_dir / "layouts" / f"source-slide-{slide_number:03d}.layout.json"
        )
        layout_path.parent.mkdir(parents=True, exist_ok=True)
        layout_path.write_text(
            json.dumps(
                {
                    "slide": slide_number,
                    "records": [
                        record for record in records if record["slide"] == slide_number
                    ],
                },
                ensure_ascii=False,
                indent=2,
            )
            + "\n",
            encoding="utf-8",
        )
        layout_paths.append(layout_path)
    manifest_path = work_dir / "template-manifest.json"
    manifest = {
        "schemaVersion": 2,
        "engine": "evoflux-direct-openxml",
        "sourcePath": str(source.resolve()),
        "sourceSha256": file_sha256(source),
        "slideCount": len(presentation.slides),
        "records": records,
        "slideArtifacts": [
            {
                "slide": index,
                "previewPath": str(previews[index - 1]),
                "layoutPath": str(layout_paths[index - 1]),
            }
            for index in range(1, len(presentation.slides) + 1)
        ],
    }
    manifest_path.write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2) + "\n", encoding="utf-8"
    )
    return TemplatePipelineResult(
        action="inspect",
        source_pptx=source,
        work_dir=work_dir,
        manifest_path=manifest_path,
        previews=previews,
        layout_paths=layout_paths,
        slide_count=len(presentation.slides),
        metadata={"engine": "evoflux-direct-openxml"},
    )


def _rewrite_relationship_ids(element: Any, mapping: dict[str, str]) -> None:
    from pptx.oxml.ns import qn

    attributes = (qn("r:id"), qn("r:embed"), qn("r:link"))
    for node in element.iter():
        for attribute in attributes:
            old = node.get(attribute)
            if old in mapping:
                node.set(attribute, mapping[old])


def _duplicate_slide(presentation: Any, source_slide: Any) -> Any:
    from pptx.shapes.shapetree import SlideShapes

    target = presentation.slides.add_slide(source_slide.slide_layout)
    duplicate = deepcopy(source_slide._element)  # noqa: SLF001 - OOXML-preserving clone
    mapping: dict[str, str] = {}
    target_layout_rel = next(
        rel for rel in target.part.rels.values() if rel.reltype.endswith("/slideLayout")
    )
    for relation in source_slide.part.rels.values():
        if relation.reltype.endswith("/slideLayout"):
            mapping[relation.rId] = target_layout_rel.rId
        elif relation.reltype.endswith("/notesSlide"):
            continue
        else:
            mapping[relation.rId] = target.part.rels._add_relationship(  # noqa: SLF001
                relation.reltype,
                relation._target,
                relation.is_external,  # noqa: SLF001
            )
    _rewrite_relationship_ids(duplicate, mapping)
    target.part._element = duplicate  # noqa: SLF001
    target._element = duplicate  # noqa: SLF001
    # ``Slide.shapes`` is cached when python-pptx creates the blank target.
    # Point that proxy at the cloned shape tree as well as replacing the XML.
    target.__dict__["shapes"] = SlideShapes(duplicate.cSld.spTree, target)
    return target


def _remove_original_slides(presentation: Any, count: int) -> None:
    slide_ids = presentation.slides._sldIdLst  # noqa: SLF001
    for slide_id in list(slide_ids)[:count]:
        presentation.part.drop_rel(slide_id.rId)
        slide_ids.remove(slide_id)


def _shape_by_id(slide: Any, shape_id: int) -> Any:
    for shape in slide.shapes:
        if shape.shape_id == shape_id:
            return shape
    raise ValueError(f"shape id {shape_id} is missing from cloned slide")


def _set_text_preserving_style(shape: Any, value: str) -> None:
    runs = [run for paragraph in shape.text_frame.paragraphs for run in paragraph.runs]
    if not runs:
        shape.text = value
        return
    runs[0].text = value
    for run in runs[1:]:
        run.text = ""


def _replace_image(slide: Any, shape: Any, source: Path) -> None:
    image_part, relation_id = slide.part.get_or_add_image_part(str(source))
    del image_part
    shape._element.blipFill.blip.rEmbed = relation_id  # noqa: SLF001


def _replace_chart_series(shape: Any, series_index: int, values: list[float]) -> None:
    from pptx.chart.data import CategoryChartData

    chart = shape.chart
    series_values = [list(series.values) for series in chart.series]
    if series_index >= len(series_values):
        raise ValueError(f"chart series index {series_index} is out of range")
    if len(values) != len(series_values[series_index]):
        raise ValueError("replacement chart series must preserve category count")
    series_values[series_index] = values
    try:
        categories = [str(category) for category in chart.plots[0].categories]
    except (AttributeError, TypeError):
        categories = [str(index + 1) for index in range(len(values))]
    data = CategoryChartData()
    data.categories = categories
    for index, series in enumerate(chart.series):
        data.add_series(str(series.name or f"Series {index + 1}"), series_values[index])
    chart.replace_data(data)


def _apply_edit(slide: Any, edit: TemplateObjectEdit, project_dir: Path) -> None:
    shape_id = int(edit.target_id.rsplit("/", 1)[1])
    shape = _shape_by_id(slide, shape_id)
    if edit.operation == "set_text":
        _set_text_preserving_style(shape, edit.text or "")
    elif edit.operation == "replace_text":
        current = shape.text
        if edit.find not in current:
            raise ValueError(f"target {edit.target_id} does not contain {edit.find!r}")
        _set_text_preserving_style(
            shape, current.replace(edit.find or "", edit.replace or "")
        )
    elif edit.operation == "replace_image":
        path = Path(edit.asset_path or "")
        if not path.is_absolute():
            path = project_dir / path
        if not path.is_file():
            raise FileNotFoundError(f"replacement image does not exist: {path}")
        _replace_image(slide, shape, path)
    elif edit.operation == "set_table_cell":
        assert edit.row is not None and edit.column is not None
        if edit.row >= len(shape.table.rows) or edit.column >= len(shape.table.columns):
            raise ValueError("table cell index is out of range")
        shape.table.cell(edit.row, edit.column).text = edit.text or ""
    else:
        assert edit.series_index is not None and edit.values is not None
        _replace_chart_series(shape, edit.series_index, edit.values)


def _build_template(
    source: Path,
    project: TemplateDeckProject,
    project_dir: Path,
    output: Path,
    work_dir: Path,
    *,
    publish: bool,
) -> TemplatePipelineResult:
    from pptx import Presentation

    presentation = Presentation(str(source))
    originals = list(presentation.slides)
    output_slides: list[Any] = []
    for plan in project.output_slides:
        output_slides.append(
            _duplicate_slide(presentation, originals[plan.source_slide - 1])
        )
    _remove_original_slides(presentation, len(originals))
    for plan, slide in zip(project.output_slides, output_slides, strict=True):
        for edit in plan.edits:
            _apply_edit(slide, edit, project_dir=project_dir)
        if plan.speaker_notes:
            slide.notes_slide.notes_text_frame.text = plan.speaker_notes
    candidate = output if publish else work_dir / "template-preview.pptx"
    candidate.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(candidate))
    previews = render_pptx_pages(candidate, work_dir / "previews")
    layout_paths: list[Path] = []
    records = _inspect_records(Presentation(str(candidate)))
    for slide_number in range(1, len(project.output_slides) + 1):
        path = work_dir / "layouts" / f"output-slide-{slide_number:03d}.layout.json"
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_text(
            json.dumps(
                {
                    "slide": slide_number,
                    "records": [
                        record for record in records if record["slide"] == slide_number
                    ],
                },
                ensure_ascii=False,
                indent=2,
            )
            + "\n",
            encoding="utf-8",
        )
        layout_paths.append(path)
    manifest_path = work_dir / "template-output-manifest.json"
    manifest_path.write_text(
        json.dumps(
            {
                "schemaVersion": 2,
                "engine": "evoflux-direct-openxml",
                "sourceSha256": project.source_sha256,
                "slideCount": len(project.output_slides),
                "records": records,
            },
            ensure_ascii=False,
            indent=2,
        )
        + "\n",
        encoding="utf-8",
    )
    return TemplatePipelineResult(
        action="compose" if publish else "render",
        source_pptx=source,
        work_dir=work_dir,
        manifest_path=manifest_path,
        output=output if publish else None,
        previews=previews,
        layout_paths=layout_paths,
        slide_count=len(project.output_slides),
        metadata={"engine": "evoflux-direct-openxml", "preserved_master_layouts": True},
    )


async def inspect_pptx_template(
    source_pptx: Path, *, workspace_root: Path, work_dir: Path
) -> TemplatePipelineResult:
    del workspace_root
    return await asyncio.to_thread(_write_inspect_artifacts, source_pptx, work_dir)


async def render_pptx_template(
    source_pptx: Path,
    project_path: Path,
    manifest_path: Path,
    *,
    workspace_root: Path,
    work_dir: Path,
) -> TemplatePipelineResult:
    del workspace_root
    project = load_template_project(project_path)
    validate_template_project(
        project, load_template_manifest(manifest_path), source_pptx=source_pptx
    )
    return await asyncio.to_thread(
        _build_template,
        source_pptx,
        project,
        project_path.parent.resolve(),
        work_dir / "template-preview.pptx",
        work_dir,
        publish=False,
    )


async def compose_pptx_template(
    source_pptx: Path,
    project_path: Path,
    manifest_path: Path,
    output: Path,
    *,
    workspace_root: Path,
    work_dir: Path,
) -> TemplatePipelineResult:
    del workspace_root
    project = load_template_project(project_path)
    validate_template_project(
        project, load_template_manifest(manifest_path), source_pptx=source_pptx
    )
    try:
        return await asyncio.to_thread(
            _build_template,
            source_pptx,
            project,
            project_path.parent.resolve(),
            output,
            work_dir,
            publish=True,
        )
    except Exception:
        output.unlink(missing_ok=True)
        raise


__all__ = [
    "TemplateDeckProject",
    "TemplateObjectEdit",
    "TemplatePipelineResult",
    "TemplateSlidePlan",
    "compose_pptx_template",
    "inspect_pptx_template",
    "load_template_manifest",
    "load_template_project",
    "pptx_sha256",
    "render_pptx_template",
    "template_catalog",
    "validate_template_project",
]
