from __future__ import annotations

import re
from types import SimpleNamespace
from zipfile import ZIP_DEFLATED, ZipFile

from docx import Document
from lxml import html as lxml_html
from openpyxl import Workbook
from openpyxl.chart import BarChart, LineChart, Reference
from openpyxl.drawing.image import Image as WorkbookImage
from PIL import Image
from pptx import Presentation
from pptx.chart.data import ChartData, XyChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_MARKER_STYLE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn as pptx_qn
from pptx.util import Inches, Pt
import pytest

from app.services.document_preview import service as preview


def _docx(path):
    document = Document()
    document.add_heading("Quarterly review", level=1)
    document.add_paragraph("Revenue increased.")
    document.save(path)


def _xlsx(path):
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Summary"
    sheet["A1"] = "Revenue"
    sheet["B1"] = 120
    sheet["B2"] = "=B1*2"
    workbook.save(path)


def _set_formula_cache(path, formula: str, value: str) -> None:
    replacement = path.with_suffix(".cached.xlsx")
    expected = f"<f>{formula}</f><v></v>".encode()
    updated = f"<f>{formula}</f><v>{value}</v>".encode()
    replaced = False
    with ZipFile(path) as source, ZipFile(replacement, "w", ZIP_DEFLATED) as target:
        for member in source.infolist():
            payload = source.read(member.filename)
            if member.filename.startswith("xl/worksheets/") and expected in payload:
                payload = payload.replace(expected, updated, 1)
                replaced = True
            target.writestr(member, payload)
    assert replaced, "formula cache target was not found in workbook XML"
    replacement.replace(path)


def _pptx(path):
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    text_box.text = "Product launch"
    text_box.text_frame.paragraphs[0].runs[0].font.name = "Aptos Display"
    presentation.save(path)


def _add_chart(slide, categories, series, *, kind, title):
    data = ChartData()
    data.categories = categories
    for name, values in series.items():
        data.add_series(name, values)
    chart = slide.shapes.add_chart(
        kind,
        Inches(0.8),
        Inches(0.8),
        Inches(11.7),
        Inches(5.8),
        data,
    ).chart
    chart.has_title = True
    chart.chart_title.text_frame.text = title
    return chart


def _replace_master_paragraph_spacing(
    paragraph_properties,
    *,
    line_percent: int,
    before_points: int,
    after_points: int,
):
    for tag in ("lnSpc", "spcBef", "spcAft"):
        for node in paragraph_properties.xpath(f"./*[local-name()='{tag}']"):
            paragraph_properties.remove(node)

    values = (
        ("lnSpc", "spcPct", line_percent),
        ("spcBef", "spcPts", before_points * 100),
        ("spcAft", "spcPts", after_points * 100),
    )
    for index, (outer_tag, inner_tag, value) in enumerate(values):
        outer = OxmlElement(f"a:{outer_tag}")
        inner = OxmlElement(f"a:{inner_tag}")
        inner.set("val", str(value))
        outer.append(inner)
        paragraph_properties.insert(index, outer)


def _set_master_level_one_style(
    presentation,
    style_name,
    *,
    size: int,
    typeface: str,
    character_spacing: int,
    line_percent: int,
    before_points: int,
    after_points: int,
    bold: bool = False,
):
    master = presentation.slide_masters[0]
    paragraph_properties = master._element.xpath(
        "./*[local-name()='txStyles']"
        f"/*[local-name()='{style_name}']"
        "/*[local-name()='lvl1pPr']"
    )[0]
    run_properties = paragraph_properties.xpath("./*[local-name()='defRPr']")[0]
    run_properties.set("sz", str(size))
    run_properties.set("spc", str(character_spacing))
    run_properties.set("b", "1" if bold else "0")
    latin = run_properties.xpath("./*[local-name()='latin']")[0]
    latin.set("typeface", typeface)
    _replace_master_paragraph_spacing(
        paragraph_properties,
        line_percent=line_percent,
        before_points=before_points,
        after_points=after_points,
    )


def _pptx_semantic_stress(path, picture_path):
    Image.new("RGB", (320, 180), "#2563eb").save(picture_path)

    presentation = Presentation()
    presentation.slide_width = Inches(40 / 3)
    presentation.slide_height = Inches(7.5)
    _set_master_level_one_style(
        presentation,
        "titleStyle",
        size=4400,
        typeface="Aptos Display",
        character_spacing=120,
        line_percent=105000,
        before_points=0,
        after_points=4,
        bold=True,
    )
    _set_master_level_one_style(
        presentation,
        "bodyStyle",
        size=3200,
        typeface="Aptos",
        character_spacing=80,
        line_percent=125000,
        before_points=12,
        after_points=6,
    )

    inherited_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    title = inherited_slide.shapes.title
    assert title is not None
    title.name = "Inherited master title placeholder"
    title.text = "Inherited master title"
    title.text_frame.margin_left = Pt(18)
    title.text_frame.margin_right = Pt(9)
    title.text_frame.margin_top = Pt(6)
    title.text_frame.margin_bottom = Pt(12)

    body = next(shape for shape in inherited_slide.placeholders if shape != title)
    body.name = "Inherited master body placeholder"
    body.text_frame.clear()
    body_paragraph = body.text_frame.paragraphs[0]
    body_paragraph.text = "Inherited master body"
    body_paragraph.level = 0
    body.text_frame.margin_left = Pt(24)
    body.text_frame.margin_right = Pt(12)
    body.text_frame.margin_top = Pt(9)
    body.text_frame.margin_bottom = Pt(9)
    inherited_slide.notes_slide.notes_text_frame.text = (
        'Presenter note <verify> & "repair" before publish.'
    )

    media_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    outer_group = media_slide.shapes.add_group_shape()
    outer_group.name = "Outer transformed group"
    inner_group = outer_group.shapes.add_group_shape()
    inner_group.name = "Inner transformed group"
    anchor = inner_group.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.1),
        Inches(0.1),
        Inches(0.4),
        Inches(0.3),
    )
    anchor.name = "Group coordinate anchor"
    anchor.fill.solid()
    anchor.fill.fore_color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    child = inner_group.shapes.add_shape(
        MSO_SHAPE.CHEVRON,
        Inches(0.55),
        Inches(0.35),
        Inches(1.6),
        Inches(0.9),
    )
    child.name = "Nested transformed chevron"
    child.text = "Grouped DNA"
    child.fill.solid()
    child.fill.fore_color.rgb = RGBColor(0x0F, 0x76, 0x6E)
    inner_group.left = Inches(0.6)
    inner_group.top = Inches(0.4)
    inner_group.width = Inches(3.4)
    inner_group.height = Inches(1.8)
    inner_group.rotation = 13
    outer_group.left = Inches(0.9)
    outer_group.top = Inches(1.5)
    outer_group.width = Inches(6.8)
    outer_group.height = Inches(3.4)
    outer_group.rotation = 27

    picture = media_slide.shapes.add_picture(
        str(picture_path),
        Inches(8.4),
        Inches(1.4),
        Inches(4.0),
        Inches(3.0),
    )
    picture.name = "Cropped flipped ellipse picture"
    picture.crop_left = 0.28
    picture.crop_right = 0.08
    picture.crop_top = 0.12
    picture.crop_bottom = 0.04
    picture.rotation = 8
    picture._element.spPr.xfrm.set("flipV", "1")
    picture._element.spPr.prstGeom.set("prst", "ellipse")

    table_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    table_frame = table_slide.shapes.add_table(
        4,
        4,
        Inches(1.0),
        Inches(1.0),
        Inches(9.0),
        Inches(3.8),
    )
    table_frame.name = "Merged rich table"
    table = table_frame.table
    for column, width in zip(
        table.columns,
        (Inches(1.2), Inches(3.6), Inches(1.8), Inches(2.4)),
        strict=True,
    ):
        column.width = width
    for row, height in zip(
        table.rows,
        (Inches(0.6), Inches(1.0), Inches(1.4), Inches(0.8)),
        strict=True,
    ):
        row.height = height
    table.cell(0, 0).merge(table.cell(0, 2))
    table.cell(0, 0).text = "Three-column heading"
    table.cell(0, 3).text = "Status"
    table.cell(1, 0).merge(table.cell(2, 0))
    table.cell(1, 0).text = "Two-row label"
    table.cell(1, 2).text = "Forecast"
    table.cell(1, 3).text = "On track"
    table.cell(2, 1).text = "Quality"
    table.cell(2, 2).text = "93%"
    table.cell(2, 3).text = "Verified"
    table.cell(3, 0).text = "Owner"
    table.cell(3, 1).text = "Platform"
    table.cell(3, 2).text = "ETA"
    table.cell(3, 3).text = "Q4"
    styled_cell = table.cell(3, 3)
    styled_cell.fill.solid()
    styled_cell.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_2
    styled_cell.margin_left = 0
    styled_cell.margin_top = 0
    styled_cell._tc.tcPr.set("anchor", "ctr")
    styled_cell.text_frame.paragraphs[0]._p.get_or_add_pPr().set("algn", "r")
    left_border = OxmlElement("a:lnL")
    left_border.set("w", "25400")
    border_fill = OxmlElement("a:solidFill")
    border_color = OxmlElement("a:schemeClr")
    border_color.set("val", "accent3")
    border_fill.append(border_color)
    left_border.append(border_fill)
    styled_cell._tc.tcPr.append(left_border)
    rich_cell = table.cell(1, 1)
    rich_cell.text_frame.clear()
    rich_paragraph = rich_cell.text_frame.paragraphs[0]
    strong_run = rich_paragraph.add_run()
    strong_run.text = "Revenue"
    strong_run.font.bold = True
    strong_run.font.color.rgb = RGBColor(0x1D, 0x4E, 0xD8)
    emphasis_run = rich_paragraph.add_run()
    emphasis_run.text = " +42%"
    emphasis_run.font.italic = True

    stacked_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    stacked_data = ChartData()
    stacked_data.categories = ["Plan", "Compose", "Render", "Repair"]
    stacked_data.add_series("Baseline", (30, -20, 85, 55))
    stacked_data.add_series("Increment", (15, -10, 60, 90))
    stacked_frame = stacked_slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED,
        Inches(0.8),
        Inches(0.7),
        Inches(11.7),
        Inches(5.9),
        stacked_data,
    )
    stacked_frame.name = "Stacked delivery chart"
    stacked_chart = stacked_frame.chart
    stacked_chart.has_title = True
    stacked_chart.chart_title.text_frame.text = "Stacked delivery"
    stacked_chart.has_legend = True
    stacked_chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    stacked_chart.value_axis.minimum_scale = -40
    stacked_chart.value_axis.maximum_scale = 100
    stacked_chart.value_axis.major_unit = 20
    stacked_chart.value_axis.has_major_gridlines = True
    stacked_chart.plots[0].has_data_labels = True
    stacked_chart.plots[0].data_labels.show_value = True

    scatter_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    scatter_data = XyChartData()
    scatter_series = scatter_data.add_series("Observed")
    for x_value, y_value in ((1, 18), (2, 42), (3, 81)):
        scatter_series.add_data_point(x_value, y_value)
    scatter_frame = scatter_slide.shapes.add_chart(
        XL_CHART_TYPE.XY_SCATTER,
        Inches(0.8),
        Inches(0.7),
        Inches(11.7),
        Inches(5.9),
        scatter_data,
    )
    scatter_frame.name = "Scatter markers only"
    scatter_chart = scatter_frame.chart
    scatter_chart.has_title = False
    scatter_chart.has_legend = False
    scatter_chart.series[0].marker.style = XL_MARKER_STYLE.CIRCLE
    scatter_chart.series[0].marker.size = 9
    scatter_chart.series[0].format.line.fill.background()

    unsupported_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    radar_data = ChartData()
    radar_data.categories = ["Quality", "Speed", "Reuse"]
    radar_data.add_series("Score", (92, 76, 88))
    radar_frame = unsupported_slide.shapes.add_chart(
        XL_CHART_TYPE.RADAR,
        Inches(0.8),
        Inches(0.7),
        Inches(11.7),
        Inches(5.9),
        radar_data,
    )
    radar_frame.name = "Unsupported radar chart"
    radar_frame.chart.has_title = True
    radar_frame.chart.chart_title.text_frame.text = "Radar source title"

    presentation.save(path)


def _replace_shape_fill(shape, fill):
    properties = shape._element.spPr
    for child in list(properties):
        if child.tag.rsplit("}", 1)[-1] in {
            "solidFill",
            "gradFill",
            "pattFill",
            "noFill",
        }:
            properties.remove(child)
    properties.insert(2, fill)


def _pptx_shape_fidelity_stress(path):
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    for index, (kind, name) in enumerate(
        (
            (MSO_SHAPE.CHEVRON, "Fidelity chevron"),
            (MSO_SHAPE.ARC, "Fidelity arc"),
            (MSO_SHAPE.CLOUD, "Fidelity cloud"),
        )
    ):
        shape = slide.shapes.add_shape(
            kind,
            Inches(0.8 + index * 2.4),
            Inches(0.8),
            Inches(2),
            Inches(1.4),
        )
        shape.name = name
        shape.text = "Actual geometry"
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x25, 0x63, 0xEB)

    gradient = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8),
        Inches(2.8),
        Inches(3),
        Inches(1.2),
    )
    gradient.name = "Fidelity gradient"
    gradient.text = "Gradient"
    gradient_fill = OxmlElement("a:gradFill")
    gradient_stops = OxmlElement("a:gsLst")
    for position, color in (("0", "0EA5E9"), ("100000", "6366F1")):
        stop = OxmlElement("a:gs")
        stop.set("pos", position)
        stop_color = OxmlElement("a:srgbClr")
        stop_color.set("val", color)
        stop.append(stop_color)
        gradient_stops.append(stop)
    gradient_fill.append(gradient_stops)
    gradient_line = OxmlElement("a:lin")
    gradient_line.set("ang", "19500000")
    gradient_fill.append(gradient_line)
    _replace_shape_fill(gradient, gradient_fill)

    patterned = slide.shapes.add_shape(
        MSO_SHAPE.HEXAGON,
        Inches(4.2),
        Inches(2.8),
        Inches(2.4),
        Inches(1.2),
    )
    patterned.name = "Fidelity pattern"
    patterned.text = "Pattern"
    pattern_fill = OxmlElement("a:pattFill")
    pattern_fill.set("prst", "diagBrick")
    for wrapper_name, color in (("fgClr", "EF4444"), ("bgClr", "FEE2E2")):
        wrapper = OxmlElement(f"a:{wrapper_name}")
        pattern_color = OxmlElement("a:srgbClr")
        pattern_color.set("val", color)
        wrapper.append(pattern_color)
        pattern_fill.append(wrapper)
    _replace_shape_fill(patterned, pattern_fill)

    flipped = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(7),
        Inches(2.8),
        Inches(2.4),
        Inches(1.2),
    )
    flipped.name = "Fidelity flipped arrow"
    flipped.text = "flipH"
    flipped._element.spPr.xfrm.set("flipH", "1")

    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.ELBOW,
        Inches(1),
        Inches(5),
        Inches(6),
        Inches(6.4),
    )
    connector.name = "Fidelity elbow connector"
    connector.line.color.rgb = RGBColor(0x25, 0x63, 0xEB)
    connector.line.width = Pt(3)
    line = connector._element.spPr.xpath("./*[local-name()='ln']")[0]
    dash = OxmlElement("a:prstDash")
    dash.set("val", "dash")
    head = OxmlElement("a:headEnd")
    head.set("type", "triangle")
    tail = OxmlElement("a:tailEnd")
    tail.set("type", "diamond")
    line.extend((dash, head, tail))

    presentation.save(path)


def _named_shape(document, name):
    matches = document.xpath(f'//*[@data-shape-name="{name}"]')
    assert len(matches) == 1, f"expected one rendered shape named {name!r}"
    return matches[0]


def _descendants_with_class(element, class_name):
    return element.xpath(
        f'.//*[contains(concat(" ", normalize-space(@class), " "), " {class_name} ")]'
    )


def _css_number(style, property_name, unit):
    match = re.search(
        rf"(?:^|;)\s*{re.escape(property_name)}\s*:\s*"
        rf"(-?\d+(?:\.\d+)?)\s*{re.escape(unit)}(?:;|$)",
        style or "",
    )
    assert match is not None, f"missing {property_name} in {style!r}"
    return float(match.group(1))


def _cqw_padding(style):
    shorthand = re.search(r"(?:^|;)\s*padding\s*:\s*([^;]+)", style or "")
    if shorthand:
        values = [
            float(value)
            for value in re.findall(r"(-?\d+(?:\.\d+)?)\s*cqw", shorthand.group(1))
        ]
        if len(values) == 1:
            return (values[0],) * 4
        if len(values) == 2:
            return (values[0], values[1], values[0], values[1])
        if len(values) == 3:
            return (values[0], values[1], values[2], values[1])
        if len(values) == 4:
            return tuple(values)
    return tuple(
        _css_number(style, property_name, "cqw")
        for property_name in (
            "padding-top",
            "padding-right",
            "padding-bottom",
            "padding-left",
        )
    )


def _crop_fraction(element, attribute):
    value = float(element.get(attribute))
    return value / 100 if value > 1 else value


@pytest.fixture
def pptx_semantic_stress_preview(monkeypatch, tmp_path):
    source = tmp_path / "pptx-semantic-stress.pptx"
    picture_path = tmp_path / "crop-source.png"
    _pptx_semantic_stress(source, picture_path)
    monkeypatch.setattr(preview.settings, "EVOFLUX_CACHE_DIR", str(tmp_path / "cache"))
    rendered = preview.render_document_preview(source).read_text(encoding="utf-8")
    return lxml_html.fromstring(rendered)


def test_render_docx_preview_uses_cache(monkeypatch, tmp_path):
    source = tmp_path / "report.docx"
    _docx(source)
    cache = tmp_path / "cache"
    monkeypatch.setattr(preview.settings, "EVOFLUX_CACHE_DIR", str(cache))

    first = preview.render_document_preview(source)
    second = preview.render_document_preview(source)

    assert first == second
    rendered = first.read_text()
    assert "Content-Security-Policy" in rendered
    assert "Quarterly review" in rendered
    assert "Revenue increased." in rendered
    assert 'class="document-flow"' in rendered
    assert "--page-width:" in rendered


def test_render_docx_preview_does_not_create_missing_header_footer(
    monkeypatch, tmp_path
):
    from docx.parts.hdrftr import FooterPart, HeaderPart

    source = tmp_path / "report-without-header-footer.docx"
    _docx(source)
    monkeypatch.setattr(preview.settings, "EVOFLUX_CACHE_DIR", str(tmp_path / "cache"))

    def fail_if_created(*_args, **_kwargs):
        raise OSError(22, "Invalid argument: extended Windows template path")

    monkeypatch.setattr(HeaderPart, "new", fail_if_created)
    monkeypatch.setattr(FooterPart, "new", fail_if_created)

    rendered = preview.render_document_preview(source).read_text(encoding="utf-8")

    assert "Quarterly review" in rendered
    assert '<div class="document-header-template"></div>' in rendered
    assert '<div class="document-footer-template"></div>' in rendered


def test_render_document_preview_cache_includes_filename(monkeypatch, tmp_path):
    source = tmp_path / "report.docx"
    renamed = tmp_path / "renamed.docx"
    _docx(source)
    renamed.write_bytes(source.read_bytes())
    monkeypatch.setattr(preview.settings, "EVOFLUX_CACHE_DIR", str(tmp_path / "cache"))

    first = preview.render_document_preview(source)
    second = preview.render_document_preview(renamed)

    assert first != second
    assert "<title>report.docx</title>" in first.read_text(encoding="utf-8")
    assert "<title>renamed.docx</title>" in second.read_text(encoding="utf-8")


def test_render_xlsx_preview(monkeypatch, tmp_path):
    source = tmp_path / "workbook.xlsx"
    _xlsx(source)
    monkeypatch.setattr(preview.settings, "EVOFLUX_CACHE_DIR", str(tmp_path / "cache"))

    rendered = preview.render_document_preview(source).read_text()

    assert "Summary" in rendered
    assert "Revenue" in rendered
    assert 'data-formula="=B1*2"' in rendered
    assert ">240</td>" in rendered
    assert "==B1*2" not in rendered
    assert 'data-cell="A1"' in rendered
    assert 'class="column-header" data-column="A"' in rendered
    assert 'class="column-header" data-column="T"' in rendered
    assert 'class="row-number" data-row="1"' in rendered
    assert 'class="row-number" data-row="40"' in rendered
    assert "<h2>Summary</h2>" not in rendered
    assert "html,body{width:100%;height:100%;overflow:hidden}" in rendered


def test_render_xlsx_preview_uses_cached_formula_results_and_visible_fallbacks(
    monkeypatch, tmp_path
):
    source = tmp_path / "formula-cache.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet["A1"] = 1
    sheet["A2"] = 2
    sheet["C1"] = 42
    sheet["C2"] = 84
    sheet["D1"] = "=XLOOKUP(1,A1:A2,C1:C2)"
    sheet["D2"] = "=UNIQUE(A1:A2)"
    workbook.save(source)
    _set_formula_cache(source, "XLOOKUP(1,A1:A2,C1:C2)", "42")
    monkeypatch.setattr(preview.settings, "EVOFLUX_CACHE_DIR", str(tmp_path / "cache"))

    rendered = preview.render_document_preview(source).read_text()

    assert 'data-formula="=XLOOKUP(1,A1:A2,C1:C2)"' in rendered
    assert 'data-display-value="42"' in rendered
    assert ">42</td>" in rendered
    assert 'data-formula="=UNIQUE(A1:A2)"' in rendered
    assert 'data-display-value="=UNIQUE(A1:A2)"' in rendered


def test_render_xlsx_preview_includes_charts_and_images(monkeypatch, tmp_path):
    source = tmp_path / "dashboard.xlsx"
    image_path = tmp_path / "badge.png"
    Image.new("RGB", (32, 24), "#2563eb").save(image_path)
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Dashboard"
    sheet.append(["Month", "Revenue"])
    sheet.append(["Jan", 120])
    sheet.append(["Feb", 180])
    chart = BarChart()
    chart.title = "Revenue trend"
    chart.add_data(
        Reference(sheet, min_col=2, min_row=1, max_row=3), titles_from_data=True
    )
    sheet.add_chart(chart, "D2")
    line_chart = LineChart()
    line_chart.title = "Revenue line"
    line_chart.add_data(
        Reference(sheet, min_col=2, min_row=1, max_row=3), titles_from_data=True
    )
    line_chart.set_categories(Reference(sheet, min_col=1, min_row=2, max_row=3))
    sheet.add_chart(line_chart, "D20")
    sheet.add_image(WorkbookImage(image_path), "D18")
    workbook.save(source)
    monkeypatch.setattr(preview.settings, "EVOFLUX_CACHE_DIR", str(tmp_path / "cache"))

    rendered = preview.render_document_preview(source).read_text()

    assert 'class="workbook-chart-svg"' in rendered
    assert "<rect " in rendered
    assert 'transform="rotate(-32 ' in rendered
    assert 'class="workbook-image"' in rendered
    assert "data:image/png;base64," in rendered


def test_render_pptx_preview(monkeypatch, tmp_path):
    source = tmp_path / "slides.pptx"
    _pptx(source)
    monkeypatch.setattr(preview.settings, "EVOFLUX_CACHE_DIR", str(tmp_path / "cache"))

    rendered = preview.render_document_preview(source).read_text()

    assert "Product launch" in rendered
    assert 'class="slide"' in rendered
    assert 'data-source-layer="slide"' in rendered
    assert "aspect-ratio:" in rendered
    assert "Arial,sans-serif" in rendered


def test_render_pptx_preview_renders_filled_picture_placeholder(monkeypatch, tmp_path):
    source_image = tmp_path / "portrait.png"
    Image.new("RGB", (200, 400), "#2563eb").save(source_image)
    source = tmp_path / "picture-placeholder.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[8])
    placeholder = next(
        shape
        for shape in slide.placeholders
        if shape.placeholder_format.type.name == "PICTURE"
    )
    picture = placeholder.insert_picture(str(source_image))
    picture.name = "Filled picture placeholder"
    presentation.save(source)
    monkeypatch.setattr(preview.settings, "EVOFLUX_CACHE_DIR", str(tmp_path / "cache"))

    rendered = preview.render_document_preview(source).read_text()

    assert "Filled picture placeholder" in rendered
    assert 'class="shape picture-frame"' in rendered
    assert "data:image/png;base64," in rendered
    assert re.search(r'data-crop-(?:top|bottom)="0\.[1-9]', rendered)


def test_render_pptx_preview_skips_picture_with_missing_relationship(
    monkeypatch, tmp_path
):
    source_image = tmp_path / "missing-relationship.png"
    Image.new("RGB", (100, 100), "#2563eb").save(source_image)
    source = tmp_path / "missing-picture-relationship.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_textbox(
        Inches(1), Inches(1), Inches(4), Inches(1)
    ).text = "The rest of the slide remains visible"
    picture = slide.shapes.add_picture(
        str(source_image), Inches(1), Inches(2), Inches(2), Inches(2)
    )
    picture._pic.blipFill.blip.attrib.pop(pptx_qn("r:embed"))
    presentation.save(source)
    monkeypatch.setattr(preview.settings, "EVOFLUX_CACHE_DIR", str(tmp_path / "cache"))

    rendered = preview.render_document_preview(source).read_text(encoding="utf-8")

    assert "The rest of the slide remains visible" in rendered
    assert 'class="picture-missing"' in rendered


def test_render_pptx_preview_keeps_titles_single_line_and_renders_area_charts(
    monkeypatch,
    tmp_path,
):
    source = tmp_path / "area-chart.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(40 / 3)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    title = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.7)
    )
    title.name = "[role:title] Delivery confidence"
    title.text = "Delivery confidence"
    title.text_frame.paragraphs[0].runs[0].font.size = Pt(18)
    _add_chart(
        slide,
        ["Plan", "Compose", "Render", "Repair"],
        {"Confidence": [42, 61, 78, 93]},
        kind=XL_CHART_TYPE.AREA,
        title="Quality gate confidence",
    )
    presentation.save(source)
    monkeypatch.setattr(preview.settings, "EVOFLUX_CACHE_DIR", str(tmp_path / "cache"))

    rendered = preview.render_document_preview(source).read_text()

    assert "white-space:nowrap" in rendered
    assert "font-size:1.8750cqw" in rendered
    assert '<polygon class="chart-area"' in rendered
    assert 'fill-opacity=".24"' in rendered
    assert ">Plan</text>" in rendered
    assert ">Repair</text>" in rendered


def test_render_pptx_preview_labels_pie_and_bar_categories(monkeypatch, tmp_path):
    source = tmp_path / "labelled-charts.pptx"
    presentation = Presentation()
    pie_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    pie = _add_chart(
        pie_slide,
        ["Automation", "Reuse", "Other"],
        {"Value": [48, 34, 18]},
        kind=XL_CHART_TYPE.DOUGHNUT,
        title="Value mix",
    )
    pie.has_legend = True
    pie.plots[0].has_data_labels = True
    bar_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _add_chart(
        bar_slide,
        ["Discover", "Deliver", "Scale"],
        {"Confidence": [42, 78, 93]},
        kind=XL_CHART_TYPE.BAR_CLUSTERED,
        title="Confidence by phase",
    )
    presentation.save(source)
    monkeypatch.setattr(preview.settings, "EVOFLUX_CACHE_DIR", str(tmp_path / "cache"))

    rendered = preview.render_document_preview(source).read_text()

    assert "Automation (48%)" in rendered
    assert "Reuse (34%)" in rendered
    assert ">Discover</text>" in rendered
    assert ">Scale</text>" in rendered


def test_render_pptx_preview_uses_literal_chart_series_names(monkeypatch, tmp_path):
    source = tmp_path / "literal-series-names.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    chart = _add_chart(
        slide,
        ["Geometry", "Text/Table", "Charts", "Median"],
        {
            "Before": [69.3, 77.1, 78.0, 75.52],
            "After": [86.66, 84.5, 84.2, 83.85],
        },
        kind=XL_CHART_TYPE.COLUMN_CLUSTERED,
        title="Fidelity score",
    )
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    for series, name in zip(chart.series, ("Before", "After"), strict=True):
        tx = series._ser.xpath("./*[local-name()='tx']")[0]
        for child in list(tx):
            tx.remove(child)
        value = OxmlElement("c:v")
        value.text = name
        tx.append(value)
    presentation.save(source)
    monkeypatch.setattr(preview.settings, "EVOFLUX_CACHE_DIR", str(tmp_path / "cache"))

    rendered = preview.render_document_preview(source).read_text()
    document = lxml_html.fromstring(rendered)
    legend = _descendants_with_class(document, "chart-legend")[0]

    assert legend.text_content() == "BeforeAfter"
    assert "Series 1" not in legend.text_content()
    assert "Series 2" not in legend.text_content()


def test_render_pptx_preview_resolves_template_theme_colors(monkeypatch, tmp_path):
    source = tmp_path / "theme.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1),
        Inches(1),
        Inches(3),
        Inches(1),
    )
    shape.fill.solid()
    shape.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
    presentation.save(source)
    monkeypatch.setattr(preview.settings, "EVOFLUX_CACHE_DIR", str(tmp_path / "cache"))

    rendered = preview.render_document_preview(source).read_text()

    assert 'data-shape-fill="solid"' in rendered
    assert 'fill="#4F81BD"' in rendered


def test_render_pptx_preview_supports_bullets_columns_and_connectors(
    monkeypatch,
    tmp_path,
):
    source = tmp_path / "rich-office.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    text_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(5), Inches(2))
    text_box.text_frame._txBody.bodyPr.set("numCol", "2")
    for index, value in enumerate(("First point", "Second point")):
        paragraph = (
            text_box.text_frame.paragraphs[0]
            if index == 0
            else text_box.text_frame.add_paragraph()
        )
        paragraph.text = value
        paragraph.level = index
        properties = paragraph._p.get_or_add_pPr()
        properties.set("marL", str(int(Inches(0.6 + index * 0.3))))
        properties.set("indent", str(-int(Inches(0.25))))
        bullet = OxmlElement("a:buChar")
        bullet.set("char", "•")
        properties.append(bullet)

    boxes = []
    for index, label in enumerate(("Discover", "Decide", "Deliver")):
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1 + index * 3.5),
            Inches(3.3),
            Inches(2.5),
            Inches(1.2),
        )
        box.text = label
        boxes.append(box)
    for left, right in zip(boxes[:-1], boxes[1:], strict=True):
        slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            left.left + left.width,
            left.top + left.height // 2,
            right.left,
            right.top + right.height // 2,
        )
    presentation.save(source)
    monkeypatch.setattr(preview.settings, "EVOFLUX_CACHE_DIR", str(tmp_path / "cache"))

    rendered = preview.render_document_preview(source).read_text()

    assert 'class="bullet-marker"' in rendered
    assert "column-count:2" in rendered
    assert "column-fill:auto" in rendered
    document = lxml_html.fromstring(rendered)
    markers = _descendants_with_class(document, "bullet-marker")
    assert [marker.text for marker in markers] == ["•", "•"]
    marker_paragraphs = [marker.xpath("ancestor::p[1]")[0] for marker in markers]
    assert _css_number(
        marker_paragraphs[1].get("style", ""), "padding-left", "cqw"
    ) > _css_number(marker_paragraphs[0].get("style", ""), "padding-left", "cqw")
    first_run = marker_paragraphs[0].xpath('.//span[@class="paragraph-content"]/span')[
        0
    ]
    assert _css_number(
        markers[0].get("style", ""), "font-size", "cqw"
    ) == pytest.approx(
        _css_number(first_run.get("style", ""), "font-size", "cqw"),
        abs=0.001,
    )
    assert rendered.count('class="shape connector"') == 2
    assert "Discover" in rendered
    assert "Deliver" in rendered
    assert '<svg class="shape-surface"' in rendered
    assert 'data-preset-geometry="roundRect"' in rendered
    assert '<rect x="0" y="0" width="1000" height="1000" rx="80"' in rendered


def test_render_pptx_preview_preserves_shape_surfaces_and_elbow_semantics(
    monkeypatch,
    tmp_path,
):
    source = tmp_path / "shape-fidelity.pptx"
    _pptx_shape_fidelity_stress(source)
    monkeypatch.setattr(preview.settings, "EVOFLUX_CACHE_DIR", str(tmp_path / "cache"))

    document = lxml_html.fromstring(
        preview.render_document_preview(source).read_text(encoding="utf-8")
    )

    chevron = _named_shape(document, "Fidelity chevron")
    chevron_surface = _descendants_with_class(chevron, "shape-surface")[0]
    assert chevron_surface.get("data-preset-geometry") == "chevron"
    assert chevron_surface.xpath(".//*[local-name()='polygon']")
    assert "drop-shadow(" in chevron_surface.get("style", "")
    assert chevron.get("data-shape-font-color") == "theme"

    arc = _named_shape(document, "Fidelity arc")
    arc_path = _descendants_with_class(arc, "shape-surface")[0].xpath(
        ".//*[local-name()='path']"
    )[0]
    assert "A500 500" in arc_path.get("d", "")
    cloud = _named_shape(document, "Fidelity cloud")
    cloud_path = _descendants_with_class(cloud, "shape-surface")[0].xpath(
        ".//*[local-name()='path']"
    )[0]
    assert cloud_path.get("d", "").count("C") >= 6

    gradient = _named_shape(document, "Fidelity gradient")
    gradient_surface = _descendants_with_class(gradient, "shape-surface")[0]
    assert gradient_surface.get("data-shape-fill") == "gradient"
    assert gradient_surface.xpath(".//*[local-name()='lineargradient']")
    assert len(gradient_surface.xpath(".//*[local-name()='stop']")) == 2

    pattern = _named_shape(document, "Fidelity pattern")
    pattern_surface = _descendants_with_class(pattern, "shape-surface")[0]
    assert pattern_surface.get("data-pattern") == "diagBrick"
    assert pattern_surface.xpath(".//*[local-name()='pattern']")

    flipped = _named_shape(document, "Fidelity flipped arrow")
    assert "scaleX(-1)" in flipped.get("style", "")
    assert flipped.get("data-text-unflip") == "horizontal"
    assert flipped.get("data-default-text-align") == "center"

    connector = _named_shape(document, "Fidelity elbow connector")
    assert connector.get("data-connector-geometry") == "bentConnector3"
    connector_path = connector.xpath("./*[local-name()='path']")[-1]
    assert re.fullmatch(r"M0 0H\d+\.\d+V1000H1000", connector_path.get("d", ""))
    assert connector_path.get("stroke-dasharray")
    assert connector_path.get("marker-start", "").startswith("url(#connector-head-")
    assert connector_path.get("marker-end", "").startswith("url(#connector-tail-")
    assert len(connector.xpath(".//*[local-name()='marker']")) == 2


def test_render_pptx_preview_resolves_inherited_master_text_responsively(
    pptx_semantic_stress_preview,
):
    document = pptx_semantic_stress_preview
    title = _named_shape(document, "Inherited master title placeholder")
    title_span = title.xpath(
        './/span[normalize-space(text())="Inherited master title"]'
    )[0]
    title_style = title_span.get("style", "")

    # A neutral wrapper with the browser's default 16px font creates an extra
    # inline strut around large native runs.  That inflates every line box and
    # clips the final line in the fixed PowerPoint text frame.
    assert title_span.getparent().tag == "p"
    assert _css_number(title_style, "font-size", "cqw") == pytest.approx(
        4.5833, abs=0.001
    )
    assert _css_number(title_style, "letter-spacing", "cqw") == pytest.approx(
        0.125, abs=0.001
    )
    assert "Aptos Display" in title_style
    assert "font-weight:700" in title_style
    title_paragraph = title_span.xpath("ancestor::p[1]")[0]
    assert "text-align:center" in title_paragraph.get("style", "")
    title_frame = _descendants_with_class(title, "text-frame")[0]
    assert _cqw_padding(title_frame.get("style", "")) == pytest.approx(
        (0.625, 0.9375, 1.25, 1.875), abs=0.001
    )

    body = _named_shape(document, "Inherited master body placeholder")
    body_span = body.xpath('.//span[normalize-space(text())="Inherited master body"]')[
        0
    ]
    body_style = body_span.get("style", "")
    assert _css_number(body_style, "font-size", "cqw") == pytest.approx(
        3.3333, abs=0.001
    )
    assert _css_number(body_style, "letter-spacing", "cqw") == pytest.approx(
        0.0833, abs=0.001
    )
    assert 'font-family:"Aptos"' in body_style
    body_frame = _descendants_with_class(body, "text-frame")[0]
    assert _cqw_padding(body_frame.get("style", "")) == pytest.approx(
        (0.9375, 1.25, 0.9375, 2.5), abs=0.001
    )
    body_paragraph = body_span.xpath("ancestor::p[1]")[0]
    paragraph_style = body_paragraph.get("style", "")
    line_height = re.search(
        r"(?:^|;)\s*line-height\s*:\s*(\d+(?:\.\d+)?)(?:em)?(?:;|$)",
        paragraph_style,
    )
    assert line_height is not None
    assert float(line_height.group(1)) == pytest.approx(1.25)
    assert _css_number(paragraph_style, "margin-top", "cqw") == pytest.approx(
        1.25, abs=0.001
    )
    assert _css_number(paragraph_style, "margin-bottom", "cqw") == pytest.approx(
        0.625, abs=0.001
    )


def test_render_pptx_preview_emits_escaped_notes_and_title_aware_slide_labels(
    pptx_semantic_stress_preview,
):
    title = _named_shape(
        pptx_semantic_stress_preview, "Inherited master title placeholder"
    )
    slide_wrap = title.xpath(
        'ancestor::*[contains(concat(" ", normalize-space(@class), " "), '
        '" slide-wrap ")][1]'
    )[0]

    label = slide_wrap.get("data-preview-label", "")
    assert "Slide 1" in label
    assert "Inherited master title" in label
    metadata = _descendants_with_class(slide_wrap, "slide-notes-metadata")
    assert len(metadata) == 1
    notes = metadata[0].get("data-preview-notes") or metadata[0].text_content()
    assert notes == 'Presenter note <verify> & "repair" before publish.'
    assert metadata[0].get("hidden") is not None
    assert not metadata[0].xpath(".//verify")


def test_render_pptx_preview_preserves_recursive_group_transform_and_identity(
    pptx_semantic_stress_preview,
):
    document = pptx_semantic_stress_preview
    outer = _named_shape(document, "Outer transformed group")
    inner = _named_shape(document, "Inner transformed group")
    child = _named_shape(document, "Nested transformed chevron")

    assert "shape-group" in outer.get("class", "").split()
    assert "shape-group" in inner.get("class", "").split()
    assert inner in outer.xpath(".//*")
    assert child in inner.xpath(".//*")
    for shape, name in (
        (outer, "Outer transformed group"),
        (inner, "Inner transformed group"),
        (child, "Nested transformed chevron"),
    ):
        assert shape.get("data-shape-id")
        assert shape.get("data-shape-name") == name
        assert shape.get("data-qa-label") == name
        assert shape.get("data-source-layer") == "slide"

    outer_style = outer.get("style", "")
    inner_style = inner.get("style", "")
    assert "transform:" in outer_style
    assert "transform:" in inner_style
    combined_transform = outer_style + inner_style
    assert any(
        transform in combined_transform
        for transform in ("matrix(", "rotate(", "scale(")
    )
    child_surface = _descendants_with_class(child, "shape-surface")[0]
    assert child_surface.get("data-preset-geometry") == "chevron"
    assert child_surface.xpath(".//*[local-name()='polygon']")


def test_render_pptx_preview_applies_picture_crop_flip_and_geometry_clip(
    pptx_semantic_stress_preview,
):
    frame = _named_shape(
        pptx_semantic_stress_preview, "Cropped flipped ellipse picture"
    )

    assert "picture-frame" in frame.get("class", "").split()
    assert _crop_fraction(frame, "data-crop-left") == pytest.approx(0.28)
    assert _crop_fraction(frame, "data-crop-right") == pytest.approx(0.08)
    assert _crop_fraction(frame, "data-crop-top") == pytest.approx(0.12)
    assert _crop_fraction(frame, "data-crop-bottom") == pytest.approx(0.04)
    frame_style = frame.get("style", "")
    assert "clip-path:ellipse(" in frame_style
    assert "scaleY(-1)" in frame_style

    image = _descendants_with_class(frame, "picture")[0]
    image_style = image.get("style", "")
    assert "object-fit:fill" in image_style
    assert _css_number(image_style, "left", "%") < 0
    assert _css_number(image_style, "top", "%") < 0
    assert _css_number(image_style, "width", "%") > 100
    assert _css_number(image_style, "height", "%") > 100


def test_render_pptx_preview_preserves_table_grid_merges_and_rich_runs(
    pptx_semantic_stress_preview,
):
    frame = _named_shape(pptx_semantic_stress_preview, "Merged rich table")
    table = frame.xpath(".//table")[0]
    assert table.get("data-table-style-id") == "5C22544A-7EE6-4342-B048-85BDC9FD1C3A"
    columns = table.xpath("./colgroup/col")
    assert len(columns) == 4
    assert [
        _css_number(column.get("style", ""), "width", "%") for column in columns
    ] == pytest.approx((13.3333, 40.0, 20.0, 26.6667), abs=0.01)

    rows = table.xpath("./tbody/tr|./tr")
    assert len(rows) == 4
    assert [
        _css_number(row.get("style", ""), "height", "%") for row in rows
    ] == pytest.approx((15.7895, 26.3158, 36.8421, 21.0526), abs=0.01)
    cells = table.xpath(".//td")
    assert len(cells) == 13
    assert [cell.get("colspan") for cell in cells].count("3") == 1
    assert [cell.get("rowspan") for cell in cells].count("2") == 1
    assert (
        "Three-column heading" in table.xpath('.//td[@colspan="3"]')[0].text_content()
    )
    assert "Two-row label" in table.xpath('.//td[@rowspan="2"]')[0].text_content()

    heading_cell = table.xpath('.//td[@colspan="3"]')[0]
    heading_style = heading_cell.get("style", "").casefold()
    assert "background:#4f81bd" in heading_style
    assert "color:#ffffff" in heading_style
    assert "font-weight:700" in heading_style
    heading_run = heading_cell.xpath('.//span[contains(text(), "Three-column")]')[0]
    assert _css_number(
        heading_run.get("style", ""), "font-size", "cqw"
    ) == pytest.approx(1.875, abs=0.001)

    styled_cell = next(cell for cell in cells if cell.text_content().strip() == "Q4")
    styled_style = styled_cell.get("style", "")
    assert "background:#C0504D" in styled_style
    assert "vertical-align:middle" in styled_style
    assert "border-left:0.2083cqw solid #9BBB59" in styled_style
    assert _cqw_padding(styled_style)[0] == pytest.approx(0, abs=0.001)
    assert _cqw_padding(styled_style)[3] == pytest.approx(0, abs=0.001)
    assert "text-align:right" in styled_cell.xpath(".//p")[0].get("style", "")

    rich_cell = next(cell for cell in cells if "Revenue" in cell.text_content())
    assert rich_cell.xpath("./p|.//p")
    revenue_run = rich_cell.xpath('.//span[normalize-space(text())="Revenue"]')[0]
    emphasis_run = rich_cell.xpath('.//span[contains(text(), "+42%") ]')[0]
    assert "font-weight:700" in revenue_run.get("style", "")
    assert "color:#1d4ed8" in revenue_run.get("style", "").casefold()
    assert "font-style:italic" in emphasis_run.get("style", "")


def test_render_pptx_preview_renders_stacked_chart_semantics(
    pptx_semantic_stress_preview,
):
    frame = _named_shape(pptx_semantic_stress_preview, "Stacked delivery chart")
    svg = _descendants_with_class(frame, "chart-svg")[0]
    classes = svg.get("class", "").split()

    assert "chart-column" in classes
    assert "chart-stacked" in classes
    assert len(_descendants_with_class(svg, "chart-axis")) >= 2
    assert len(_descendants_with_class(svg, "chart-gridline")) >= 6
    legends = _descendants_with_class(svg, "chart-legend")
    assert len(legends) == 1
    assert legends[0].get("data-legend-position") == "bottom"
    assert "Baseline" in legends[0].text_content()
    assert "Increment" in legends[0].text_content()
    axis_labels = [
        label.text_content()
        for label in _descendants_with_class(svg, "chart-axis-label")
    ]
    assert axis_labels == ["100", "80", "60", "40", "20", "0", "-20", "-40"]
    gridlines = _descendants_with_class(svg, "chart-gridline")
    assert all(float(line.get("y1")) == float(line.get("y2")) for line in gridlines)

    columns = _descendants_with_class(svg, "chart-column")
    assert len(columns) == 8
    for column in columns:
        top = float(column.get("y"))
        bottom = top + float(column.get("height"))
        assert 28 <= top <= 248
        assert 28 <= bottom <= 248

    data_labels = _descendants_with_class(svg, "chart-data-label")
    assert len(data_labels) == 8
    assert all(28 <= float(label.get("y")) <= 248 for label in data_labels)
    assert {label.text_content() for label in data_labels} >= {"85", "60", "90"}
    assert "-40" in svg.text_content()
    assert "100" in svg.text_content()


def test_render_pptx_preview_renders_scatter_markers_without_invented_title(
    pptx_semantic_stress_preview,
):
    frame = _named_shape(pptx_semantic_stress_preview, "Scatter markers only")
    svg = _descendants_with_class(frame, "chart-svg")[0]

    assert "chart-scatter" in svg.get("class", "").split()
    markers = _descendants_with_class(svg, "chart-marker")
    assert len(markers) == 3
    assert all(marker.tag == "ellipse" for marker in markers)
    assert all(marker.get("data-marker-style") == "circle" for marker in markers)
    assert not svg.xpath(".//polyline")
    x_axis_labels = [
        label.text_content()
        for label in _descendants_with_class(svg, "chart-axis-label-x")
    ]
    y_axis_labels = [
        label.text_content()
        for label in _descendants_with_class(svg, "chart-axis-label-y")
    ]
    assert x_axis_labels == ["0", "1", "2", "3", "4"]
    assert y_axis_labels == ["90", "80", "70", "60", "50", "40", "30", "20", "10", "0"]
    assert not _descendants_with_class(frame, "chart-title")
    assert not frame.xpath(".//strong")
    assert "Chart" not in frame.text_content()


def test_chart_auto_ticks_and_default_marker_match_office_density():
    axis = SimpleNamespace(major_unit=None)
    chart = SimpleNamespace(value_axis=axis)

    assert preview._chart_tick_values(chart, -40, 100) == [
        -40,
        -20,
        0,
        20,
        40,
        60,
        80,
        100,
    ]

    series = SimpleNamespace(marker=SimpleNamespace(size=None, style=None))
    marker = lxml_html.fragment_fromstring(
        preview._chart_marker_svg(
            series,
            x=100,
            y=80,
            color="#4472c4",
            x_compensation=1.75,
            default_style="DIAMOND",
        )
    )
    assert marker.tag == "polygon"
    assert marker.get("data-marker-style") == "diamond"


def test_render_pptx_preview_marks_unsupported_chart_instead_of_drawing_bars(
    pptx_semantic_stress_preview,
):
    frame = _named_shape(pptx_semantic_stress_preview, "Unsupported radar chart")
    unsupported = _descendants_with_class(frame, "chart-unsupported")

    assert len(unsupported) == 1
    assert unsupported[0].get("data-chart-type") == "radar"
    assert "unsupported" in unsupported[0].text_content().casefold()
    assert not _descendants_with_class(frame, "chart-svg")
    assert not frame.xpath(".//rect")


def test_render_document_preview_invalidates_when_source_changes(
    monkeypatch,
    tmp_path,
):
    source = tmp_path / "workbook.xlsx"
    _xlsx(source)
    monkeypatch.setattr(preview.settings, "EVOFLUX_CACHE_DIR", str(tmp_path / "cache"))

    first = preview.render_document_preview(source)
    workbook = Workbook()
    workbook.active["A1"] = "Changed"
    workbook.save(source)
    second = preview.render_document_preview(source)

    assert first != second


def test_render_document_preview_rejects_unsupported_file(tmp_path):
    source = tmp_path / "legacy.xls"
    source.write_bytes(b"legacy")
    with pytest.raises(preview.DocumentPreviewUnsupportedError):
        preview.render_document_preview(source)


def test_render_document_preview_surfaces_parse_error(monkeypatch, tmp_path):
    source = tmp_path / "broken.pptx"
    source.write_bytes(b"not an OpenXML package")
    monkeypatch.setattr(preview.settings, "EVOFLUX_CACHE_DIR", str(tmp_path / "cache"))

    with pytest.raises(preview.DocumentPreviewError, match="Could not render"):
        preview.render_document_preview(source)


def test_render_document_preview_rejects_oversized_generated_html(
    monkeypatch,
    tmp_path,
):
    source = tmp_path / "report.docx"
    _docx(source)
    monkeypatch.setattr(preview.settings, "EVOFLUX_CACHE_DIR", str(tmp_path / "cache"))
    monkeypatch.setattr(preview, "MAX_DOCUMENT_PREVIEW_HTML_BYTES", 16)
    monkeypatch.setattr(preview, "_render_source", lambda _source: "x" * 17)

    with pytest.raises(
        preview.DocumentPreviewUnsupportedError,
        match="too large for the in-app viewer",
    ):
        preview.render_document_preview(source)


def test_render_pdf_preview_is_bounded_and_cleans_rasters(monkeypatch, tmp_path):
    from app.services.document_preview import pdf

    source = tmp_path / "large.pdf"
    source.write_bytes(b"pdf fixture")
    monkeypatch.setattr(preview.settings, "EVOFLUX_CACHE_DIR", str(tmp_path / "cache"))
    captured = {}

    def fake_render(source_path, render_dir, **kwargs):
        captured.update(kwargs)
        render_dir.mkdir(parents=True)
        page = render_dir / "page-001.png"
        page.write_bytes(b"preview")
        return [page]

    monkeypatch.setattr(pdf, "count_pdf_pages", lambda _source: 5)
    monkeypatch.setattr(pdf, "render_pdf_pages", fake_render)

    rendered = preview._render_pdf(source)

    assert "Preview shows 1 of 5 pages" in rendered
    assert captured == {
        "dpi": 120,
        "max_pages": preview.MAX_PDF_PREVIEW_PAGES,
        "max_total_bytes": preview.MAX_PDF_PREVIEW_RASTER_BYTES,
        "max_pixels_per_page": preview.MAX_PDF_PREVIEW_PIXELS_PER_PAGE,
    }
    assert not list((tmp_path / "cache").rglob("*-pages"))
