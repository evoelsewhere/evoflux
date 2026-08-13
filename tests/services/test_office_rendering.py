from __future__ import annotations

from copy import deepcopy
from pathlib import Path
from types import SimpleNamespace

from docx import Document
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.shapes.placeholder import PicturePlaceholder
from pptx.util import Inches, Pt

from app.agent.builtin_plugins.documents import rendering
from app.agent.builtin_plugins.documents.rendering import internal


def test_internal_renderer_is_always_available() -> None:
    assert rendering.renderer_available() is True


def test_internal_renderer_uses_bundled_fonts_at_requested_size() -> None:
    assert internal._FONT_ROOT.is_dir()
    regular = internal._font(18)
    bold = internal._font(22, bold=True)
    assert isinstance(regular, ImageFont.FreeTypeFont)
    assert isinstance(bold, ImageFont.FreeTypeFont)
    assert regular.size == 18
    assert bold.size == 22


def test_internal_renderer_maps_common_office_fonts_to_unicode_safe_fallback() -> None:
    regular = internal._font(18, family="Arial")
    bold = internal._font(18, family="Helvetica Neue", bold=True)

    assert isinstance(regular, ImageFont.FreeTypeFont)
    assert isinstance(bold, ImageFont.FreeTypeFont)
    assert Path(regular.path).name == "WorkSans-Regular.ttf"
    assert Path(bold.path).name == "WorkSans-Bold.ttf"
    assert regular.getbbox("Từ tín hiệu người dùng") is not None


def test_render_pages_creates_docx_preview_without_external_binary(
    tmp_path: Path,
) -> None:
    source = tmp_path / "letter.docx"
    document = Document()
    document.add_heading("Portable preview", level=1)
    document.add_paragraph("No office suite is installed or launched.")
    document.save(str(source))

    pages, issues = rendering.render_pages(
        source, tmp_path / "previews", code_prefix="docx"
    )

    assert issues == []
    assert [path.name for path in pages] == ["page-001.png"]
    assert pages[0].stat().st_size > 0


def test_render_pages_reports_unsupported_format(tmp_path: Path) -> None:
    source = tmp_path / "notes.txt"
    source.write_text("plain text", encoding="utf-8")

    pages, issues = rendering.render_pages(
        source, tmp_path / "previews", code_prefix="document"
    )

    assert pages == []
    assert issues[0]["code"] == "document-render-failed"
    assert "unsupported render format" in issues[0]["message"]


def _solid_rectangle(slide, *, x, y, width, height, color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*color)
    shape.line.fill.background()
    return shape


def test_pptx_renderer_paints_master_and_layout_shapes_behind_slide_shapes(
    tmp_path: Path,
) -> None:
    source = tmp_path / "template-layers.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(10)
    presentation.slide_height = Inches(5)
    layout = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(layout)

    master_shape = _solid_rectangle(
        slide,
        x=0,
        y=0,
        width=10,
        height=5,
        color=(255, 0, 0),
    )
    layout_shape = _solid_rectangle(
        slide,
        x=1,
        y=1,
        width=4,
        height=3,
        color=(0, 255, 0),
    )
    _solid_rectangle(
        slide,
        x=2,
        y=2,
        width=1,
        height=1,
        color=(0, 0, 255),
    )
    layout.slide_master.shapes._spTree.insert_element_before(  # noqa: SLF001
        deepcopy(master_shape._element),  # noqa: SLF001
        "p:extLst",
    )
    layout.shapes._spTree.insert_element_before(  # noqa: SLF001
        deepcopy(layout_shape._element),  # noqa: SLF001
        "p:extLst",
    )
    slide.shapes._spTree.remove(master_shape._element)  # noqa: SLF001
    slide.shapes._spTree.remove(layout_shape._element)  # noqa: SLF001
    presentation.save(source)

    [rendered] = internal.render_pptx_pages(source, tmp_path / "layers", width=1000)

    with Image.open(rendered).convert("RGB") as image:
        assert image.getpixel((50, 50)) == (255, 0, 0)
        assert image.getpixel((150, 150)) == (0, 255, 0)
        assert image.getpixel((250, 250)) == (0, 0, 255)
        assert image.getpixel((900, 450)) == (255, 0, 0)


def test_pptx_renderer_deduplicates_inherited_placeholders(tmp_path: Path) -> None:
    source = tmp_path / "placeholder-deduplication.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(10)
    presentation.slide_height = Inches(5)
    layout = presentation.slide_layouts[5]
    layout_title = layout.placeholders.get(0)
    assert layout_title is not None
    layout_title.text = "LAYOUT GHOST"
    layout_run = layout_title.text_frame.paragraphs[0].runs[0]
    layout_run.font.size = Pt(28)
    layout_run.font.color.rgb = RGBColor(255, 0, 0)

    slide = presentation.slides.add_slide(layout)
    slide_title = slide.shapes.title
    assert slide_title is not None
    slide_title.text = "SLIDE TITLE"
    slide_run = slide_title.text_frame.paragraphs[0].runs[0]
    slide_run.font.size = Pt(28)
    slide_run.font.color.rgb = RGBColor(0, 0, 255)
    presentation.save(source)

    [rendered] = internal.render_pptx_pages(source, tmp_path / "dedup", width=1000)

    with Image.open(rendered).convert("RGB") as image:
        pixels = list(image.getdata())
    red_pixels = sum(
        red > 180 and green < 90 and blue < 90 for red, green, blue in pixels
    )
    blue_pixels = sum(
        blue > 150 and red < 90 and green < 90 for red, green, blue in pixels
    )
    assert red_pixels == 0
    assert blue_pixels > 100


def test_pptx_renderer_applies_picture_crop_flip_and_rotation(tmp_path: Path) -> None:
    source_image = tmp_path / "picture-source.png"
    picture = Image.new("RGB", (100, 50), "#ff0000")
    ImageDraw.Draw(picture).rectangle((50, 0, 99, 49), fill="#0000ff")
    picture.save(source_image)
    green_image = tmp_path / "rotation-source.png"
    Image.new("RGB", (100, 50), "#00ff00").save(green_image)

    source = tmp_path / "picture-transforms.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(10)
    presentation.slide_height = Inches(5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    cropped = slide.shapes.add_picture(
        str(source_image), Inches(0.5), Inches(0.5), Inches(2), Inches(1)
    )
    cropped.crop_left = 0.5
    flipped = slide.shapes.add_picture(
        str(source_image), Inches(3), Inches(0.5), Inches(2), Inches(1)
    )
    flipped._element.spPr.xfrm.set("flipH", "1")  # noqa: SLF001
    rotated = slide.shapes.add_picture(
        str(green_image), Inches(6), Inches(2), Inches(2), Inches(1)
    )
    rotated.rotation = 90
    presentation.save(source)

    [rendered] = internal.render_pptx_pages(
        source, tmp_path / "picture-transforms", width=1000
    )

    with Image.open(rendered).convert("RGB") as image:
        assert image.getpixel((100, 100)) == (0, 0, 255)
        assert image.getpixel((325, 100)) == (0, 0, 255)
        assert image.getpixel((475, 100)) == (255, 0, 0)
        assert image.getpixel((700, 160)) == (0, 255, 0)
        assert image.getpixel((610, 250)) == (255, 255, 255)


def test_pptx_renderer_renders_filled_picture_placeholder_crop(tmp_path: Path) -> None:
    source_image = tmp_path / "placeholder-source.png"
    picture = Image.new("RGB", (100, 50), "#ff0000")
    ImageDraw.Draw(picture).rectangle((50, 0, 99, 49), fill="#0000ff")
    picture.save(source_image)

    source = tmp_path / "filled-picture-placeholder.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(10)
    presentation.slide_height = Inches(5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[8])
    picture_placeholder = next(
        placeholder
        for placeholder in slide.placeholders
        if isinstance(placeholder, PicturePlaceholder)
    )
    filled = picture_placeholder.insert_picture(str(source_image))
    filled.crop_left = 0.5
    presentation.save(source)

    [rendered] = internal.render_pptx_pages(
        source, tmp_path / "filled-picture-placeholder", width=1000
    )

    with Image.open(rendered).convert("RGB") as image:
        center_x = round(
            (filled.left + filled.width / 2) / presentation.slide_width * 1000
        )
        center_y = round(
            (filled.top + filled.height / 2) / presentation.slide_height * 500
        )
        assert image.getpixel((center_x, center_y)) == (0, 0, 255)


def test_pptx_renderer_accepts_shapes_without_explicit_transform() -> None:
    shape = SimpleNamespace(
        _element=SimpleNamespace(spPr=SimpleNamespace(xfrm=None)),
        rotation=0,
    )

    transformed = internal._pptx_transform_image(
        Image.new("RGBA", (20, 10), "#2563eb"), shape
    )

    assert transformed.size == (20, 10)


def test_pptx_renderer_preserves_rich_text_layout_and_bullets(tmp_path: Path) -> None:
    source = tmp_path / "rich-text.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(10)
    presentation.slide_height = Inches(5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(3))
    frame = text_box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.5)
    frame.margin_right = Inches(0.25)
    frame.margin_top = Inches(0.25)
    frame.margin_bottom = Inches(0.25)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    heading = frame.paragraphs[0]
    heading.alignment = PP_ALIGN.RIGHT
    red_run = heading.add_run()
    red_run.text = "RED"
    red_run.font.size = Pt(30)
    red_run.font.bold = True
    red_run.font.color.rgb = RGBColor(255, 0, 0)
    blue_run = heading.add_run()
    blue_run.text = " BLUE"
    blue_run.font.size = Pt(18)
    blue_run.font.italic = True
    blue_run.font.underline = True
    blue_run.font.color.rgb = RGBColor(0, 0, 255)

    bullet = frame.add_paragraph()
    bullet_run = bullet.add_run()
    bullet_run.text = "Bullet"
    bullet_run.font.size = Pt(20)
    bullet_run.font.color.rgb = RGBColor(0, 128, 0)
    properties = bullet._p.get_or_add_pPr()  # noqa: SLF001
    properties.set("marL", str(int(Inches(0.5))))
    properties.set("indent", str(-int(Inches(0.25))))
    marker = OxmlElement("a:buChar")
    marker.set("char", "•")
    properties.append(marker)
    presentation.save(source)

    [rendered] = internal.render_pptx_pages(source, tmp_path / "rich-text", width=1000)

    with Image.open(rendered).convert("RGB") as image:
        color_points: dict[str, list[tuple[int, int]]] = {
            "red": [],
            "blue": [],
            "green": [],
        }
        for y in range(image.height):
            for x in range(image.width):
                red, green, blue = image.getpixel((x, y))
                if red > 180 and green < 90 and blue < 90:
                    color_points["red"].append((x, y))
                if blue > 150 and red < 90 and green < 90:
                    color_points["blue"].append((x, y))
                if green > 60 and red < 90 and blue < 90:
                    color_points["green"].append((x, y))

    assert len(color_points["red"]) > 1_000
    assert min(x for x, _ in color_points["red"]) > 500
    assert min(y for _, y in color_points["red"]) > 200
    assert len(color_points["blue"]) > 300
    assert max(x for x, _ in color_points["blue"]) >= 670
    assert len(color_points["green"]) > 350
    assert min(x for x, _ in color_points["green"]) < 190
    assert max(x for x, _ in color_points["green"]) > 250
