"""The rendered preview must stay machine-readable for layout inspection."""

from __future__ import annotations

import os
import stat
from pathlib import Path

import pytest

from app.services.document_preview.inspection import (
    ELEMENT_LABEL_ATTRIBUTE,
    PREVIEW_ITEM_ATTRIBUTE,
    PREVIEW_LABEL_ATTRIBUTE,
    summarize_document_preview,
)
from app.services.document_preview.security import mark_cached_preview_used

pytest.importorskip("pptx", reason="office-preview extra not installed")


def _write(tmp_path: Path, body: str) -> Path:
    path = tmp_path / "preview.html"
    path.write_text(f"<html><body>{body}</body></html>", encoding="utf-8")
    return path


def test_reads_items_elements_text_and_geometry(tmp_path):
    preview = _write(
        tmp_path,
        f'<article {PREVIEW_ITEM_ATTRIBUTE} {PREVIEW_LABEL_ATTRIBUTE}="Slide 1 — Intro">'
        f'<div {ELEMENT_LABEL_ATTRIBUTE}="Title" '
        'style="left:6.0%;top:16.0%;width:45.0%;height:10.0%">'
        "<span>Hello   world</span></div>"
        "</article>",
    )

    summary = summarize_document_preview(preview)

    assert [item.label for item in summary.items] == ["Slide 1 — Intro"]
    element = summary.items[0].elements[0]
    assert element.label == "Title"
    assert element.text == "Hello world"
    assert (element.left, element.top, element.width, element.height) == (
        6.0,
        16.0,
        45.0,
        10.0,
    )
    assert element.out_of_bounds() is False


@pytest.mark.parametrize(
    "style",
    [
        "left:110.0%;top:5.0%;width:20.0%;height:10.0%",  # past the right edge
        "left:5.0%;top:95.0%;width:20.0%;height:20.0%",  # past the bottom edge
        "left:-20.0%;top:5.0%;width:30.0%;height:10.0%",  # before the left edge
    ],
)
def test_flags_elements_outside_the_page_box(tmp_path, style):
    preview = _write(
        tmp_path,
        f'<article {PREVIEW_ITEM_ATTRIBUTE} {PREVIEW_LABEL_ATTRIBUTE}="Slide 1">'
        f'<div {ELEMENT_LABEL_ATTRIBUTE}="Stray" style="{style}">text</div>'
        "</article>",
    )

    summary = summarize_document_preview(preview)

    assert [element.label for element in summary.items[0].out_of_bounds()] == ["Stray"]


def test_elements_without_geometry_are_never_flagged(tmp_path):
    preview = _write(
        tmp_path,
        f'<article {PREVIEW_ITEM_ATTRIBUTE} {PREVIEW_LABEL_ATTRIBUTE}="Sheet1">'
        f'<td {ELEMENT_LABEL_ATTRIBUTE}="A1">42</td>'
        "</article>",
    )

    summary = summarize_document_preview(preview)

    assert summary.items[0].out_of_bounds() == []
    assert summary.items[0].elements[0].text == "42"


def test_budgets_bound_items_elements_and_text(tmp_path):
    items = "".join(
        f'<article {PREVIEW_ITEM_ATTRIBUTE} {PREVIEW_LABEL_ATTRIBUTE}="Slide {index}">'
        f'<div {ELEMENT_LABEL_ATTRIBUTE}="A">{"x" * 50}</div>'
        f'<div {ELEMENT_LABEL_ATTRIBUTE}="B">second</div>'
        "</article>"
        for index in range(4)
    )
    preview = _write(tmp_path, items)

    summary = summarize_document_preview(
        preview, max_items=2, max_elements=1, max_text_chars=10
    )

    assert len(summary.items) == 2
    assert summary.items_truncated is True
    assert summary.items[0].elements_truncated is True
    assert len(summary.items[0].elements) == 1
    assert summary.items[0].elements[0].text.endswith("…")
    assert len(summary.items[0].elements[0].text) <= 11


def test_real_renderer_output_matches_the_inspection_contract(tmp_path):
    """Guard against the renderer and its reader drifting apart."""

    from pptx import Presentation
    from pptx.util import Inches

    from app.services.document_preview import render_document_preview

    source = tmp_path / "deck.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    inside = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    inside.text_frame.text = "inside the slide"
    stray = slide.shapes.add_textbox(Inches(11), Inches(6.5), Inches(4), Inches(2))
    stray.text_frame.text = "hangs off the edge"
    presentation.save(source)

    summary = summarize_document_preview(render_document_preview(source))

    assert len(summary.items) == 1
    texts = [element.text for element in summary.items[0].elements]
    assert "inside the slide" in texts
    assert [element.text for element in summary.items[0].out_of_bounds()] == [
        "hangs off the edge"
    ]


def test_marking_a_cached_preview_survives_a_second_use(tmp_path):
    """A cache hit must not fail where ``follow_symlinks`` is unsupported."""

    cached = tmp_path / "cached.html"
    cached.write_text("<html></html>", encoding="utf-8")

    mark_cached_preview_used(cached)
    mark_cached_preview_used(cached)

    assert cached.exists()


def test_marking_ignores_a_directory(tmp_path):
    target = tmp_path / "a-directory"
    target.mkdir()

    mark_cached_preview_used(target)

    assert stat.S_ISDIR(os.stat(target).st_mode)


def test_marking_ignores_a_missing_path(tmp_path):
    mark_cached_preview_used(tmp_path / "absent.html")
