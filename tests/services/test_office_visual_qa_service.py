"""Tests for LibreOffice-free Office visual QA."""

from __future__ import annotations

from pathlib import Path

import pytest
from docx import Document
from openpyxl import Workbook
from openpyxl.chart import BarChart, Reference
from openpyxl.drawing.image import Image as WorkbookImage
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from app.agent.builtin_skills.pptx.scripts import qa as pptx_qa
from app.agent.builtin_skills.pptx.scripts import stylekit as pptx_stylekit
from app.services import office_preview_service as preview
from app.services import office_visual_qa_service as visual


def _pptx(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    text_box.text = "Browser-rendered visual QA"
    presentation.save(path)


def _docx(path: Path) -> None:
    document = Document()
    document.sections[0].header.paragraphs[0].text = "QA header"
    document.add_heading("Browser-rendered document", level=1)
    document.add_paragraph("First page")
    document.add_page_break()
    document.add_paragraph("Second page")
    document.save(path)


def _xlsx(path: Path) -> None:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "QA"
    sheet["A1"] = "Revenue"
    sheet["B1"] = 1234.5
    sheet["B1"].number_format = "$#,##0.00"
    sheet["A2"] = "Forecast"
    sheet["B2"] = 1600
    chart = BarChart()
    chart.title = "Revenue"
    chart.add_data(Reference(sheet, min_col=2, min_row=1, max_row=2))
    sheet.add_chart(chart, "D2")
    image_path = path.with_suffix(".png")
    Image.new("RGB", (32, 24), "#0f766e").save(image_path)
    sheet.add_image(WorkbookImage(image_path), "D18")
    workbook.save(path)


def test_find_chromium_prefers_configured_path(monkeypatch, tmp_path: Path) -> None:
    executable = tmp_path / "chromium"
    executable.write_bytes(b"browser")
    monkeypatch.setattr(
        visual.settings,
        "EVOFLUX_CHROMIUM_PATH",
        str(executable),
    )

    assert visual.find_chromium() == executable.resolve()


def test_render_reports_structural_fallback_without_browser(
    monkeypatch,
    tmp_path: Path,
) -> None:
    source = tmp_path / "slides.pptx"
    _pptx(source)
    monkeypatch.setattr(visual, "find_chromium", lambda: None)

    report = visual.render_office_images(source, tmp_path / "render")

    assert report["status"] == "unavailable"
    assert report["engine"] == "structural-only"
    assert report["confidence"] == "none"


def test_render_all_office_formats_with_chromium_when_available(
    monkeypatch,
    tmp_path: Path,
) -> None:
    if visual.find_chromium() is None:
        pytest.skip("Chrome/Chromium is unavailable in this test environment")
    monkeypatch.setattr(preview.settings, "EVOFLUX_CACHE_DIR", str(tmp_path / "cache"))
    fixtures = (
        ("slides.pptx", _pptx, "medium", 1),
        ("document.docx", _docx, "approximate", 2),
        ("workbook.xlsx", _xlsx, "medium", 1),
    )

    for filename, builder, confidence, image_count in fixtures:
        source = tmp_path / filename
        builder(source)
        report = visual.render_office_images(
            source,
            tmp_path / f"{source.stem}-render",
        )

        assert report["status"] == "rendered"
        assert report["engine"] == "chromium-openxml"
        assert report["confidence"] == confidence
        assert report["errors"] == []
        assert len(report["images"]) == image_count
        assert all(Path(image).is_file() for image in report["images"])


def test_compare_rendered_images_reports_pixel_changes(tmp_path: Path) -> None:
    before = tmp_path / "before.png"
    after = tmp_path / "after.png"
    Image.new("RGB", (20, 20), "white").save(before)
    Image.new("RGB", (20, 20), "black").save(after)

    report = visual.compare_rendered_images([str(before)], [str(after)])

    assert report["errors"] == []
    assert report["items"][0]["pixel_channel_change_ratio"] == 1
    assert report["items"][0]["changed_bounds"] == (0, 0, 20, 20)


def test_compare_rendered_images_rejects_item_count_change() -> None:
    report = visual.compare_rendered_images(["before.png"], [])

    assert report["errors"] == ["Rendered item count changed from 1 to 0"]


def test_docx_long_table_paginates_without_overflow(
    monkeypatch,
    tmp_path: Path,
) -> None:
    if visual.find_chromium() is None:
        pytest.skip("Chrome/Chromium is unavailable in this test environment")
    source = tmp_path / "long-table.docx"
    document = Document()
    table = document.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "Item"
    table.rows[0].cells[1].text = "Description"
    for index in range(80):
        cells = table.add_row().cells
        cells[0].text = str(index + 1)
        cells[1].text = f"Quality assurance row {index + 1}"
    document.save(source)
    monkeypatch.setattr(preview.settings, "EVOFLUX_CACHE_DIR", str(tmp_path / "cache"))

    report = visual.render_office_images(source, tmp_path / "render")

    assert report["status"] == "rendered"
    assert report["errors"] == []
    assert len(report["images"]) >= 2


def test_pptx_visual_lint_detects_shape_collisions(
    monkeypatch,
    tmp_path: Path,
) -> None:
    if visual.find_chromium() is None:
        pytest.skip("Chrome/Chromium is unavailable in this test environment")
    source = tmp_path / "overlap.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    first = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(5), Inches(1.2))
    first.text = "First message"
    second = slide.shapes.add_textbox(Inches(4), Inches(2.3), Inches(5), Inches(1.2))
    second.text = "Second message"
    presentation.save(source)
    monkeypatch.setattr(preview.settings, "EVOFLUX_CACHE_DIR", str(tmp_path / "cache"))

    report = visual.render_office_images(source, tmp_path / "render")

    assert report["status"] == "rendered"
    assert any("visually overlap" in error for error in report["errors"])


def test_layout_first_pptx_passes_structural_and_visual_qa(
    monkeypatch,
    tmp_path: Path,
) -> None:
    if visual.find_chromium() is None:
        pytest.skip("Chrome/Chromium is unavailable in this test environment")
    source = tmp_path / "layout-first.pptx"
    image_path = tmp_path / "visual.png"
    Image.new("RGB", (1200, 800), "#d8e7e4").save(image_path)
    presentation = pptx_stylekit.new_wide_presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    plan = pptx_stylekit.layout_plan(presentation, "split")
    guard = pptx_stylekit.LayoutGuard(plan)
    pptx_stylekit.set_background(slide, "F7F5F0")
    pptx_stylekit.add_title(
        slide,
        "One composition clarifies the story",
        guard=guard,
    )
    text_region = plan.region("text")
    pptx_stylekit.add_text(
        slide,
        "Short copy fits its declared region. The visual carries the evidence.",
        left=text_region.left,
        top=text_region.top,
        width=text_region.width,
        height=text_region.height,
        font="Aptos",
        size=22,
        color="20303C",
        max_lines=6,
        guard=guard,
    )
    visual_region = plan.region("visual")
    pptx_stylekit.add_image_cover(
        slide,
        image_path,
        left=visual_region.left,
        top=visual_region.top,
        width=visual_region.width,
        height=visual_region.height,
        guard=guard,
    )
    presentation.save(source)
    monkeypatch.setattr(preview.settings, "EVOFLUX_CACHE_DIR", str(tmp_path / "cache"))

    structural = pptx_qa.inspect_pptx(source)
    rendered = visual.render_office_images(source, tmp_path / "render")

    assert structural["errors"] == []
    assert rendered["status"] == "rendered"
    assert rendered["errors"] == []
