from __future__ import annotations

import hashlib
import json
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor
from pptx.util import Inches
import pytest
from pydantic import ValidationError

from app.agent.builtin_plugins.documents.engines import pptx_template as template_engine
from app.agent.builtin_plugins.documents.engines.pptx_template import (
    TemplateDeckProject,
    compose_pptx_template,
    inspect_pptx_template,
    load_template_manifest,
    template_catalog,
    validate_template_project,
)


def _source(tmp_path: Path) -> tuple[Path, str]:
    source = tmp_path / "template.pptx"
    source.write_bytes(b"pptx-template-fixture")
    digest = hashlib.sha256(source.read_bytes()).hexdigest()
    return source, digest


def _manifest(tmp_path: Path, digest: str) -> Path:
    manifest = tmp_path / "template-manifest.json"
    manifest.write_text(
        json.dumps(
            {
                "schemaVersion": 2,
                "sourceSha256": digest,
                "slideCount": 3,
                "records": [
                    {"id": "sh/1/2", "slide": 1, "kind": "textbox"},
                    {"id": "im/1/3", "slide": 1, "kind": "image"},
                    {"id": "tb/2/4", "slide": 2, "kind": "table"},
                    {"id": "ch/2/5", "slide": 2, "kind": "chart"},
                ],
            }
        ),
        encoding="utf-8",
    )
    return manifest


def _project(digest: str) -> TemplateDeckProject:
    return TemplateDeckProject.model_validate(
        {
            "schema_version": 2,
            "title": "Inherited deck",
            "source_sha256": digest,
            "template_confirmed": True,
            "output_slides": [
                {
                    "output_slide": 1,
                    "source_slide": 1,
                    "narrative_role": "opening",
                    "edits": [
                        {
                            "operation": "replace_text",
                            "target_id": "sh/1/2",
                            "find": "Old",
                            "replace": "New",
                        },
                        {
                            "operation": "replace_image",
                            "target_id": "im/1/3",
                            "asset_path": "hero.png",
                        },
                    ],
                },
                {
                    "output_slide": 2,
                    "source_slide": 2,
                    "narrative_role": "evidence",
                    "edits": [
                        {
                            "operation": "set_table_cell",
                            "target_id": "tb/2/4",
                            "row": 1,
                            "column": 2,
                            "text": "$4.2M",
                        },
                        {
                            "operation": "set_chart_series",
                            "target_id": "ch/2/5",
                            "series_index": 0,
                            "values": [1, 2, 3],
                        },
                    ],
                },
            ],
            "omitted_source_slides": [3],
        }
    )


def _picture_alt_text(shape) -> str:
    properties = shape._element.xpath(
        "./*[local-name()='nvPicPr']/*[local-name()='cNvPr']"
    )
    assert properties
    return str(properties[0].get("descr") or properties[0].get("title") or "")


def _set_picture_alt_text(shape, value: str) -> None:
    properties = shape._element.xpath(
        "./*[local-name()='nvPicPr']/*[local-name()='cNvPr']"
    )
    assert properties
    properties[0].set("descr", value)


def _image_partname(shape) -> str:
    relation_id = shape._pic.blip_rId
    return str(shape.part.related_part(relation_id).partname)


def test_template_catalog_declares_direct_openxml_preservation() -> None:
    catalog = template_catalog()
    assert catalog["workflow"] == "direct-openxml-pptx-template"
    assert catalog["style_behavior"]["uploaded_template_is_style_confirmation"]
    assert "replace_image" in catalog["supported_edits"]
    assert "set_image" in catalog["supported_placeholder_fills"]


def test_template_project_validates_complete_typed_mapping(tmp_path: Path) -> None:
    source, digest = _source(tmp_path)
    manifest = load_template_manifest(_manifest(tmp_path, digest))

    result = validate_template_project(_project(digest), manifest, source_pptx=source)

    assert result["valid"] is True
    assert result["edit_count"] == 4


def test_template_project_requires_sequential_output_slides() -> None:
    with pytest.raises(ValidationError, match="sequentially"):
        TemplateDeckProject.model_validate(
            {
                "schema_version": 2,
                "title": "Bad map",
                "source_sha256": "0" * 64,
                "template_confirmed": True,
                "output_slides": [
                    {
                        "output_slide": 2,
                        "source_slide": 1,
                        "narrative_role": "opening",
                    }
                ],
            }
        )


def test_template_project_rejects_changed_source(tmp_path: Path) -> None:
    source, digest = _source(tmp_path)
    manifest = load_template_manifest(_manifest(tmp_path, digest))
    project = _project(digest)
    source.write_bytes(b"changed")

    with pytest.raises(ValueError, match="changed after inspection"):
        validate_template_project(project, manifest, source_pptx=source)


def test_template_project_rejects_non_raster_placeholder_asset() -> None:
    with pytest.raises(ValidationError, match="supported raster"):
        TemplateDeckProject.model_validate(
            {
                "schema_version": 2,
                "title": "Unsafe image fill",
                "source_sha256": "0" * 64,
                "template_confirmed": True,
                "output_slides": [
                    {
                        "output_slide": 1,
                        "source_slide": 1,
                        "narrative_role": "image proof",
                        "reuse_mode": "use-layout",
                        "placeholder_fills": [
                            {
                                "operation": "set_image",
                                "placeholder_idx": 1,
                                "placeholder_type": "PICTURE",
                                "asset_path": "payload.svg",
                            }
                        ],
                    }
                ],
                "omitted_source_slides": [],
            }
        )


def test_template_replacement_asset_cannot_escape_project(tmp_path: Path) -> None:
    project_dir = tmp_path / "project"
    project_dir.mkdir()
    outside = tmp_path / "outside.png"
    Image.new("RGB", (8, 8), "#2563eb").save(outside)

    with pytest.raises(ValueError, match="inside the project directory"):
        template_engine._project_asset(project_dir, "../outside.png")


@pytest.mark.asyncio
async def test_inspect_and_compose_preserve_source_shapes(tmp_path: Path) -> None:
    source = tmp_path / "source.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    textbox.text = "Old title"
    presentation.save(source)

    inspected = await inspect_pptx_template(
        source, workspace_root=tmp_path, work_dir=tmp_path / "inspect"
    )
    manifest = load_template_manifest(inspected.manifest_path)
    target = next(
        record["id"] for record in manifest["records"] if record["kind"] == "textbox"
    )
    project_path = tmp_path / "project.json"
    project_path.write_text(
        json.dumps(
            {
                "schema_version": 2,
                "title": "Edited",
                "source_sha256": manifest["sourceSha256"],
                "template_confirmed": True,
                "output_slides": [
                    {
                        "output_slide": 1,
                        "source_slide": 1,
                        "narrative_role": "opening",
                        "edits": [
                            {
                                "operation": "replace_text",
                                "target_id": target,
                                "find": "Old",
                                "replace": "New",
                            }
                        ],
                    }
                ],
                "omitted_source_slides": [],
            }
        ),
        encoding="utf-8",
    )
    output = tmp_path / "output.pptx"

    result = await compose_pptx_template(
        source,
        project_path,
        inspected.manifest_path,
        output,
        workspace_root=tmp_path,
        work_dir=tmp_path / "compose",
    )

    assert result.passed is True
    reopened = Presentation(output)
    assert len(reopened.slides) == 1
    assert reopened.slides[0].shapes[0].text == "New title"


@pytest.mark.asyncio
async def test_replace_text_preserves_rich_run_structure(tmp_path: Path) -> None:
    source = tmp_path / "rich-text-template.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    paragraph = textbox.text_frame.paragraphs[0]
    first = paragraph.add_run()
    first.text = "Revenue "
    first.font.bold = True
    first.font.color.rgb = RGBColor(0xDC, 0x26, 0x26)
    middle = paragraph.add_run()
    middle.text = "Old"
    middle.font.italic = True
    middle.font.color.rgb = RGBColor(0x25, 0x63, 0xEB)
    last = paragraph.add_run()
    last.text = " outlook"
    last.font.color.rgb = RGBColor(0x16, 0xA3, 0x4A)
    presentation.save(source)

    inspected = await inspect_pptx_template(
        source, workspace_root=tmp_path, work_dir=tmp_path / "inspect-rich-text"
    )
    manifest = load_template_manifest(inspected.manifest_path)
    target = next(
        record["id"] for record in manifest["records"] if record["kind"] == "textbox"
    )
    project_path = tmp_path / "rich-text-project.json"
    project_path.write_text(
        json.dumps(
            {
                "schema_version": 2,
                "title": "Rich text replacement",
                "source_sha256": manifest["sourceSha256"],
                "template_confirmed": True,
                "output_slides": [
                    {
                        "output_slide": 1,
                        "source_slide": 1,
                        "narrative_role": "opening",
                        "edits": [
                            {
                                "operation": "replace_text",
                                "target_id": target,
                                "find": "Old",
                                "replace": "New",
                            }
                        ],
                    }
                ],
                "omitted_source_slides": [],
            }
        ),
        encoding="utf-8",
    )
    output = tmp_path / "rich-text-output.pptx"

    result = await compose_pptx_template(
        source,
        project_path,
        inspected.manifest_path,
        output,
        workspace_root=tmp_path,
        work_dir=tmp_path / "compose-rich-text",
    )

    assert result.passed is True
    reopened_runs = (
        Presentation(output).slides[0].shapes[0].text_frame.paragraphs[0].runs
    )
    assert [run.text for run in reopened_runs] == ["Revenue ", "New", " outlook"]
    assert reopened_runs[0].font.bold is True
    assert reopened_runs[0].font.color.rgb == RGBColor(0xDC, 0x26, 0x26)
    assert reopened_runs[1].font.italic is True
    assert reopened_runs[1].font.color.rgb == RGBColor(0x25, 0x63, 0xEB)
    assert reopened_runs[2].font.color.rgb == RGBColor(0x16, 0xA3, 0x4A)


@pytest.mark.asyncio
async def test_inspect_records_lineage_and_filled_picture_placeholder(
    tmp_path: Path,
) -> None:
    source_image = tmp_path / "portrait.png"
    Image.new("RGB", (200, 400), "#dc2626").save(source_image)
    source = tmp_path / "picture-template.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[8])
    picture_placeholder = next(
        shape
        for shape in slide.placeholders
        if shape.placeholder_format.type.name == "PICTURE"
    )
    picture = picture_placeholder.insert_picture(str(source_image))
    _set_picture_alt_text(picture, "Portrait source")
    presentation.save(source)

    inspected = await inspect_pptx_template(
        source, workspace_root=tmp_path, work_dir=tmp_path / "inspect-picture"
    )
    manifest = load_template_manifest(inspected.manifest_path)

    image_record = next(
        record
        for record in manifest["records"]
        if record["shapeId"] == picture.shape_id
    )
    assert image_record["kind"] == "image"
    assert image_record["placeholder"] == {"idx": 1, "type": "PICTURE"}
    assert image_record["image"]["altText"] == "Portrait source"
    assert image_record["image"]["pixelWidth"] == 200
    assert image_record["image"]["pixelHeight"] == 400
    assert image_record["image"]["crop"]["top"] > 0
    lineage = manifest["slideLineage"][0]
    assert lineage["layout"]["name"] == slide.slide_layout.name
    assert lineage["layout"]["partName"] == str(slide.slide_layout.part.partname)
    assert len(lineage["layout"]["sha256"]) == 64
    assert lineage["master"]["partName"] == str(
        slide.slide_layout.slide_master.part.partname
    )
    assert len(lineage["master"]["sha256"]) == 64
    assert any(
        item["placeholder"]["type"] == "PICTURE"
        for item in lineage["layout"]["placeholders"]
    )


@pytest.mark.asyncio
async def test_inspect_flags_group_shapes_and_children_as_non_addressable(
    tmp_path: Path,
) -> None:
    source = tmp_path / "group-template.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    group = slide.shapes.add_group_shape()
    group.name = "Grouped composition"
    child = group.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    child.name = "Grouped label"
    child.text = "Nested editable-looking text"
    presentation.save(source)

    inspected = await inspect_pptx_template(
        source, workspace_root=tmp_path, work_dir=tmp_path / "inspect-group"
    )
    manifest = load_template_manifest(inspected.manifest_path)

    group_record = next(
        record for record in manifest["records"] if record["kind"] == "group"
    )
    assert group_record["id"].startswith("gr/1/")
    assert group_record["addressable"] is False
    assert len(group_record["children"]) == 1
    assert group_record["children"][0]["kind"] == "textbox"
    assert group_record["children"][0]["text"] == "Nested editable-looking text"
    assert group_record["children"][0]["addressable"] is False


@pytest.mark.asyncio
async def test_replace_filled_picture_placeholder_preserves_frame_and_applies_alt(
    tmp_path: Path,
) -> None:
    source_image = tmp_path / "source-portrait.png"
    replacement_image = tmp_path / "replacement-landscape.png"
    Image.new("RGB", (200, 400), "#dc2626").save(source_image)
    Image.new("RGB", (400, 200), "#2563eb").save(replacement_image)
    source = tmp_path / "picture-template.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[8])
    picture_placeholder = next(
        shape
        for shape in slide.placeholders
        if shape.placeholder_format.type.name == "PICTURE"
    )
    picture = picture_placeholder.insert_picture(str(source_image))
    _set_picture_alt_text(picture, "Old description")
    source_crop = (
        picture.crop_left,
        picture.crop_right,
        picture.crop_top,
        picture.crop_bottom,
    )
    source_layout = str(slide.slide_layout.part.partname)
    slide.notes_slide.notes_text_frame.text = "Preserve this source note"
    presentation.save(source)

    inspected = await inspect_pptx_template(
        source, workspace_root=tmp_path, work_dir=tmp_path / "inspect-replace"
    )
    manifest = load_template_manifest(inspected.manifest_path)
    target = next(
        record["id"] for record in manifest["records"] if record["kind"] == "image"
    )
    project_path = tmp_path / "replace-project.json"
    project_path.write_text(
        json.dumps(
            {
                "schema_version": 2,
                "title": "Picture replacement",
                "source_sha256": manifest["sourceSha256"],
                "template_confirmed": True,
                "output_slides": [
                    {
                        "output_slide": 1,
                        "source_slide": 1,
                        "narrative_role": "image proof",
                        "edits": [
                            {
                                "operation": "replace_image",
                                "target_id": target,
                                "asset_path": replacement_image.name,
                                "alt": "Replacement landscape",
                            }
                        ],
                    }
                ],
                "omitted_source_slides": [],
            }
        ),
        encoding="utf-8",
    )
    output = tmp_path / "picture-output.pptx"

    result = await compose_pptx_template(
        source,
        project_path,
        inspected.manifest_path,
        output,
        workspace_root=tmp_path,
        work_dir=tmp_path / "compose-replace",
    )

    assert result.passed is True
    assert result.metadata["lineage_verified"] is True
    assert result.metadata["preserved_master_layouts"] is True
    reopened = Presentation(output)
    output_slide = reopened.slides[0]
    output_picture = next(
        shape for shape in output_slide.shapes if hasattr(shape, "image")
    )
    assert output_picture.image.size == (400, 200)
    assert _picture_alt_text(output_picture) == "Replacement landscape"
    assert (
        output_picture.crop_left,
        output_picture.crop_right,
        output_picture.crop_top,
        output_picture.crop_bottom,
    ) == pytest.approx(source_crop)
    assert str(output_slide.slide_layout.part.partname) == source_layout
    assert output_slide.notes_slide.notes_text_frame.text == "Preserve this source note"


@pytest.mark.asyncio
async def test_duplicate_slide_isolates_chart_and_image_parts_and_preserves_notes(
    tmp_path: Path,
) -> None:
    source_image = tmp_path / "shared.png"
    Image.new("RGB", (320, 180), "#0f766e").save(source_image)
    source = tmp_path / "chart-template.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(
        str(source_image), Inches(7), Inches(0.5), Inches(4), Inches(2.25)
    )
    data = CategoryChartData()
    data.categories = ["A", "B"]
    data.add_series("Series", [1, 2])
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1),
        Inches(1),
        Inches(5),
        Inches(3),
        data,
    )
    slide.notes_slide.notes_text_frame.text = "Source evidence note"
    presentation.save(source)

    inspected = await inspect_pptx_template(
        source, workspace_root=tmp_path, work_dir=tmp_path / "inspect-chart"
    )
    manifest = load_template_manifest(inspected.manifest_path)
    chart_target = next(
        record["id"] for record in manifest["records"] if record["kind"] == "chart"
    )
    project_path = tmp_path / "duplicate-project.json"
    project_path.write_text(
        json.dumps(
            {
                "schema_version": 2,
                "title": "Isolated duplicates",
                "source_sha256": manifest["sourceSha256"],
                "template_confirmed": True,
                "output_slides": [
                    {
                        "output_slide": 1,
                        "source_slide": 1,
                        "narrative_role": "edited chart",
                        "edits": [
                            {
                                "operation": "set_chart_series",
                                "target_id": chart_target,
                                "series_index": 0,
                                "values": [9, 8],
                            }
                        ],
                    },
                    {
                        "output_slide": 2,
                        "source_slide": 1,
                        "narrative_role": "preserved chart",
                        "edits": [],
                        "speaker_notes": "Explicit replacement note",
                    },
                ],
                "omitted_source_slides": [],
            }
        ),
        encoding="utf-8",
    )
    output = tmp_path / "duplicate-output.pptx"

    result = await compose_pptx_template(
        source,
        project_path,
        inspected.manifest_path,
        output,
        workspace_root=tmp_path,
        work_dir=tmp_path / "compose-chart",
    )

    assert result.passed is True
    reopened = Presentation(output)
    first_chart = next(
        shape
        for shape in reopened.slides[0].shapes
        if getattr(shape, "has_chart", False)
    )
    second_chart = next(
        shape
        for shape in reopened.slides[1].shapes
        if getattr(shape, "has_chart", False)
    )
    assert list(first_chart.chart.series[0].values) == [9.0, 8.0]
    assert list(second_chart.chart.series[0].values) == [1.0, 2.0]
    assert first_chart.chart.part.partname != second_chart.chart.part.partname
    first_image = next(
        shape for shape in reopened.slides[0].shapes if hasattr(shape, "image")
    )
    second_image = next(
        shape for shape in reopened.slides[1].shapes if hasattr(shape, "image")
    )
    assert _image_partname(first_image) != _image_partname(second_image)
    assert (
        reopened.slides[0].notes_slide.notes_text_frame.text == "Source evidence note"
    )
    assert reopened.slides[1].notes_slide.notes_text_frame.text == (
        "Explicit replacement note"
    )


@pytest.mark.asyncio
async def test_compose_fails_when_cloned_slide_layout_lineage_drifts(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "lineage-template.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Expected layout"
    presentation.save(source)
    inspected = await inspect_pptx_template(
        source, workspace_root=tmp_path, work_dir=tmp_path / "inspect-lineage"
    )
    manifest = load_template_manifest(inspected.manifest_path)
    project_path = tmp_path / "lineage-project.json"
    project_path.write_text(
        json.dumps(
            {
                "schema_version": 2,
                "title": "Lineage check",
                "source_sha256": manifest["sourceSha256"],
                "template_confirmed": True,
                "output_slides": [
                    {
                        "output_slide": 1,
                        "source_slide": 1,
                        "narrative_role": "preserve layout",
                        "edits": [],
                    }
                ],
                "omitted_source_slides": [],
            }
        ),
        encoding="utf-8",
    )
    original_duplicate = template_engine._duplicate_slide

    def duplicate_with_wrong_layout(target_presentation, source_slide):
        target = original_duplicate(target_presentation, source_slide)
        layout_relation = next(
            relation
            for relation in target.part.rels.values()
            if relation.reltype.endswith("/slideLayout")
        )
        # Two Content retains compatible title/body placeholder positions while
        # still proving that a different layout part is rejected by lineage QA.
        wrong_layout = target_presentation.slide_layouts[3]
        assert wrong_layout.part.partname != source_slide.slide_layout.part.partname
        layout_relation._target = wrong_layout.part
        return target

    monkeypatch.setattr(
        template_engine, "_duplicate_slide", duplicate_with_wrong_layout
    )

    result = await compose_pptx_template(
        source,
        project_path,
        inspected.manifest_path,
        tmp_path / "lineage-output.pptx",
        workspace_root=tmp_path,
        work_dir=tmp_path / "compose-lineage",
    )

    assert result.passed is False
    assert result.metadata["lineage_verified"] is False
    assert result.metadata["preserved_master_layouts"] is False
    assert any(issue["code"] == "template-lineage-drift" for issue in result.issues)


@pytest.mark.asyncio
async def test_use_layout_fills_native_text_and_picture_placeholders_after_reopen(
    tmp_path: Path,
) -> None:
    source_picture = tmp_path / "source-picture.png"
    replacement_picture = tmp_path / "layout-picture.png"
    Image.new("RGB", (240, 480), "#dc2626").save(source_picture)
    Image.new("RGB", (640, 360), "#2563eb").save(replacement_picture)
    source = tmp_path / "layout-source.pptx"
    presentation = Presentation()
    wrong_layout_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    wrong_layout_slide.shapes.add_textbox(
        Inches(1), Inches(1), Inches(4), Inches(1)
    ).text = "Wrong source layout"
    source_slide = presentation.slides.add_slide(presentation.slide_layouts[8])
    source_slide.shapes.title.text = "Populated source title"
    source_body = next(
        shape
        for shape in source_slide.placeholders
        if shape.placeholder_format.type.name == "BODY"
    )
    source_body.text = "Populated source body"
    source_picture_placeholder = next(
        shape
        for shape in source_slide.placeholders
        if shape.placeholder_format.type.name == "PICTURE"
    )
    source_picture_placeholder.insert_picture(str(source_picture))
    source_slide.shapes.add_textbox(
        Inches(0.2), Inches(0.2), Inches(2), Inches(0.4)
    ).text = "Do not copy this local object"
    source_slide.notes_slide.notes_text_frame.text = "Do not copy source notes"
    source_layout_part = str(source_slide.slide_layout.part.partname)
    source_master_part = str(source_slide.slide_layout.slide_master.part.partname)
    presentation.save(source)

    inspected = await inspect_pptx_template(
        source, workspace_root=tmp_path, work_dir=tmp_path / "inspect-use-layout"
    )
    manifest = load_template_manifest(inspected.manifest_path)
    layout_placeholders = manifest["slideLineage"][1]["layout"]["placeholders"]
    placeholder_idx = {
        record["placeholder"]["type"]: record["placeholder"]["idx"]
        for record in layout_placeholders
    }
    project_path = tmp_path / "use-layout-project.json"
    project_path.write_text(
        json.dumps(
            {
                "schema_version": 2,
                "title": "Native layout authoring",
                "source_sha256": manifest["sourceSha256"],
                "template_confirmed": True,
                "output_slides": [
                    {
                        "output_slide": 1,
                        "source_slide": 2,
                        "narrative_role": "fresh native layout",
                        "reuse_mode": "use-layout",
                        "placeholder_fills": [
                            {
                                "operation": "set_text",
                                "placeholder_idx": placeholder_idx["TITLE"],
                                "placeholder_type": "TITLE",
                                "text": "Authored from the inspected layout",
                            },
                            {
                                "operation": "set_image",
                                "placeholder_idx": placeholder_idx["PICTURE"],
                                "placeholder_type": "PICTURE",
                                "asset_path": replacement_picture.name,
                                "alt": "Blue landscape evidence",
                            },
                            {
                                "operation": "set_text",
                                "placeholder_idx": placeholder_idx["BODY"],
                                "placeholder_type": "BODY",
                                "text": "Native body placeholder remains editable.",
                            },
                        ],
                    }
                ],
                "omitted_source_slides": [1],
            }
        ),
        encoding="utf-8",
    )
    output = tmp_path / "use-layout-output.pptx"

    result = await compose_pptx_template(
        source,
        project_path,
        inspected.manifest_path,
        output,
        workspace_root=tmp_path,
        work_dir=tmp_path / "compose-use-layout",
    )

    assert result.passed is True
    assert result.metadata["lineage_verified"] is True
    assert result.metadata["use_layout_slide_count"] == 1
    assert result.metadata["placeholder_fill_count"] == 3
    reopened = Presentation(output)
    assert len(reopened.slides) == 1
    slide = reopened.slides[0]
    assert str(slide.slide_layout.part.partname) == source_layout_part
    assert str(slide.slide_layout.slide_master.part.partname) == source_master_part
    placeholders = {
        (shape.placeholder_format.idx, shape.placeholder_format.type.name): shape
        for shape in slide.placeholders
    }
    title = placeholders[(placeholder_idx["TITLE"], "TITLE")]
    picture = placeholders[(placeholder_idx["PICTURE"], "PICTURE")]
    body = placeholders[(placeholder_idx["BODY"], "BODY")]
    assert title.text == "Authored from the inspected layout"
    assert body.text == "Native body placeholder remains editable."
    assert picture.image.size == (640, 360)
    assert _picture_alt_text(picture) == "Blue landscape evidence"
    assert all(
        "Do not copy" not in getattr(shape, "text", "") for shape in slide.shapes
    )
    assert slide.has_notes_slide is False


@pytest.mark.asyncio
async def test_use_layout_rejects_placeholder_type_not_in_inspected_lineage(
    tmp_path: Path,
) -> None:
    source = tmp_path / "layout-validation-source.pptx"
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[8])
    presentation.save(source)
    inspected = await inspect_pptx_template(
        source,
        workspace_root=tmp_path,
        work_dir=tmp_path / "inspect-layout-validation",
    )
    manifest = load_template_manifest(inspected.manifest_path)
    picture_record = next(
        record
        for record in manifest["slideLineage"][0]["layout"]["placeholders"]
        if record["placeholder"]["type"] == "PICTURE"
    )
    project = TemplateDeckProject.model_validate(
        {
            "schema_version": 2,
            "title": "Invalid layout target",
            "source_sha256": manifest["sourceSha256"],
            "template_confirmed": True,
            "output_slides": [
                {
                    "output_slide": 1,
                    "source_slide": 1,
                    "narrative_role": "invalid target",
                    "reuse_mode": "use-layout",
                    "placeholder_fills": [
                        {
                            "operation": "set_text",
                            "placeholder_idx": picture_record["placeholder"]["idx"],
                            "placeholder_type": "BODY",
                            "text": "This must not target a picture slot.",
                        }
                    ],
                }
            ],
            "omitted_source_slides": [],
        }
    )

    with pytest.raises(ValueError, match="has type.*PICTURE.*not BODY"):
        validate_template_project(project, manifest, source_pptx=source)
