from __future__ import annotations

import base64
from io import BytesIO
import json
from pathlib import Path
import shutil
from typing import Any

from PIL import Image
from pptx import Presentation
import pytest

from app.agent.builtin_plugins.documents.engines import pptx_html as pptx_html_engine
from app.agent.builtin_plugins.documents.engines.html_slide_broker import (
    get_html_slide_render_broker,
)
from app.agent.builtin_plugins.documents.engines.pptx_html import (
    compose_html_pptx_project,
    html_pptx_catalog,
    load_html_pptx_project,
    validate_html_pptx_project,
)


def _png_base64(width: int, height: int, *, blank: bool = False) -> str:
    image = (
        Image.new("RGB", (width, height), "white")
        if blank
        else Image.effect_noise((width, height), 48).convert("RGB")
    )
    stream = BytesIO()
    image.save(stream, format="PNG")
    return base64.b64encode(stream.getvalue()).decode("ascii")


def _project(tmp_path: Path) -> Path:
    (tmp_path / "slide.html").write_text(
        """
        <section data-slide-root class="relative h-full w-full">
          <h1>HTML owns the composition</h1>
          <img src="asset://hero" alt="Hero" />
        </section>
        """,
        encoding="utf-8",
    )
    (tmp_path / "slide.css").write_text(
        "[data-slide-root]{background:linear-gradient(135deg,#07111f,#3355ff)}",
        encoding="utf-8",
    )
    Image.new("RGB", (120, 80), "#22d3ee").save(tmp_path / "hero.png")
    path = tmp_path / "deck.json"
    path.write_text(
        json.dumps(
            {
                "schema_version": 7,
                "title": "HTML deck",
                "width": 1280,
                "height": 720,
                "slides": [
                    {
                        "id": "opening",
                        "html_path": "slide.html",
                        "style_paths": ["slide.css"],
                        "assets": {"hero": "hero.png"},
                    }
                ],
            }
        ),
        encoding="utf-8",
    )
    return path


def _text_value(element: dict[str, Any]) -> str:
    if "paragraphs" in element:
        return "\n".join(
            "".join(str(run.get("text") or "") for run in paragraph["runs"])
            for paragraph in element["paragraphs"]
        )
    return str(element.get("text") or "")


def _coverage(
    elements: list[dict[str, Any]],
    *,
    flattened: list[dict[str, Any]] | None = None,
) -> dict[str, Any]:
    native = [element for element in elements if element.get("kind") == "text"]
    flattened = flattened or []
    return {
        "visible_blocks": len(native) + len(flattened),
        "visible_characters": sum(len(_text_value(element)) for element in native)
        + sum(int(item["characters"]) for item in flattened),
        "native_blocks": len(native),
        "native_characters": sum(len(_text_value(element)) for element in native),
        "flattened": flattened,
    }


def _response(
    *,
    elements: list[dict[str, Any]] | None = None,
    coverage: dict[str, Any] | None = None,
    blank: bool = False,
    shell_width: int = 2560,
    shell_height: int = 1440,
    issues: list[dict[str, Any]] | None = None,
) -> dict[str, Any]:
    elements = elements or []
    return {
        "preview_png_base64": _png_base64(1280, 720, blank=blank),
        "shell_png_base64": _png_base64(shell_width, shell_height, blank=blank),
        "editable_elements": elements,
        "text_coverage": coverage or _coverage(elements),
        "issues": issues or [],
    }


def _copy_source_preview(output: Path, tmp_path: Path, work_name: str = "work"):
    def render(source: Path, render_dir: Path, *, width: int = 1280) -> list[Path]:
        assert source == output
        assert width == 1280
        render_dir.mkdir(parents=True, exist_ok=True)
        destination = render_dir / "slide-001.png"
        shutil.copyfile(
            tmp_path / work_name / "previews" / "slide-001.png", destination
        )
        return [destination]

    return render


def test_html_catalog_and_validation_has_no_slide_dna(tmp_path: Path) -> None:
    path = _project(tmp_path)
    project = load_html_pptx_project(path)

    result = validate_html_pptx_project(project, path)
    catalog = html_pptx_catalog()

    assert result["valid"] is True
    assert result["representation"] == "html-shell-editable-text"
    assert result["quality_policy"]["evidence"] == "runtime-render-only"
    assert catalog["workflow"] == "html-shell-editable-text-pptx"
    assert catalog["schema_version"] == 7
    assert "slide_dna" not in catalog
    assert "qa_ledger" not in catalog
    assert catalog["html_contract"]["text"]["art_opt_out"] == (
        'data-pptx-text-mode="art"'
    )


def test_html_schema_rejects_removed_dna_fields(tmp_path: Path) -> None:
    path = _project(tmp_path)
    value = json.loads(path.read_text(encoding="utf-8"))
    value["dna_path"] = "slide-dna.json"
    path.write_text(json.dumps(value), encoding="utf-8")

    with pytest.raises(ValueError, match="dna_path"):
        load_html_pptx_project(path)


@pytest.mark.parametrize(
    ("html", "message"),
    [
        ("<section data-slide-root><script>x()</script></section>", "unsafe HTML"),
        (
            '<section data-slide-root style="background:url(https://bad.test/x)"></section>',
            "network or executable URL",
        ),
        (
            '<section data-slide-root><img src="asset://missing"></section>',
            "undeclared assets",
        ),
    ],
)
def test_html_validation_rejects_unsafe_or_missing_resources(
    tmp_path: Path, html: str, message: str
) -> None:
    path = _project(tmp_path)
    (tmp_path / "slide.html").write_text(html, encoding="utf-8")

    with pytest.raises(ValueError, match=message):
        validate_html_pptx_project(load_html_pptx_project(path), path)


def test_html_validation_rejects_css_that_escapes_style_element(
    tmp_path: Path,
) -> None:
    path = _project(tmp_path)
    (tmp_path / "slide.css").write_text(
        "</style><img src=x onerror=alert(1)>", encoding="utf-8"
    )

    with pytest.raises(ValueError, match="unsafe CSS"):
        validate_html_pptx_project(load_html_pptx_project(path), path)


def test_html_project_paths_cannot_escape_project(tmp_path: Path) -> None:
    project_dir = tmp_path / "project"
    project_dir.mkdir()
    path = _project(project_dir)
    value = json.loads(path.read_text(encoding="utf-8"))
    value["slides"][0]["html_path"] = "../outside.html"
    (tmp_path / "outside.html").write_text(
        "<section data-slide-root></section>", encoding="utf-8"
    )
    path.write_text(json.dumps(value), encoding="utf-8")

    with pytest.raises(ValueError, match="inside the project directory"):
        validate_html_pptx_project(load_html_pptx_project(path), path)


@pytest.mark.asyncio
async def test_html_pipeline_keeps_shell_and_exports_rich_editable_text(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    path = _project(tmp_path)
    elements = [
        {
            "kind": "text",
            "name": "Title",
            "role": "title",
            "x": 80,
            "y": 90,
            "width": 800,
            "height": 150,
            "padding": {"left": 4, "right": 4, "top": 2, "bottom": 2},
            "text_align": "left",
            "vertical_align": "top",
            "line_height_ratio": 1.05,
            "rotation": 0,
            "paragraphs": [
                {
                    "runs": [
                        {
                            "text": "HTML owns ",
                            "font_family": "Arial",
                            "font_size": 48,
                            "bold": True,
                            "color": "#FFFFFF",
                            "letter_spacing": 1.5,
                        },
                        {
                            "text": "the composition",
                            "font_family": "Arial",
                            "font_size": 48,
                            "italic": True,
                            "color": "#22D3EE",
                        },
                    ]
                },
                {
                    "runs": [
                        {
                            "text": "Editable after export",
                            "font_family": "Arial",
                            "font_size": 24,
                            "color": "#FFFFFF",
                        }
                    ],
                    "bullet": {"kind": "bullet", "marker": "•", "level": 0},
                },
                {
                    "runs": [
                        {
                            "text": "Fourth ordered item",
                            "font_family": "Arial",
                            "font_size": 24,
                            "color": "#FFFFFF",
                        }
                    ],
                    "bullet": {"kind": "number", "level": 0, "start": 4},
                },
            ],
        }
    ]

    async def fake_request(*args: Any, **kwargs: Any) -> dict[str, Any]:
        del args, kwargs
        return _response(elements=elements)

    monkeypatch.setattr(get_html_slide_render_broker(), "request", fake_request)
    output = tmp_path / "output.pptx"
    monkeypatch.setattr(
        pptx_html_engine,
        "render_pptx_pages",
        _copy_source_preview(output, tmp_path),
    )

    result = await compose_html_pptx_project(
        path,
        output,
        workspace_root=tmp_path,
        work_dir=tmp_path / "work",
        session_id="session-1",
    )

    assert result.passed is True
    assert result.output == output
    assert result.metadata["editable_text_object_count"] == 1
    assert result.metadata["text_conversion"][0]["native_character_ratio"] == 1.0
    assert result.metadata["reopened_parity"]["median"] == 1.0
    assert "slide_dna" not in result.metadata
    manifest = json.loads(result.manifest_path.read_text(encoding="utf-8"))
    assert manifest["schemaVersion"] == 7
    assert manifest["qualityPolicy"]["evidence"] == "runtime-render-only"
    assert manifest["structuralCounts"]["editable_text_objects"] == 1
    deck = Presentation(str(output))
    assert len(deck.slides) == 1
    assert len(deck.slides[0].shapes) == 2  # HTML shell + editable text.
    title = deck.slides[0].shapes[1]
    assert title.has_text_frame is True
    assert len(title.text_frame.paragraphs) == 3
    assert title.text_frame.paragraphs[0].runs[0].text == "HTML owns "
    assert title.text_frame.paragraphs[0].runs[1].text == "the composition"
    assert title.text_frame.paragraphs[0].runs[1].font.italic is True
    first_properties = title.text_frame.paragraphs[0].runs[0]._r.get_or_add_rPr()  # noqa: SLF001
    assert first_properties.get("spc") == "112"
    assert title.text_frame.paragraphs[1].runs[0].text == "Editable after export"
    numbering = title.text_frame.paragraphs[2]._p.xpath(  # noqa: SLF001
        "./*[local-name()='pPr']/*[local-name()='buAutoNum']"
    )
    assert len(numbering) == 1
    assert numbering[0].get("type") == "arabicPeriod"
    assert numbering[0].get("startAt") == "4"


@pytest.mark.asyncio
async def test_html_pipeline_accepts_legacy_flat_text_manifest(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    path = _project(tmp_path)
    elements = [
        {
            "kind": "text",
            "name": "Legacy title",
            "x": 80,
            "y": 90,
            "width": 800,
            "height": 90,
            "text": "Still editable",
            "font_family": "Arial",
            "font_size": 48,
            "bold": True,
            "color": "#FFFFFF",
        }
    ]

    async def fake_request(*args: Any, **kwargs: Any) -> dict[str, Any]:
        return _response(elements=elements)

    monkeypatch.setattr(get_html_slide_render_broker(), "request", fake_request)
    output = tmp_path / "legacy.pptx"
    monkeypatch.setattr(
        pptx_html_engine,
        "render_pptx_pages",
        _copy_source_preview(output, tmp_path),
    )

    result = await compose_html_pptx_project(
        path,
        output,
        workspace_root=tmp_path,
        work_dir=tmp_path / "work",
        session_id="session-1",
    )

    assert result.passed is True
    assert Presentation(output).slides[0].shapes[1].text == "Still editable"


@pytest.mark.asyncio
async def test_html_pipeline_rejects_unaccounted_text(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    path = _project(tmp_path)

    async def fake_request(*args: Any, **kwargs: Any) -> dict[str, Any]:
        return _response(
            coverage={
                "visible_blocks": 1,
                "visible_characters": 25,
                "native_blocks": 0,
                "native_characters": 0,
                "flattened": [],
            }
        )

    monkeypatch.setattr(get_html_slide_render_broker(), "request", fake_request)
    output = tmp_path / "unaccounted.pptx"

    result = await compose_html_pptx_project(
        path,
        output,
        workspace_root=tmp_path,
        work_dir=tmp_path / "work",
        session_id="session-1",
    )

    assert result.passed is False
    assert result.output is None
    assert not output.exists()
    assert any(issue["code"] == "unaccounted-html-text" for issue in result.issues)


@pytest.mark.asyncio
async def test_html_pipeline_requires_renderer_text_coverage(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    path = _project(tmp_path)

    async def fake_request(*args: Any, **kwargs: Any) -> dict[str, Any]:
        value = _response()
        value.pop("text_coverage")
        return value

    monkeypatch.setattr(get_html_slide_render_broker(), "request", fake_request)

    with pytest.raises(ValueError, match="text_coverage"):
        await compose_html_pptx_project(
            path,
            tmp_path / "missing-coverage.pptx",
            workspace_root=tmp_path,
            work_dir=tmp_path / "work",
            session_id="session-1",
        )


@pytest.mark.asyncio
async def test_html_pipeline_rejects_low_reopened_visual_parity(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    path = _project(tmp_path)

    async def fake_request(*args: Any, **kwargs: Any) -> dict[str, Any]:
        return _response()

    def fake_reopened_render(
        source: Path, render_dir: Path, *, width: int = 1280
    ) -> list[Path]:
        render_dir.mkdir(parents=True, exist_ok=True)
        destination = render_dir / "slide-001.png"
        Image.new("RGB", (width, 720), "white").save(destination)
        return [destination]

    monkeypatch.setattr(get_html_slide_render_broker(), "request", fake_request)
    monkeypatch.setattr(pptx_html_engine, "render_pptx_pages", fake_reopened_render)
    output = tmp_path / "low-parity.pptx"

    result = await compose_html_pptx_project(
        path,
        output,
        workspace_root=tmp_path,
        work_dir=tmp_path / "work",
        session_id="session-1",
    )

    assert result.passed is False
    assert result.output is None
    assert result.metadata["reopened_parity"]["median"] < 0.9
    assert any(issue["code"] == "render-parity-below-target" for issue in result.issues)


@pytest.mark.asyncio
async def test_html_pipeline_rejects_blank_webview_preview(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    path = _project(tmp_path)

    async def fake_request(*args: Any, **kwargs: Any) -> dict[str, Any]:
        return _response(blank=True)

    monkeypatch.setattr(get_html_slide_render_broker(), "request", fake_request)
    output = tmp_path / "blank.pptx"

    result = await compose_html_pptx_project(
        path,
        output,
        workspace_root=tmp_path,
        work_dir=tmp_path / "work",
        session_id="session-1",
    )

    assert result.passed is False
    assert result.output is None
    assert any(issue["code"] == "invalid-html-slide-preview" for issue in result.issues)


@pytest.mark.asyncio
async def test_html_pipeline_rejects_wrong_shell_dimensions(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    path = _project(tmp_path)

    async def fake_request(*args: Any, **kwargs: Any) -> dict[str, Any]:
        return _response(shell_width=1280, shell_height=720)

    monkeypatch.setattr(get_html_slide_render_broker(), "request", fake_request)

    with pytest.raises(ValueError, match="shell must be 2560x1440"):
        await compose_html_pptx_project(
            path,
            tmp_path / "wrong-shell.pptx",
            workspace_root=tmp_path,
            work_dir=tmp_path / "work",
            session_id="session-1",
        )
