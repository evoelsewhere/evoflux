from __future__ import annotations

import json
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
import pytest
from pydantic import ValidationError

from app.services import pptx_native_pipeline as pipeline


def _native_project() -> dict[str, object]:
    return {
        "schema_version": 3,
        "mode": "new",
        "quality_profile": "native",
        "title": "Native deck",
        "slides": [
            {
                "id": "title",
                "elements": [
                    {
                        "type": "text",
                        "position": {
                            "left": 80,
                            "top": 80,
                            "width": 800,
                            "height": 120,
                        },
                        "text": "Artifact Fabric",
                        "font_size": 48,
                    },
                    {
                        "type": "chart",
                        "position": {
                            "left": 80,
                            "top": 240,
                            "width": 800,
                            "height": 360,
                        },
                        "categories": ["DOCX", "XLSX", "PPTX", "PDF"],
                        "series": [{"name": "Native", "values": [1, 1, 1, 1]}],
                    },
                ],
            }
        ],
    }


def test_native_schema_exposes_editable_openxml_types(tmp_path: Path) -> None:
    project_path = tmp_path / "deck.json"
    project_path.write_text(json.dumps(_native_project()), encoding="utf-8")

    project = pipeline.load_native_pptx_project(project_path)
    result = pipeline.validate_native_pptx_project(project, project_path)

    assert result["element_count"] == 2
    assert pipeline.native_pptx_catalog()["workflow"] == ("evoflux-openxml-svg-pptx")


def test_native_pptx_rejects_out_of_slide_geometry(tmp_path: Path) -> None:
    raw = _native_project()
    raw["width"] = 1280
    slides = raw["slides"]
    assert isinstance(slides, list)
    slides[0]["elements"][0]["position"] = {  # type: ignore[index]
        "left": 1200,
        "top": 10,
        "width": 200,
        "height": 40,
    }
    path = tmp_path / "bad.json"
    path.write_text(json.dumps(raw), encoding="utf-8")

    with pytest.raises(ValidationError, match="right slide boundary"):
        pipeline.load_native_pptx_project(path)


def test_fidelity_profile_accepts_only_project_local_svg(tmp_path: Path) -> None:
    (tmp_path / "slide.svg").write_text(
        '<svg xmlns="http://www.w3.org/2000/svg" width="1280" height="720">'
        '<rect width="1280" height="720" fill="#073b82"/></svg>',
        encoding="utf-8",
    )
    path = tmp_path / "deck.json"
    path.write_text(
        json.dumps(
            {
                "schema_version": 3,
                "mode": "new",
                "quality_profile": "fidelity",
                "title": "SVG deck",
                "slides": [
                    {
                        "id": "exact",
                        "visual_shell": {
                            "svg_path": "slide.svg",
                            "alt": "Exact title slide",
                        },
                    }
                ],
            }
        ),
        encoding="utf-8",
    )

    project = pipeline.load_native_pptx_project(path)
    assert pipeline.validate_native_pptx_project(project, path)["valid"] is True


def test_visual_parity_rejects_structural_displacement(tmp_path: Path) -> None:
    reference_path = tmp_path / "reference.png"
    preview_path = tmp_path / "preview.png"
    reference = Image.new("RGB", (1280, 720), "white")
    preview = reference.copy()
    ImageDraw.Draw(reference).rectangle((120, 120, 760, 520), fill="#073b82")
    ImageDraw.Draw(preview).rectangle((220, 120, 860, 520), fill="#073b82")
    reference.save(reference_path)
    preview.save(preview_path)

    metric = pipeline._visual_parity(
        reference_path,
        preview_path,
        max_changed_pixel_ratio=0.02,
        max_mean_absolute_error=0.01,
    )

    assert metric["passed"] is False
    assert metric["changed_pixel_ratio"] > 0.02


@pytest.mark.asyncio
async def test_fidelity_svg_round_trips_without_browser(tmp_path: Path) -> None:
    (tmp_path / "slide.svg").write_text(
        '<svg xmlns="http://www.w3.org/2000/svg" width="1280" height="720">'
        '<rect width="1280" height="720" fill="#071426"/>'
        '<text x="80" y="180" fill="white" font-size="64">EvoFlux</text>'
        "</svg>",
        encoding="utf-8",
    )
    project_path = tmp_path / "deck.json"
    project_path.write_text(
        json.dumps(
            {
                "schema_version": 3,
                "mode": "new",
                "quality_profile": "fidelity",
                "title": "Fidelity smoke",
                "slides": [
                    {
                        "id": "fidelity",
                        "visual_shell": {
                            "svg_path": "slide.svg",
                            "alt": "Dark title slide",
                        },
                    }
                ],
            }
        ),
        encoding="utf-8",
    )
    output = tmp_path / "fidelity.pptx"

    result = await pipeline.compose_native_pptx_project(
        project_path,
        output,
        workspace_root=tmp_path,
        work_dir=tmp_path / "work",
    )

    assert result.passed is True
    assert result.metadata["visual_parity"][0]["passed"] is True
    presentation = Presentation(output)
    assert len(presentation.slides) == 1
    assert len(presentation.slides[0].shapes) == 1
