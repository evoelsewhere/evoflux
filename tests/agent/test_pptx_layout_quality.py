"""Regression tests for layout-first PPTX generation and QA."""

from __future__ import annotations

import zipfile
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from app.agent.builtin_skills.pptx.scripts import icons as pptx_icons
from app.agent.builtin_skills.pptx.scripts import office_features
from app.agent.builtin_skills.pptx.scripts import qa as pptx_qa
from app.agent.builtin_skills.pptx.scripts import stylekit


def _slide(presentation):
    return presentation.slides.add_slide(presentation.slide_layouts[6])


def _textbox(slide, text: str, left: float, top: float, width: float, height: float):
    shape = slide.shapes.add_textbox(
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.text = text
    shape.text_frame.paragraphs[0].runs[0].font.size = Pt(20)
    return shape


def test_layout_plan_produces_safe_non_overlapping_regions() -> None:
    presentation = stylekit.new_wide_presentation()

    plan = stylekit.layout_plan(presentation, "split")

    text = plan.region("text")
    visual = plan.region("visual")
    assert plan.safe_canvas.contains(text)
    assert plan.safe_canvas.contains(visual)
    assert not text.intersects(visual)


def test_layout_guard_rejects_collisions() -> None:
    presentation = stylekit.new_wide_presentation()
    plan = stylekit.layout_plan(presentation, "hero")
    guard = stylekit.LayoutGuard(plan)
    canvas = plan.region("canvas")
    guard.reserve(
        "first",
        left=canvas.left,
        top=canvas.top,
        width=Inches(4),
        height=Inches(2),
    )

    with pytest.raises(stylekit.LayoutError, match="collides"):
        guard.reserve(
            "second",
            left=canvas.left + Inches(3),
            top=canvas.top + Inches(1),
            width=Inches(4),
            height=Inches(2),
        )


def test_title_and_kicker_fit_above_guarded_content() -> None:
    presentation = stylekit.new_wide_presentation()
    slide = _slide(presentation)
    plan = stylekit.layout_plan(presentation, "hero")
    guard = stylekit.LayoutGuard(plan)

    stylekit.add_title(
        slide,
        "One clear takeaway",
        kicker="Context",
        guard=guard,
    )
    guard.reserve_region(plan.region("canvas"))

    assert len(slide.shapes) == 2


def test_add_text_rejects_copy_that_cannot_fit() -> None:
    presentation = stylekit.new_wide_presentation()
    slide = _slide(presentation)

    with pytest.raises(stylekit.TextOverflowError, match="change layout"):
        stylekit.add_text(
            slide,
            "This sentence is intentionally far too long for the tiny text box.",
            left=Inches(1),
            top=Inches(2),
            width=Inches(1.5),
            height=Inches(0.35),
            font="Aptos",
            size=18,
            color="20303C",
        )


def test_quality_ledger_reports_all_text_issues_after_one_build() -> None:
    presentation = stylekit.new_wide_presentation()
    slide = _slide(presentation)
    ledger = stylekit.QualityLedger()

    stylekit.add_text(
        slide,
        "This body copy is too small.",
        left=Inches(1),
        top=Inches(2),
        width=Inches(3),
        height=Inches(0.5),
        font="Aptos",
        size=12,
        color="20303C",
        ledger=ledger,
    )
    stylekit.add_text(
        slide,
        "This sentence is intentionally far too long for the tiny text box.",
        left=Inches(1),
        top=Inches(3),
        width=Inches(1.5),
        height=Inches(0.35),
        font="Aptos",
        size=18,
        color="20303C",
        ledger=ledger,
    )

    assert len(slide.shapes) == 2
    assert len(ledger.errors) == 2
    with pytest.raises(stylekit.LayoutError, match="2 issue"):
        ledger.raise_if_errors()


def test_curated_icon_search_resolves_semantic_aliases() -> None:
    assert len(pptx_icons.list_icons()) == 47
    assert pptx_icons.resolve_icon("growth") == "trending-up"
    assert any(
        match.name == "chart-line" for match in pptx_icons.search_icons("analytics")
    )


def test_add_icon_embeds_themeable_svg_and_qa_metadata(tmp_path: Path) -> None:
    source = tmp_path / "icons.pptx"
    presentation = stylekit.new_wide_presentation()
    slide = _slide(presentation)
    picture = pptx_icons.add_icon(
        slide,
        "agent",
        left=Inches(1),
        top=Inches(2),
        size=Inches(0.7),
        color="2F6D68",
    )
    presentation.save(source)

    assert picture.name == "[icon:lucide:bot]"
    with zipfile.ZipFile(source) as package:
        svg_parts = [
            name
            for name in package.namelist()
            if name.startswith("ppt/media/") and name.endswith(".svg")
        ]
        assert len(svg_parts) == 1
        assert b'stroke="#2F6D68"' in package.read(svg_parts[0])

    report = pptx_qa.inspect_pptx(source)
    metrics = report["layout"]["metrics"][0]
    assert report["errors"] == []
    assert metrics["icon_count"] == 1
    assert metrics["icon_families"] == ["lucide"]


def test_icon_qa_flags_tiny_or_mixed_family_icons(tmp_path: Path) -> None:
    source = tmp_path / "mixed-icons.pptx"
    presentation = stylekit.new_wide_presentation()
    slide = _slide(presentation)
    pptx_icons.add_icon(
        slide,
        "rocket",
        left=Inches(1),
        top=Inches(2),
        size=Inches(0.2),
    )
    other = pptx_icons.add_icon(
        slide,
        "target",
        left=Inches(2),
        top=Inches(2),
        size=Inches(0.6),
    )
    other.name = "[icon:other:target]"
    presentation.save(source)

    report = pptx_qa.inspect_pptx(source)

    assert any("only 0.20in" in warning for warning in report["warnings"])
    assert any("families lucide, other" in warning for warning in report["warnings"])


def test_native_chart_and_table_remain_editable_office_objects(tmp_path: Path) -> None:
    source = tmp_path / "native-features.pptx"
    presentation = stylekit.new_wide_presentation()
    chart_slide = _slide(presentation)
    chart = office_features.add_native_chart(
        chart_slide,
        ["Q1", "Q2", "Q3"],
        {"Actual": [12, 18, 25], "Plan": [14, 19, 23]},
        left=Inches(0.8),
        top=Inches(1.3),
        width=Inches(7),
        height=Inches(4.8),
        kind="line",
        title="Actual moves ahead of plan",
    )
    office_features.set_accessibility(
        chart,
        title="Actual versus plan",
        description="Quarterly line chart.",
    )
    table_slide = _slide(presentation)
    office_features.add_native_table(
        table_slide,
        ["Workstream", "Owner", "Status"],
        [["Platform", "Mai", "On track"], ["Launch", "Linh", "Ready"]],
        left=Inches(0.8),
        top=Inches(1.3),
        width=Inches(7),
        height=Inches(3.2),
        column_weights=[2, 1, 1.2],
        alignments=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.CENTER],
    )
    presentation.save(source)

    with zipfile.ZipFile(source) as package:
        names = package.namelist()
        assert any(name.startswith("ppt/charts/chart") for name in names)
        assert any(name.startswith("ppt/embeddings/") for name in names)

    report = pptx_qa.inspect_pptx(source)
    summary = report["layout"]["office_feature_summary"]
    assert report["errors"] == []
    assert summary["native_charts"] == 1
    assert summary["native_tables"] == 1
    assert summary["accessible_shapes"] == 1


def test_rich_text_columns_bullets_theme_and_morph_are_native_ooxml(
    tmp_path: Path,
) -> None:
    source = tmp_path / "rich-morph.pptx"
    presentation = stylekit.new_wide_presentation()
    first = _slide(presentation)
    second = _slide(presentation)
    theme = office_features.theme_from_presentation(presentation)
    rich_text = office_features.add_rich_text(
        first,
        [
            office_features.RichParagraph(
                runs=(
                    office_features.RichTextRun(
                        "Evidence-led",
                        bold=True,
                        color="2F6D68",
                    ),
                    office_features.RichTextRun(" argument with editable runs."),
                ),
                bullet=True,
            ),
            office_features.RichParagraph(
                runs=(
                    office_features.RichTextRun(
                        "Second column-ready point",
                        italic=True,
                    ),
                ),
                bullet=True,
                level=1,
            ),
        ],
        left=Inches(0.8),
        top=Inches(1.4),
        width=Inches(5.4),
        height=Inches(3.4),
        columns=2,
        theme=theme,
    )
    first_object = first.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(7),
        Inches(2),
        Inches(2),
        Inches(1.2),
    )
    second_object = second.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(8),
        Inches(2),
        Inches(2.8),
        Inches(1.6),
    )
    office_features.set_morph_identity(first_object, "hero-object")
    office_features.set_morph_identity(second_object, "hero-object")
    office_features.set_slide_transition(second, "morph")
    presentation.save(source)

    assert theme.title_font
    assert theme.body_font
    assert rich_text.text == (
        "Evidence-led argument with editable runs.\nSecond column-ready point"
    )
    with zipfile.ZipFile(source) as package:
        first_xml = package.read("ppt/slides/slide1.xml")
        second_xml = package.read("ppt/slides/slide2.xml")
        assert b'numCol="2"' in first_xml
        assert b"<a:buChar" in first_xml
        assert b'b="1"' in first_xml
        assert b'i="1"' in first_xml
        assert b'name="!!hero-object"' in first_xml
        assert b'name="!!hero-object"' in second_xml
        assert b":morph" in second_xml

    reopened = Presentation(source)
    assert reopened.slides[0].shapes[0].text == rich_text.text
    report = pptx_qa.inspect_pptx(source)
    assert report["errors"] == []
    assert report["powerpoint_features"]["morph_transitions"] == 1
    assert report["powerpoint_features"]["placeholders"] == 0


def test_grouped_process_keeps_nodes_labels_and_connectors_editable(
    tmp_path: Path,
) -> None:
    source = tmp_path / "grouped-process.pptx"
    presentation = stylekit.new_wide_presentation()
    slide = _slide(presentation)

    group = office_features.add_grouped_process(
        slide,
        ["Discover", "Decide", "Deliver"],
        left=Inches(1),
        top=Inches(1.7),
        width=Inches(10),
        height=Inches(3),
    )
    presentation.save(source)

    child_types = [
        getattr(getattr(shape, "shape_type", None), "name", "")
        for shape in group.shapes
    ]
    assert child_types.count("LINE") == 2
    assert [shape.text for shape in group.shapes if shape.has_text_frame] == [
        "1",
        "Discover",
        "2",
        "Decide",
        "3",
        "Deliver",
    ]
    report = pptx_qa.inspect_pptx(source)
    metrics = report["layout"]["metrics"][0]["office_features"]
    assert report["errors"] == []
    assert metrics["groups"] == 1
    assert metrics["connectors"] == 2
    assert report["powerpoint_features"]["groups"] == 1


def test_transition_hyperlink_gradient_and_shadow_use_native_ooxml(
    tmp_path: Path,
) -> None:
    source = tmp_path / "native-ooxml.pptx"
    presentation = stylekit.new_wide_presentation()
    slide = _slide(presentation)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1),
        Inches(1.5),
        Inches(4),
        Inches(2.5),
    )
    office_features.apply_gradient_fill(
        shape,
        [(0, "F7F5F0", 1), (1, "2F6D68", 0.85)],
        angle=25,
    )
    office_features.apply_soft_shadow(shape)
    office_features.set_accessibility(
        shape,
        title="Feature panel",
        description="Gradient panel with a restrained shadow.",
    )
    office_features.set_shape_hyperlink(shape, "https://example.com/details")
    office_features.set_slide_transition(slide, "fade")
    presentation.save(source)

    with zipfile.ZipFile(source) as package:
        slide_xml = package.read("ppt/slides/slide1.xml")
        assert b"<a:gradFill" in slide_xml
        assert b"<a:outerShdw" in slide_xml
        assert b"<p:transition" in slide_xml
        assert b"<p:fade" in slide_xml
        assert b"<a:hlinkClick" in slide_xml
        assert b'descr="Gradient panel with a restrained shadow."' in slide_xml


def test_image_cover_uses_native_focal_crop_and_alt_text(tmp_path: Path) -> None:
    image_path = tmp_path / "wide.png"
    Image.new("RGB", (1200, 600), "white").save(image_path)
    presentation = stylekit.new_wide_presentation()
    slide = _slide(presentation)

    picture = stylekit.add_image_cover(
        slide,
        image_path,
        left=Inches(1),
        top=Inches(1),
        width=Inches(3),
        height=Inches(3),
        focal_x=0.9,
        alt_text="Product subject positioned on the right.",
    )

    assert picture.crop_left > picture.crop_right
    properties = picture._element.xpath(".//*[local-name()='cNvPr']")
    assert properties[0].get("descr") == "Product subject positioned on the right."


def test_qa_warns_when_long_deck_is_shape_only(tmp_path: Path) -> None:
    source = tmp_path / "shape-only.pptx"
    presentation = stylekit.new_wide_presentation()
    for slide_number in range(5):
        slide = _slide(presentation)
        _textbox(slide, f"Slide {slide_number + 1}", 1, 2, 5, 1)
    presentation.save(source)

    report = pptx_qa.inspect_pptx(source)

    assert any("shape-only treatment" in warning for warning in report["warnings"])


def test_structural_qa_rejects_overlapping_text_items(tmp_path: Path) -> None:
    source = tmp_path / "overlap.pptx"
    presentation = stylekit.new_wide_presentation()
    slide = _slide(presentation)
    _textbox(slide, "First message", 1, 2, 5, 1.2)
    _textbox(slide, "Second message", 4, 2.3, 5, 1.2)
    presentation.save(source)

    report = pptx_qa.inspect_pptx(source)

    assert any("overlap by" in error for error in report["errors"])
    assert report["layout"]["metrics"][0]["overlaps"]


def test_structural_qa_rejects_excessive_rounded_cards(tmp_path: Path) -> None:
    source = tmp_path / "cards.pptx"
    presentation = stylekit.new_wide_presentation()
    slide = _slide(presentation)
    for index in range(4):
        slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8 + index * 3.05),
            Inches(2),
            Inches(2.6),
            Inches(2.2),
        )
    presentation.save(source)

    report = pptx_qa.inspect_pptx(source)

    assert any("rounded rectangles" in error for error in report["errors"])


def test_template_baseline_suppresses_unchanged_design_debt(tmp_path: Path) -> None:
    source = tmp_path / "template.pptx"
    presentation = stylekit.new_wide_presentation()
    slide = _slide(presentation)
    for index in range(4):
        slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8 + index * 3.05),
            Inches(2),
            Inches(2.6),
            Inches(2.2),
        )
    presentation.save(source)

    report = pptx_qa.inspect_pptx(source, reference=source)

    assert not any("rounded rectangles" in error for error in report["errors"])
    assert report["layout"]["findings"] == []
