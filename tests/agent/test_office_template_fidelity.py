"""Regression tests for template-first Office edits."""

from __future__ import annotations

import zipfile
from io import BytesIO
from pathlib import Path

import pytest
from docx import Document
from openpyxl import Workbook, load_workbook
from openpyxl.styles import PatternFill
from pptx import Presentation
from pptx.util import Inches

from app.agent.builtin_skills.docx.scripts import accept_changes as docx_accept
from app.agent.builtin_skills.docx.scripts import template as docx_template
from app.agent.builtin_skills.pptx.scripts import office_features
from app.agent.builtin_skills.pptx.scripts import template as pptx_template
from app.agent.builtin_skills.template_fidelity import package_hashes, patch_package
from app.agent.builtin_skills.xlsx.scripts import template as xlsx_template


def _protected_hashes(path: Path, prefixes: tuple[str, ...]) -> dict[str, str]:
    return {
        part: digest
        for part, digest in package_hashes(path).items()
        if part.startswith(prefixes)
    }


def test_docx_accept_changes_uses_ooxml_and_preserves_other_parts(
    tmp_path: Path,
) -> None:
    pristine = tmp_path / "pristine.docx"
    source = tmp_path / "tracked.docx"
    output = tmp_path / "accepted.docx"
    document = Document()
    document.add_paragraph("Keep")
    document.add_paragraph("Base")
    document.save(pristine)
    with zipfile.ZipFile(pristine) as package:
        document_xml = package.read("word/document.xml")
        settings_xml = package.read("word/settings.xml")
    old_run = b"<w:r><w:t>Base</w:t></w:r>"
    tracked_run = (
        b'<w:del w:id="1" w:author="Reviewer"><w:r>'
        b"<w:delText>Deleted</w:delText></w:r></w:del>"
        b'<w:ins w:id="2" w:author="Reviewer"><w:r>'
        b"<w:t>Accepted</w:t></w:r></w:ins>"
    )
    assert old_run in document_xml
    settings_xml = settings_xml.replace(
        b"</w:settings>",
        b"<w:trackRevisions/></w:settings>",
    )
    patch_package(
        pristine,
        source,
        {
            "word/document.xml": document_xml.replace(old_run, tracked_run),
            "word/settings.xml": settings_xml,
        },
    )
    styles_before = package_hashes(source)["word/styles.xml"]

    _, message = docx_accept.accept_changes(str(source), str(output))

    assert message.startswith("Accepted ")
    accepted = Document(output)
    assert [paragraph.text for paragraph in accepted.paragraphs] == [
        "Keep",
        "Accepted",
    ]
    with zipfile.ZipFile(output) as package:
        accepted_xml = package.read("word/document.xml")
        accepted_settings = package.read("word/settings.xml")
    assert b"<w:ins" not in accepted_xml
    assert b"<w:del" not in accepted_xml
    assert b"trackRevisions" not in accepted_settings
    assert package_hashes(output)["word/styles.xml"] == styles_before


def test_pptx_template_edit_preserves_master_layout_and_theme(tmp_path: Path) -> None:
    source = tmp_path / "template.pptx"
    output = tmp_path / "edited.pptx"
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Original title"
    table = slide.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(5), Inches(2))
    table.table.cell(0, 0).text = "Old cell"
    deck.save(source)

    manifest = pptx_template.inspect(source)
    shapes = manifest["slides"][0]["shapes"]
    title_id = next(
        shape["shape_id"] for shape in shapes if shape["text"] == "Original title"
    )
    table_id = next(
        shape["shape_id"] for shape in shapes if shape["kind"] == "graphicFrame"
    )
    plan = {
        "edits": [
            {
                "slide": 1,
                "shape_id": title_id,
                "action": "replace_text",
                "text": "Edited title",
            },
            {
                "slide": 1,
                "shape_id": table_id,
                "action": "replace_table_cell",
                "row": 0,
                "column": 0,
                "text": "Edited cell",
            },
        ]
    }

    report = pptx_template.apply(source, output, plan)

    assert report["errors"] == []
    assert report["changed_parts"] == ["ppt/slides/slide1.xml"]
    assert _protected_hashes(
        source, ("ppt/slideMasters/", "ppt/slideLayouts/", "ppt/theme/")
    ) == _protected_hashes(
        output, ("ppt/slideMasters/", "ppt/slideLayouts/", "ppt/theme/")
    )
    edited = Presentation(output)
    assert edited.slides[0].shapes.title.text == "Edited title"
    edited_table = next(
        shape.table for shape in edited.slides[0].shapes if shape.has_table
    )
    assert edited_table.cell(0, 0).text == "Edited cell"
    assert pptx_template.verify(source, output, plan)["errors"] == []

    tampered = tmp_path / "tampered.pptx"
    with zipfile.ZipFile(output) as package:
        slide_xml = package.read("ppt/slides/slide1.xml").replace(
            b"Edited title", b"Unplanned title"
        )
    patch_package(output, tampered, {"ppt/slides/slide1.xml": slide_xml})
    mismatch = pptx_template.verify(source, tampered, plan)
    assert mismatch["plan_mismatches"] == ["ppt/slides/slide1.xml"]
    assert mismatch["errors"]


def test_pptx_template_edits_chart_workbook_without_rebuilding_slide(
    tmp_path: Path,
) -> None:
    source = tmp_path / "chart-template.pptx"
    output = tmp_path / "chart-edited.pptx"
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    office_features.add_native_chart(
        slide,
        ["Q1", "Q2"],
        {"Actual": [10, 12], "Plan": [11, 13]},
        left=Inches(1),
        top=Inches(1),
        width=Inches(7),
        height=Inches(4),
        kind="line",
    )
    deck.save(source)
    manifest = pptx_template.inspect(source)
    chart_shape = next(
        shape for shape in manifest["slides"][0]["shapes"] if shape["has_chart"]
    )
    plan = {
        "edits": [
            {
                "slide": 1,
                "shape_id": chart_shape["shape_id"],
                "action": "replace_chart_data",
                "categories": ["H1", "H2", "H3"],
                "series": [
                    {"name": "Actual", "values": [21, 25, 30]},
                    {"name": "Plan", "values": [22, 24, 28]},
                ],
            }
        ]
    }

    report = pptx_template.apply(source, output, plan)

    assert report["errors"] == []
    assert len(report["changed_parts"]) == 2
    assert any(part.startswith("ppt/charts/chart") for part in report["changed_parts"])
    assert any(part.startswith("ppt/embeddings/") for part in report["changed_parts"])
    assert "ppt/slides/slide1.xml" not in report["changed_parts"]
    edited = Presentation(output)
    chart = next(shape.chart for shape in edited.slides[0].shapes if shape.has_chart)
    assert [series.name for series in chart.series] == ["Actual", "Plan"]
    assert [list(series.values) for series in chart.series] == [
        [21.0, 25.0, 30.0],
        [22.0, 24.0, 28.0],
    ]
    assert list(chart.plots[0].categories.flattened_labels) == [
        ("H1",),
        ("H2",),
        ("H3",),
    ]
    with zipfile.ZipFile(output) as package:
        workbook_part = next(
            name
            for name in package.namelist()
            if name.startswith("ppt/embeddings/") and name.endswith(".xlsx")
        )
        workbook = load_workbook(BytesIO(package.read(workbook_part)))
    assert workbook.active["A2"].value == "H1"
    assert workbook.active["C4"].value == 28
    assert pptx_template.verify(source, output, plan)["errors"] == []


def test_pptx_template_fills_placeholder_and_replaces_rich_text_runs(
    tmp_path: Path,
) -> None:
    source = tmp_path / "rich-template.pptx"
    output = tmp_path / "rich-edited.pptx"
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[5])
    slide.shapes.title.text = "Template title"
    body = slide.shapes.add_textbox(
        Inches(1),
        Inches(2),
        Inches(6),
        Inches(2),
    )
    body.text = "Template body"
    deck.save(source)
    manifest = pptx_template.inspect(source)
    shapes = manifest["slides"][0]["shapes"]
    title_shape = next(shape for shape in shapes if shape["placeholder"])
    body_shape = next(shape for shape in shapes if shape["text"] == "Template body")
    plan = {
        "edits": [
            {
                "slide": 1,
                "shape_id": title_shape["shape_id"],
                "action": "fill_placeholder",
                "text": "Filled title",
            },
            {
                "slide": 1,
                "shape_id": body_shape["shape_id"],
                "action": "replace_rich_text",
                "paragraphs": [
                    {
                        "bullet": True,
                        "runs": [
                            {
                                "text": "Lead",
                                "bold": True,
                                "color": "2F6D68",
                                "size": 22,
                            },
                            {"text": " with supporting evidence."},
                        ],
                    },
                    {
                        "bullet": True,
                        "level": 1,
                        "runs": [{"text": "Editable detail", "italic": True}],
                    },
                ],
            },
        ]
    }

    report = pptx_template.apply(source, output, plan)

    assert report["errors"] == []
    assert report["changed_parts"] == ["ppt/slides/slide1.xml"]
    edited = Presentation(output)
    assert edited.slides[0].shapes.title.text == "Filled title"
    edited_body = next(
        shape
        for shape in edited.slides[0].shapes
        if shape.has_text_frame and shape.text.startswith("Lead")
    )
    assert edited_body.text == ("Lead with supporting evidence.\nEditable detail")
    assert edited_body.text_frame.paragraphs[0].runs[0].font.bold
    assert edited_body.text_frame.paragraphs[0].runs[0].font.size.pt == 22
    assert edited_body.text_frame.paragraphs[1].runs[0].font.italic
    with zipfile.ZipFile(output) as package:
        slide_xml = package.read("ppt/slides/slide1.xml")
    assert slide_xml.count(b":buChar") == 2
    assert pptx_template.verify(source, output, plan)["errors"] == []


def test_docx_template_edit_preserves_styles_and_headers(tmp_path: Path) -> None:
    source = tmp_path / "template.docx"
    output = tmp_path / "edited.docx"
    document = Document()
    document.add_heading("Original heading", level=1)
    table = document.add_table(rows=1, cols=1)
    table.cell(0, 0).text = "Original cell"
    document.sections[0].header.paragraphs[0].text = "Brand header"
    document.save(source)
    plan = {
        "edits": [
            {
                "action": "replace_paragraph",
                "paragraph": 0,
                "text": "Edited heading",
            },
            {
                "action": "replace_table_cell",
                "table": 0,
                "row": 0,
                "column": 0,
                "text": "Edited cell",
            },
        ]
    }

    report = docx_template.apply(source, output, plan)

    assert report["errors"] == []
    assert report["changed_parts"] == ["word/document.xml"]
    before = package_hashes(source)
    after = package_hashes(output)
    assert before["word/styles.xml"] == after["word/styles.xml"]
    header_part = next(part for part in before if part.startswith("word/header"))
    assert before[header_part] == after[header_part]
    edited = Document(output)
    assert edited.paragraphs[0].text == "Edited heading"
    assert edited.tables[0].cell(0, 0).text == "Edited cell"
    assert edited.sections[0].header.paragraphs[0].text == "Brand header"
    assert docx_template.verify(source, output, plan)["errors"] == []


def test_xlsx_template_edit_preserves_style_formula_and_charts_parts(
    tmp_path: Path,
) -> None:
    source = tmp_path / "template.xlsx"
    output = tmp_path / "edited.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Model"
    sheet["A1"] = "Original"
    sheet["A1"].fill = PatternFill("solid", fgColor="00AAFF")
    sheet["B1"] = 10
    sheet["C1"] = "=B1*2"
    workbook.save(source)
    original_style_id = load_workbook(source)["Model"]["A1"].style_id
    plan = {
        "edits": [
            {
                "action": "set_cell",
                "sheet": "Model",
                "cell": "A1",
                "kind": "string",
                "value": "Edited",
            },
            {
                "action": "set_cell",
                "sheet": "Model",
                "cell": "B1",
                "kind": "number",
                "value": 25,
            },
        ]
    }

    report = xlsx_template.apply(source, output, plan)

    assert report["errors"] == []
    assert report["changed_parts"] == ["xl/worksheets/sheet1.xml"]
    before = package_hashes(source)
    after = package_hashes(output)
    assert before["xl/styles.xml"] == after["xl/styles.xml"]
    edited = load_workbook(output, data_only=False)
    assert edited["Model"]["A1"].value == "Edited"
    assert edited["Model"]["A1"].style_id == original_style_id
    assert edited["Model"]["B1"].value == 25
    assert edited["Model"]["C1"].value == "=B1*2"
    assert xlsx_template.verify(source, output, plan)["errors"] == []


def test_package_patch_refuses_to_overwrite_source(tmp_path: Path) -> None:
    source = tmp_path / "source.zip"
    with zipfile.ZipFile(source, "w") as package:
        package.writestr("part.xml", b"before")

    with pytest.raises(ValueError, match="overwrite"):
        patch_package(source, source, {"part.xml": b"after"})
