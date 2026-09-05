"""The document_preview tool exposes the host viewer engine to the agent."""

from __future__ import annotations

from pathlib import Path

import pytest

from app.agent.sandbox import SandboxConfig, get_sandbox, set_sandbox
from app.agent.tools.builtin.document_preview import (
    _document_preview,
    document_preview,
)

pytest.importorskip("pptx", reason="office-preview extra not installed")


@pytest.fixture
def workspace(tmp_path, monkeypatch):
    root = tmp_path / "workspace"
    root.mkdir()
    token = set_sandbox(SandboxConfig(workspace=str(root)))
    yield root
    try:
        from app.agent import sandbox as sandbox_module

        sandbox_module._sandbox.reset(token)
    except Exception:  # pragma: no cover - contextvar reset is best effort
        pass


def _deck(path: Path, *, stray: bool = False) -> Path:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "a title that belongs here"
    if stray:
        off = slide.shapes.add_textbox(Inches(11), Inches(6.5), Inches(4), Inches(2))
        off.text_frame.text = "hangs off the edge"
    presentation.save(path)
    return path


@pytest.mark.asyncio
async def test_reports_pages_text_and_geometry(workspace):
    _deck(workspace / "deck.pptx")

    report = await _document_preview("deck.pptx")

    assert "[Document preview: deck.pptx]" in report
    assert "Pages: 1" in report
    assert "a title that belongs here" in report
    assert "Every laid-out element sits inside its page box." in report


@pytest.mark.asyncio
async def test_flags_an_element_outside_the_page(workspace):
    _deck(workspace / "broken.pptx", stray=True)

    report = await _document_preview("broken.pptx")

    assert "Elements outside the page box:" in report
    assert "OUTSIDE PAGE" in report
    assert "hangs off the edge" in report


@pytest.mark.asyncio
async def test_verbose_includes_elements_without_text(workspace):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    labelled = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    labelled.text_frame.text = "has text"
    slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1), Inches(3), Inches(2), Inches(1)
    )
    presentation.save(workspace / "mixed.pptx")

    quiet = await _document_preview("mixed.pptx")
    loud = await _document_preview("mixed.pptx", verbose=True)

    assert loud.count("\n") > quiet.count("\n")
    assert "has text" in quiet


@pytest.mark.asyncio
async def test_unsupported_format_names_the_supported_ones(workspace):
    (workspace / "notes.txt").write_text("plain", encoding="utf-8")

    report = await _document_preview("notes.txt")

    assert "[Cannot preview: notes.txt]" in report
    for suffix in (".docx", ".pdf", ".pptx", ".xlsx"):
        assert suffix in report


@pytest.mark.asyncio
async def test_missing_and_directory_paths_raise(workspace):
    (workspace / "folder").mkdir()

    with pytest.raises(FileNotFoundError):
        await _document_preview("absent.pptx")
    with pytest.raises(IsADirectoryError):
        await _document_preview("folder")


@pytest.mark.asyncio
async def test_path_outside_the_workspace_is_denied(workspace, tmp_path):
    outside = _deck(tmp_path / "outside.pptx")

    with pytest.raises(PermissionError):
        await _document_preview(str(outside))


@pytest.mark.asyncio
async def test_reports_an_absolute_rendered_path(workspace):
    """The cache root may be configured relative to the process directory."""

    _deck(workspace / "deck.pptx")

    report = await _document_preview("deck.pptx")
    rendered = Path(report.split("engine: ", 1)[1].split(" (", 1)[0])

    assert rendered.is_absolute()
    assert rendered.is_file()


@pytest.mark.asyncio
async def test_rendering_does_not_widen_the_sandbox(workspace):
    """The report carries the content, so the cache stays out of reach."""

    _deck(workspace / "deck.pptx")
    sandbox = get_sandbox()

    report = await _document_preview("deck.pptx")
    rendered = Path(report.split("engine: ", 1)[1].split(" (", 1)[0])

    with pytest.raises(PermissionError):
        sandbox.validate_path(str(rendered))


@pytest.mark.asyncio
async def test_repeated_calls_reuse_the_cache_without_failing(workspace):
    _deck(workspace / "deck.pptx")

    first = await _document_preview("deck.pptx")
    second = await _document_preview("deck.pptx")

    assert first == second


def test_tool_is_registered_read_only_with_the_supported_formats():
    assert document_preview.name == "document_preview"
    assert document_preview.read_only is True
    for suffix in (".docx", ".pdf", ".pptx", ".xlsx"):
        assert suffix in document_preview.description


def test_tool_is_wired_into_the_agent_registry():
    from app.agent.loader import _default_tool_registry

    assert "document_preview" in _default_tool_registry()
